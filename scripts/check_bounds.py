#!/usr/bin/env python3
"""Check that every shape in a .pptx lies within the 16:9 widescreen canvas."""
import argparse
import json
import subprocess
import sys

try:
    from pptx import Presentation
except ImportError:
    subprocess.run(
        [sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"],
        check=True,
    )
    from pptx import Presentation

CANVAS_W = 12192000  # 13.333" in EMU
CANVAS_H = 6858000   # 7.5" in EMU


def shape_label(shape):
    name = getattr(shape, "name", None)
    if name:
        return name
    if getattr(shape, "is_placeholder", False):
        try:
            return f"placeholder:{shape.placeholder_format.type}"
        except Exception:
            return "placeholder"
    return "shape"


def check(path):
    prs = Presentation(path)
    violations = []
    for sidx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            L = shape.left
            T = shape.top
            W = shape.width
            H = shape.height
            if L is None or T is None or W is None or H is None:
                continue
            edges = []
            if L < 0:
                edges.append("left<0")
            if T < 0:
                edges.append("top<0")
            if L + W > CANVAS_W:
                edges.append("right>canvas")
            if T + H > CANVAS_H:
                edges.append("bottom>canvas")
            if edges:
                violations.append({
                    "slide": sidx,
                    "shape": shape_label(shape),
                    "bounds": [L, T, W, H],
                    "violation": ",".join(edges),
                })
    return violations


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()
    violations = check(args.pptx)
    if args.json:
        print(json.dumps({"violations": violations}, indent=2))
    else:
        for v in violations:
            L, T, W, H = v["bounds"]
            print(
                f"slide={v['slide']} shape={v['shape']} "
                f"bounds=({L},{T},{W},{H}) violation={v['violation']}"
            )
    sys.exit(1 if violations else 0)


if __name__ == "__main__":
    main()
