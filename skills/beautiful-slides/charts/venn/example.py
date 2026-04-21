"""Render the venn diagram in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_2SET = {
    "title": "Engineering skill overlap",
    "sets": [
        {"label": "Frontend", "items": ["React", "CSS", "A11y"]},
        {"label": "Backend", "items": ["Go", "SQL", "K8s"]},
    ],
    "intersections": {
        "ab": "REST APIs",
    },
}

DATA_3SET_SKILLS = {
    "title": "Cross-functional team capabilities",
    "sets": [
        {"label": "Design", "items": ["Figma", "UX research"]},
        {"label": "Engineering", "items": ["Code", "CI/CD"]},
        {"label": "Product", "items": ["Roadmap", "Metrics"]},
    ],
    "intersections": {
        "ab": "Design systems",
        "ac": "User stories",
        "bc": "Analytics",
        "abc": "Ship great products",
    },
}

DATA_3SET_STRATEGY = {
    "title": "Strategic sweet spot",
    "transparency": 60,
    "sets": [
        {"label": "Desirable", "items": ["User need", "Market pull"]},
        {"label": "Feasible", "items": ["Tech ready", "Team skills"]},
        {"label": "Viable", "items": ["Revenue", "Unit economics"]},
    ],
    "intersections": {
        "ab": "Prototype",
        "ac": "Pitch deck",
        "bc": "Infra play",
        "abc": "Build this",
    },
}


MODE_DATA = {
    "sv-keynote": DATA_3SET_SKILLS,
    "editorial-magazine": DATA_2SET,
    "playful-marketing": DATA_3SET_STRATEGY,
    "consulting-boardroom": DATA_3SET_SKILLS,
    "craft-minimal": DATA_2SET,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
