"""Venn diagram renderer — native python-pptx shapes only.

2-3 overlapping circles with semi-transparent fills showing set relationships.
Labels placed in each distinct region (A-only, B-only, intersection, etc).

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn


# ---- helpers ---------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 word_wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Emu(int(Pt(2).emu))
    tf.margin_right = Emu(int(Pt(2).emu))
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_textbox_multiline(slide, x, y, w, h, lines, font_name, sizes_pt,
                           colors_hex, bolds, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE):
    """Add a textbox with multiple paragraphs (one per line)."""
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(int(Pt(2).emu))
    tf.margin_right = Emu(int(Pt(2).emu))
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = str(line)
        run.font.name = font_name
        run.font.size = Pt(sizes_pt[i] if i < len(sizes_pt) else sizes_pt[-1])
        run.font.bold = bolds[i] if i < len(bolds) else bolds[-1]
        c = colors_hex[i] if i < len(colors_hex) else colors_hex[-1]
        run.font.color.rgb = _rgb(c)
    return tb


def _add_oval_transparent(slide, x, y, w, h, fill_hex, alpha_pct=50,
                          line_hex=None, line_width_pt=1.5):
    """Add an oval with transparent fill. alpha_pct is 0 (opaque) to 100 (invisible)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    # Set transparency via XML (python-pptx doesn't expose fill transparency)
    sp = shp._element
    solid = sp.find('.//' + qn('a:solidFill'))
    if solid is not None:
        srgb = solid.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_val = int((100 - alpha_pct) * 1000)  # 0-100000 scale
            alpha_elem = srgb.makeelement(qn('a:alpha'), {'val': str(alpha_val)})
            srgb.append(alpha_elem)

    if line_hex is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Pt(line_width_pt)
    shp.shadow.inherit = False
    return shp


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# ---- main render -----------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a Venn diagram onto *slide* inside *bounds*, styled by *tokens*."""
    x0, y0, w0, h0 = bounds

    # --- unpack tokens -------------------------------------------------------
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = int(tokens["font_size_base_pt"])

    # --- unpack data ---------------------------------------------------------
    title = data.get("title")
    sets = list(data.get("sets") or [])[:3]  # max 3 circles
    # Each set: {label: str, color: str (optional)}
    # Intersection labels
    intersections = data.get("intersections") or {}
    # For 2-set: {"ab": "Overlap text"}
    # For 3-set: {"ab": ..., "ac": ..., "bc": ..., "abc": ...}
    transparency = data.get("transparency", 65)  # 0-100

    if len(sets) < 2:
        return

    n_sets = min(len(sets), 3)

    # --- background ----------------------------------------------------------
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    # --- outer padding -------------------------------------------------------
    pad = int(min(w0, h0) * 0.035)
    ix, iy = x0 + pad, y0 + pad
    iw, ih = w0 - 2 * pad, h0 - 2 * pad

    # --- title ---------------------------------------------------------------
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.55))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, ix, iy, iw, title_h, title,
                     font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + int(pad * 0.3)
        ih = (y0 + h0 - pad) - iy

    # --- compute circle layout -----------------------------------------------
    # Available drawing area
    draw_x = ix
    draw_y = iy
    draw_w = iw
    draw_h = ih

    # Circle colors: use set-specific or fall back to primary/accent/muted
    default_colors = [primary, accent, muted]

    label_pt = max(9, int(base_pt * 0.9))
    region_label_pt = max(8, int(base_pt * 0.78))
    set_label_pt = max(10, int(base_pt * 1.0))
    label_h = int(Pt(label_pt).emu * 2.0)

    if n_sets == 2:
        # --- 2-circle Venn: side-by-side with overlap -----------------------
        # Overlap fraction controls how much the circles intersect
        overlap_frac = 0.35

        # Size circles to fit the drawing area
        # Two overlapping circles: total width = 2*d - overlap*d = d*(2-overlap)
        max_d_from_w = draw_w / (2.0 - overlap_frac)
        max_d_from_h = draw_h * 0.85  # leave room for labels
        d = min(max_d_from_w, max_d_from_h)
        r = d / 2

        # Center the pair vertically and horizontally
        total_w = d * (2.0 - overlap_frac)
        cx_start = draw_x + (draw_w - total_w) / 2
        cy = draw_y + draw_h * 0.48  # slightly above center for labels below

        c1x = cx_start + r
        c2x = cx_start + total_w - r
        c1y = c2y = cy

        colors = [
            sets[0].get("color") or default_colors[0],
            sets[1].get("color") or default_colors[1],
        ]

        # Draw circles (draw in order so overlap shows)
        _add_oval_transparent(slide,
                              c1x - r, c1y - r, d, d,
                              colors[0], transparency,
                              line_hex=colors[0], line_width_pt=1.5)
        _add_oval_transparent(slide,
                              c2x - r, c2y - r, d, d,
                              colors[1], transparency,
                              line_hex=colors[1], line_width_pt=1.5)

        # Labels for A-only region (left side of circle A)
        region_w = int(r * 0.9)
        region_h = int(Pt(region_label_pt).emu * 3.0)

        def _trunc_items_2set(items_list, max_chars=20):
            result = []
            for item in items_list[:4]:
                s = str(item)
                if len(s) > max_chars:
                    s = s[:max_chars - 3] + "..."
                result.append(s)
            return result

        if sets[0].get("items"):
            items_text = "\n".join(_trunc_items_2set(sets[0]["items"]))
            _add_textbox(slide,
                         c1x - r * 0.75, c1y - region_h / 2,
                         region_w, region_h,
                         items_text, font_body, region_label_pt, text_c,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        if sets[1].get("items"):
            items_text = "\n".join(_trunc_items_2set(sets[1]["items"]))
            _add_textbox(slide,
                         c2x + r * 0.75 - region_w, c2y - region_h / 2,
                         region_w, region_h,
                         items_text, font_body, region_label_pt, text_c,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Intersection label
        ab_label = intersections.get("ab") or intersections.get("AB") or ""
        if ab_label:
            mid_x = (c1x + c2x) / 2
            int_w = int(d * overlap_frac * 0.7)
            int_h = int(Pt(region_label_pt).emu * 3.0)
            _add_textbox(slide,
                         mid_x - int_w / 2, cy - int_h / 2,
                         int_w, int_h,
                         ab_label, font_body, region_label_pt, text_c,
                         bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)

        # Set name labels (above or below circles) — clamped to bounds
        set_label_h = int(Pt(set_label_pt).emu * 1.8)
        set_label_w = int(r * 1.6)

        sl_x = max(x0, min(x0 + w0 - set_label_w, c1x - set_label_w / 2))
        sl_y = max(y0, c1y - r - set_label_h - int(pad * 0.2))
        _add_textbox(slide, sl_x, sl_y, set_label_w, set_label_h,
                     sets[0].get("label", "A"), font_display, set_label_pt,
                     colors[0],
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        sl_x = max(x0, min(x0 + w0 - set_label_w, c2x - set_label_w / 2))
        sl_y = max(y0, c2y - r - set_label_h - int(pad * 0.2))
        _add_textbox(slide, sl_x, sl_y, set_label_w, set_label_h,
                     sets[1].get("label", "B"), font_display, set_label_pt,
                     colors[1],
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

    elif n_sets == 3:
        # --- 3-circle Venn: triangular arrangement --------------------------
        # Place 3 circle centers in an equilateral-triangle pattern
        overlap_frac = 0.35

        # The three centers form a triangle; the bounding width is roughly
        # 2*r + separation. We compute so they fit in draw area.
        # Triangle height = separation * sin(60) = separation * 0.866
        # We want: 2*r + sep <= draw_w  and  2*r + sep*0.866 <= draw_h*0.85
        # Also: sep = d * (1 - overlap_frac)

        max_d_from_w = draw_w / (2.0 - overlap_frac)
        max_d_from_h = draw_h * 0.80 / (1.0 + (1 - overlap_frac) * 0.866 / 2 + 0.5)
        d = min(max_d_from_w, max_d_from_h)
        r = d / 2
        sep = d * (1 - overlap_frac)  # center-to-center distance

        # Center of the triangle arrangement
        center_x = draw_x + draw_w / 2
        center_y = draw_y + draw_h * 0.48

        # Three centers: top, bottom-left, bottom-right
        tri_h = sep * math.sin(math.radians(60))
        c1x = center_x
        c1y = center_y - tri_h / 3
        c2x = center_x - sep / 2
        c2y = center_y + tri_h * 2 / 3
        c3x = center_x + sep / 2
        c3y = center_y + tri_h * 2 / 3

        colors = [
            sets[0].get("color") or default_colors[0],
            sets[1].get("color") or default_colors[1],
            sets[2].get("color") or default_colors[2],
        ]

        centers = [(c1x, c1y), (c2x, c2y), (c3x, c3y)]

        # Draw circles
        for i, (cx, cy) in enumerate(centers):
            _add_oval_transparent(slide,
                                  cx - r, cy - r, d, d,
                                  colors[i], transparency,
                                  line_hex=colors[i], line_width_pt=1.5)

        # Set name labels — clamped to stay within slide bounds
        set_label_h = int(Pt(set_label_pt).emu * 1.8)
        set_label_w = int(r * 1.6)

        # Helper to clamp label position within bounds
        def _clamp_lbl(lx, ly, lw, lh):
            lx = max(x0, min(x0 + w0 - lw, lx))
            ly = max(y0, min(y0 + h0 - lh, ly))
            return lx, ly

        # Top circle label: above
        sl_x = c1x - set_label_w / 2
        sl_y = c1y - r - set_label_h - int(pad * 0.2)
        sl_x, sl_y = _clamp_lbl(sl_x, sl_y, set_label_w, set_label_h)
        _add_textbox(slide, sl_x, sl_y, set_label_w, set_label_h,
                     sets[0].get("label", "A"), font_display, set_label_pt,
                     colors[0],
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        # Bottom-left label: below-left
        sl_x = c2x - set_label_w - int(pad * 0.1)
        sl_y = c2y + r + int(pad * 0.1)
        sl_x, sl_y = _clamp_lbl(sl_x, sl_y, set_label_w, set_label_h)
        _add_textbox(slide, sl_x, sl_y, set_label_w, set_label_h,
                     sets[1].get("label", "B"), font_display, set_label_pt,
                     colors[1],
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # Bottom-right label: below-right
        sl_x = c3x + int(pad * 0.1)
        sl_y = c3y + r + int(pad * 0.1)
        sl_x, sl_y = _clamp_lbl(sl_x, sl_y, set_label_w, set_label_h)
        _add_textbox(slide, sl_x, sl_y, set_label_w, set_label_h,
                     sets[2].get("label", "C"), font_display, set_label_pt,
                     colors[2],
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        # Region labels in each exclusive area
        # Auto-scale font when all 3 sets have items to reduce overlap
        any_items = any(s.get("items") for s in sets[:3])
        if any_items:
            region_label_pt = max(7, int(region_label_pt * 0.85))
        region_h = int(Pt(region_label_pt).emu * 2.5)
        region_w = int(r * 0.55)  # narrower than before to avoid overlap

        def _truncate_items(items_list, max_chars=20):
            result = []
            for item in items_list[:3]:
                s = str(item)
                if len(s) > max_chars:
                    s = s[:max_chars - 3] + "..."
                result.append(s)
            return result

        # A-only: above center, pushed further up
        if sets[0].get("items"):
            items_text = "\n".join(_truncate_items(sets[0]["items"]))
            _add_textbox(slide,
                         c1x - region_w / 2, c1y - r * 0.70,
                         region_w, region_h,
                         items_text, font_body, region_label_pt, text_c,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # B-only: bottom-left, pushed further left
        if sets[1].get("items"):
            items_text = "\n".join(_truncate_items(sets[1]["items"]))
            _add_textbox(slide,
                         c2x - r * 0.70, c2y + r * 0.05,
                         region_w, region_h,
                         items_text, font_body, region_label_pt, text_c,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # C-only: bottom-right, pushed further right
        if sets[2].get("items"):
            items_text = "\n".join(_truncate_items(sets[2]["items"]))
            _add_textbox(slide,
                         c3x + r * 0.70 - region_w, c3y + r * 0.05,
                         region_w, region_h,
                         items_text, font_body, region_label_pt, text_c,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Pairwise intersections
        int_w = int(r * 0.7)
        int_h = int(Pt(region_label_pt).emu * 2.5)

        def _trunc(s, n=20):
            return s if len(s) <= n else s[:n - 3] + "..."

        ab_label = intersections.get("ab") or intersections.get("AB") or ""
        if ab_label:
            ab_label = _trunc(ab_label)
            mx = (c1x + c2x) / 2
            my = (c1y + c2y) / 2
            _add_textbox(slide, mx - int_w / 2, my - int_h / 2,
                         int_w, int_h,
                         ab_label, font_body, region_label_pt, text_c,
                         bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)

        ac_label = intersections.get("ac") or intersections.get("AC") or ""
        if ac_label:
            ac_label = _trunc(ac_label)
            mx = (c1x + c3x) / 2
            my = (c1y + c3y) / 2
            _add_textbox(slide, mx - int_w / 2, my - int_h / 2,
                         int_w, int_h,
                         ac_label, font_body, region_label_pt, text_c,
                         bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)

        bc_label = intersections.get("bc") or intersections.get("BC") or ""
        if bc_label:
            bc_label = _trunc(bc_label)
            mx = (c2x + c3x) / 2
            my = (c2y + c3y) / 2
            _add_textbox(slide, mx - int_w / 2, my - int_h / 2,
                         int_w, int_h,
                         bc_label, font_body, region_label_pt, text_c,
                         bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)

        # Triple intersection (center of all three)
        abc_label = intersections.get("abc") or intersections.get("ABC") or ""
        if abc_label:
            abc_label = _trunc(abc_label)
            mx = (c1x + c2x + c3x) / 3
            my = (c1y + c2y + c3y) / 3
            _add_textbox(slide, mx - int_w / 2, my - int_h / 2,
                         int_w, int_h,
                         abc_label, font_body, region_label_pt, text_c,
                         bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)
