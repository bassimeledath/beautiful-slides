import os
import sys
from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


def _data():
    return {
        "title": "Q2 program status",
        "rows": [
            "Auth service",
            "Payment gateway",
            "Analytics pipeline",
            "Mobile app",
            "Infrastructure",
            "Documentation",
        ],
        "columns": ["Schedule", "Budget", "Quality", "Risk", "Dependencies"],
        "statuses": [
            ["G", "G", "A", "G", "G"],
            ["A", "R", "G", "A", "R"],
            ["G", "G", "G", "G", "A"],
            ["R", "A", "A", "R", "G"],
            ["G", "G", "G", "A", "G"],
            ["A", "G", "R", "G", "G"],
        ],
        "show_labels": True,
    }


def _make_pptx(mode_name, tokens, data, out_dir):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg_shape = slide.shapes.add_shape(
        1,
        0, 0,
        prs.slide_width, prs.slide_height,
    )
    bg_shape.shadow.inherit = False
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(tokens["bg"])
    bg_shape.line.fill.background()
    bg_shape.text_frame.text = ""

    margin = Inches(0.6)
    bounds = (
        margin,
        margin,
        prs.slide_width - 2 * margin,
        prs.slide_height - 2 * margin,
    )

    render(slide, data, tokens, bounds)

    out_path = os.path.join(out_dir, f"example-{mode_name}.pptx")
    prs.save(out_path)
    return out_path


def main():
    data = _data()
    written = []
    for name, tokens in MODES.items():
        path = _make_pptx(name, tokens, data, HERE)
        written.append(path)
        print(f"wrote {path}")
    print(f"OK: {len(written)} files")


if __name__ == "__main__":
    main()
