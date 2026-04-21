from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    t = max(0.0, min(1.0, t))
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _luminance(hex_):
    h = hex_.lstrip("#")
    r = int(h[0:2], 16) / 255.0
    g = int(h[2:4], 16) / 255.0
    b = int(h[4:6], 16) / 255.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _set_text(tf, text, font_name, size_pt, color_hex, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        for r in list(p.runs):
            r.text = ""
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)


def _add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_w_emu=0, radius_emu=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_emu > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    if radius_emu > 0:
        try:
            short = min(w, h)
            adj = max(0.0, min(0.5, (radius_emu / short)))
            shp.adjustments[0] = adj
        except Exception:
            pass
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None or line_w_emu <= 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Emu(int(line_w_emu))
    shp.text_frame.text = ""
    tf = shp.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return shp


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    _set_text(tb.text_frame, text, font_name, size_pt, color_hex, bold=bold, align=align, anchor=anchor)
    return tb


def render(slide, data, tokens, bounds):
    """Render a RAG (Red/Amber/Green) status matrix.

    Rows = workstreams/items, columns = criteria/milestones.
    Each cell contains a colored status indicator (R/A/G) derived from token
    colors. Statuses: "R" (red/primary), "A" (amber/lerp), "G" (green/accent),
    or "grey"/"none" for not-applicable.
    """
    x, y, w, h = bounds

    rows = list(data.get("rows", []))
    columns = list(data.get("columns", []))
    statuses = data.get("statuses", [])
    title = data.get("title")
    show_labels = bool(data.get("show_labels", True))

    n_rows = len(rows)
    n_cols = len(columns)
    if n_rows == 0 or n_cols == 0:
        return None

    # Extract tokens
    bg_hex = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(bg_hex.lstrip("#"))
    primary_hex = tokens["primary"]
    accent_hex = tokens["accent"]
    text_hex = tokens["text"]
    muted_hex = tokens["muted"]
    font_body = tokens["font_body"]
    font_display = tokens["font_display"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg_hex)
    bg_shape.line.fill.background()

    radius_emu = radius_px * 9525

    # RAG colors — hardcoded semantic colors so that R=red, A=amber, G=green
    # regardless of the slide theme. Dark-mode variants are slightly brighter.
    is_dark = _luminance(bg_hex) < 0.35
    if is_dark:
        rag_colors = {
            "R": "#EF5350",  # red (bright for dark bg)
            "A": "#FFB74D",  # amber (bright for dark bg)
            "G": "#66BB6A",  # green (bright for dark bg)
        }
    else:
        rag_colors = {
            "R": "#E53935",  # red
            "A": "#F5A623",  # amber
            "G": "#43A047",  # green
        }
    grey_hex = _lerp_hex(bg_hex, muted_hex, 0.3)

    # Font sizes
    label_pt = max(7, int(round(base_pt * 0.72)))
    title_pt = max(base_pt + 2, int(round(base_pt * 1.15)))
    status_label_pt = max(7, int(round(base_pt * 0.65)))
    header_pt = max(7, int(round(base_pt * 0.65)))

    # Layout
    title_h = 0
    if title:
        title_h = int(Pt(title_pt * 1.6))

    # Row labels on left
    row_label_w = int(Pt(base_pt * 3.0))
    for lbl in rows:
        est = int(Pt(base_pt * 0.55 * max(3, len(str(lbl)))))
        if est > row_label_w:
            row_label_w = est
    row_label_w = min(row_label_w, int(w * 0.30))

    col_label_h = int(Pt(base_pt * 1.8))

    # Legend at bottom
    legend_h = int(Pt(base_pt * 1.8))

    inner_x = x
    inner_y = y + title_h
    inner_w = w
    inner_h = h - title_h - legend_h

    grid_x = inner_x + row_label_w
    grid_y = inner_y + col_label_h
    grid_w = inner_w - row_label_w
    grid_h = inner_h - col_label_h

    gap = max(1, int(Emu(0.5 * 9525)))
    cell_w = (grid_w - gap * (n_cols - 1)) / n_cols
    cell_h = (grid_h - gap * (n_rows - 1)) / n_rows

    # Title
    if title:
        _add_textbox(
            slide,
            x, y, w, title_h,
            title,
            font_display,
            title_pt,
            text_hex,
            bold=True,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )

    # Column headers
    for ci, cl in enumerate(columns):
        cx = grid_x + ci * (cell_w + gap)
        _add_textbox(
            slide,
            int(cx), int(inner_y), int(cell_w), int(col_label_h),
            cl,
            font_body,
            header_pt,
            text_hex,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.BOTTOM,
        )

    # Row labels — truncate with ellipsis if too long for the allocated width
    max_label_chars = max(10, int(row_label_w / (Pt(label_pt) * 0.55)))
    for ri, rl in enumerate(rows):
        ry = grid_y + ri * (cell_h + gap)
        pad_right = int(Pt(base_pt * 0.3))
        display_label = str(rl)
        if len(display_label) > max_label_chars:
            display_label = display_label[:max_label_chars - 1] + "\u2026"
        _add_textbox(
            slide,
            int(inner_x), int(ry), int(row_label_w - pad_right), int(cell_h),
            display_label,
            font_body,
            label_pt,
            text_hex,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Draw cells with RAG indicators
    for ri in range(n_rows):
        row_statuses = statuses[ri] if ri < len(statuses) else []
        for ci in range(n_cols):
            cx = grid_x + ci * (cell_w + gap)
            cy = grid_y + ri * (cell_h + gap)

            # Get status for this cell
            status = row_statuses[ci] if ci < len(row_statuses) else None
            status_key = str(status).upper().strip() if status else ""

            # Draw cell background
            _add_rect(
                slide,
                cx, cy, cell_w, cell_h,
                bg_hex,
                line_hex=muted_hex,
                line_w_emu=int(Emu(0.25 * 9525)),
                radius_emu=radius_emu,
            )

            if status_key not in ("R", "A", "G"):
                # Empty or N/A — just the background cell
                if status_key and status_key not in ("", "NONE", "N/A", "-"):
                    # Show custom text
                    _add_textbox(
                        slide,
                        int(cx), int(cy), int(cell_w), int(cell_h),
                        str(status),
                        font_mono,
                        status_label_pt,
                        text_hex,
                        align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE,
                    )
                continue

            # Draw colored circle indicator
            indicator_color = rag_colors[status_key]
            circle_size = min(cell_w, cell_h) * 0.55
            circle_x = cx + (cell_w - circle_size) / 2
            circle_y = cy + (cell_h - circle_size) / 2

            if show_labels:
                # With label: circle on top portion, label below
                circle_size = min(cell_w, cell_h) * 0.42
                circle_x = cx + (cell_w - circle_size) / 2
                circle_y = cy + cell_h * 0.15
            else:
                circle_y = cy + (cell_h - circle_size) / 2

            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Emu(int(circle_x)), Emu(int(circle_y)),
                Emu(int(circle_size)), Emu(int(circle_size)),
            )
            circle.shadow.inherit = False
            circle.fill.solid()
            circle.fill.fore_color.rgb = _rgb(indicator_color)
            circle.line.fill.background()
            circle.text_frame.text = ""

            if show_labels:
                label_y = circle_y + circle_size
                label_h = cell_h - (circle_y - cy) - circle_size
                _add_textbox(
                    slide,
                    int(cx), int(label_y), int(cell_w), int(max(1, label_h)),
                    status_key,
                    font_mono,
                    status_label_pt,
                    text_hex,
                    align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.TOP,
                )

    # Legend at bottom
    legend_y = inner_y + inner_h
    legend_items = [("R", "Red / At Risk"), ("A", "Amber / Caution"), ("G", "Green / On Track")]
    legend_item_w = grid_w / len(legend_items)

    for idx, (key, label) in enumerate(legend_items):
        li_x = grid_x + idx * legend_item_w
        dot_size = int(Pt(base_pt * 0.6))
        dot_y = legend_y + (legend_h - dot_size) / 2

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(int(li_x)), Emu(int(dot_y)),
            Emu(dot_size), Emu(dot_size),
        )
        dot.shadow.inherit = False
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(rag_colors[key])
        dot.line.fill.background()
        dot.text_frame.text = ""

        _add_textbox(
            slide,
            int(li_x + dot_size + Pt(base_pt * 0.3)), int(legend_y),
            int(legend_item_w - dot_size - Pt(base_pt * 0.3)), int(legend_h),
            label,
            font_body,
            max(6, int(round(base_pt * 0.55))),
            text_hex,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    return None
