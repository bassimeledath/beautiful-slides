"""Data table — clean styled table with header row, aligned columns,
optional row striping, and optional highlight cells.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    """Linearly interpolate between two hex colors."""
    t = max(0.0, min(1.0, t))
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _luminance(hex_):
    h = hex_.lstrip("#")
    r = int(h[0:2], 16) / 255.0
    g = int(h[2:4], 16) / 255.0
    b = int(h[4:6], 16) / 255.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_w_emu=0,
              radius_emu=0):
    shape_type = (MSO_SHAPE.ROUNDED_RECTANGLE if radius_emu > 0
                  else MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(shape_type, Emu(int(x)), Emu(int(y)),
                                 Emu(int(w)), Emu(int(h)))
    if radius_emu > 0:
        try:
            short = min(w, h)
            adj = max(0.0, min(0.5, radius_emu / short))
            shp.adjustments[0] = adj
        except Exception:
            pass
    try:
        shp.shadow.inherit = False
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None or line_w_emu <= 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Emu(int(line_w_emu))
    return shp


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def render(slide, data, tokens, bounds):
    """Render a styled data table.

    Parameters
    ----------
    slide : pptx.slide.Slide
    data : dict
        {
            "title": "Optional table title",
            "columns": [
                {"label": "Company", "align": "left"},
                {"label": "Revenue", "align": "right"},
                ...
            ],
            "rows": [
                ["Acme Corp", "$4.2M", ...],
                ...
            ],
            "highlight_cells": [             # optional list of (row, col)
                [0, 1], [2, 3],
            ],
            "row_striping": true,            # optional, default true
        }
    tokens : dict
    bounds : tuple (x_emu, y_emu, w_emu, h_emu)
    """
    x, y, w, h = bounds

    columns = list(data.get("columns") or [])
    rows = list(data.get("rows") or [])
    title = data.get("title")
    highlight_cells = set()
    for pair in (data.get("highlight_cells") or []):
        if len(pair) >= 2:
            highlight_cells.add((int(pair[0]), int(pair[1])))
    row_striping = data.get("row_striping", True)

    if not columns or not rows:
        return None

    n_cols = len(columns)
    n_rows = len(rows)

    # Token extraction
    bg_hex = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(bg_hex.lstrip("#"))
    primary_hex = tokens["primary"]
    accent_hex = tokens["accent"]
    text_hex = tokens["text"]
    muted_hex = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg_hex)
    bg_shape.line.fill.background()

    radius_emu = radius_px * 9525

    # Font sizes
    title_pt = max(base_pt + 2, int(round(base_pt * 1.15)))
    header_pt = max(8, int(round(base_pt * 0.85)))
    cell_pt = max(8, int(round(base_pt * 0.82)))

    # Title space — dynamic height for multi-line wrapping
    title_h = 0
    if title:
        char_w_emu = Pt(title_pt).emu * 0.55
        chars_per_line = max(1, int(w / char_w_emu))
        title_lines = max(1, -(-len(title) // chars_per_line))  # ceil division
        title_h = int(Pt(title_pt).emu * 2.0 * title_lines)
        _add_textbox(slide, x, y, w, title_h, title,
                     font_display, title_pt, text_hex,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Table area
    table_x = x
    table_y = y + title_h
    table_w = w
    table_h = h - title_h

    # Row heights — compact: sized to fit text, NOT to fill available space
    max_col_label_len = max((len(col.get("label") or col.get("name") or "") for col in columns), default=5)
    header_mult = 3.5 if max_col_label_len > 12 else 2.8
    header_h = int(Pt(header_pt).emu * header_mult)
    body_avail = table_h - header_h
    # Compact row height: font size * 2.5 gives comfortable padding around text
    row_h = int(Pt(cell_pt).emu * 2.5)
    # Safety cap so rows don't exceed bounds
    if row_h * n_rows > body_avail:
        row_h = body_avail // max(n_rows, 1)

    # Column widths: weighted by content length
    # Scan content to compute a weight for each column
    col_weights = []
    for ci in range(n_cols):
        col_def = columns[ci] if ci < len(columns) else {}
        header_len = len(col_def.get("label") or col_def.get("name") or "")
        # Average cell length in this column
        cell_lens = [len(str(row[ci])) if ci < len(row) else 0 for row in rows]
        avg_len = (sum(cell_lens) / max(len(cell_lens), 1)) if cell_lens else 5
        max_len = max(cell_lens) if cell_lens else 5
        # Use a blend of avg and max for the weight
        content_len = max(header_len, (avg_len + max_len) / 2)
        # Right-aligned (numeric) columns need less space
        a = (col_def.get("align") or "left").lower()
        if a == "right":
            content_len *= 0.7
        col_weights.append(max(3, content_len))
    total_weight = sum(col_weights)
    col_widths = [int(table_w * cw / total_weight) for cw in col_weights]
    # Distribute rounding remainder to the first column
    remainder = table_w - sum(col_widths)
    col_widths[0] += remainder

    # Precompute column x-offsets
    col_x_offsets = []
    cx_accum = table_x
    for cw in col_widths:
        col_x_offsets.append(cx_accum)
        cx_accum += cw

    # Padding inside each cell
    pad_x = int(Pt(base_pt * 0.5).emu)

    # Estimate max chars per column for truncation
    char_w_emu = Pt(cell_pt).emu * 0.55
    col_max_chars = [max(3, int((cw - 2 * pad_x) / char_w_emu)) for cw in col_widths]

    # Stripe color: subtle mix toward primary
    bg_lum = _luminance(bg_hex)
    if bg_lum < 0.3:
        # Dark mode: lighten slightly
        stripe_hex = _lerp_hex(bg_hex, muted_hex, 0.10)
    else:
        # Light mode: darken slightly
        stripe_hex = _lerp_hex(bg_hex, muted_hex, 0.08)

    # Highlight fill: subtle accent tint
    highlight_hex = _lerp_hex(bg_hex, primary_hex, 0.15)

    # Header separator line
    hairline = int(0.75 * 9525)  # ~0.75pt

    # Determine alignment per column
    aligns = []
    for col in columns:
        a = (col.get("align") or "left").lower()
        if a == "right":
            aligns.append(PP_ALIGN.RIGHT)
        elif a == "center":
            aligns.append(PP_ALIGN.CENTER)
        else:
            aligns.append(PP_ALIGN.LEFT)

    # -- Draw header row --
    for ci, col in enumerate(columns):
        cx = col_x_offsets[ci]
        cw = col_widths[ci]
        # Header background
        _add_rect(slide, cx, table_y, cw, header_h, bg_hex)

        # Header text
        label = col.get("label") or col.get("name") or ""
        _add_textbox(slide, cx + pad_x, table_y, cw - 2 * pad_x, header_h,
                     label.upper(),
                     font_body, header_pt, text_hex,
                     bold=True, align=aligns[ci], anchor=MSO_ANCHOR.MIDDLE)

    # Header bottom border
    sep_y = table_y + header_h
    sep_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(table_x)), Emu(int(sep_y)),
        Emu(int(table_w)), Emu(int(hairline * 2)),
    )
    sep_shape.fill.solid()
    sep_shape.fill.fore_color.rgb = _rgb(primary_hex)
    sep_shape.line.fill.background()
    try:
        sep_shape.shadow.inherit = False
    except Exception:
        pass

    # -- Draw body rows --
    body_y = sep_y + hairline * 2
    for ri, row in enumerate(rows):
        ry = body_y + ri * row_h

        # Row background (striping)
        if row_striping and ri % 2 == 1:
            row_bg = stripe_hex
        else:
            row_bg = bg_hex

        for ci in range(n_cols):
            cx = col_x_offsets[ci]
            cw = col_widths[ci]
            cell_val = row[ci] if ci < len(row) else ""

            # Check highlight
            is_highlight = (ri, ci) in highlight_cells
            cell_bg = highlight_hex if is_highlight else row_bg

            # Cell background
            _add_rect(slide, cx, ry, cw, row_h, cell_bg)

            # Cell text — truncate with ellipsis if too long
            cell_text = str(cell_val)
            mc = col_max_chars[ci]
            if len(cell_text) > mc + 3:
                cell_text = cell_text[:mc] + "..."

            cell_font = font_mono if aligns[ci] == PP_ALIGN.RIGHT else font_body
            cell_color = text_hex
            if is_highlight:
                cell_color = primary_hex

            _add_textbox(slide, cx + pad_x, ry, cw - 2 * pad_x, row_h,
                         cell_text,
                         cell_font, cell_pt, cell_color,
                         bold=is_highlight, align=aligns[ci],
                         anchor=MSO_ANCHOR.MIDDLE)

        # Subtle row separator
        if ri < n_rows - 1:
            row_sep_y = ry + row_h
            row_sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(int(table_x)), Emu(int(row_sep_y)),
                Emu(int(table_w)), Emu(int(hairline)),
            )
            row_sep.fill.solid()
            row_sep.fill.fore_color.rgb = _rgb(muted_hex)
            row_sep.line.fill.background()
            try:
                row_sep.shadow.inherit = False
            except Exception:
                pass

    return None
