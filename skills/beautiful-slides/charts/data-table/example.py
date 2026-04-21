"""Render data tables in all 5 modes. Generates one .pptx per mode."""

import os
import sys

from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


SAMPLE_DATA = {
    "title": "Q1 2025 Revenue by Region",
    "columns": [
        {"label": "Region", "align": "left"},
        {"label": "Revenue", "align": "right"},
        {"label": "Growth", "align": "right"},
        {"label": "Margin", "align": "right"},
        {"label": "Headcount", "align": "right"},
    ],
    "rows": [
        ["North America", "$12.4M", "+18%", "72%", "142"],
        ["EMEA",          "$8.1M",  "+12%", "68%", "89"],
        ["APAC",          "$5.6M",  "+34%", "65%", "64"],
        ["LATAM",         "$2.3M",  "+22%", "61%", "31"],
        ["Middle East",   "$1.1M",  "+8%",  "59%", "12"],
    ],
    "highlight_cells": [[0, 1], [2, 2]],
    "row_striping": True,
}


def _paint_slide_bg(slide, hex_):
    from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        slide.part.package.presentation_part.presentation.slide_width,
        slide.part.package.presentation_part.presentation.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(hex_)
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def build(mode, tokens):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _paint_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    bounds = (
        margin,
        margin,
        prs.slide_width - 2 * margin,
        prs.slide_height - 2 * margin,
    )

    render(slide, SAMPLE_DATA, tokens, bounds)

    out = os.path.join(HERE, f"example-{mode}.pptx")
    prs.save(out)
    return out


def main():
    for mode, tokens in MODES.items():
        path = build(mode, tokens)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
