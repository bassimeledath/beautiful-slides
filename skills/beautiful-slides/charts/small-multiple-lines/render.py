"""Small-multiple lines — grid of mini line charts sharing a common y-scale.

4-12 panels in 2x2, 2x3, 3x3, or 3x4 grid. Each panel has a title and
draws a line using native python-pptx connector shapes.
Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Emu(int(x1)), Emu(int(y1)),
                                      Emu(int(x2)), Emu(int(y2)))
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(weight_pt)
    return conn


def _add_filled_circle(slide, cx, cy, r_emu, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
                                 Emu(int(r_emu * 2)), Emu(int(r_emu * 2)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    try:
        shp.shadow.inherit = False
    except Exception:
        pass
    return shp


def _choose_grid(n):
    """Pick a (rows, cols) grid that fits n panels."""
    if n <= 2:
        return (1, 2)
    if n <= 4:
        return (2, 2)
    if n <= 6:
        return (2, 3)
    if n <= 9:
        return (3, 3)
    return (3, 4)  # max 12


def _nice_ticks(vmin, vmax, target=4):
    """Return ticks, lo, hi for the shared y-axis."""
    if vmax <= vmin:
        vmax = vmin + 1
    span = vmax - vmin
    raw = span / max(target, 1)
    mag = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    for mult in (1, 2, 2.5, 5, 10):
        step = mult * mag
        if span / step <= target * 1.5:
            break
    lo = math.floor(vmin / step) * step
    hi = math.ceil(vmax / step) * step
    ticks = []
    v = lo
    while v <= hi + 1e-9:
        ticks.append(round(v, 6))
        v += step
    return ticks, lo, hi


def _fmt_num(v):
    """Format number with compact notation (K/M) for large values."""
    av = abs(v)
    if av >= 1_000_000:
        s = f"{v / 1_000_000:.1f}M"
        return s.replace(".0M", "M")
    if av >= 10_000:
        s = f"{v / 1_000:.0f}K"
        return s
    if av >= 1_000:
        s = f"{v / 1_000:.1f}K"
        return s.replace(".0K", "K")
    if abs(v - round(v)) < 1e-6:
        return f"{int(round(v))}"
    return f"{v:.1f}"


def _estimate_lines(text, font_pt, avail_w):
    """Estimate how many lines text will wrap to given available width."""
    char_w = Pt(font_pt).emu * 0.55
    chars_per_line = max(1, int(avail_w / char_w))
    return max(1, math.ceil(len(text) / chars_per_line))


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 1].rstrip() + "\u2026"


def render(slide, data, tokens, bounds):
    """Render a grid of small-multiple line charts with shared y-scale.

    Parameters
    ----------
    slide : pptx.slide.Slide
    data : dict
        panels: list of {title, values: list[float]}
        x_labels: list[str] — shared x-axis labels (optional)
        title: str — overall title (optional)
        y_label: str — shared y-axis label (optional)
    tokens : dict
    bounds : tuple (x, y, w, h) in EMU
    """
    x0, y0, w0, h0 = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x0)), Emu(int(y0)),
        Emu(int(w0)), Emu(int(h0)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    panels = list(data.get("panels", []))
    x_labels = list(data.get("x_labels", []))
    title = data.get("title")
    y_label = data.get("y_label")

    if not panels:
        return

    # Limit to 12 panels
    panels = panels[:12]
    n = len(panels)
    rows, cols = _choose_grid(n)

    # Layout
    pad = int(min(w0, h0) * 0.03)
    cursor_y = y0 + pad

    # Overall title -- dynamic height for multi-line wrapping
    title_h = 0
    if title:
        # Reduce title font for dense grids
        if n >= 9:
            title_pt = int(base_pt * 1.1)
        else:
            title_pt = int(base_pt * 1.4)
        title_avail_w = w0 - 2 * pad
        title_lines = _estimate_lines(title, title_pt, title_avail_w)
        title_h = int(Pt(title_pt).emu * 1.3 * title_lines + Pt(title_pt).emu * 0.4)
        title_h = min(title_h, int(h0 * 0.15))  # cap at 15% of height
        _add_textbox(slide, x0 + pad, cursor_y, title_avail_w, title_h,
                     title, font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += title_h + int(Pt(base_pt * 0.15).emu)

    # Y-label subtitle
    if y_label:
        sub_pt = max(8, int(base_pt * 0.8))
        sub_h = int(Pt(sub_pt).emu * 1.5)
        _add_textbox(slide, x0 + pad, cursor_y, w0 - 2 * pad, sub_h,
                     y_label, font_body, sub_pt, text_c,
                     bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += sub_h

    # Grid area
    grid_x = x0 + pad
    grid_y = cursor_y + int(Pt(base_pt * 0.2).emu)
    grid_w = w0 - 2 * pad
    grid_h = (y0 + h0 - pad) - grid_y

    # Cell sizing — tighter gutters for dense grids
    if n >= 9:
        gutter_x = int(grid_w * 0.02)
        gutter_y = int(grid_h * 0.025)
    else:
        gutter_x = int(grid_w * 0.03)
        gutter_y = int(grid_h * 0.04)
    cell_w = (grid_w - gutter_x * (cols - 1)) // cols
    cell_h = (grid_h - gutter_y * (rows - 1)) // rows

    # Compute shared y-range across all panels
    all_vals = []
    for p in panels:
        for v in p.get("values", []):
            if v is not None:
                all_vals.append(float(v))
    if not all_vals:
        return
    vmin = min(all_vals)
    vmax = max(all_vals)
    span = vmax - vmin if vmax > vmin else max(abs(vmax), 1.0)
    vmin_p = vmin - span * 0.05
    vmax_p = vmax + span * 0.05
    ticks, lo, hi = _nice_ticks(vmin_p, vmax_p, target=4)
    if hi == lo:
        hi = lo + 1

    # Font sizes for panels -- reduce when grid is dense (9+ panels)
    if n >= 9:
        panel_title_pt = max(6, int(base_pt * 0.55))
        tick_pt = max(5, int(base_pt * 0.45))
    elif n >= 7:
        panel_title_pt = max(7, int(base_pt * 0.65))
        tick_pt = max(6, int(base_pt * 0.5))
    else:
        panel_title_pt = max(8, int(base_pt * 0.85))
        tick_pt = max(7, int(base_pt * 0.65))
    panel_title_h = int(Pt(panel_title_pt).emu * 1.3)

    # Estimate max chars that fit in one line for panel titles
    panel_title_char_w = Pt(panel_title_pt).emu * 0.55
    panel_title_max_chars = max(8, int(cell_w / panel_title_char_w))

    # Layout for each panel: title on top, then plot area with y-ticks on left
    # Reduce y-axis label width for dense grids to prevent overflow
    if n >= 9:
        left_tick_w = int(cell_w * 0.22)
    else:
        left_tick_w = int(cell_w * 0.18)
    x_label_h = int(Pt(tick_pt).emu * 1.4) if x_labels else 0

    for idx, panel in enumerate(panels):
        if idx >= rows * cols:
            break

        row = idx // cols
        col = idx % cols
        cx = grid_x + col * (cell_w + gutter_x)
        cy = grid_y + row * (cell_h + gutter_y)

        panel_title = panel.get("title", "")
        # Truncate panel title to fit in one line
        panel_title = _truncate(panel_title, panel_title_max_chars)
        values = [float(v) if v is not None else None
                  for v in panel.get("values", [])]

        # Panel title
        _add_textbox(slide, cx, cy, cell_w, panel_title_h,
                     panel_title, font_body, panel_title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

        # Plot area within cell
        plot_x = cx + left_tick_w
        plot_y = cy + panel_title_h
        plot_w = cell_w - left_tick_w - int(cell_w * 0.03)
        plot_h = cell_h - panel_title_h - x_label_h - int(Pt(tick_pt).emu * 0.3)

        if plot_w <= 0 or plot_h <= 0:
            continue

        def y_to_emu(v):
            frac = (v - lo) / (hi - lo)
            return plot_y + plot_h - frac * plot_h

        # Subtle horizontal gridlines (only bottom and top ticks for cleanliness)
        for tv in [ticks[0], ticks[-1]]:
            yy = y_to_emu(tv)
            _add_line(slide, plot_x, yy, plot_x + plot_w, yy, muted, 0.4)

        # Y-tick labels: only show first and last tick for compactness
        tick_label_w = left_tick_w - int(Pt(tick_pt).emu * 0.2)
        tick_label_h = int(Pt(tick_pt).emu * 1.3)
        for tv in [ticks[0], ticks[-1]]:
            yy = y_to_emu(tv)
            _add_textbox(slide,
                         cx, yy - tick_label_h // 2,
                         tick_label_w, tick_label_h,
                         _fmt_num(tv), font_mono, tick_pt, text_c,
                         bold=False, align=PP_ALIGN.RIGHT,
                         anchor=MSO_ANCHOR.MIDDLE)

        # X positions
        n_pts = len(values)
        if n_pts <= 1:
            def x_to_emu(i):
                return plot_x + plot_w / 2
        else:
            step_x = plot_w / (n_pts - 1)
            def x_to_emu(i):
                return plot_x + i * step_x

        # Build valid points
        pts = []
        for i, v in enumerate(values):
            if v is not None:
                pts.append((x_to_emu(i), y_to_emu(v), i, v))

        # Draw line segments
        for a, b in zip(pts, pts[1:]):
            _add_line(slide, a[0], a[1], b[0], b[1], primary, 1.75)

        # Endpoint dot
        if pts:
            last = pts[-1]
            dot_r = int(Pt(2.5).emu)
            _add_filled_circle(slide, last[0], last[1], dot_r, accent)

        # X-axis labels: first and last only (if provided)
        if x_labels and n_pts > 0:
            shown = set()
            for show_i in [0, min(n_pts - 1, len(x_labels) - 1)]:
                if show_i in shown or show_i >= len(x_labels):
                    continue
                shown.add(show_i)
                lx = x_to_emu(show_i)
                lbl_w = int(plot_w * 0.4)
                align = PP_ALIGN.LEFT if show_i == 0 else PP_ALIGN.RIGHT
                label_x = lx if show_i == 0 else lx - lbl_w
                _add_textbox(slide, label_x,
                             plot_y + plot_h + int(Pt(tick_pt).emu * 0.15),
                             lbl_w, x_label_h,
                             str(x_labels[show_i]), font_body, tick_pt, text_c,
                             bold=False, align=align, anchor=MSO_ANCHOR.TOP)

    return None
