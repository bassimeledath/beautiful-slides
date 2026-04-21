"""Run `render()` 5 times -- once per mode -- writing example-<mode>.pptx per run."""

import os
import sys

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module
from tokens import MODES  # noqa: E402
from render import render  # noqa: E402


DATA = {
    "title": "Monthly Active Users by Region",
    "y_label": "MAU (thousands)",
    "x_labels": ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
    "panels": [
        {"title": "North America", "values": [120, 125, 130, 128, 135, 142]},
        {"title": "Europe",        "values": [80, 82, 85, 88, 90, 94]},
        {"title": "APAC",          "values": [45, 50, 55, 62, 70, 78]},
        {"title": "LATAM",         "values": [30, 32, 31, 35, 38, 40]},
        {"title": "MEA",           "values": [15, 16, 18, 20, 22, 25]},
        {"title": "ANZ",           "values": [12, 13, 14, 15, 16, 18]},
    ],
}


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    for mode_name, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        margin = Inches(0.4)
        bounds = (
            int(margin),
            int(margin),
            int(prs.slide_width - 2 * margin),
            int(prs.slide_height - 2 * margin),
        )
        render(slide, DATA, tokens, bounds)

        out_path = os.path.join(out_dir, f"example-{mode_name}.pptx")
        prs.save(out_path)
        print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
