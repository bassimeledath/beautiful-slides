"""Render the scatter chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_BASIC = {
    "title": "Ad spend vs. revenue",
    "x_label": "Ad spend ($K)",
    "y_label": "Revenue ($K)",
    "show_labels": True,
    "show_trend_line": False,
    "points": [
        {"x": 12, "y":  95, "label": "Acme"},
        {"x": 28, "y": 160, "label": "Beta"},
        {"x": 41, "y": 205, "label": "Gamma"},
        {"x": 55, "y": 240, "label": "Delta"},
        {"x": 68, "y": 290, "label": "Echo"},
        {"x": 82, "y": 340, "label": "Foxtrot"},
    ],
}

DATA_TREND = {
    "title": "R&D investment vs. patent output",
    "x_label": "R&D ($M)",
    "y_label": "Patents filed",
    "show_labels": True,
    "show_trend_line": True,
    "points": [
        {"x":  5, "y": 12, "label": "StartupA"},
        {"x": 15, "y": 28, "label": "StartupB"},
        {"x": 22, "y": 35, "label": "MidCo"},
        {"x": 38, "y": 55, "label": "BigCorp"},
        {"x": 50, "y": 48, "label": "Outlier"},
        {"x": 60, "y": 72, "label": "TechGiant"},
        {"x": 75, "y": 90, "label": "MegaCorp"},
    ],
}

DATA_DIVIDERS = {
    "title": "Customer satisfaction vs. retention",
    "x_label": "Satisfaction score",
    "y_label": "Retention rate (%)",
    "show_labels": True,
    "show_trend_line": False,
    "x_divider": 7,
    "y_divider": 80,
    "points": [
        {"x": 3.2, "y": 45, "label": "Seg A"},
        {"x": 4.5, "y": 62, "label": "Seg B"},
        {"x": 6.1, "y": 55, "label": "Seg C"},
        {"x": 7.8, "y": 88, "label": "Seg D"},
        {"x": 8.5, "y": 92, "label": "Seg E"},
        {"x": 9.1, "y": 95, "label": "Seg F"},
        {"x": 5.0, "y": 85, "label": "Seg G"},
    ],
}


MODE_DATA = {
    "sv-keynote": DATA_BASIC,
    "editorial-magazine": DATA_TREND,
    "playful-marketing": DATA_DIVIDERS,
    "consulting-boardroom": DATA_TREND,
    "craft-minimal": DATA_BASIC,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
