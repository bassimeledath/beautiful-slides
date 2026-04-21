"""Scatter chart renderer — native python-pptx shapes only.

Points on x/y axes showing relationship between two continuous variables.
Supports optional trend line, optional quadrant divider lines, and
optional direct labels per point.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ---- helpers ---------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap=False):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_circle(slide, cx, cy, r_emu, fill_hex, line_hex=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
        Emu(int(r_emu * 2)), Emu(int(r_emu * 2)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
    shp.shadow.inherit = False
    return shp


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt, dash=None):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(int(x1)), Emu(int(y1)),
        Emu(int(x2)), Emu(int(y2)),
    )
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(weight_pt)
    if dash is not None:
        try:
            conn.line.dash_style = dash
        except Exception:
            pass
    return conn


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _nice_ticks(vmin, vmax, target=5):
    if vmax <= vmin:
        vmax = vmin + 1
    span = vmax - vmin
    raw = span / max(target, 1)
    mag = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    for mult in (1, 2, 2.5, 5, 10):
        step = mult * mag
        if span / step <= target * 1.5:
            break
    lo = math.floor(vmin / step) * step
    hi = math.ceil(vmax / step) * step
    ticks = []
    v = lo
    while v <= hi + 1e-9:
        ticks.append(round(v, 6))
        v += step
    return ticks, lo, hi


def _fmt(v):
    if abs(v - round(v)) < 1e-6:
        return f"{int(round(v))}"
    return f"{v:.1f}"


# ---- main render -----------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a scatter chart onto *slide* inside *bounds*, styled by *tokens*."""
    x0, y0, w0, h0 = bounds

    # --- unpack tokens -------------------------------------------------------
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # --- unpack data ---------------------------------------------------------
    title = data.get("title")
    points = list(data.get("points") or [])
    x_label = data.get("x_label")
    y_label = data.get("y_label")
    show_labels = bool(data.get("show_labels", True))
    show_trend_line = bool(data.get("show_trend_line", False))
    x_divider = data.get("x_divider")  # optional vertical line at this x
    y_divider = data.get("y_divider")  # optional horizontal line at this y
    point_radius_pt = data.get("point_radius_pt", 5)

    if not points:
        return

    # --- background ----------------------------------------------------------
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    # --- outer padding -------------------------------------------------------
    pad = int(min(w0, h0) * 0.035)
    ix, iy = x0 + pad, y0 + pad
    iw, ih = w0 - 2 * pad, h0 - 2 * pad

    # --- title ---------------------------------------------------------------
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.55))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, ix, iy, iw, title_h, title,
                     font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + int(pad * 0.3)
        ih = (y0 + h0 - pad) - iy

    # --- compute plot area ---------------------------------------------------
    tick_pt = max(8, int(base_pt * 0.8))
    label_pt = max(8, int(base_pt * 0.85))

    left_margin = int(iw * 0.08)
    bottom_margin = int(Pt(tick_pt).emu * 1.8) + (
        int(Pt(label_pt).emu * 1.6) if x_label else 0
    )
    top_margin = int(Pt(label_pt).emu * 2.0) if y_label else int(Pt(tick_pt).emu * 0.6)
    right_margin = int(iw * 0.04)

    plot_x = ix + left_margin
    plot_y = iy + top_margin
    plot_w = max(1, iw - left_margin - right_margin)
    plot_h = max(1, ih - top_margin - bottom_margin)

    if plot_w <= 0 or plot_h <= 0:
        return

    # --- compute ranges ------------------------------------------------------
    xs = [p["x"] for p in points]
    ys = [p["y"] for p in points]
    x_span = max(xs) - min(xs) if max(xs) > min(xs) else max(abs(max(xs)), 1)
    y_span = max(ys) - min(ys) if max(ys) > min(ys) else max(abs(max(ys)), 1)

    x_ticks, x_lo, x_hi = _nice_ticks(min(xs) - x_span * 0.08,
                                        max(xs) + x_span * 0.08)
    y_ticks, y_lo, y_hi = _nice_ticks(min(ys) - y_span * 0.08,
                                        max(ys) + y_span * 0.08)

    def to_px(vx, vy):
        fx = (vx - x_lo) / (x_hi - x_lo) if x_hi != x_lo else 0.5
        fy = (vy - y_lo) / (y_hi - y_lo) if y_hi != y_lo else 0.5
        return plot_x + fx * plot_w, plot_y + plot_h - fy * plot_h

    # --- gridlines -----------------------------------------------------------
    hairline_pt = 0.5
    for tv in y_ticks:
        _, yy = to_px(0, tv)
        _add_line(slide, plot_x, yy, plot_x + plot_w, yy, muted, hairline_pt)
    for tv in x_ticks:
        xx, _ = to_px(tv, 0)
        _add_line(slide, xx, plot_y, xx, plot_y + plot_h, muted, hairline_pt)

    # --- axes ----------------------------------------------------------------
    _add_line(slide, plot_x, plot_y + plot_h, plot_x + plot_w,
              plot_y + plot_h, muted, 0.75)
    _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, muted, 0.75)

    # --- tick labels ---------------------------------------------------------
    tick_label_h = int(Pt(tick_pt).emu * 1.4)
    # y-axis
    y_tick_w = int(left_margin * 0.92)
    for tv in y_ticks:
        _, yy = to_px(0, tv)
        _add_textbox(slide,
                     plot_x - y_tick_w - int(pad * 0.15),
                     yy - tick_label_h / 2,
                     y_tick_w, tick_label_h,
                     _fmt(tv), font_mono, tick_pt, text_c,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # x-axis
    x_label_w = int(plot_w / max(len(x_ticks), 1) * 1.6)
    right_edge = x0 + w0  # absolute right boundary of bounds
    for tv in x_ticks:
        xx, _ = to_px(tv, 0)
        lbl_x = xx - x_label_w / 2
        # Clamp so the textbox doesn't exceed the right boundary
        if lbl_x + x_label_w > right_edge:
            lbl_x = right_edge - x_label_w
        # Also clamp left
        if lbl_x < x0:
            lbl_x = x0
        _add_textbox(slide,
                     lbl_x,
                     plot_y + plot_h + int(pad * 0.15),
                     x_label_w, tick_label_h,
                     _fmt(tv), font_mono, tick_pt, text_c,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # --- axis labels ---------------------------------------------------------
    if x_label:
        lbl_h = int(Pt(label_pt).emu * 1.6)
        _add_textbox(slide,
                     plot_x, plot_y + plot_h + tick_label_h + int(pad * 0.2),
                     plot_w, lbl_h,
                     x_label, font_body, label_pt, text_c,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    if y_label:
        # Render as a dedicated subtitle row above the plot area (left-aligned
        # next to the y-axis) so it doesn't overlap tick labels.
        y_lbl_display = y_label if len(y_label) <= 50 else y_label[:47] + "..."
        lbl_h = int(Pt(label_pt).emu * 1.6)
        _add_textbox(slide, plot_x, plot_y - lbl_h - int(pad * 0.1),
                     plot_w, lbl_h,
                     y_lbl_display, font_body, label_pt, text_c,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)

    # --- optional divider lines ----------------------------------------------
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        dash = MSO_LINE_DASH_STYLE.DASH
    except Exception:
        dash = None

    if x_divider is not None:
        dx, _ = to_px(x_divider, 0)
        if plot_x <= dx <= plot_x + plot_w:
            _add_line(slide, dx, plot_y, dx, plot_y + plot_h, muted, 1.0,
                      dash=dash)
    if y_divider is not None:
        _, dy = to_px(0, y_divider)
        if plot_y <= dy <= plot_y + plot_h:
            _add_line(slide, plot_x, dy, plot_x + plot_w, dy, muted, 1.0,
                      dash=dash)

    # --- trend line (OLS) ----------------------------------------------------
    if show_trend_line and len(points) >= 2:
        n = len(points)
        sx = sum(p["x"] for p in points)
        sy = sum(p["y"] for p in points)
        sxx = sum(p["x"] ** 2 for p in points)
        sxy = sum(p["x"] * p["y"] for p in points)
        denom = n * sxx - sx * sx
        if abs(denom) > 1e-12:
            slope = (n * sxy - sx * sy) / denom
            intercept = (sy - slope * sx) / n
            # Draw from x_lo to x_hi
            ty1 = slope * x_lo + intercept
            ty2 = slope * x_hi + intercept
            tx1, ty1e = to_px(x_lo, ty1)
            tx2, ty2e = to_px(x_hi, ty2)
            # Clamp to plot bounds
            tx1 = max(plot_x, min(plot_x + plot_w, tx1))
            tx2 = max(plot_x, min(plot_x + plot_w, tx2))
            ty1e = max(plot_y, min(plot_y + plot_h, ty1e))
            ty2e = max(plot_y, min(plot_y + plot_h, ty2e))
            _add_line(slide, tx1, ty1e, tx2, ty2e, accent, 1.5, dash=dash)

    # --- points --------------------------------------------------------------
    r_emu = int(Pt(point_radius_pt).emu)
    point_label_pt = max(7, int(base_pt * 0.7))
    point_label_h = int(Pt(point_label_pt).emu * 1.4)

    # Collision detection: track occupied label rectangles
    occupied_rects = []  # list of (lx, ly, lx+w, ly+h)
    gap = int(Pt(1).emu)  # small gap between labels

    def _rects_overlap(r1, r2):
        return r1[0] < r2[2] and r1[2] > r2[0] and r1[1] < r2[3] and r1[3] > r2[1]

    def _any_collision(rect):
        for occ in occupied_rects:
            if _rects_overlap(rect, occ):
                return True
        return False

    for p in points:
        cx, cy = to_px(p["x"], p["y"])
        color = p.get("color") or primary
        # Ensure point stays within plot area (clamp center to keep circle inside bounds)
        cx = max(plot_x + r_emu, min(plot_x + plot_w - r_emu, cx))
        cy = max(plot_y + r_emu, min(plot_y + plot_h - r_emu, cy))
        _add_circle(slide, cx, cy, r_emu, color)

        if show_labels and p.get("label"):
            lbl = str(p["label"])
            # Truncate long labels
            if len(lbl) > 25:
                lbl = lbl[:22] + "..."
            est_w = int(Pt(point_label_pt).emu * 0.55 * len(lbl))
            est_w = max(est_w, int(Pt(point_label_pt).emu * 2))

            # Try multiple positions: right, left, above, below
            spacing = r_emu + int(Pt(2).emu)
            candidates = [
                (cx + spacing, cy - point_label_h / 2),                          # right
                (cx - spacing - est_w, cy - point_label_h / 2),                  # left
                (cx - est_w / 2, cy - spacing - point_label_h),                  # above
                (cx - est_w / 2, cy + spacing),                                  # below
            ]
            placed = False
            for lx, ly in candidates:
                # Clamp within plot area
                lx = max(plot_x, min(plot_x + plot_w - est_w, lx))
                ly = max(plot_y, min(plot_y + plot_h - point_label_h, ly))
                rect = (lx - gap, ly - gap, lx + est_w + gap, ly + point_label_h + gap)
                if not _any_collision(rect):
                    _add_textbox(slide, lx, ly, est_w, point_label_h,
                                 lbl, font_body, point_label_pt, text_c,
                                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
                    occupied_rects.append(rect)
                    placed = True
                    break
            # If no position works, skip this label to avoid overlap
