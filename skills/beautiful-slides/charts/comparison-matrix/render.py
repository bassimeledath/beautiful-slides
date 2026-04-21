"""Comparison matrix — rows=items, cols=criteria, cells=checkmarks/crosses/partial.

Classic "us vs them" competitive grid. Cells contain check marks, crosses,
or partial/half indicators drawn as native shapes.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    """Linearly interpolate between two hex colors."""
    t = max(0.0, min(1.0, t))
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _luminance(hex_):
    h = hex_.lstrip("#")
    r = int(h[0:2], 16) / 255.0
    g = int(h[2:4], 16) / 255.0
    b = int(h[4:6], 16) / 255.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_w_emu=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y)),
                                 Emu(int(w)), Emu(int(h)))
    try:
        shp.shadow.inherit = False
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None or line_w_emu <= 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Emu(int(line_w_emu))
    return shp


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_check(slide, cx, cy, size, color_hex):
    """Draw a checkmark as a Unicode glyph in a textbox centered at (cx, cy)."""
    # We use a simple Unicode checkmark rendered in the body font
    tb = slide.shapes.add_textbox(
        Emu(int(cx - size // 2)), Emu(int(cy - size // 2)),
        Emu(int(size)), Emu(int(size)),
    )
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u2713"  # Check mark
    run.font.size = Pt(int(size / 12700 * 0.55))
    run.font.bold = True
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_cross(slide, cx, cy, size, color_hex):
    """Draw a cross as a Unicode glyph."""
    tb = slide.shapes.add_textbox(
        Emu(int(cx - size // 2)), Emu(int(cy - size // 2)),
        Emu(int(size)), Emu(int(size)),
    )
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u2715"  # Multiplication X
    run.font.size = Pt(int(size / 12700 * 0.50))
    run.font.bold = True
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_partial(slide, cx, cy, size, color_hex):
    """Draw a partial/half indicator using a half-filled circle (Unicode)."""
    tb = slide.shapes.add_textbox(
        Emu(int(cx - size // 2)), Emu(int(cy - size // 2)),
        Emu(int(size)), Emu(int(size)),
    )
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u25D0"  # Circle with left half black
    run.font.size = Pt(int(size / 12700 * 0.50))
    run.font.bold = False
    run.font.color.rgb = _rgb(color_hex)
    return tb


def render(slide, data, tokens, bounds):
    """Render a comparison matrix.

    Parameters
    ----------
    slide : pptx.slide.Slide
    data : dict
        {
            "title": "Feature Comparison",           # optional
            "row_labels": ["Us", "Competitor A", "Competitor B"],
            "col_labels": ["SSO", "API", "Mobile", "Analytics"],
            "values": [                              # row-major
                ["check", "check", "check", "partial"],
                ["check", "cross", "check", "cross"],
                ["cross", "cross", "partial", "cross"],
            ],
            "highlight_row": 0,                      # optional: "our" row index
        }
    tokens : dict
    bounds : tuple (x_emu, y_emu, w_emu, h_emu)

    Cell values: "check" | "cross" | "partial" | true | false | 1 | 0 | 0.5
    """
    x, y, w, h = bounds

    row_labels = list(data.get("row_labels") or [])
    col_labels = list(data.get("col_labels") or [])
    values = list(data.get("values") or [])
    title = data.get("title")
    highlight_row = data.get("highlight_row")

    n_rows = len(row_labels)
    n_cols = len(col_labels)
    if n_rows == 0 or n_cols == 0:
        return None

    # Token extraction
    bg_hex = tokens["bg"]
    primary_hex = tokens["primary"]
    accent_hex = tokens["accent"]
    text_hex = tokens["text"]
    muted_hex = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg_hex)
    bg_shape.line.fill.background()


    # Font sizes
    title_pt = max(base_pt + 2, int(round(base_pt * 1.15)))
    # Adaptive header font: scale down further when many columns
    if n_cols > 7:
        header_pt = max(6, int(round(base_pt * 0.80 * 7 / n_cols)))
    else:
        header_pt = max(7, int(round(base_pt * 0.80)))
    row_label_pt = max(8, int(round(base_pt * 0.88)))

    # Title — dynamic height for multi-line wrapping
    title_h = 0
    if title:
        char_w_emu = Pt(title_pt).emu * 0.55
        chars_per_line = max(1, int(w / char_w_emu))
        title_lines = max(1, -(-len(title) // chars_per_line))  # ceil division
        title_h = int(Pt(title_pt).emu * 2.0 * title_lines)
        _add_textbox(slide, x, y, w, title_h, title,
                     font_display, title_pt, text_hex,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Grid area
    grid_x = x
    grid_y = y + title_h
    grid_w = w
    grid_h = h - title_h

    # Row label column width — based on longest label
    max_label_chars = max((len(str(lbl)) for lbl in row_labels), default=5)
    row_label_w = int(Pt(row_label_pt).emu * 0.55 * (max_label_chars + 3))
    row_label_w = min(row_label_w, int(grid_w * 0.28))

    # Column header row height — taller when many columns (headers need to wrap)
    col_header_mult = 4.5 if n_cols > 6 else 3.0
    col_header_h = int(Pt(header_pt).emu * col_header_mult)

    # Cell dimensions
    cell_area_x = grid_x + row_label_w
    cell_area_y = grid_y + col_header_h
    cell_area_w = grid_w - row_label_w
    cell_area_h = grid_h - col_header_h

    cell_w = cell_area_w // max(n_cols, 1)
    cell_h = cell_area_h // max(n_rows, 1)

    # Gap between cells
    gap = max(1, int(0.5 * 9525))  # ~0.5px

    # Glyph size for check/cross/partial
    glyph_size = int(min(cell_w, cell_h) * 0.50)

    # Highlight row background color
    bg_lum = _luminance(bg_hex)
    if highlight_row is not None:
        highlight_bg = _lerp_hex(bg_hex, primary_hex, 0.08)
    else:
        highlight_bg = bg_hex

    # Stripe color for alternating rows
    stripe_hex = _lerp_hex(bg_hex, muted_hex, 0.06)

    # -- Draw column headers --
    # Estimate max chars that fit per column for truncation
    hdr_char_w_emu = Pt(header_pt).emu * 0.55
    hdr_max_chars = max(5, int(cell_w / hdr_char_w_emu))
    for ci, cl in enumerate(col_labels):
        cx = cell_area_x + ci * cell_w
        hdr_text = str(cl)
        # Truncate long headers with ellipsis
        if len(hdr_text) > hdr_max_chars + 3:
            hdr_text = hdr_text[:hdr_max_chars] + "..."
        _add_textbox(slide, int(cx), int(grid_y), int(cell_w), int(col_header_h),
                     hdr_text,
                     font_body, header_pt, text_hex,
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # -- Draw rows --
    for ri, rl in enumerate(row_labels):
        ry = cell_area_y + ri * cell_h
        is_highlight = (highlight_row is not None and ri == highlight_row)

        # Row background
        if is_highlight:
            row_bg = highlight_bg
        elif ri % 2 == 1:
            row_bg = stripe_hex
        else:
            row_bg = bg_hex

        # Full row background rectangle (including label area)
        _add_rect(slide, grid_x, ry, grid_w, cell_h, row_bg)

        # Row label
        label_color = primary_hex if is_highlight else text_hex
        pad_x = int(Pt(base_pt * 0.4).emu)
        _add_textbox(slide, int(grid_x + pad_x), int(ry),
                     int(row_label_w - 2 * pad_x), int(cell_h),
                     str(rl),
                     font_body, row_label_pt, label_color,
                     bold=is_highlight, align=PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)

        # Cell values
        row_vals = values[ri] if ri < len(values) else []
        for ci in range(n_cols):
            cx = cell_area_x + ci * cell_w
            cell_cx = cx + cell_w // 2
            cell_cy = ry + cell_h // 2

            v = row_vals[ci] if ci < len(row_vals) else "cross"

            # Normalize value
            if v is True or v == 1 or str(v).lower() in ("check", "yes", "true", "1"):
                cell_type = "check"
            elif v is False or v == 0 or str(v).lower() in ("cross", "no", "false", "0", "x"):
                cell_type = "cross"
            else:
                cell_type = "partial"

            if cell_type == "check":
                color = accent_hex
                _draw_check(slide, cell_cx, cell_cy, glyph_size, color)
            elif cell_type == "cross":
                color = text_hex
                _draw_cross(slide, cell_cx, cell_cy, glyph_size, color)
            else:  # partial
                color = primary_hex
                _draw_partial(slide, cell_cx, cell_cy, glyph_size, color)

    # Bottom border line
    bottom_y = cell_area_y + n_rows * cell_h
    hairline = int(0.5 * 9525)
    if bottom_y <= y + h:
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(grid_x)), Emu(int(bottom_y)),
            Emu(int(grid_w)), Emu(int(hairline)),
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = _rgb(muted_hex)
        sep.line.fill.background()
        try:
            sep.shadow.inherit = False
        except Exception:
            pass

    return None
