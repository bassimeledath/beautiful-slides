"""Render comparison matrices in all 5 modes. Generates one .pptx per mode."""

import os
import sys

from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


SAMPLE_DATA = {
    "title": "Platform Feature Comparison",
    "row_labels": ["Our Platform", "Competitor A", "Competitor B", "Competitor C"],
    "col_labels": ["SSO/SAML", "REST API", "Mobile App", "Analytics",
                   "Custom Reports", "24/7 Support"],
    "values": [
        ["check", "check", "check", "check", "check", "check"],
        ["check", "check", "cross", "partial", "cross", "check"],
        ["check", "cross", "check", "cross", "cross", "partial"],
        ["cross", "cross", "partial", "cross", "cross", "cross"],
    ],
    "highlight_row": 0,
}


def _paint_slide_bg(slide, hex_):
    from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        slide.part.package.presentation_part.presentation.slide_width,
        slide.part.package.presentation_part.presentation.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(hex_)
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def build(mode, tokens):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _paint_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    bounds = (
        margin,
        margin,
        prs.slide_width - 2 * margin,
        prs.slide_height - 2 * margin,
    )

    render(slide, SAMPLE_DATA, tokens, bounds)

    out = os.path.join(HERE, f"example-{mode}.pptx")
    prs.save(out)
    return out


def main():
    for mode, tokens in MODES.items():
        path = build(mode, tokens)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
