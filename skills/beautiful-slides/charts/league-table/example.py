"""Render the league-table chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render  # noqa: E402


DATA = {
    "title": "Sales leaderboard, Q1 2025",
    "columns": [
        {"name": "Rep",      "key": "rep",    "type": "text",   "width_pct": 2.0},
        {"name": "Revenue",  "key": "rev",    "type": "bar",    "suffix": "K"},
        {"name": "Deals",    "key": "deals",  "type": "number"},
        {"name": "Trend",    "key": "trend",  "type": "arrow"},
        {"name": "Tier",     "key": "tier",   "type": "badge"},
    ],
    "rows": [
        {"rep": "Alice Chen",     "rev": 420, "deals": 18, "trend": "up",   "tier": "Gold"},
        {"rep": "Bob Martinez",   "rev": 380, "deals": 15, "trend": "up",   "tier": "Gold"},
        {"rep": "Carlos Diaz",    "rev": 310, "deals": 12, "trend": "flat", "tier": "Silver"},
        {"rep": "Dana Kim",       "rev": 280, "deals": 10, "trend": "down", "tier": "Silver"},
        {"rep": "Eve Okonkwo",    "rev": 250, "deals": 9,  "trend": "up",   "tier": "Silver"},
        {"rep": "Frank Rossi",    "rev": 210, "deals": 8,  "trend": "flat", "tier": "Bronze"},
        {"rep": "Grace Tanaka",   "rev": 190, "deals": 7,  "trend": "down", "tier": "Bronze"},
        {"rep": "Hiro Patel",     "rev": 140, "deals": 5,  "trend": "down", "tier": "Bronze"},
    ],
    "highlight_top": 2,
    "highlight_bottom": 2,
    "show_rank": True,
}

DATA_COMPACT = {
    "title": "Product NPS rankings",
    "columns": [
        {"name": "Product",  "key": "product", "type": "text",   "width_pct": 2.5},
        {"name": "NPS",      "key": "nps",     "type": "bar"},
        {"name": "Change",   "key": "change",  "type": "arrow"},
        {"name": "Status",   "key": "status",  "type": "badge"},
    ],
    "rows": [
        {"product": "Cloud Platform", "nps": 72, "change": "up",   "status": "Strong"},
        {"product": "Mobile App",     "nps": 65, "change": "up",   "status": "Strong"},
        {"product": "Desktop Suite",  "nps": 48, "change": "flat", "status": "OK"},
        {"product": "API Gateway",    "nps": 41, "change": "down", "status": "At risk"},
        {"product": "Legacy Portal",  "nps": 22, "change": "down", "status": "Critical"},
    ],
    "highlight_top": 1,
    "highlight_bottom": 1,
    "show_rank": True,
}

MODE_DATA = {
    "sv-keynote": DATA,
    "editorial-magazine": DATA_COMPACT,
    "playful-marketing": DATA,
    "consulting-boardroom": DATA_COMPACT,
    "craft-minimal": DATA,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
