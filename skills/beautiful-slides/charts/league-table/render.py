"""League table renderer -- native python-pptx shapes only.

Ranked table with inline visual elements: bars showing magnitude,
directional arrows, and optional row highlighting for top/bottom entries.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


_EMU_PER_PX = 9525


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _set_fill(shape, hex_):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_)


def _no_fill(shape):
    shape.fill.background()


def _no_line(shape):
    shape.line.fill.background()


def _add_text(slide, x, y, w, h, text, font_name, size_pt, hex_color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold = bold
    return tb


def _add_rect(slide, x, y, w, h, fill_hex, radius_px=0):
    w_i = max(1, int(w))
    h_i = max(1, int(h))
    if radius_px and radius_px > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i)
        )
        short_emu = min(w_i, h_i)
        radius_emu = radius_px * _EMU_PER_PX
        ratio = max(0.0, min(0.5, radius_emu / short_emu / 2.0))
        try:
            shape.adjustments[0] = ratio
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i)
        )
    _set_fill(shape, fill_hex)
    _no_line(shape)
    return shape


def _add_line(slide, x1, y1, x2, y2, hex_color, width_emu):
    ln = slide.shapes.add_connector(1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    ln.line.color.rgb = _rgb(hex_color)
    ln.line.width = Emu(int(width_emu))
    return ln


def _draw_arrow(slide, cx, cy, size, direction, color_hex):
    """Draw a simple up/down/flat arrow using a triangle or dash shape.

    direction: "up", "down", or "flat"
    """
    half = size // 2
    if direction == "up":
        # Upward triangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Emu(int(cx - half)), Emu(int(cy - half)),
            Emu(int(size)), Emu(int(size))
        )
        _set_fill(shape, color_hex)
        _no_line(shape)
    elif direction == "down":
        # Downward triangle (flip via rotation)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Emu(int(cx - half)), Emu(int(cy - half)),
            Emu(int(size)), Emu(int(size))
        )
        shape.rotation = 180.0
        _set_fill(shape, color_hex)
        _no_line(shape)
    else:
        # Flat: small horizontal dash
        dash_h = max(1, size // 5)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(cx - half)), Emu(int(cy - dash_h // 2)),
            Emu(int(size)), Emu(int(dash_h))
        )
        _set_fill(shape, color_hex)
        _no_line(shape)
    return shape


def _fmt(v, suffix=""):
    if isinstance(v, float):
        if abs(v - round(v)) < 1e-9:
            s = f"{int(round(v))}"
        else:
            s = f"{v:.1f}"
    else:
        s = str(v)
    return f"{s}{suffix}" if suffix else s


def render(slide, data, tokens, bounds):
    """Render a league table into *slide* within *bounds*.

    data keys
    ---------
    title : str | None           Optional chart title.
    columns : list[dict]         Column definitions. Each dict:
        name : str               Header text.
        key : str                Key into each row dict.
        type : str               "text" | "number" | "bar" | "arrow" | "badge"
        suffix : str             Optional suffix for number/bar values.
        width_pct : float        Relative width weight (default 1.0).
    rows : list[dict]            Data rows. Each is a dict keyed by column keys.
    highlight_top : int          Number of top rows to highlight (default 0).
    highlight_bottom : int       Number of bottom rows to highlight (default 0).
    show_rank : bool             Prepend a rank column (default True).
    value_suffix : str           Global suffix (used if column-level not set).
    """
    x, y, w, h = bounds

    title = data.get("title")
    columns = list(data.get("columns") or [])
    rows = list(data.get("rows") or [])
    highlight_top = int(data.get("highlight_top") or 0)
    highlight_bottom = int(data.get("highlight_bottom") or 0)
    show_rank = bool(data.get("show_rank", True))
    global_suffix = data.get("value_suffix") or ""

    if not columns or not rows:
        return

    # Tokens
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h))
    )
    _set_fill(bg_shape, bg)
    _no_line(bg_shape)

    # Padding
    pad = int(min(w, h) * 0.04)
    ix = x + pad
    iy = y + pad
    iw = w - 2 * pad
    ih = h - 2 * pad

    # Title — allocate enough height for multi-line wrapping
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.6))
        est_chars_per_line = max(1, int(iw / (Pt(title_pt).emu * 0.55)))
        n_title_lines = max(1, math.ceil(len(title) / est_chars_per_line))
        title_h = int(Pt(title_pt).emu * 1.8 * n_title_lines)
        _add_text(
            slide, ix, iy, iw, title_h, title,
            font_name=font_display, size_pt=title_pt, hex_color=text_c,
            bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        iy += title_h + int(Pt(base_pt).emu * 0.5)
        ih = (y + h - pad) - iy

    # Sizing
    header_pt = max(int(base_pt * 0.85), 8)
    row_pt = max(int(base_pt * 0.92), 9)
    header_h = int(Pt(header_pt).emu * 2.4)
    n_rows = len(rows)
    available_h = ih - header_h
    row_h = min(int(Pt(row_pt).emu * 3.0), available_h // max(n_rows, 1))
    row_h = max(row_h, int(Pt(row_pt).emu * 1.8))

    # Column widths
    rank_w = int(iw * 0.06) if show_rank else 0
    table_w = iw - rank_w
    total_weight = sum(c.get("width_pct", 1.0) for c in columns)
    col_widths = []
    for c in columns:
        w_pct = c.get("width_pct", 1.0)
        col_widths.append(int(table_w * w_pct / total_weight))

    # Find max value for bar columns
    bar_maxes = {}
    for ci, col in enumerate(columns):
        if col.get("type") == "bar":
            key = col["key"]
            vals = [abs(float(r.get(key) or 0)) for r in rows]
            bar_maxes[key] = max(vals) if vals else 1

    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    # Header row
    hx = ix + rank_w
    if show_rank:
        _add_text(
            slide, ix, iy, rank_w, header_h, "#",
            font_name=font_body, size_pt=header_pt, hex_color=text_c,
            bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
    for ci, col in enumerate(columns):
        col_type = col.get("type", "text")
        align = PP_ALIGN.LEFT if col_type == "text" else PP_ALIGN.RIGHT
        if col_type in ("bar", "arrow"):
            align = PP_ALIGN.LEFT
        _add_text(
            slide, hx, iy, col_widths[ci], header_h,
            col["name"],
            font_name=font_body, size_pt=header_pt, hex_color=text_c,
            bold=True, align=align, anchor=MSO_ANCHOR.MIDDLE,
        )
        hx += col_widths[ci]

    # Header underline
    header_base_y = iy + header_h
    _add_line(slide, ix, header_base_y, ix + iw, header_base_y, muted, hairline * 2)

    # Rows
    for ri, row in enumerate(rows):
        ry = header_base_y + ri * row_h

        # Ensure we don't draw beyond bounds
        if ry + row_h > y + h - pad:
            break

        # Row highlight
        is_top = ri < highlight_top
        is_bottom = ri >= n_rows - highlight_bottom
        if is_top or is_bottom:
            highlight_color = primary if is_top else accent
            # Semi-transparent highlight via low-opacity rectangle
            hl = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(int(ix)), Emu(int(ry)),
                Emu(int(iw)), Emu(int(row_h))
            )
            hl.fill.solid()
            hl.fill.fore_color.rgb = _rgb(highlight_color)
            _no_line(hl)
            # Apply transparency via XML
            try:
                from lxml import etree
                _ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                srgb = hl._element.find(f".//{{{_ns}}}srgbClr")
                if srgb is not None:
                    alpha = etree.SubElement(srgb, f"{{{_ns}}}alpha")
                    alpha.set("val", "12000")  # 12% opacity
            except Exception:
                pass

        # Row separator
        if ri > 0:
            _add_line(slide, ix, ry, ix + iw, ry, muted, hairline)

        # Rank
        rx = ix
        if show_rank:
            _add_text(
                slide, rx, ry, rank_w, row_h, str(ri + 1),
                font_name=font_mono, size_pt=row_pt, hex_color=text_c,
                bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            rx += rank_w

        # Columns
        for ci, col in enumerate(columns):
            col_type = col.get("type", "text")
            key = col["key"]
            val = row.get(key)
            suffix = col.get("suffix") or global_suffix
            cw = col_widths[ci]
            cx = rx
            cell_pad = int(Pt(row_pt).emu * 0.3)

            if col_type == "text":
                _add_text(
                    slide, cx + cell_pad, ry, cw - 2 * cell_pad, row_h,
                    str(val or ""),
                    font_name=font_body, size_pt=row_pt, hex_color=text_c,
                    bold=(ri < highlight_top), align=PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE,
                )

            elif col_type == "number":
                _add_text(
                    slide, cx + cell_pad, ry, cw - 2 * cell_pad, row_h,
                    _fmt(val, suffix),
                    font_name=font_mono, size_pt=row_pt, hex_color=text_c,
                    bold=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                )

            elif col_type == "bar":
                # Inline bar showing magnitude
                bar_max = bar_maxes.get(key, 1)
                v = float(val or 0)
                frac = abs(v) / bar_max if bar_max > 0 else 0
                bar_region_w = int(cw * 0.55)
                label_w = cw - bar_region_w - 2 * cell_pad
                bar_h = int(row_h * 0.35)
                bar_w = int(bar_region_w * frac)
                bar_y = ry + (row_h - bar_h) // 2

                # Value label on left side
                _add_text(
                    slide, cx + cell_pad, ry, label_w, row_h,
                    _fmt(val, suffix),
                    font_name=font_mono, size_pt=row_pt, hex_color=text_c,
                    bold=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                )

                # Inline bar
                if bar_w > 0:
                    bar_x = cx + cell_pad + label_w + cell_pad
                    _add_rect(slide, bar_x, bar_y, bar_w, bar_h, primary, radius_px)

            elif col_type == "arrow":
                # Directional arrow
                direction = "flat"
                arrow_color = muted
                if val is not None:
                    if isinstance(val, str):
                        direction = val.lower()
                    elif isinstance(val, (int, float)):
                        if val > 0:
                            direction = "up"
                        elif val < 0:
                            direction = "down"
                        else:
                            direction = "flat"

                if direction == "up":
                    arrow_color = accent
                elif direction == "down":
                    arrow_color = primary

                arrow_size = int(Pt(row_pt).emu * 0.8)
                arrow_cx = cx + cw // 2
                arrow_cy = ry + row_h // 2
                _draw_arrow(slide, arrow_cx, arrow_cy, arrow_size, direction, arrow_color)

            elif col_type == "badge":
                # Small badge / label
                badge_text = str(val or "")
                if badge_text:
                    badge_pt = max(int(row_pt * 0.8), 7)
                    badge_w = int(Pt(badge_pt).emu * 0.6 * max(len(badge_text), 2)) + cell_pad * 2
                    badge_h = int(Pt(badge_pt).emu * 1.6)
                    badge_x = cx + (cw - badge_w) // 2
                    badge_y = ry + (row_h - badge_h) // 2
                    _add_rect(slide, badge_x, badge_y, badge_w, badge_h, muted, radius_px)
                    _add_text(
                        slide, badge_x, badge_y, badge_w, badge_h,
                        badge_text,
                        font_name=font_mono, size_pt=badge_pt, hex_color=bg,
                        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    )

            rx += cw
