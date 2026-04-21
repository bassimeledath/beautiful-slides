"""Value-chain (Porter) — chevron arrows + support bars + margin bar.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _luminance(hex_):
    """Return relative luminance (0-1) of a hex color."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255.0, int(h[2:4], 16) / 255.0, int(h[4:6], 16) / 255.0
    r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
    g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
    b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_color(fill_hex, light="#FFFFFF", dark="#000000"):
    """Return light or dark hex depending on fill luminance."""
    return light if _luminance(fill_hex) < 0.4 else dark


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(max(1, int(w))), Emu(max(1, int(h))))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(2)))
    tf.margin_right = Emu(int(Pt(2)))
    tf.margin_top = Emu(int(Pt(1)))
    tf.margin_bottom = Emu(int(Pt(1)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_chevron(slide, cx, cy, cw, ch, notch, fill_hex):
    """Draw a single chevron (pentagon / arrow shape).

    Parameters
    ----------
    cx, cy : left-top corner of the bounding box
    cw, ch : width and height
    notch  : width of the V-notch on the left side (0 for the first chevron)
    fill_hex : fill color
    """
    # Points (clockwise from top-left):
    #   For first chevron (notch=0): flat left edge
    #     TL -> TR-tip-base -> arrow-tip -> BR-tip-base -> BL
    #   For subsequent (notch>0): V-notch on left
    #     TL -> TR-tip-base -> arrow-tip -> BR-tip-base -> BL -> notch-point

    arrow_point = int(cw * 0.15)  # how far the arrow tip extends
    left = int(cx)
    right = int(cx + cw - arrow_point)
    tip = int(cx + cw)
    top = int(cy)
    bottom = int(cy + ch)
    mid_y = int(cy + ch / 2)

    if notch <= 0:
        # First chevron: flat left, arrow right
        pts = [
            (left, top),
            (right, top),
            (tip, mid_y),
            (right, bottom),
            (left, bottom),
        ]
    else:
        # Subsequent: V-notch left, arrow right
        notch_x = int(cx + notch)
        pts = [
            (left, top),
            (right, top),
            (tip, mid_y),
            (right, bottom),
            (left, bottom),
            (notch_x, mid_y),
        ]

    ff = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
    ff.add_line_segments(pts[1:], close=True)
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    return shape


def render(slide, data, tokens, bounds):
    """Render a Porter value chain.

    data:
        title          - optional string
        primary        - list of str (primary activity labels)
        support        - list of str (support activity labels)
        margin_label   - str (default "Margin")

    Also accepts legacy keys primary_activities / support_activities
    (list of {"label": str}).
    """
    x, y, w, h = bounds

    bg = tokens["bg"]
    primary_c = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = int(tokens["font_size_base_pt"])

    # --- Background ---
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()

    # --- Parse data (support both new flat and legacy dict formats) ---
    raw_primary = data.get("primary") or data.get("primary_activities") or []
    raw_support = data.get("support") or data.get("support_activities") or []

    primary_labels = []
    for item in raw_primary:
        if isinstance(item, str):
            primary_labels.append(item)
        elif isinstance(item, dict):
            primary_labels.append(item.get("label", ""))
    support_labels = []
    for item in raw_support:
        if isinstance(item, str):
            support_labels.append(item)
        elif isinstance(item, dict):
            support_labels.append(item.get("label", ""))

    margin_label = data.get("margin_label", "Margin")
    title = data.get("title")

    if not primary_labels:
        return

    # --- Layout constants ---
    pad = int(min(w, h) * 0.035)
    ix = x + pad
    iy = y + pad
    iw = w - 2 * pad
    ih = h - 2 * pad

    # Title
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.5))
        title_h = int(Pt(title_pt) * 1.8)
        _add_textbox(
            slide, ix, iy, iw, title_h,
            title, font_display, title_pt, text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        iy += title_h + int(Pt(base_pt) * 0.4)
        ih = (y + h - pad) - iy

    # Margin bar dimensions
    margin_w = int(iw * 0.06)
    right_pad = int(iw * 0.01)
    content_w = iw - margin_w - right_pad - int(iw * 0.02)  # gap before margin

    # Section heights
    section_label_h = int(Pt(base_pt) * 1.6)
    section_label_pt = max(int(base_pt * 0.7), 8)
    gap_between = int(ih * 0.03)

    has_support = len(support_labels) > 0
    if has_support:
        primary_zone_h = int(ih * 0.45)
        support_zone_h = int(ih * 0.40)
    else:
        primary_zone_h = int(ih * 0.75)
        support_zone_h = 0

    # --- Primary Activities ---
    # Section label
    _add_textbox(
        slide, ix, iy, content_w, section_label_h,
        "Primary Activities", font_body, section_label_pt, muted,
        align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.BOTTOM,
    )

    chevron_y = iy + section_label_h
    chevron_h = primary_zone_h - section_label_h

    n_primary = len(primary_labels)
    gap_pct = 0.025
    total_gap = content_w * gap_pct * (n_primary - 1) if n_primary > 1 else 0
    chevron_w = int((content_w - total_gap) / n_primary)
    chevron_gap = int(content_w * gap_pct) if n_primary > 1 else 0
    notch_depth = int(chevron_w * 0.15)

    # Font size for chevron labels — scale down if many items
    chevron_label_pt = max(int(base_pt * 0.85), 9)
    if n_primary > 6:
        chevron_label_pt = max(int(base_pt * 0.7), 8)

    for i, label in enumerate(primary_labels):
        t = i / max(n_primary - 1, 1)
        # Interpolate primary -> accent, then desaturate 25% toward bg
        raw_fill = _lerp_hex(primary_c, accent, t)
        fill = _lerp_hex(raw_fill, bg, 0.25)

        cx = ix + i * (chevron_w + chevron_gap)
        notch = 0 if i == 0 else notch_depth

        _draw_chevron(slide, cx, chevron_y, chevron_w, chevron_h, notch, fill)

        # Text inside chevron — use luminance-based contrast
        text_color = _contrast_color(fill)
        # Offset text area to account for notch and arrow tip
        text_inset_left = notch + int(chevron_w * 0.04)
        arrow_tip = int(chevron_w * 0.15)
        text_area_w = chevron_w - text_inset_left - arrow_tip
        _add_textbox(
            slide,
            cx + text_inset_left, chevron_y, max(1, text_area_w), chevron_h,
            label, font_body, chevron_label_pt, text_color,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )

    # --- Support Activities ---
    if has_support:
        support_y = iy + primary_zone_h + gap_between

        _add_textbox(
            slide, ix, support_y, content_w, section_label_h,
            "Support Activities", font_body, section_label_pt, muted,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.BOTTOM,
        )

        bars_y = support_y + section_label_h
        bars_h = support_zone_h - section_label_h

        n_support = len(support_labels)
        bar_gap = int(bars_h * 0.04)
        bar_h = max(int((bars_h - bar_gap * (n_support - 1)) / n_support), int(Pt(base_pt) * 1.2))
        bar_label_pt = max(int(base_pt * 0.8), 9)

        # Subtle tint: primary blended 85% toward bg
        bar_fill = _lerp_hex(primary_c, bg, 0.85)

        for i, label in enumerate(support_labels):
            by = bars_y + i * (bar_h + bar_gap)
            # Bar rect
            bar_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(int(ix)), Emu(int(by)),
                Emu(int(content_w)), Emu(int(bar_h)),
            )
            bar_shape.fill.solid()
            bar_shape.fill.fore_color.rgb = _rgb(bar_fill)
            bar_shape.line.fill.background()

            # Text left-aligned inside bar
            _add_textbox(
                slide,
                ix + int(Pt(base_pt) * 0.6), by,
                content_w - int(Pt(base_pt) * 1.2), bar_h,
                label, font_body, bar_label_pt, text_c,
                align=PP_ALIGN.LEFT, bold=False, anchor=MSO_ANCHOR.MIDDLE,
            )

    # --- Margin bar (right side, spans both sections) ---
    margin_x = ix + content_w + int(iw * 0.02)
    margin_top = chevron_y
    if has_support:
        margin_bottom = support_y + support_zone_h
    else:
        margin_bottom = chevron_y + chevron_h
    margin_h = margin_bottom - margin_top

    # Ensure margin bar stays within bounds
    max_right = x + w - int(iw * 0.01)
    if margin_x + margin_w > max_right:
        margin_w = max(1, max_right - margin_x)

    margin_fill = _lerp_hex(accent, bg, 0.40)

    margin_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(margin_x)), Emu(int(margin_top)),
        Emu(int(margin_w)), Emu(int(margin_h)),
    )
    margin_shape.fill.solid()
    margin_shape.fill.fore_color.rgb = _rgb(margin_fill)
    margin_shape.line.fill.background()

    # Margin text — horizontal, centered vertically
    margin_text_color = _contrast_color(margin_fill)
    margin_pt = max(int(base_pt * 0.75), 9)
    _add_textbox(
        slide,
        margin_x, margin_top, margin_w, margin_h,
        margin_label, font_body, margin_pt, margin_text_color,
        align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
    )

    return None
