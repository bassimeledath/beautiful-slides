"""Render the histogram chart in all 5 modes as separate 16:9 slides."""

import random
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render  # noqa: E402


# Synthetic response-time data (log-normal-ish)
random.seed(42)
_RAW = sorted([round(random.gauss(120, 40), 1) for _ in range(200)])

DATA_BASIC = {
    "title": "API response time distribution",
    "values": _RAW,
    "bins": 12,
    "x_label": "Latency (ms)",
    "y_label": "Frequency",
    "show_counts": False,
}

DATA_COUNTS = {
    "title": "Deal size distribution, Q1 pipeline",
    "values": [round(random.uniform(5, 250), 0) for _ in range(150)],
    "bins": 10,
    "x_label": "Deal size ($K)",
    "show_counts": True,
    "value_suffix": "K",
}

MODE_DATA = {
    "sv-keynote": DATA_BASIC,
    "editorial-magazine": DATA_COUNTS,
    "playful-marketing": DATA_BASIC,
    "consulting-boardroom": DATA_COUNTS,
    "craft-minimal": DATA_BASIC,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
