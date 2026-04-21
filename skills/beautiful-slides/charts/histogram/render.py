"""Histogram chart renderer -- native python-pptx shapes only.

Adjacent bars (no gaps) representing frequency of values within binned ranges.
X-axis is continuous and binned. Configurable bin count.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


_EMU_PER_PX = 9525


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _set_fill(shape, hex_):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_)


def _no_line(shape):
    shape.line.fill.background()


def _add_text(slide, x, y, w, h, text, font_name, size_pt, hex_color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold = bold
    return tb


def _add_bar(slide, x, y, w, h, fill_hex):
    """Adjacent histogram bar -- always a plain rectangle (no rounding, no gaps)."""
    w_i = max(1, int(w))
    h_i = max(1, int(h))
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i)
    )
    _set_fill(shape, fill_hex)
    _no_line(shape)
    return shape


def _add_line(slide, x1, y1, x2, y2, hex_color, width_emu):
    ln = slide.shapes.add_connector(1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    ln.line.color.rgb = _rgb(hex_color)
    ln.line.width = Emu(int(width_emu))
    return ln


def _nice_ticks(vmax, target=5):
    """Return (ticks_list, ceiling) for y-axis frequency."""
    if vmax <= 0:
        return [0, 1], 1
    raw = vmax / target
    magnitude = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    residual = raw / magnitude
    if residual < 1.5:
        step = 1 * magnitude
    elif residual < 3:
        step = 2 * magnitude
    elif residual < 7:
        step = 5 * magnitude
    else:
        step = 10 * magnitude
    top = step * math.ceil(vmax / step)
    ticks = []
    v = 0.0
    while v <= top + 1e-9:
        ticks.append(v)
        v += step
    return ticks, top


def _fmt(v):
    if abs(v - round(v)) < 1e-9:
        return f"{int(round(v))}"
    return f"{v:.1f}"


def _bin_values(values, bins):
    """Bin raw values into *bins* equal-width buckets.

    Returns (edges, counts) where len(edges) == bins + 1 and len(counts) == bins.
    """
    if not values:
        return [], []
    vmin = min(values)
    vmax = max(values)
    if vmax == vmin:
        vmax = vmin + 1
    width = (vmax - vmin) / bins
    edges = [vmin + i * width for i in range(bins + 1)]
    counts = [0] * bins
    for v in values:
        idx = int((v - vmin) / width)
        if idx >= bins:
            idx = bins - 1
        counts[idx] += 1
    return edges, counts


def render(slide, data, tokens, bounds):
    """Render a histogram into *slide* within *bounds*.

    data keys
    ---------
    values : list[float]       Raw data points to bin.
    bins : int                 Number of bins (default 10).
    title : str | None         Optional chart title.
    x_label : str | None       Optional x-axis label.
    y_label : str | None       Optional y-axis label (defaults to "Frequency").
    value_suffix : str         Suffix for x-axis tick labels.
    show_counts : bool         Show count above each bar (default False).
    """
    x, y, w, h = bounds

    values = list(data.get("values") or [])
    n_bins = int(data.get("bins") or 10)
    title = data.get("title")
    x_label = data.get("x_label")
    y_label = data.get("y_label")
    value_suffix = data.get("value_suffix") or ""
    show_counts = bool(data.get("show_counts", False))

    if not values:
        return

    edges, counts = _bin_values(values, n_bins)
    if not counts:
        return

    # Tokens
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h))
    )
    _set_fill(bg_shape, bg)
    _no_line(bg_shape)

    # Padding
    pad = int(min(w, h) * 0.04)
    ix = x + pad
    iy = y + pad
    iw = w - 2 * pad
    ih = h - 2 * pad

    # Title
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.6))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_text(
            slide, ix, iy, iw, title_h, title,
            font_name=font_display, size_pt=title_pt, hex_color=text_c,
            bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        iy += title_h + int(Pt(base_pt).emu * 0.6)
        ih = (y + h - pad) - iy

    # Y-label (subtitle style)
    if y_label:
        sub_pt = max(int(base_pt * 0.85), 8)
        sub_h = int(Pt(sub_pt).emu * 1.6)
        _add_text(
            slide, ix, iy, iw, sub_h, y_label,
            font_name=font_body, size_pt=sub_pt, hex_color=text_c,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        iy += sub_h
        ih -= sub_h

    # Sizing
    tick_pt = max(int(base_pt * 0.78), 8)
    cat_pt = max(int(base_pt * 0.85), 8)
    val_pt = max(int(base_pt * 0.78), 8)

    # Margins for axes
    freq_max = max(counts)
    freq_ticks, freq_top = _nice_ticks(freq_max)
    tick_labels = [_fmt(t) for t in freq_ticks]
    max_tick_chars = max(len(s) for s in tick_labels) if tick_labels else 1
    left_margin = int(Pt(tick_pt).emu * 0.65 * (max_tick_chars + 1))
    bottom_margin = int(Pt(cat_pt).emu * 2.2)
    if x_label:
        bottom_margin += int(Pt(cat_pt).emu * 1.6)
    top_margin = int(Pt(val_pt).emu * 1.6) if show_counts else int(Pt(val_pt).emu * 0.4)
    right_margin = int(Pt(val_pt).emu * 1.2)

    plot_x = ix + left_margin
    plot_y = iy + top_margin
    plot_w = max(1, iw - left_margin - right_margin)
    plot_h = max(1, ih - top_margin - bottom_margin)

    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    # Y-axis gridlines and tick labels
    for i, t in enumerate(freq_ticks):
        ty = plot_y + plot_h - int(plot_h * (t / freq_top)) if freq_top > 0 else plot_y + plot_h
        lbl_w = left_margin - int(Pt(tick_pt).emu * 0.3)
        lbl_h = int(Pt(tick_pt).emu * 1.4)
        _add_text(
            slide, ix, ty - lbl_h // 2, lbl_w, lbl_h,
            _fmt(t),
            font_name=font_mono, size_pt=tick_pt, hex_color=text_c,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )
        if i > 0:
            _add_line(slide, plot_x, ty, plot_x + plot_w, ty, muted, hairline)

    # Baseline (x-axis)
    base_y = plot_y + plot_h
    _add_line(slide, plot_x, base_y, plot_x + plot_w, base_y, muted, hairline * 2)

    # Left axis
    _add_line(slide, plot_x, plot_y, plot_x, base_y, muted, hairline * 2)

    # Bars -- adjacent (no gaps)
    bar_w = plot_w / n_bins
    for bi in range(n_bins):
        count = counts[bi]
        bh = int(plot_h * (count / freq_top)) if freq_top > 0 else 0
        bx = plot_x + bi * bar_w
        by = base_y - bh
        if bh > 0:
            _add_bar(slide, bx, by, bar_w, bh, primary)

        if show_counts and count > 0:
            vh = int(Pt(val_pt).emu * 1.3)
            _add_text(
                slide, bx, by - vh - int(Pt(val_pt).emu * 0.1),
                bar_w, vh, str(count),
                font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
            )

    # X-axis tick labels at bin edges
    # Show a subset to avoid crowding: first, last, and every N-th
    max_x_labels = 8
    label_skip = max(1, (n_bins + max_x_labels - 1) // max_x_labels)
    for ei in range(0, n_bins + 1):
        if ei % label_skip != 0 and ei != n_bins:
            continue
        ex = plot_x + ei * bar_w
        lbl_w = int(bar_w * 1.6)
        lbl_h = int(Pt(cat_pt).emu * 1.6)
        _add_text(
            slide, ex - lbl_w // 2,
            base_y + int(Pt(cat_pt).emu * 0.3),
            lbl_w, lbl_h,
            _fmt(edges[ei]) + value_suffix,
            font_name=font_mono, size_pt=cat_pt, hex_color=text_c,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )

    # X-axis label
    if x_label:
        xl_pt = max(int(base_pt * 0.85), 8)
        xl_h = int(Pt(xl_pt).emu * 1.6)
        _add_text(
            slide, plot_x,
            base_y + int(Pt(cat_pt).emu * 2.0),
            plot_w, xl_h,
            x_label,
            font_name=font_body, size_pt=xl_pt, hex_color=text_c,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )
