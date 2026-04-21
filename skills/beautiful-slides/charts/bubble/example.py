"""Render the bubble chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_BASIC = {
    "title": "Market opportunity map",
    "x_label": "Market growth (%)",
    "y_label": "Market share (%)",
    "size_label": "Revenue ($M)",
    "show_labels": True,
    "show_size_legend": True,
    "bubbles": [
        {"x": 12, "y": 35, "size": 50,  "label": "Alpha"},
        {"x": 28, "y": 60, "size": 120, "label": "Beta"},
        {"x": 45, "y": 22, "size": 30,  "label": "Gamma"},
        {"x": 60, "y": 75, "size": 200, "label": "Delta"},
        {"x": 38, "y": 50, "size": 80,  "label": "Epsilon"},
    ],
}

DATA_PORTFOLIO = {
    "title": "Product portfolio view",
    "x_label": "Customer satisfaction",
    "y_label": "Revenue contribution (%)",
    "size_label": "Headcount",
    "show_labels": True,
    "show_size_legend": True,
    "bubbles": [
        {"x": 8.5, "y": 40, "size": 200, "label": "Core Platform"},
        {"x": 6.2, "y": 15, "size": 50,  "label": "Legacy API"},
        {"x": 9.1, "y": 25, "size": 80,  "label": "Mobile App"},
        {"x": 7.0, "y": 10, "size": 30,  "label": "Admin Tools"},
        {"x": 8.0, "y": 30, "size": 120, "label": "Analytics"},
        {"x": 5.5, "y":  5, "size": 15,  "label": "Experiments"},
    ],
}

DATA_DENSE = {
    "title": "Startup ecosystem snapshot",
    "x_label": "Funding round ($M)",
    "y_label": "Monthly growth (%)",
    "size_label": "Team size",
    "show_labels": True,
    "show_size_legend": True,
    "bubbles": [
        {"x":  2, "y": 15, "size": 8,   "label": "Seed1"},
        {"x":  5, "y": 22, "size": 12,  "label": "Seed2"},
        {"x": 10, "y": 30, "size": 25,  "label": "SeriesA1"},
        {"x": 15, "y": 18, "size": 35,  "label": "SeriesA2"},
        {"x": 25, "y": 40, "size": 60,  "label": "SeriesB1"},
        {"x": 35, "y": 28, "size": 80,  "label": "SeriesB2"},
        {"x": 50, "y": 35, "size": 150, "label": "SeriesC"},
        {"x": 70, "y": 20, "size": 200, "label": "Late"},
    ],
}


MODE_DATA = {
    "sv-keynote": DATA_BASIC,
    "editorial-magazine": DATA_PORTFOLIO,
    "playful-marketing": DATA_DENSE,
    "consulting-boardroom": DATA_PORTFOLIO,
    "craft-minimal": DATA_BASIC,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
