"""Bubble chart renderer — native python-pptx shapes only.

Scatter plot where point size (area-scaled circles) encodes a third variable.
Optional size legend. Max 15-20 bubbles. Direct-labeled.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ---- helpers ---------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap=False):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_circle(slide, cx, cy, r_emu, fill_hex, line_hex=None, line_width_pt=0):
    d = int(r_emu * 2)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
        Emu(d), Emu(d),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Pt(line_width_pt)
    shp.shadow.inherit = False
    return shp


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(int(x1)), Emu(int(y1)),
        Emu(int(x2)), Emu(int(y2)),
    )
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(weight_pt)
    return conn


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _nice_ticks(vmin, vmax, target=5):
    if vmax <= vmin:
        vmax = vmin + 1
    span = vmax - vmin
    raw = span / max(target, 1)
    mag = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    for mult in (1, 2, 2.5, 5, 10):
        step = mult * mag
        if span / step <= target * 1.5:
            break
    lo = math.floor(vmin / step) * step
    hi = math.ceil(vmax / step) * step
    ticks = []
    v = lo
    while v <= hi + 1e-9:
        ticks.append(round(v, 6))
        v += step
    return ticks, lo, hi


def _fmt(v):
    if abs(v - round(v)) < 1e-6:
        return f"{int(round(v))}"
    return f"{v:.1f}"


# ---- main render -----------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a bubble chart onto *slide* inside *bounds*, styled by *tokens*."""
    x0, y0, w0, h0 = bounds

    # --- unpack tokens -------------------------------------------------------
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # --- unpack data ---------------------------------------------------------
    title = data.get("title")
    bubbles = list(data.get("bubbles") or [])[:20]  # cap at 20
    x_label = data.get("x_label")
    y_label = data.get("y_label")
    size_label = data.get("size_label")  # legend label for size variable
    show_labels = bool(data.get("show_labels", True))
    show_size_legend = bool(data.get("show_size_legend", True))
    min_radius_pt = data.get("min_radius_pt", 4)
    max_radius_pt = data.get("max_radius_pt", 28)

    if not bubbles:
        return

    # --- background ----------------------------------------------------------
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    # --- outer padding -------------------------------------------------------
    pad = int(min(w0, h0) * 0.035)
    ix, iy = x0 + pad, y0 + pad
    iw, ih = w0 - 2 * pad, h0 - 2 * pad

    # --- title ---------------------------------------------------------------
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.55))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, ix, iy, iw, title_h, title,
                     font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + int(pad * 0.3)
        ih = (y0 + h0 - pad) - iy

    # --- compute plot area ---------------------------------------------------
    tick_pt = max(8, int(base_pt * 0.8))
    label_pt = max(8, int(base_pt * 0.85))

    left_margin = int(iw * 0.08)
    bottom_margin = int(Pt(tick_pt).emu * 1.8) + (
        int(Pt(label_pt).emu * 1.6) if x_label else 0
    )
    top_margin = int(Pt(label_pt).emu * 2.0) if y_label else int(Pt(tick_pt).emu * 0.6)
    # Extra right margin for size legend
    legend_w = int(iw * 0.22) if (show_size_legend and size_label) else 0
    right_margin = int(iw * 0.05) + legend_w

    plot_x = ix + left_margin
    plot_y = iy + top_margin
    plot_w = max(1, iw - left_margin - right_margin)
    plot_h = max(1, ih - top_margin - bottom_margin)

    if plot_w <= 0 or plot_h <= 0:
        return

    # --- compute ranges ------------------------------------------------------
    xs = [b["x"] for b in bubbles]
    ys = [b["y"] for b in bubbles]
    sizes = [b.get("size", 1) for b in bubbles]
    x_span = max(xs) - min(xs) if max(xs) > min(xs) else max(abs(max(xs)), 1)
    y_span = max(ys) - min(ys) if max(ys) > min(ys) else max(abs(max(ys)), 1)

    # Add extra padding to accommodate largest bubbles
    max_r_emu = int(Pt(max_radius_pt).emu)
    x_pad_frac = max_r_emu / max(plot_w, 1) * (x_span if x_span else 1) * 1.2
    y_pad_frac = max_r_emu / max(plot_h, 1) * (y_span if y_span else 1) * 1.2

    x_ticks, x_lo, x_hi = _nice_ticks(min(xs) - x_pad_frac,
                                        max(xs) + x_pad_frac)
    y_ticks, y_lo, y_hi = _nice_ticks(min(ys) - y_pad_frac,
                                        max(ys) + y_pad_frac)

    def to_px(vx, vy):
        fx = (vx - x_lo) / (x_hi - x_lo) if x_hi != x_lo else 0.5
        fy = (vy - y_lo) / (y_hi - y_lo) if y_hi != y_lo else 0.5
        return plot_x + fx * plot_w, plot_y + plot_h - fy * plot_h

    # --- size scaling (area-proportional) ------------------------------------
    s_min = min(sizes) if sizes else 1
    s_max = max(sizes) if sizes else 1
    if s_max == s_min:
        s_max = s_min + 1

    def size_to_r_emu(s):
        # Area scales linearly with value; radius scales with sqrt
        frac = (s - s_min) / (s_max - s_min)
        min_area = math.pi * min_radius_pt ** 2
        max_area = math.pi * max_radius_pt ** 2
        area = min_area + frac * (max_area - min_area)
        r_pt = math.sqrt(area / math.pi)
        return int(Pt(r_pt).emu)

    # --- gridlines -----------------------------------------------------------
    hairline_pt = 0.5
    for tv in y_ticks:
        _, yy = to_px(0, tv)
        _add_line(slide, plot_x, yy, plot_x + plot_w, yy, muted, hairline_pt)
    for tv in x_ticks:
        xx, _ = to_px(tv, 0)
        _add_line(slide, xx, plot_y, xx, plot_y + plot_h, muted, hairline_pt)

    # --- axes ----------------------------------------------------------------
    _add_line(slide, plot_x, plot_y + plot_h, plot_x + plot_w,
              plot_y + plot_h, muted, 0.75)
    _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, muted, 0.75)

    # --- tick labels ---------------------------------------------------------
    tick_label_h = int(Pt(tick_pt).emu * 1.4)
    y_tick_w = int(left_margin * 0.92)
    for tv in y_ticks:
        _, yy = to_px(0, tv)
        _add_textbox(slide,
                     plot_x - y_tick_w - int(pad * 0.15),
                     yy - tick_label_h / 2,
                     y_tick_w, tick_label_h,
                     _fmt(tv), font_mono, tick_pt, text_c,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    x_label_w = int(plot_w / max(len(x_ticks), 1) * 1.6)
    for tv in x_ticks:
        xx, _ = to_px(tv, 0)
        _add_textbox(slide,
                     xx - x_label_w / 2,
                     plot_y + plot_h + int(pad * 0.15),
                     x_label_w, tick_label_h,
                     _fmt(tv), font_mono, tick_pt, text_c,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # --- axis labels ---------------------------------------------------------
    if x_label:
        lbl_h = int(Pt(label_pt).emu * 1.6)
        _add_textbox(slide,
                     plot_x, plot_y + plot_h + tick_label_h + int(pad * 0.2),
                     plot_w, lbl_h,
                     x_label, font_body, label_pt, text_c,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    if y_label:
        # Render as a dedicated subtitle row above the plot area so it
        # doesn't overlap tick labels.
        y_lbl_display = y_label if len(y_label) <= 50 else y_label[:47] + "..."
        lbl_h = int(Pt(label_pt).emu * 1.6)
        _add_textbox(slide, plot_x, plot_y - lbl_h - int(pad * 0.1),
                     plot_w, lbl_h,
                     y_lbl_display, font_body, label_pt, text_c,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)

    # --- bubbles (draw largest first so small ones overlay) -------------------
    sorted_bubbles = sorted(bubbles, key=lambda b: b.get("size", 1), reverse=True)
    # Scale label font down for dense charts to reduce overlap
    n_bubbles = len(bubbles)
    if n_bubbles > 12:
        label_scale = 0.48
    elif n_bubbles > 8:
        label_scale = 0.55
    else:
        label_scale = 0.62
    point_label_pt = max(6, int(base_pt * label_scale))
    point_label_h = int(Pt(point_label_pt).emu * 1.3)

    # Collision detection: track occupied label rectangles
    # Also register each bubble circle as an occupied zone so labels
    # don't overlap other bubbles.
    occupied_rects = []  # list of (lx, ly, lx+w, ly+h)
    lbl_gap = int(Pt(1).emu)

    def _rects_overlap(r1, r2):
        return r1[0] < r2[2] and r1[2] > r2[0] and r1[1] < r2[3] and r1[3] > r2[1]

    def _any_collision(rect):
        for occ in occupied_rects:
            if _rects_overlap(rect, occ):
                return True
        return False

    for b in sorted_bubbles:
        cx, cy = to_px(b["x"], b["y"])
        r = size_to_r_emu(b.get("size", 1))
        color = b.get("color") or primary

        # Clamp center so the full circle stays inside plot area
        cx = max(plot_x + r, min(plot_x + plot_w - r, cx))
        cy = max(plot_y + r, min(plot_y + plot_h - r, cy))

        _add_circle(slide, cx, cy, r, color, line_hex=bg, line_width_pt=1)

        # Register the bubble itself as an occupied rect for label collision
        occupied_rects.append((cx - r, cy - r, cx + r, cy + r))

        if show_labels and b.get("label"):
            lbl = str(b["label"])
            # Truncate long labels
            if len(lbl) > 25:
                lbl = lbl[:22] + "..."
            est_w = int(Pt(point_label_pt).emu * 0.55 * len(lbl))
            est_w = max(est_w, int(Pt(point_label_pt).emu * 2))

            # Try to center label inside bubble if it fits
            if est_w < r * 1.6 and point_label_h < r * 1.6:
                lx = cx - est_w / 2
                ly = cy - point_label_h / 2
                rect = (lx - lbl_gap, ly - lbl_gap,
                        lx + est_w + lbl_gap, ly + point_label_h + lbl_gap)
                if not _any_collision(rect):
                    _add_textbox(slide, lx, ly, est_w, point_label_h,
                                 lbl, font_body, point_label_pt, bg,
                                 bold=True,
                                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                    occupied_rects.append(rect)
                    continue  # placed inside bubble

            # Try multiple external positions: right, left, above, below,
            # and diagonal offsets for dense charts
            spacing = r + int(Pt(2).emu)
            diag = int(spacing * 0.71)  # ~cos(45deg) * spacing
            candidates = [
                (cx + spacing, cy - point_label_h / 2),         # right
                (cx - spacing - est_w, cy - point_label_h / 2), # left
                (cx - est_w / 2, cy - spacing - point_label_h), # above
                (cx - est_w / 2, cy + spacing),                 # below
                (cx + diag, cy - diag - point_label_h),         # top-right
                (cx - diag - est_w, cy - diag - point_label_h), # top-left
                (cx + diag, cy + diag),                          # bottom-right
                (cx - diag - est_w, cy + diag),                  # bottom-left
            ]
            for lx, ly in candidates:
                lx = max(plot_x, min(plot_x + plot_w - est_w, lx))
                ly = max(plot_y, min(plot_y + plot_h - point_label_h, ly))
                rect = (lx - lbl_gap, ly - lbl_gap,
                        lx + est_w + lbl_gap, ly + point_label_h + lbl_gap)
                if not _any_collision(rect):
                    _add_textbox(slide, lx, ly, est_w, point_label_h,
                                 lbl, font_body, point_label_pt, text_c,
                                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
                    occupied_rects.append(rect)
                    break
            # If no position works, skip this label to avoid overlap

    # --- size legend ---------------------------------------------------------
    if show_size_legend and size_label and legend_w > 0:
        legend_x = plot_x + plot_w + int(pad * 0.5)
        legend_y = plot_y + int(plot_h * 0.1)
        legend_inner_w = legend_w - int(pad * 0.5)
        # Hard limit: legend must not extend beyond bounds
        legend_bottom_limit = y0 + h0 - pad

        # Title for legend
        leg_title_h = int(Pt(tick_pt).emu * 1.6)
        _add_textbox(slide, legend_x, legend_y, legend_inner_w, leg_title_h,
                     size_label, font_body, tick_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                     word_wrap=True)
        legend_y += leg_title_h + int(pad * 0.5)

        # Show 3 reference bubbles: min, mid, max
        # Cap the drawn radius so circles fit in the legend column
        ref_vals = [s_min, (s_min + s_max) / 2, s_max]
        max_r_legend = int(legend_inner_w * 0.18)

        for rv in ref_vals:
            r = size_to_r_emu(rv)
            r_draw = min(r, max_r_legend)

            if legend_y + r_draw * 2 > legend_bottom_limit:
                break

            circle_cx = legend_x + r_draw + int(pad * 0.3)
            circle_cy = legend_y + r_draw
            _add_circle(slide, circle_cx, circle_cy, r_draw, muted)

            lbl_x = circle_cx + r_draw + int(pad * 0.4)
            lbl_w = max(1, (legend_x + legend_inner_w) - lbl_x)
            if lbl_w > 0:
                _add_textbox(slide, lbl_x, circle_cy - tick_label_h // 2,
                             lbl_w, tick_label_h,
                             _fmt(rv), font_mono, tick_pt, text_c,
                             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

            # Increase vertical spacing between reference circles
            legend_y += r_draw * 2 + int(pad * 1.2)
