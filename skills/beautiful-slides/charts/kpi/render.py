"""KPI tile — big-number, label, delta, footnote."""

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


_LIGHT_BG_HEXES = {"F6F1E8", "FFF4EB", "FFFFFF", "FCFBF8"}


def _is_dark_mode(bg_hex):
    return bg_hex.lstrip("#").upper() not in _LIGHT_BG_HEXES


def _add_textbox(slide, x, y, w, h, text, *, font, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 letter_spacing=None):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = _rgb(color_hex)
    return tb


def render(slide, data, tokens, bounds):
    x, y, w, h = bounds

    label = data.get("label", "")
    value = data.get("value", "")
    delta = data.get("delta")
    direction = data.get("delta_direction")
    footnote = data.get("footnote")

    bg_hex = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(bg_hex.lstrip("#"))
    text_hex = tokens["text"]
    muted_hex = tokens["muted"]
    primary_hex = tokens["primary"]
    accent_hex = tokens["accent"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = tokens["font_size_base_pt"]
    radius_px = tokens.get("radius_px", 0) or 0

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg_hex)
    bg_shape.line.fill.background()

    dark = _is_dark_mode(bg_hex)

    # Tile background — subtle rounded rect with hairline border.
    radius_emu = Emu(int(radius_px * 9525))  # px → EMU (approx 9525 EMU/px)
    if radius_px > 0:
        tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Emu(int(x)), Emu(int(y)),
                                      Emu(int(w)), Emu(int(h)))
        # python-pptx rounded-rect adjustment is relative (0.0–~0.5).
        try:
            short_side = min(w, h)
            adj = max(0.0, min(0.5, float(radius_emu) / float(short_side)))
            tile.adjustments[0] = adj
        except Exception:
            pass
    else:
        tile = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Emu(int(x)), Emu(int(y)),
                                      Emu(int(w)), Emu(int(h)))

    tile.fill.solid()
    tile.fill.fore_color.rgb = _rgb(bg_hex)
    tile.line.color.rgb = _rgb(muted_hex)
    tile.line.width = Emu(6350)  # 0.5 pt hairline
    # Remove default shadow.
    try:
        sp = tile.shadow
        sp.inherit = False
    except Exception:
        pass

    # Padding inside tile.
    pad_x = int(w * 0.07)
    pad_y = int(h * 0.09)
    inner_x = x + pad_x
    inner_y = y + pad_y
    inner_w = w - 2 * pad_x
    inner_h = h - 2 * pad_y

    # Label (top, small caps-ish via font).
    label_size = max(8, int(base_pt * 0.78))
    label_h = Pt(label_size).emu * 1.5
    _add_textbox(slide, inner_x, inner_y, inner_w, label_h,
                 label.upper() if label else "",
                 font=font_body, size_pt=label_size,
                 color_hex=text_hex, bold=True, align=PP_ALIGN.LEFT)

    # Value — hero number. Size based on tile height.
    value_pt_from_h = (h / 12700.0) * 0.22
    value_pt = max(36, min(int(value_pt_from_h), 96))
    # Auto-size: reduce font when value string is long (> 6 chars)
    val_len = max(1, len(value))
    if val_len > 6:
        value_pt = max(18, int(value_pt * 6.0 / val_len))
    # Further shrink if estimated width still exceeds available width
    est_char_w_emu = Pt(value_pt).emu * 0.55
    est_w = est_char_w_emu * val_len
    if est_w > inner_w:
        scale = inner_w / est_w
        value_pt = max(14, int(value_pt * scale))

    value_h = Pt(value_pt).emu * 1.2
    value_y = inner_y + label_h + Pt(base_pt * 0.35).emu
    _add_textbox(slide, inner_x, value_y, inner_w, value_h,
                 value,
                 font=font_display, size_pt=value_pt,
                 color_hex=text_hex, bold=True, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP)

    # Neon underline under the hero number in dark mode (sv-keynote signature).
    if dark:
        underline_w = int(inner_w * 0.18)
        underline_h = Emu(int(Pt(3).emu))
        underline_y = value_y + value_h + Pt(base_pt * 0.15).emu
        ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Emu(int(inner_x)), Emu(int(underline_y)),
                                    Emu(int(underline_w)), underline_h)
        ul.fill.solid()
        ul.fill.fore_color.rgb = _rgb(primary_hex)
        ul.line.fill.background()
    else:
        underline_y = value_y + value_h
        underline_h = 0

    # Delta — under the value.
    delta_y = (underline_y + (Pt(base_pt * 0.45).emu if dark else Pt(base_pt * 0.25).emu)
               + (int(underline_h) if dark else 0))
    if delta:
        if direction == "up":
            delta_color = accent_hex
            glyph = "\u25B2 "  # ▲
        elif direction == "down":
            delta_color = primary_hex if not dark else accent_hex
            # If primary is the same hue family as up-accent we already got,
            # fall back to muted to signal "attention" without clashing.
            if delta_color == accent_hex:
                delta_color = text_hex
            glyph = "\u25BC "  # ▼
        else:
            delta_color = text_hex
            glyph = ""

        delta_pt = max(10, int(base_pt * 0.95))
        delta_str = f"{glyph}{delta}"
        # Cap delta text: reduce font if string is long, truncate if extreme
        if len(delta_str) > 25:
            delta_str = delta_str[:22] + "..."
        if len(delta_str) > 15:
            delta_pt = max(8, int(delta_pt * 15.0 / len(delta_str)))
        delta_h = Pt(delta_pt).emu * 1.5
        _add_textbox(slide, inner_x, delta_y, inner_w, delta_h,
                     delta_str,
                     font=font_mono, size_pt=delta_pt,
                     color_hex=delta_color, bold=False, align=PP_ALIGN.LEFT)
        delta_y_end = delta_y + delta_h
    else:
        delta_y_end = delta_y

    # Footnote at bottom.
    if footnote:
        foot_pt = max(7, int(base_pt * 0.7))
        foot_h = Pt(foot_pt).emu * 1.4
        foot_y = y + h - pad_y - foot_h
        if foot_y < delta_y_end:
            foot_y = delta_y_end + Pt(base_pt * 0.2).emu
        _add_textbox(slide, inner_x, foot_y, inner_w, foot_h,
                     footnote,
                     font=font_body, size_pt=foot_pt,
                     color_hex=text_hex, bold=False, align=PP_ALIGN.LEFT)

    return tile
