from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.LEFT, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_trapezoid(slide, tl_x, tl_y, tr_x, bl_x, br_x, bottom_y, fill_hex):
    ff = slide.shapes.build_freeform(int(tl_x), int(tl_y), scale=1.0)
    ff.add_line_segments(
        [
            (int(tr_x), int(tl_y)),
            (int(br_x), int(bottom_y)),
            (int(bl_x), int(bottom_y)),
        ],
        close=True,
    )
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    return shape


def render(slide, data, tokens, bounds):
    x, y, w, h = bounds

    bg = tokens["bg"]
    primary = tokens["primary"]
    muted = tokens["muted"]
    text_c = tokens["text"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = tokens["font_size_base_pt"]

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()

    stages = data.get("stages", [])
    if not stages:
        return
    title = data.get("title")
    show_conversion = data.get("show_conversion", True)
    value_format = data.get("value_format", "{:,}")

    n = len(stages)
    max_val = max(s["value"] for s in stages) or 1

    # Layout regions
    cur_y = y
    pad = Emu(Pt(6)).emu if False else int(Pt(6))  # not used directly; keep simple
    title_h = 0
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide,
            x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.6))

    # Reserve label columns on left and right
    label_col_w = int(w * 0.22)
    value_col_w = int(w * 0.18)
    gap_labels = int(w * 0.015)
    funnel_x = x + label_col_w + gap_labels
    funnel_w = w - label_col_w - value_col_w - 2 * gap_labels
    if funnel_w < int(w * 0.3):
        funnel_w = int(w * 0.3)
        label_col_w = int((w - funnel_w) * 0.55)
        value_col_w = int((w - funnel_w) * 0.45) - 2 * gap_labels
        funnel_x = x + label_col_w + gap_labels

    funnel_center = funnel_x + funnel_w // 2
    max_funnel_w = int(funnel_w * 0.98)
    min_funnel_w = max(int(max_funnel_w * 0.18), int(Emu(Pt(8)).emu))

    # Vertical layout: n stages + (n-1) gaps; if show_conversion, the gap holds text
    avail_h = (y + h) - cur_y
    inter_gap = int(Pt(base_pt * (1.6 if show_conversion else 0.55)))
    total_gaps = inter_gap * (n - 1)
    stage_h = max(int((avail_h - total_gaps) / n), int(Pt(base_pt * 1.4)))

    def width_for(val):
        t = val / max_val
        return int(min_funnel_w + (max_funnel_w - min_funnel_w) * t)

    stage_y = cur_y
    prev_bottom_w = None
    for i, stage in enumerate(stages):
        # Top width = width for this stage's value
        top_w = width_for(stage["value"])
        # Bottom width = width for next stage (or slightly narrower if last)
        if i < n - 1:
            bot_w = width_for(stages[i + 1]["value"])
        else:
            bot_w = max(int(top_w * 0.45), min_funnel_w // 2)

        tl_x = funnel_center - top_w // 2
        tr_x = funnel_center + top_w // 2
        bl_x = funnel_center - bot_w // 2
        br_x = funnel_center + bot_w // 2
        tl_y = stage_y
        bottom_y = stage_y + stage_h

        t = i / (n - 1) if n > 1 else 0.0
        fill_hex = _lerp_hex(primary, muted, t)
        _draw_trapezoid(slide, tl_x, tl_y, tr_x, bl_x, br_x, bottom_y, fill_hex)

        # Label on left
        label_x = x
        label_w = label_col_w
        _add_textbox(
            slide,
            label_x, stage_y, label_w, stage_h,
            stage["label"], font_body, base_pt, text_c,
            align=PP_ALIGN.RIGHT, bold=False, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Value on right
        val_x = x + label_col_w + gap_labels + funnel_w + gap_labels
        val_w = (x + w) - val_x
        if val_w < int(w * 0.1):
            val_w = int(w * 0.1)
        val_text = value_format.format(stage["value"])
        _add_textbox(
            slide,
            val_x, stage_y, val_w, stage_h,
            val_text, font_mono, base_pt, text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Conversion between this stage and next
        if show_conversion and i < n - 1:
            next_val = stages[i + 1]["value"]
            if stage["value"] > 0:
                pct = (next_val / stage["value"]) * 100
            else:
                pct = 0
            conv_text = f"{pct:.0f}% \u2193"
            conv_y = bottom_y
            conv_h = inter_gap
            _add_textbox(
                slide,
                funnel_x, conv_y, funnel_w, conv_h,
                conv_text, font_mono, max(int(base_pt * 0.75), 8), text_c,
                align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE,
            )

        stage_y = bottom_y + inter_gap

    return None
