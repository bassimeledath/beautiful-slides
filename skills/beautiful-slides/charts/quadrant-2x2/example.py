"""Render the quadrant-2x2 chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_EFFORT_IMPACT = {
    "title": "Effort vs. impact prioritization",
    "x_label": "Effort",
    "y_label": "Impact",
    "x_low_label": "Low",
    "x_high_label": "High",
    "y_low_label": "Low",
    "y_high_label": "High",
    "quadrant_labels": [
        "Rethink",
        "Big bets",
        "Fill-ins",
        "Quick wins",
    ],
    "items": [
        {"x": 0.2, "y": 0.85, "label": "AI assistant"},
        {"x": 0.75, "y": 0.9, "label": "Platform rebuild"},
        {"x": 0.8, "y": 0.25, "label": "Logo refresh"},
        {"x": 0.15, "y": 0.2, "label": "Email footer"},
        {"x": 0.6, "y": 0.7, "label": "API v2"},
        {"x": 0.35, "y": 0.45, "label": "Dashboard"},
    ],
}

DATA_BCG = {
    "title": "Growth-share matrix",
    "x_label": "Relative market share",
    "y_label": "Market growth rate",
    "x_low_label": "Low",
    "x_high_label": "High",
    "y_low_label": "Low",
    "y_high_label": "High",
    "quadrant_labels": [
        "Question marks",
        "Stars",
        "Cash cows",
        "Dogs",
    ],
    "items": [
        {"x": 0.15, "y": 0.8, "label": "New SaaS"},
        {"x": 0.3,  "y": 0.7, "label": "Mobile app"},
        {"x": 0.8,  "y": 0.85, "label": "Cloud platform"},
        {"x": 0.75, "y": 0.2, "label": "Enterprise suite"},
        {"x": 0.85, "y": 0.15, "label": "Legacy product"},
        {"x": 0.2,  "y": 0.3, "label": "Side project"},
        {"x": 0.6,  "y": 0.6, "label": "Analytics tool"},
    ],
}

DATA_SKILL_WILL = {
    "title": "Team skill-will matrix",
    "x_label": "Skill",
    "y_label": "Will",
    "x_low_label": "Low",
    "x_high_label": "High",
    "y_low_label": "Low",
    "y_high_label": "High",
    "quadrant_labels": [
        "Guide",
        "Delegate",
        "Direct",
        "Excite",
    ],
    "items": [
        {"x": 0.25, "y": 0.8, "label": "Alice"},
        {"x": 0.7,  "y": 0.85, "label": "Bob"},
        {"x": 0.8,  "y": 0.3, "label": "Carol"},
        {"x": 0.15, "y": 0.15, "label": "Dave"},
        {"x": 0.55, "y": 0.6, "label": "Eve"},
    ],
}


MODE_DATA = {
    "sv-keynote": DATA_EFFORT_IMPACT,
    "editorial-magazine": DATA_BCG,
    "playful-marketing": DATA_SKILL_WILL,
    "consulting-boardroom": DATA_BCG,
    "craft-minimal": DATA_EFFORT_IMPACT,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
