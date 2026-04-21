"""Quadrant 2x2 chart renderer — native python-pptx shapes only.

The iconic consulting strategy slide: four-quadrant grid with labeled axes
and quadrant labels. Items placed as dots or labeled markers.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ---- helpers ---------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_circle(slide, cx, cy, r_emu, fill_hex, line_hex=None):
    d = int(r_emu * 2)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
        Emu(d), Emu(d),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
    shp.shadow.inherit = False
    return shp


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(int(x1)), Emu(int(y1)),
        Emu(int(x2)), Emu(int(y2)),
    )
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(weight_pt)
    return conn


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# ---- main render -----------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a 2x2 quadrant chart onto *slide* inside *bounds*."""
    x0, y0, w0, h0 = bounds

    # --- unpack tokens -------------------------------------------------------
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # --- unpack data ---------------------------------------------------------
    title = data.get("title")
    x_label = data.get("x_label")           # e.g., "Effort"
    y_label = data.get("y_label")           # e.g., "Impact"
    x_low_label = data.get("x_low_label")   # e.g., "Low"
    x_high_label = data.get("x_high_label") # e.g., "High"
    y_low_label = data.get("y_low_label")   # e.g., "Low"
    y_high_label = data.get("y_high_label") # e.g., "High"
    # Quadrant labels (clockwise from top-left)
    quadrant_labels = data.get("quadrant_labels") or []
    # Items to place: list of {x: 0-1, y: 0-1, label: str}
    items = list(data.get("items") or [])
    point_radius_pt = data.get("point_radius_pt", 5)

    # --- background ----------------------------------------------------------
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    # --- outer padding -------------------------------------------------------
    pad = int(min(w0, h0) * 0.035)
    ix, iy = x0 + pad, y0 + pad
    iw, ih = w0 - 2 * pad, h0 - 2 * pad

    # --- title ---------------------------------------------------------------
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.55))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, ix, iy, iw, title_h, title,
                     font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + int(pad * 0.3)
        ih = (y0 + h0 - pad) - iy

    # --- compute grid area ---------------------------------------------------
    label_pt = max(8, int(base_pt * 0.85))
    axis_label_pt = max(9, int(base_pt * 0.9))

    # Reserve space for axis labels around the grid
    # Dynamically adjust left_margin based on y_label length
    if y_label:
        _y_lbl_chars = min(len(y_label), 40)
        left_margin = int(Pt(axis_label_pt).emu * max(3.0, min(5.0, _y_lbl_chars * 0.12)))
    else:
        left_margin = int(Pt(axis_label_pt).emu * 1.5)
    bottom_margin = int(Pt(axis_label_pt).emu * 2.5) if x_label else int(Pt(axis_label_pt).emu * 1.0)
    top_margin = int(Pt(axis_label_pt).emu * 1.5)
    right_margin = int(Pt(axis_label_pt).emu * 1.5)

    grid_x = ix + left_margin
    grid_y = iy + top_margin
    grid_w = max(1, iw - left_margin - right_margin)
    grid_h = max(1, ih - top_margin - bottom_margin)

    if grid_w <= 0 or grid_h <= 0:
        return

    mid_x = grid_x + grid_w / 2
    mid_y = grid_y + grid_h / 2

    # --- draw quadrant background fills (very subtle) ------------------------
    # Use muted color with the background to give a subtle differentiation
    # We'll just draw the cross lines and quadrant labels — no background tints
    # to keep it clean like a consulting slide.

    # --- cross lines (the 2x2 grid) -----------------------------------------
    line_weight = 1.5
    # Vertical center line
    _add_line(slide, mid_x, grid_y, mid_x, grid_y + grid_h, muted, line_weight)
    # Horizontal center line
    _add_line(slide, grid_x, mid_y, grid_x + grid_w, mid_y, muted, line_weight)

    # --- outer border (subtle) -----------------------------------------------
    border_weight = 0.75
    _add_line(slide, grid_x, grid_y, grid_x + grid_w, grid_y, muted, border_weight)
    _add_line(slide, grid_x + grid_w, grid_y, grid_x + grid_w, grid_y + grid_h, muted, border_weight)
    _add_line(slide, grid_x, grid_y + grid_h, grid_x + grid_w, grid_y + grid_h, muted, border_weight)
    _add_line(slide, grid_x, grid_y, grid_x, grid_y + grid_h, muted, border_weight)

    # --- quadrant labels (centered in each quadrant, ~20% smaller) -----------
    q_label_pt = max(9, int(base_pt * 0.88))
    q_label_h = int(Pt(q_label_pt).emu * 1.8)
    half_w = grid_w / 2
    half_h = grid_h / 2

    # Order: [top-left, top-right, bottom-right, bottom-left]
    # Following the standard 2x2 convention
    q_positions = [
        (grid_x, grid_y, half_w, half_h),                           # top-left
        (grid_x + half_w, grid_y, half_w, half_h),                  # top-right
        (grid_x + half_w, grid_y + half_h, half_w, half_h),         # bottom-right
        (grid_x, grid_y + half_h, half_w, half_h),                  # bottom-left
    ]
    q_pad = int(min(half_w, half_h) * 0.08)

    # Collision detection helpers
    occupied_rects = []  # list of (x1, y1, x2, y2)
    lbl_gap = int(Pt(1).emu)

    def _rects_overlap(r1, r2):
        return r1[0] < r2[2] and r1[2] > r2[0] and r1[1] < r2[3] and r1[3] > r2[1]

    def _any_collision(rect):
        for occ in occupied_rects:
            if _rects_overlap(rect, occ):
                return True
        return False

    for i, ql in enumerate(quadrant_labels[:4]):
        if not ql:
            continue
        qx, qy, qw, qh = q_positions[i]
        _add_textbox(slide,
                     qx + q_pad, qy + q_pad,
                     qw - 2 * q_pad, q_label_h,
                     ql, font_display, q_label_pt, text_c,
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # Record quadrant label as exclusion zone for item labels
        occupied_rects.append((qx + q_pad, qy + q_pad,
                               qx + qw - q_pad, qy + q_pad + q_label_h))

    # --- axis labels ---------------------------------------------------------
    ax_label_h = int(Pt(axis_label_pt).emu * 1.6)

    # x-axis label (centered below grid)
    if x_label:
        _add_textbox(slide,
                     grid_x, grid_y + grid_h + int(Pt(axis_label_pt).emu * 0.5),
                     grid_w, ax_label_h,
                     x_label, font_body, axis_label_pt, text_c,
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # y-axis label (to the left, vertically centered — rendered as horizontal text)
    if y_label:
        y_lbl_display = y_label if len(y_label) <= 40 else y_label[:37] + "..."
        y_lbl_w = left_margin - int(pad * 0.3)
        y_lbl_h = int(Pt(axis_label_pt).emu * 3.0)  # allow wrapping
        _add_textbox(slide,
                     ix, grid_y + (grid_h - y_lbl_h) / 2,
                     y_lbl_w, y_lbl_h,
                     y_lbl_display, font_body, axis_label_pt, text_c,
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Low/High endpoint labels on axes
    endpoint_pt = max(8, int(base_pt * 0.78))
    endpoint_h = int(Pt(endpoint_pt).emu * 1.4)

    if x_low_label:
        _add_textbox(slide,
                     grid_x, grid_y + grid_h + int(Pt(endpoint_pt).emu * 0.15),
                     half_w, endpoint_h,
                     x_low_label, font_body, endpoint_pt, text_c,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    if x_high_label:
        _add_textbox(slide,
                     grid_x + half_w, grid_y + grid_h + int(Pt(endpoint_pt).emu * 0.15),
                     half_w, endpoint_h,
                     x_high_label, font_body, endpoint_pt, text_c,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)
    if y_low_label:
        y_end_w = left_margin - int(pad * 0.3)
        _add_textbox(slide,
                     ix, grid_y + grid_h - endpoint_h,
                     y_end_w, endpoint_h,
                     y_low_label, font_body, endpoint_pt, text_c,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    if y_high_label:
        y_end_w = left_margin - int(pad * 0.3)
        _add_textbox(slide,
                     ix, grid_y,
                     y_end_w, endpoint_h,
                     y_high_label, font_body, endpoint_pt, text_c,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # --- items (placed as dots with labels) ----------------------------------
    r_emu = int(Pt(point_radius_pt).emu)
    item_label_pt = max(6, int(base_pt * 0.60))
    item_label_h = int(Pt(item_label_pt).emu * 1.3)

    for item in items:
        # x and y are 0..1 fractions within the grid (0=left/bottom, 1=right/top)
        fx = float(item.get("x", 0.5))
        fy = float(item.get("y", 0.5))
        fx = max(0.0, min(1.0, fx))
        fy = max(0.0, min(1.0, fy))

        cx = grid_x + fx * grid_w
        cy = grid_y + (1.0 - fy) * grid_h  # invert y: 1 = top

        color = item.get("color") or primary

        # Clamp so dot stays inside grid
        cx = max(grid_x + r_emu, min(grid_x + grid_w - r_emu, cx))
        cy = max(grid_y + r_emu, min(grid_y + grid_h - r_emu, cy))

        _add_circle(slide, cx, cy, r_emu, color)

        if item.get("label"):
            lbl = str(item["label"])
            # Truncate long labels
            if len(lbl) > 25:
                lbl = lbl[:22] + "..."
            est_w = int(Pt(item_label_pt).emu * 0.55 * len(lbl))
            est_w = max(est_w, int(Pt(item_label_pt).emu * 2))

            # Try multiple positions: right, left, above, below
            spacing = r_emu + int(Pt(2).emu)
            candidates = [
                (cx + spacing, cy - item_label_h / 2),                           # right
                (cx - spacing - est_w, cy - item_label_h / 2),                   # left
                (cx - est_w / 2, cy - spacing - item_label_h),                   # above
                (cx - est_w / 2, cy + spacing),                                  # below
            ]
            placed = False
            for lx, ly in candidates:
                lx = max(grid_x, min(grid_x + grid_w - est_w, lx))
                ly = max(grid_y, min(grid_y + grid_h - item_label_h, ly))
                rect = (lx - lbl_gap, ly - lbl_gap,
                        lx + est_w + lbl_gap, ly + item_label_h + lbl_gap)
                if not _any_collision(rect):
                    _add_textbox(slide, lx, ly, est_w, item_label_h,
                                 lbl, font_body, item_label_pt, text_c,
                                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                                 word_wrap=False)
                    occupied_rects.append(rect)
                    placed = True
                    break
            # If no position works, skip this label to avoid overlap
