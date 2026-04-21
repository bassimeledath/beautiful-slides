"""Stacked-bar chart renderer -- native python-pptx shapes only.

Bars split into additive colored segments showing part-to-whole within
each category.  Supports vertical / horizontal orientation and an optional
percentage-stacked variant.  2-5 segments (series) max.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# ── constants ──────────────────────────────────────────────────────────
_EMU_PER_PX = 9525
_EMU_PER_PT = 12700

# Palette slots beyond primary/accent: derive from muted + text.
_EXTRA_COLORS = ["primary", "accent", "muted", "text", "bg"]


# ── helpers ────────────────────────────────────────────────────────────
def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _set_fill(shape, hex_):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_)


def _no_line(shape):
    shape.line.fill.background()


def _style_run(run, font_name, size_pt, hex_color, bold=False):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold = bold


def _add_text(slide, x, y, w, h, text, font_name, size_pt, hex_color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _style_run(run, font_name, size_pt, hex_color, bold=bold)
    return tb


def _add_bar(slide, x, y, w, h, fill_hex, radius_px):
    w_i = max(1, int(w))
    h_i = max(1, int(h))
    if radius_px and radius_px > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i),
        )
        short_emu = min(w_i, h_i)
        radius_emu = radius_px * _EMU_PER_PX
        ratio = max(0.0, min(0.5, radius_emu / short_emu / 2.0))
        try:
            shape.adjustments[0] = ratio
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i),
        )
    _set_fill(shape, fill_hex)
    _no_line(shape)
    return shape


def _add_line(slide, x1, y1, x2, y2, hex_color, width_emu):
    ln = slide.shapes.add_connector(
        1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)),
    )
    ln.line.color.rgb = _rgb(hex_color)
    ln.line.width = Emu(int(width_emu))
    return ln


def _nice_ticks(vmax, target=5):
    if vmax <= 0:
        return [0, 1], 1
    import math
    raw = vmax / target
    magnitude = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    residual = raw / magnitude
    if residual < 1.5:
        step = 1 * magnitude
    elif residual < 3:
        step = 2 * magnitude
    elif residual < 7:
        step = 5 * magnitude
    else:
        step = 10 * magnitude
    top = step * math.ceil(vmax / step)
    ticks = []
    v = 0.0
    while v <= top + 1e-9:
        ticks.append(v)
        v += step
    return ticks, top


def _fmt(v, suffix=""):
    if abs(v - round(v)) < 1e-9:
        s = f"{int(round(v))}"
    else:
        s = f"{v:.1f}"
    return f"{s}{suffix}" if suffix else s


def _series_colors(tokens, n):
    """Return *n* hex colors from the token palette, cycling if needed."""
    # Build a palette list: primary, accent, then blend muted toward primary/accent
    palette = [tokens["primary"], tokens["accent"]]
    if n > 2:
        palette.append(tokens["muted"])
    if n > 3:
        palette.append(tokens["text"])
    if n > 4:
        # A fifth color derived by lightening primary
        palette.append(tokens.get("muted", tokens["primary"]))
    return palette[:n]


# ── public API ─────────────────────────────────────────────────────────
def render(slide, data, tokens, bounds):
    x, y, w, h = bounds

    orientation = (data.get("orientation") or "vertical").lower()
    title = data.get("title")
    categories = list(data.get("categories") or [])
    series = list(data.get("series") or [])
    value_suffix = data.get("value_suffix") or ""
    show_values = bool(data.get("show_values", True))
    percent = bool(data.get("percent", False))

    if not categories or not series:
        return

    series = series[:5]  # cap at 5 segments

    bg = tokens["bg"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    colors = _series_colors(tokens, len(series))

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    _set_fill(bg_shape, bg)
    _no_line(bg_shape)

    # Padding
    pad = int(min(w, h) * 0.04)
    ix, iy = x + pad, y + pad
    iw, ih = w - 2 * pad, h - 2 * pad

    # Title
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.6))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_text(slide, ix, iy, iw, title_h, title,
                  font_name=font_display, size_pt=title_pt, hex_color=text_c,
                  bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + int(Pt(base_pt).emu * 0.6)
        ih = (y + h - pad) - iy

    # Legend
    legend_h = 0
    if len(series) >= 2:
        legend_pt = max(int(base_pt * 0.85), 8)
        legend_h = int(Pt(legend_pt).emu * 2.0)
        _draw_legend(slide, ix, iy, iw, legend_h, series, colors,
                     font_body, legend_pt, text_c, muted, radius_px)
        iy += legend_h
        ih -= legend_h

    # Compute stacked totals per category
    n_cat = len(categories)
    totals = []
    for ci in range(n_cat):
        t = 0
        for s in series:
            vals = s.get("values") or []
            t += vals[ci] if ci < len(vals) else 0
        totals.append(t)

    if percent:
        # percentage-stacked: all bars go to 100%
        vmax_display = 100
        ticks, _ = _nice_ticks(100)
        ticks = [0, 25, 50, 75, 100]
    else:
        vmax_raw = max(totals) if totals else 1
        ticks, vmax_display = _nice_ticks(max(vmax_raw, 0.0))

    tick_pt = max(int(base_pt * 0.78), 8)
    cat_pt = max(int(base_pt * 0.92), 9)
    val_pt = max(int(base_pt * 0.70), 7)

    if orientation == "horizontal":
        _draw_horizontal(
            slide, ix, iy, iw, ih, categories, series, colors,
            ticks, vmax_display, totals, percent, text_c, muted,
            font_body, font_mono, cat_pt, tick_pt, val_pt,
            radius_px, value_suffix, show_values,
        )
    else:
        _draw_vertical(
            slide, ix, iy, iw, ih, categories, series, colors,
            ticks, vmax_display, totals, percent, text_c, muted,
            font_body, font_mono, cat_pt, tick_pt, val_pt,
            radius_px, value_suffix, show_values,
        )


# ── legend ─────────────────────────────────────────────────────────────
def _draw_legend(slide, x, y, w, h, series, colors, font_body, pt,
                 text_c, muted, radius_px):
    swatch = int(Pt(pt).emu * 0.9)
    gap = int(Pt(pt).emu * 0.4)
    item_gap = int(Pt(pt).emu * 1.2)

    widths = []
    for s in series:
        name = s.get("name") or ""
        est_text_w = int(Pt(pt).emu * 0.55 * max(len(name), 1))
        widths.append(swatch + gap + est_text_w)
    total = sum(widths) + item_gap * (len(widths) - 1)
    cursor = x + max(0, w - total)
    cy = y + (h - swatch) // 2
    ty = y
    for (s, c), wd in zip(zip(series, colors), widths):
        sw = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(int(cursor)), Emu(int(cy)),
            Emu(swatch), Emu(swatch),
        )
        _set_fill(sw, c)
        _no_line(sw)
        tx = cursor + swatch + gap
        tw = wd - swatch - gap
        _add_text(slide, tx, ty, tw, h, s.get("name") or "",
                  font_name=font_body, size_pt=pt, hex_color=text_c,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        cursor += wd + item_gap


# ── vertical ───────────────────────────────────────────────────────────
def _draw_vertical(slide, x, y, w, h, categories, series, colors,
                   ticks, vmax, totals, percent, text_c, muted,
                   font_body, font_mono, cat_pt, tick_pt, val_pt,
                   radius_px, suffix, show_values):
    tick_labels = [_fmt(t) for t in ticks]
    max_tick_chars = max(len(s) for s in tick_labels) if tick_labels else 1
    left_margin = int(Pt(tick_pt).emu * 0.65 * (max_tick_chars + 1))
    bottom_margin = int(Pt(cat_pt).emu * 2.2)
    top_margin = int(Pt(val_pt).emu * 1.6)
    right_margin = int(Pt(val_pt).emu * 1.2)

    plot_x = x + left_margin
    plot_y = y + top_margin
    plot_w = max(1, w - left_margin - right_margin)
    plot_h = max(1, h - top_margin - bottom_margin)

    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    # Grid / tick labels
    for i, t in enumerate(ticks):
        ty = plot_y + plot_h - int(plot_h * (t / vmax)) if vmax > 0 else plot_y + plot_h
        lbl_w = left_margin - int(Pt(tick_pt).emu * 0.3)
        lbl_h = int(Pt(tick_pt).emu * 1.4)
        label = _fmt(t) + ("%" if percent and t != 0 and i == len(ticks) - 1 else
                           (suffix if t != 0 and i == len(ticks) - 1 else ""))
        _add_text(slide, x, ty - lbl_h // 2, lbl_w, lbl_h, label,
                  font_name=font_mono, size_pt=tick_pt, hex_color=text_c,
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        if i > 0:
            _add_line(slide, plot_x, ty, plot_x + plot_w, ty, muted, hairline)

    # Baseline
    base_y = plot_y + plot_h
    _add_line(slide, plot_x, base_y, plot_x + plot_w, base_y, muted, hairline * 2)

    n_cat = len(categories)
    group_w = plot_w / n_cat
    inner_pad = group_w * 0.18
    bar_w = group_w - 2 * inner_pad

    for ci, cat in enumerate(categories):
        group_left = plot_x + ci * group_w
        # Category label
        cat_y = base_y + int(Pt(cat_pt).emu * 0.4)
        cat_h = int(Pt(cat_pt).emu * 1.6)
        _add_text(slide, group_left, cat_y, group_w, cat_h, str(cat),
                  font_name=font_body, size_pt=cat_pt, hex_color=text_c,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        total = totals[ci] if ci < len(totals) else 0
        cursor_y = base_y  # bottom of stack

        for si, (s, color) in enumerate(zip(series, colors)):
            vals = s.get("values") or []
            v = vals[ci] if ci < len(vals) else 0
            if percent and total > 0:
                seg_ratio = v / total
                seg_h = int(plot_h * seg_ratio)
            else:
                seg_h = int(plot_h * (max(v, 0) / vmax)) if vmax > 0 else 0

            if seg_h < 1:
                continue
            bx = group_left + inner_pad
            by = cursor_y - seg_h
            # Only use radius on the topmost visible segment
            is_top = (si == len(series) - 1) or all(
                (s2.get("values") or [ci] if ci < len(s2.get("values") or []) else [0])[0 if ci >= len(s2.get("values") or []) else ci] == 0
                for s2 in series[si + 1:]
            )
            r = radius_px if is_top else 0
            _add_bar(slide, bx, by, bar_w, seg_h, color, r)

            # Segment value label (inside segment if tall enough)
            if show_values and seg_h > int(Pt(val_pt).emu * 1.5):
                if percent and total > 0:
                    vt = f"{v / total * 100:.0f}%"
                else:
                    vt = _fmt(v, suffix)
                vh = int(Pt(val_pt).emu * 1.3)
                _add_text(slide, bx, by + (seg_h - vh) // 2, bar_w, vh, vt,
                          font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            cursor_y -= seg_h

        # Total label above the full bar
        if show_values and not percent:
            total_label = _fmt(total, suffix)
            vh = int(Pt(val_pt).emu * 1.3)
            _add_text(
                slide, group_left + inner_pad,
                cursor_y - vh - int(Pt(val_pt).emu * 0.1),
                bar_w, vh, total_label,
                font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
            )


# ── horizontal ─────────────────────────────────────────────────────────
def _draw_horizontal(slide, x, y, w, h, categories, series, colors,
                     ticks, vmax, totals, percent, text_c, muted,
                     font_body, font_mono, cat_pt, tick_pt, val_pt,
                     radius_px, suffix, show_values):
    max_cat_chars = max((len(str(c)) for c in categories), default=4)
    left_margin = int(Pt(cat_pt).emu * 0.55 * (max_cat_chars + 2))
    bottom_margin = int(Pt(tick_pt).emu * 2.0)
    top_margin = int(Pt(val_pt).emu * 0.6)
    right_margin = int(Pt(val_pt).emu * 3.0)

    plot_x = x + left_margin
    plot_y = y + top_margin
    plot_w = max(1, w - left_margin - right_margin)
    plot_h = max(1, h - top_margin - bottom_margin)

    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    # X-axis ticks
    for i, t in enumerate(ticks):
        tx = plot_x + int(plot_w * (t / vmax)) if vmax > 0 else plot_x
        tlh = int(Pt(tick_pt).emu * 1.4)
        label = _fmt(t) + ("%" if percent and t != 0 and i == len(ticks) - 1 else
                           (suffix if t != 0 and i == len(ticks) - 1 else ""))
        _add_text(
            slide, tx - int(Pt(tick_pt).emu * 2),
            plot_y + plot_h + int(Pt(tick_pt).emu * 0.3),
            int(Pt(tick_pt).emu * 4), tlh, label,
            font_name=font_mono, size_pt=tick_pt, hex_color=text_c,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )
        if i > 0:
            _add_line(slide, tx, plot_y, tx, plot_y + plot_h, muted, hairline)

    # Baseline (left axis)
    _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, muted, hairline * 2)

    n_cat = len(categories)
    row_h = plot_h / n_cat
    inner_pad = row_h * 0.18
    bar_h = row_h - 2 * inner_pad

    for ci, cat in enumerate(categories):
        row_top = plot_y + ci * row_h
        cat_h_em = int(Pt(cat_pt).emu * 1.4)
        _add_text(
            slide, x, row_top + (row_h - cat_h_em) / 2,
            left_margin - int(Pt(cat_pt).emu * 0.4), cat_h_em, str(cat),
            font_name=font_body, size_pt=cat_pt, hex_color=text_c,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )

        total = totals[ci] if ci < len(totals) else 0
        cursor_x = plot_x

        for si, (s, color) in enumerate(zip(series, colors)):
            vals = s.get("values") or []
            v = vals[ci] if ci < len(vals) else 0
            if percent and total > 0:
                seg_ratio = v / total
                seg_w = int(plot_w * seg_ratio)
            else:
                seg_w = int(plot_w * (max(v, 0) / vmax)) if vmax > 0 else 0

            if seg_w < 1:
                continue
            by = row_top + inner_pad
            # Only use radius on the rightmost visible segment
            is_right = (si == len(series) - 1) or all(
                (s2.get("values") or [])[ci] == 0 if ci < len(s2.get("values") or []) else True
                for s2 in series[si + 1:]
            )
            r = radius_px if is_right else 0
            _add_bar(slide, cursor_x, by, seg_w, bar_h, color, r)

            # Segment label inside segment
            if show_values and seg_w > int(Pt(val_pt).emu * 3.0):
                if percent and total > 0:
                    vt = f"{v / total * 100:.0f}%"
                else:
                    vt = _fmt(v, suffix)
                vh = int(Pt(val_pt).emu * 1.3)
                _add_text(
                    slide, cursor_x + (seg_w - int(Pt(val_pt).emu * 3)) // 2,
                    by + (bar_h - vh) // 2,
                    int(Pt(val_pt).emu * 3), vh, vt,
                    font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                )

            cursor_x += seg_w

        # Total label after bar
        if show_values and not percent:
            total_label = _fmt(total, suffix)
            vh = int(Pt(val_pt).emu * 1.3)
            vw_em = int(Pt(val_pt).emu * 4)
            _add_text(
                slide, cursor_x + int(Pt(val_pt).emu * 0.3),
                row_top + inner_pad + (bar_h - vh) / 2,
                vw_em, vh, total_label,
                font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
            )
