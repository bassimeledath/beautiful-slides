"""Render the stacked-bar chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_VERTICAL = {
    "orientation": "vertical",
    "title": "Revenue by product, per region",
    "categories": ["AMER", "EMEA", "APAC", "LATAM"],
    "series": [
        {"name": "Product A", "values": [8.2, 5.1, 3.4, 1.2]},
        {"name": "Product B", "values": [4.1, 3.9, 2.8, 0.9]},
        {"name": "Product C", "values": [2.0, 1.5, 1.1, 0.6]},
    ],
    "value_suffix": "M",
    "show_values": True,
}

DATA_HORIZONTAL = {
    "orientation": "horizontal",
    "title": "Headcount by function",
    "categories": ["Engineering", "Sales", "Marketing", "Support", "G&A"],
    "series": [
        {"name": "Full-time", "values": [120, 85, 40, 55, 30]},
        {"name": "Contract",  "values": [30, 15, 10, 20, 5]},
    ],
    "show_values": True,
}

DATA_PERCENT = {
    "orientation": "vertical",
    "title": "Channel mix by quarter (% of total)",
    "categories": ["Q1", "Q2", "Q3", "Q4"],
    "series": [
        {"name": "Direct",   "values": [45, 42, 38, 35]},
        {"name": "Partner",  "values": [30, 33, 37, 40]},
        {"name": "Self-serve", "values": [25, 25, 25, 25]},
    ],
    "show_values": True,
    "percent": True,
}


MODE_DATA = {
    "sv-keynote": DATA_VERTICAL,
    "editorial-magazine": DATA_HORIZONTAL,
    "playful-marketing": DATA_PERCENT,
    "consulting-boardroom": DATA_VERTICAL,
    "craft-minimal": DATA_HORIZONTAL,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
