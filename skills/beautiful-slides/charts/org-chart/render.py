"""Org Chart — top-down tree hierarchy of people/teams connected by lines."""

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _luminance(hex_):
    """Return relative luminance (0-1) of a hex color."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255.0, int(h[2:4], 16) / 255.0, int(h[4:6], 16) / 255.0
    r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
    g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
    b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_text(fill_hex, light="#FFFFFF", dark="#1A1A1A"):
    """Return dark text if fill is light, light text if fill is dark."""
    return dark if _luminance(fill_hex) > 0.35 else light


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "\u2026"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(3)))
    tf.margin_right = Emu(int(Pt(3)))
    tf.margin_top = Emu(int(Pt(2)))
    tf.margin_bottom = Emu(int(Pt(2)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_multiline_textbox(slide, x, y, w, h, lines, font_name, sizes,
                           colors, bolds, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE):
    """Add a textbox with multiple paragraphs (one per line)."""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(4)))
    tf.margin_right = Emu(int(Pt(4)))
    tf.margin_top = Emu(int(Pt(3)))
    tf.margin_bottom = Emu(int(Pt(3)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line_text
        run.font.name = font_name
        run.font.size = Pt(sizes[i] if i < len(sizes) else sizes[-1])
        run.font.bold = bolds[i] if i < len(bolds) else bolds[-1]
        run.font.color.rgb = _rgb(colors[i] if i < len(colors) else colors[-1])
    return tb


def _draw_connector(slide, x1, y1, x2, y2, color_hex, thickness_pt):
    """Draw an L-shaped connector: vertical from (x1,y1) down to y2's level, then horizontal."""
    mid_y = (y1 + y2) // 2
    line_w = max(int(Pt(thickness_pt)), 1)

    # Vertical segment from parent bottom to midpoint
    v1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(x1 - line_w // 2), Emu(y1),
        Emu(line_w), Emu(mid_y - y1),
    )
    v1.fill.solid()
    v1.fill.fore_color.rgb = _rgb(color_hex)
    v1.line.fill.background()

    # Horizontal segment at midpoint
    left_x = min(x1, x2)
    right_x = max(x1, x2)
    if right_x - left_x > 0:
        hz = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(left_x - line_w // 2), Emu(mid_y - line_w // 2),
            Emu(right_x - left_x + line_w), Emu(line_w),
        )
        hz.fill.solid()
        hz.fill.fore_color.rgb = _rgb(color_hex)
        hz.line.fill.background()

    # Vertical segment from midpoint down to child top
    v2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(x2 - line_w // 2), Emu(mid_y),
        Emu(line_w), Emu(y2 - mid_y),
    )
    v2.fill.solid()
    v2.fill.fore_color.rgb = _rgb(color_hex)
    v2.line.fill.background()


def _count_leaves(node):
    """Count leaf nodes (nodes with no children) in the subtree."""
    children = node.get("children", [])
    if not children:
        return 1
    return sum(_count_leaves(c) for c in children)


def _tree_depth(node):
    """Return the depth of the tree (1 for a single node)."""
    children = node.get("children", [])
    if not children:
        return 1
    return 1 + max(_tree_depth(c) for c in children)


def _layout_tree(node, x_start, y, node_w, node_h, level_gap, sibling_gap):
    """Recursively compute (cx, cy, w, h) for each node. Returns list of positioned nodes."""
    result = []
    children = node.get("children", [])

    if not children:
        cx = x_start + node_w // 2
        cy = y
        result.append({"node": node, "cx": cx, "cy": cy, "children_pos": []})
        return result

    # Layout children first
    child_positions = []
    child_x = x_start
    for child in children:
        leaves = _count_leaves(child)
        child_span = max(leaves * (node_w + sibling_gap) - sibling_gap, node_w)
        child_result = _layout_tree(
            child, child_x, y + node_h + level_gap,
            node_w, node_h, level_gap, sibling_gap,
        )
        child_positions.append(child_result)
        child_x += child_span + sibling_gap

    # Parent centered over its children
    first_child_cx = child_positions[0][0]["cx"]
    last_child_cx = child_positions[-1][0]["cx"]
    cx = (first_child_cx + last_child_cx) // 2
    cy = y

    children_top_level = [cp[0] for cp in child_positions]
    result.append({"node": node, "cx": cx, "cy": cy, "children_pos": children_top_level})

    for cp in child_positions:
        result.extend(cp)

    return result


def render(slide, data, tokens, bounds):
    """Render an org chart (tree hierarchy).

    data:
        title    - optional string
        root     - dict with:
            name     - str (person/team name)
            role     - optional str (title/role)
            children - optional list of same shape (recursive)
        Supports 5-20 nodes total.
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]
    radius_px = tokens.get("radius_px", 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    root = data.get("root")
    if not root:
        return
    title = data.get("title")

    # --- layout ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_h = (y + h) - cur_y
    avail_w = w

    # Compute tree dimensions
    depth = _tree_depth(root)
    leaves = _count_leaves(root)

    # Node sizing — adapt to available space
    max_node_w = int(avail_w * 0.18)
    min_node_w = int(avail_w * 0.08)
    # Scale node width based on number of leaves
    node_w = max(min_node_w, min(max_node_w, int(avail_w / (leaves * 1.3))))

    node_h = int(node_w * 0.65)
    level_gap = max(int((avail_h - depth * node_h) / depth), int(Pt(base_pt)))
    sibling_gap = max(int(node_w * 0.25), int(Pt(base_pt * 0.5)))

    # Count total nodes for adaptive font scaling
    total_nodes = len(_layout_tree(root, 0, 0, node_w, node_h, level_gap, sibling_gap))
    # For large trees, reduce font size more aggressively instead of truncating
    if total_nodes > 15:
        font_scale = max(0.65, 10 / total_nodes)
    elif total_nodes > 12:
        font_scale = 0.80
    else:
        font_scale = 1.0
    # Estimate max chars that fit in a node — at least 12 for names, 15 for roles
    max_name_chars = max(12, int(node_w / (Pt(max(int(base_pt * 0.75 * font_scale), 7)) * 0.55)))
    max_role_chars = max(15, int(node_w / (Pt(max(int(base_pt * 0.55 * font_scale), 6)) * 0.48)))

    # Layout the tree
    positioned = _layout_tree(root, x, cur_y, node_w, node_h, level_gap, sibling_gap)

    # Calculate actual tree width and re-center
    all_cx = [p["cx"] for p in positioned]
    tree_left = min(all_cx) - node_w // 2
    tree_right = max(all_cx) + node_w // 2
    tree_width = tree_right - tree_left
    offset_x = x + (avail_w - tree_width) // 2 - tree_left

    # Apply offset to center tree
    for p in positioned:
        p["cx"] += offset_x

    # Recalculate after centering — clamp to bounds
    for p in positioned:
        p["cx"] = max(x + node_w // 2, min(p["cx"], x + w - node_w // 2))

    # Also clamp vertical
    for p in positioned:
        p["cy"] = max(cur_y, min(p["cy"], y + h - node_h))

    # Draw connectors first (behind nodes)
    connector_color = _lerp_hex(muted, primary, 0.3)
    connector_thickness = max(base_pt * 0.1, 1.0)
    for p in positioned:
        parent_cx = p["cx"]
        parent_bottom = p["cy"] + node_h
        for child_p in p["children_pos"]:
            child_cx = child_p["cx"]
            child_top = child_p["cy"]
            _draw_connector(
                slide, parent_cx, parent_bottom,
                child_cx, child_top,
                connector_color, connector_thickness,
            )

    # Draw nodes
    for i, p in enumerate(positioned):
        node = p["node"]
        cx = p["cx"]
        cy = p["cy"]
        node_x = cx - node_w // 2
        node_y = cy

        # Depth-based color: root gets primary, deeper nodes blend toward muted
        depth_level = (cy - cur_y) / max(avail_h - node_h, 1)
        node_color = _lerp_hex(primary, accent, min(depth_level * 1.2, 1.0))

        # Draw node rectangle
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE
        node_shape = slide.shapes.add_shape(
            shape_type,
            Emu(node_x), Emu(node_y),
            Emu(node_w), Emu(node_h),
        )
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = _rgb(node_color)
        node_shape.line.fill.background()

        # Text content -- truncate long names/roles
        name = _truncate(node.get("name", ""), max_name_chars)
        role = _truncate(node.get("role", ""), max_role_chars)

        # Luminance-based text color: dark text on light fills, light on dark
        node_text_color = _contrast_text(node_color)
        # Role text slightly muted relative to the name
        node_role_color = _lerp_hex(node_text_color, node_color, 0.25)

        if role:
            name_size = max(int(base_pt * 0.75 * font_scale), 7)
            role_size = max(int(base_pt * 0.55 * font_scale), 6)
            # Ensure role text stays readable (minimum 6pt)
            role_size = max(role_size, 6)
            _add_multiline_textbox(
                slide, node_x, node_y, node_w, node_h,
                [name, role],
                font_body,
                [name_size, role_size],
                [node_text_color, node_role_color],
                [True, False],
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
        else:
            name_size = max(int(base_pt * 0.7 * font_scale), 7)
            _add_textbox(
                slide, node_x, node_y, node_w, node_h,
                name, font_body, name_size, node_text_color,
                align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )

    return None
