import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


DATA = {
    "title": "Engineering Organization",
    "root": {
        "name": "Alice Chen",
        "role": "CEO",
        "children": [
            {
                "name": "Bob Park",
                "role": "VP Engineering",
                "children": [
                    {"name": "Carol Wu", "role": "Backend Lead"},
                    {"name": "Dan Lee", "role": "Frontend Lead"},
                    {"name": "Grace Kim", "role": "Infra Lead"},
                ],
            },
            {
                "name": "Eve Rao",
                "role": "VP Product",
                "children": [
                    {"name": "Frank Li", "role": "Sr. PM"},
                    {"name": "Hana Sato", "role": "PM"},
                ],
            },
            {
                "name": "Ivan Petrov",
                "role": "VP Design",
                "children": [
                    {"name": "Jade Obi", "role": "UX Lead"},
                ],
            },
        ],
    },
}


DATA_MAX = {
    "title": "Global Enterprise Organization",
    "root": {
        "name": "Dr. Alexandra Worthington-Smith",
        "role": "Chief Executive Officer",
        "children": [
            {
                "name": "Michael Christopher Johnson",
                "role": "VP Engineering & Technology",
                "children": [
                    {"name": "Sarah Elizabeth Chen", "role": "Director of Backend"},
                    {"name": "James Robert Wilson", "role": "Director of Frontend"},
                    {"name": "Patricia Anne Lee", "role": "Director of Infrastructure"},
                    {"name": "William Thomas Davis", "role": "Director of Security"},
                ],
            },
            {
                "name": "Jennifer Marie Anderson",
                "role": "VP Product Management",
                "children": [
                    {"name": "Christopher Paul Brown", "role": "Senior Product Manager"},
                    {"name": "Jessica Lynn Martinez", "role": "Product Manager"},
                    {"name": "Daniel James Taylor", "role": "Product Analyst"},
                ],
            },
            {
                "name": "Robert Andrew Thompson",
                "role": "VP Design & UX Research",
                "children": [
                    {"name": "Amanda Nicole Garcia", "role": "UX Design Lead"},
                    {"name": "Matthew Ryan Hernandez", "role": "Visual Design Lead"},
                ],
            },
            {
                "name": "Katharine Elizabeth Moore",
                "role": "Chief Financial Officer",
                "children": [
                    {"name": "Steven Patrick Clark", "role": "Finance Director"},
                ],
            },
        ],
    },
}


def _set_slide_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))


def build(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA, tokens, (x, y, w, h))

    prs.save(out_path)


def build_max(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA_MAX, tokens, (x, y, w, h))

    prs.save(out_path)


def main():
    for name, tokens in MODES.items():
        path = os.path.join(HERE, f"example-{name}.pptx")
        build(name, tokens, path)
        print(f"wrote {path}")
        path_max = os.path.join(HERE, f"example-max-{name}.pptx")
        build_max(name, tokens, path_max)
        print(f"wrote {path_max}")


if __name__ == "__main__":
    main()
