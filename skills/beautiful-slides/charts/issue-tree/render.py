"""Issue tree — MECE branching decomposition (consulting-style).

Root question on the left, branches fan out to the right with sub-branches.
Each level decomposes the parent into mutually exclusive, collectively
exhaustive components.
"""

import math

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "\u2026"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.LEFT, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(4)))
    tf.margin_right = Emu(int(Pt(4)))
    tf.margin_top = Emu(int(Pt(2)))
    tf.margin_bottom = Emu(int(Pt(2)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _count_leaves(node):
    """Count leaf nodes in the subtree."""
    children = node.get("children", [])
    if not children:
        return 1
    return sum(_count_leaves(c) for c in children)


def _max_depth(node, d=0):
    """Find the maximum depth of the tree."""
    children = node.get("children", [])
    if not children:
        return d
    return max(_max_depth(c, d + 1) for c in children)


def _draw_horizontal_bracket(slide, x1, y_top, y_bottom, color_hex, thickness):
    """Draw a vertical line with short horizontal ticks at top and bottom.

    This forms the right-side bracket that visually groups children.
    """
    tick_w = int(thickness * 4)
    bar_w = max(int(thickness), int(Pt(1.5)))

    # Vertical bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(x1), Emu(y_top),
        Emu(bar_w), Emu(y_bottom - y_top),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(color_hex)
    bar.line.fill.background()

    return bar


def render(slide, data, tokens, bounds):
    """Render an issue tree (MECE decomposition).

    data:
        title    - optional string
        root     - nested tree: {"label": str, "children": [same structure]}
                   Typically 2-4 levels deep, 2-4 children per node.
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    root = data.get("root")
    if not root:
        return
    title = data.get("title")

    # --- layout title ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_w = w
    avail_h = (y + h) - cur_y
    area_x = x
    area_y = cur_y

    md = _max_depth(root)
    num_levels = md + 1
    total_leaves = _count_leaves(root)

    # Column widths: leaf columns get more space for long descriptive text
    # Reserve some space for connector lines between columns
    connector_gap = int(avail_w * 0.03)
    usable_w = avail_w - connector_gap * num_levels

    # Reverse bias: give leaf columns more weight since they have longest text
    col_weights = []
    for d in range(num_levels):
        if d == num_levels - 1:
            col_weights.append(1.5)  # leaf column gets most space
        elif d == 0:
            col_weights.append(1.2)  # root column gets moderate space
        else:
            col_weights.append(0.9)  # interior columns get less
    total_weight = sum(col_weights)
    col_widths = [int(usable_w * cw / total_weight) for cw in col_weights]

    # Column x-positions
    col_x = []
    cx = area_x
    for d in range(num_levels):
        col_x.append(cx)
        cx += col_widths[d] + connector_gap

    # Node height based on available vertical space and leaf count
    row_h = int(avail_h / total_leaves) if total_leaves > 0 else avail_h
    node_h = max(int(row_h * 0.65), int(Pt(base_pt) * 1.5))
    node_h = min(node_h, int(avail_h * 0.08))

    # Adaptive font scaling based on total leaf count
    font_scale = min(1.0, 8 / max(total_leaves, 1))
    adaptive_pt = max(int(base_pt * 0.65 * font_scale + base_pt * 0.35), 7)
    leaf_font = max(int(adaptive_pt * 0.85), 7)

    # Estimate max chars per column for truncation
    max_chars_per_col = []
    for d in range(num_levels):
        est_chars = max(15, int(col_widths[d] / (Pt(adaptive_pt) * 0.55)))
        max_chars_per_col.append(est_chars)

    # Depth-based colors
    depth_colors = []
    for d in range(num_levels):
        t = d / max(num_levels - 1, 1)
        depth_colors.append(_lerp_hex(primary, accent, t))

    # Connector line color and thickness
    line_color = _lerp_hex(muted, primary, 0.25)
    line_thick = max(int(Pt(1.5)), int(Pt(base_pt * 0.1)))

    # --- recursive layout and draw ---
    def draw_node(node, depth, y_start, y_avail, parent_right_x=None, parent_cy=None):
        """Draw a node and its children recursively.

        Returns the centre-y of this node for connector purposes.
        """
        children = node.get("children", [])
        label = node.get("label", "")

        col_idx = min(depth, num_levels - 1)
        nx = col_x[col_idx]
        nw = col_widths[col_idx]

        if not children:
            # Leaf node: centre in allocated vertical space
            cy = y_start + y_avail // 2
            ny = cy - node_h // 2
            # Clamp
            ny = max(area_y, min(ny, area_y + avail_h - node_h))

            # Draw leaf box
            fill = _lerp_hex(bg, depth_colors[col_idx], 0.12)
            border = depth_colors[col_idx]
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(nx), Emu(ny), Emu(nw), Emu(node_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.color.rgb = _rgb(border)
            shape.line.width = Pt(1)

            trunc_label = _truncate(label, max_chars_per_col[col_idx])
            _add_textbox(
                slide, nx, ny, nw, node_h,
                trunc_label, font_body, leaf_font, text_c,
                align=PP_ALIGN.LEFT, bold=False, anchor=MSO_ANCHOR.MIDDLE,
            )

            # Connector from parent
            if parent_right_x is not None:
                _draw_connector_line(slide, parent_right_x, parent_cy, nx, cy,
                                     line_color, line_thick)

            return cy

        # Interior node: position based on children
        my_leaves = _count_leaves(node)
        child_y = y_start
        child_centres = []

        for child in children:
            child_leaves = _count_leaves(child)
            child_avail = int(y_avail * child_leaves / my_leaves)
            child_avail = max(child_avail, node_h + int(Pt(base_pt * 0.3)))

            child_cy = draw_node(
                child, depth + 1, child_y, child_avail,
                parent_right_x=nx + nw,
                parent_cy=None,  # will be set after we know this node's cy
            )
            child_centres.append(child_cy)
            child_y += child_avail

        # This node's vertical centre = average of children
        cy = int(sum(child_centres) / len(child_centres))
        ny = cy - node_h // 2
        # Clamp
        ny = max(area_y, min(ny, area_y + avail_h - node_h))
        cy = ny + node_h // 2

        # Draw this node
        trunc_label = _truncate(label, max_chars_per_col[col_idx])
        if depth == 0:
            # Root node: prominent fill
            fill = primary
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(nx), Emu(ny), Emu(nw), Emu(node_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.fill.background()
            _add_textbox(
                slide, nx, ny, nw, node_h,
                trunc_label, font_body, adaptive_pt, bg,
                align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )
        else:
            # Interior node: bordered
            fill = _lerp_hex(bg, depth_colors[col_idx], 0.06)
            border = depth_colors[col_idx]
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(nx), Emu(ny), Emu(nw), Emu(node_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.color.rgb = _rgb(border)
            shape.line.width = Pt(1.5)
            _add_textbox(
                slide, nx, ny, nw, node_h,
                trunc_label, font_body, adaptive_pt, text_c,
                align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )

        # Draw connector from parent to this node
        if parent_right_x is not None:
            _draw_connector_line(slide, parent_right_x, parent_cy, nx, cy,
                                 line_color, line_thick)

        # Draw bracket + horizontal connectors from this node to children
        right_x = nx + nw
        bracket_x = right_x + connector_gap // 3

        if len(child_centres) > 1:
            # Vertical bracket line spanning children
            top_child_cy = min(child_centres)
            bot_child_cy = max(child_centres)
            _draw_horizontal_bracket(
                slide, bracket_x, top_child_cy, bot_child_cy,
                line_color, line_thick,
            )

        # Now draw horizontal lines from bracket to each child (retro-fix parent_cy)
        # We drew children with parent_cy=None; draw connectors now
        next_col_x = col_x[min(depth + 1, num_levels - 1)] if depth + 1 < num_levels else right_x + connector_gap

        for child_cy_val in child_centres:
            # Horizontal line from this node's right edge to bracket
            bar_w = max(int(line_thick), int(Pt(1.5)))
            line_seg_w = bracket_x - right_x
            if line_seg_w > 0:
                seg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Emu(right_x), Emu(child_cy_val - bar_w // 2),
                    Emu(line_seg_w), Emu(bar_w),
                )
                seg.fill.solid()
                seg.fill.fore_color.rgb = _rgb(line_color)
                seg.line.fill.background()

            # Horizontal line from bracket to child node
            child_col_x = next_col_x
            line_seg_w2 = child_col_x - bracket_x
            if line_seg_w2 > 0:
                seg2 = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Emu(bracket_x), Emu(child_cy_val - bar_w // 2),
                    Emu(line_seg_w2), Emu(bar_w),
                )
                seg2.fill.solid()
                seg2.fill.fore_color.rgb = _rgb(line_color)
                seg2.line.fill.background()

        return cy

    def _draw_connector_line(slide, x1, y1, x2, y2, color_hex, thickness):
        """Draw a simple horizontal connector line."""
        # For issue trees we use the bracket pattern instead; this is a fallback
        # for parent->this node connectors that aren't handled by the bracket logic.
        pass

    draw_node(root, 0, area_y, avail_h)

    return None
