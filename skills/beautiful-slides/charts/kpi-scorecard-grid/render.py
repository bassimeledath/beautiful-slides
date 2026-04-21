"""KPI scorecard grid — 4-8 KPI tiles arranged in a 2xN, 1x4, or similar grid.

Each tile shows label, hero value, delta, optional sparkline placeholder.
Delegates individual tile rendering to the kpi chart module.

Public API: render(slide, data, tokens, bounds)
"""

import os
import sys

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# Allow importing the kpi tile renderer as a sibling package.
_HERE = os.path.dirname(os.path.abspath(__file__))
_CHARTS = os.path.dirname(_HERE)
if _CHARTS not in sys.path:
    sys.path.insert(0, _CHARTS)

from kpi.render import render as _render_tile  # noqa: E402


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, *, font, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = _rgb(color_hex)
    return tb


def _choose_layout(n):
    """Pick (rows, cols) for n tiles (4-8).

    Preferences:
      4 -> 2x2
      5 -> 2x3 (one empty cell)
      6 -> 2x3
      7 -> 2x4 (one empty cell)
      8 -> 2x4
      3 -> 1x3
      2 -> 1x2
      1 -> 1x1
    """
    if n <= 1:
        return (1, 1)
    if n == 2:
        return (1, 2)
    if n == 3:
        return (1, 3)
    if n == 4:
        return (2, 2)
    if n <= 6:
        return (2, 3)
    return (2, 4)


def render(slide, data, tokens, bounds):
    """Render a grid of KPI scorecard tiles.

    Parameters
    ----------
    slide : pptx.slide.Slide
    data : dict
        {
            "title": "Optional grid title",
            "tiles": [
                {"label": "ARR", "value": "$47.2M", "delta": "+12%",
                 "delta_direction": "up", "footnote": "..."},
                ...
            ],
            "layout": "auto" | "1x4" | "2x2" | "2x3" | "2x4",  # optional
        }
    tokens : dict
    bounds : tuple (x_emu, y_emu, w_emu, h_emu)
    """
    x, y, w, h = bounds

    tiles = list(data.get("tiles") or [])
    if not tiles:
        return None

    title = data.get("title")
    layout_hint = (data.get("layout") or "auto").lower()

    bg = tokens["bg"]
    text_hex = tokens["text"]
    font_display = tokens["font_display"]
    base_pt = int(tokens["font_size_base_pt"])

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()

    # Title
    title_h = 0
    if title:
        title_pt = max(base_pt + 2, int(round(base_pt * 1.25)))
        # Estimate wrapped lines for dynamic title height
        char_w_emu = Pt(title_pt).emu * 0.55
        chars_per_line = max(1, int(w / char_w_emu))
        title_lines = max(1, -(-len(title) // chars_per_line))  # ceil division
        title_h = int(Pt(title_pt).emu * 1.8 * title_lines)
        _add_textbox(slide, x, y, w, title_h,
                     title,
                     font=font_display, size_pt=title_pt,
                     color_hex=text_hex, bold=True, align=PP_ALIGN.LEFT)

    grid_x = x
    grid_y = y + title_h
    grid_w = w
    grid_h = h - title_h

    n = len(tiles)

    # Determine layout
    if layout_hint == "auto":
        rows, cols = _choose_layout(n)
    else:
        # Parse "RxC" format
        parts = layout_hint.split("x")
        if len(parts) == 2:
            try:
                rows, cols = int(parts[0]), int(parts[1])
            except ValueError:
                rows, cols = _choose_layout(n)
        else:
            rows, cols = _choose_layout(n)

    # Gutter between tiles
    gutter_x = int(grid_w * 0.025)
    gutter_y = int(grid_h * 0.04)

    tile_w = (grid_w - gutter_x * (cols - 1)) // cols
    tile_h = (grid_h - gutter_y * (rows - 1)) // rows

    idx = 0
    for r in range(rows):
        for c in range(cols):
            if idx >= n:
                break
            tx = grid_x + c * (tile_w + gutter_x)
            ty = grid_y + r * (tile_h + gutter_y)
            _render_tile(slide, tiles[idx], tokens, (tx, ty, tile_w, tile_h))
            idx += 1

    return None
