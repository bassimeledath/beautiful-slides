"""Render KPI scorecard grids in all 5 modes. Generates one .pptx per mode."""

import os
import sys

from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402


SAMPLE_DATA = {
    "title": "Q1 2025 Executive Scorecard",
    "layout": "auto",
    "tiles": [
        {
            "label": "ARR",
            "value": "$47.2M",
            "delta": "+12.4% vs plan",
            "delta_direction": "up",
            "footnote": "Source: NetSuite, Apr 10",
        },
        {
            "label": "Gross Margin",
            "value": "72.8%",
            "delta": "-1.1 pts QoQ",
            "delta_direction": "down",
            "footnote": "Non-GAAP, ex. one-time",
        },
        {
            "label": "NPS",
            "value": "64",
            "delta": "+6 vs Q4",
            "delta_direction": "up",
            "footnote": "n=1,842, mixed segment",
        },
        {
            "label": "Burn Multiple",
            "value": "1.2x",
            "delta": "-0.3x QoQ",
            "delta_direction": "down",
            "footnote": "Net burn / net new ARR",
        },
        {
            "label": "Headcount",
            "value": "284",
            "delta": "+18 hires",
            "delta_direction": "up",
            "footnote": "Eng: 142, GTM: 89",
        },
        {
            "label": "Pipeline",
            "value": "$12.8M",
            "delta": "+34% QoQ",
            "delta_direction": "up",
            "footnote": "Stage 2+ weighted",
        },
    ],
}


def _paint_slide_bg(slide, hex_):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        slide.part.package.presentation_part.presentation.slide_width,
        slide.part.package.presentation_part.presentation.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(hex_)
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def build(mode, tokens):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _paint_slide_bg(slide, tokens["bg"])

    sw = prs.slide_width
    sh = prs.slide_height

    margin = Inches(0.6)
    bounds = (
        margin,
        margin,
        sw - 2 * margin,
        sh - 2 * margin,
    )

    render(slide, SAMPLE_DATA, tokens, bounds)

    out = os.path.join(HERE, f"example-{mode}.pptx")
    prs.save(out)
    return out


def main():
    for mode, tokens in MODES.items():
        path = build(mode, tokens)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
