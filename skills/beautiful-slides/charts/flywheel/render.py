"""Flywheel — circular arrangement of steps with directional arrows forming a loop."""

import math

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "\u2026"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(3)))
    tf.margin_right = Emu(int(Pt(3)))
    tf.margin_top = Emu(int(Pt(2)))
    tf.margin_bottom = Emu(int(Pt(2)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_arrow(slide, x1, y1, x2, y2, color_hex, thickness):
    """Draw a directional arrow from (x1,y1) to (x2,y2) as a triangle at end + line."""
    dx = x2 - x1
    dy = y2 - y1
    length = math.sqrt(dx * dx + dy * dy)
    if length < 1:
        return

    # Unit vector
    ux = dx / length
    uy = dy / length
    # Perpendicular
    px = -uy
    py = ux

    # Arrow head (triangle) — 20% of the line length but capped
    head_len = min(length * 0.25, thickness * 5)
    head_w = head_len * 0.7

    # Arrow head tip at (x2, y2)
    # Base of arrow head
    base_x = x2 - ux * head_len
    base_y = y2 - uy * head_len

    # Triangle points
    p1x, p1y = int(x2), int(y2)
    p2x, p2y = int(base_x + px * head_w / 2), int(base_y + py * head_w / 2)
    p3x, p3y = int(base_x - px * head_w / 2), int(base_y - py * head_w / 2)

    ff = slide.shapes.build_freeform(p1x, p1y, scale=1.0)
    ff.add_line_segments([(p2x, p2y), (p3x, p3y)], close=True)
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_hex)
    shape.line.fill.background()

    # Shaft (thin rect from start to base of arrowhead)
    shaft_len = length - head_len
    if shaft_len > 0:
        # Draw shaft as a thin freeform quad
        sw = thickness / 2
        s1x = int(x1 + px * sw)
        s1y = int(y1 + py * sw)
        s2x = int(x1 - px * sw)
        s2y = int(y1 - py * sw)
        s3x = int(base_x - px * sw)
        s3y = int(base_y - py * sw)
        s4x = int(base_x + px * sw)
        s4y = int(base_y + py * sw)

        ff2 = slide.shapes.build_freeform(s1x, s1y, scale=1.0)
        ff2.add_line_segments([(s2x, s2y), (s3x, s3y), (s4x, s4y)], close=True)
        shaft = ff2.convert_to_shape()
        shaft.fill.solid()
        shaft.fill.fore_color.rgb = _rgb(color_hex)
        shaft.line.fill.background()


def render(slide, data, tokens, bounds):
    """Render a flywheel (circular loop of steps with arrows).

    data:
        title    - optional string
        center   - optional string (text in the center of the wheel)
        steps    - list of {"label": str}  (3-6 items)
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]
    radius_px = tokens.get("radius_px", 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    steps = data.get("steps", [])
    if not steps:
        return
    title = data.get("title")
    center_text = data.get("center")
    n = len(steps)

    # --- layout ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_h = (y + h) - cur_y
    avail_w = w

    # The flywheel is inscribed in a circle that fits in avail bounds
    cx = x + avail_w // 2
    cy = cur_y + avail_h // 2
    # Radius of the circle on which step nodes sit
    max_radius = min(avail_w, avail_h) // 2
    # Reserve space for labels outside the node circles
    node_r = int(max_radius * 0.13)
    orbit_r = int(max_radius * 0.55)
    label_pad = int(max_radius * 0.28)

    # Draw center circle (decorative ring)
    center_ring_r = int(orbit_r * 0.38)
    ring_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(cx - center_ring_r), Emu(cy - center_ring_r),
        Emu(center_ring_r * 2), Emu(center_ring_r * 2),
    )
    ring_shape.fill.solid()
    ring_shape.fill.fore_color.rgb = _rgb(_lerp_hex(bg, primary, 0.1))
    ring_shape.line.color.rgb = _rgb(muted)
    ring_shape.line.width = Pt(1)

    # Center text
    if center_text:
        ct_w = int(center_ring_r * 1.8)
        ct_h = int(center_ring_r * 1.4)
        center_pt = max(int(base_pt * 0.75), 7) if len(center_text) > 10 else max(int(base_pt * 0.85), 8)
        _add_textbox(
            slide,
            cx - ct_w // 2, cy - ct_h // 2, ct_w, ct_h,
            _truncate(center_text, 20), font_display, center_pt, text_c,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )

    # Position steps around the circle
    # Start at top (-pi/2) and go clockwise
    positions = []
    for i in range(n):
        angle = -math.pi / 2 + (2 * math.pi * i / n)
        nx = cx + int(orbit_r * math.cos(angle))
        ny = cy + int(orbit_r * math.sin(angle))
        positions.append((nx, ny, angle))

    # Draw arrows between consecutive nodes
    arrow_thickness = max(int(Pt(base_pt * 0.2)), int(node_r * 0.2))
    for i in range(n):
        j = (i + 1) % n
        sx, sy, _ = positions[i]
        ex, ey, _ = positions[j]

        # Shorten arrow: start from edge of source node, end before target node
        dx = ex - sx
        dy = ey - sy
        seg_len = math.sqrt(dx * dx + dy * dy)
        if seg_len < 1:
            continue
        ux = dx / seg_len
        uy = dy / seg_len

        # Offset start and end from node centres
        offset = node_r + int(node_r * 0.4)
        a_sx = sx + ux * offset
        a_sy = sy + uy * offset
        a_ex = ex - ux * offset
        a_ey = ey - uy * offset

        t = i / n
        arrow_color = _lerp_hex(muted, primary, 0.4)
        _draw_arrow(slide, a_sx, a_sy, a_ex, a_ey, arrow_color, arrow_thickness)

    # Draw step nodes and labels
    for i, (nx, ny, angle) in enumerate(positions):
        t = i / (n - 1) if n > 1 else 0.0
        node_color = _lerp_hex(primary, accent, t)

        # Node circle
        node_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(nx - node_r), Emu(ny - node_r),
            Emu(node_r * 2), Emu(node_r * 2),
        )
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = _rgb(node_color)
        node_shape.line.fill.background()

        # Step number inside node
        num_size = max(int(base_pt * 0.65), 7)
        num_w = node_r * 2
        num_h = node_r * 2
        _add_textbox(
            slide,
            nx - num_w // 2, ny - num_h // 2, num_w, num_h,
            str(i + 1), font_body, num_size, bg,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label outside the node -- scale with step count
        label_w_factor = max(0.25, 0.45 - 0.04 * n)
        label_w = int(max_radius * label_w_factor)
        label_h = int(Pt(base_pt) * min(3, max(1.8, 3.5 - 0.3 * n)))

        # Push labels further from center for readability
        label_dist = node_r + int(label_pad * 0.7)
        lx = nx + int(label_dist * math.cos(angle)) - label_w // 2
        ly = ny + int(label_dist * math.sin(angle)) - label_h // 2

        # Clamp label to bounds
        lx = max(x, min(lx, x + w - label_w))
        ly = max(cur_y, min(ly, y + h - label_h))

        # Adaptive font for labels
        label_font = max(int(base_pt * 0.75), 7) if n >= 5 else base_pt
        # Truncate long step descriptions
        max_label_chars = max(20, int(label_w / (Pt(label_font) * 0.55)))
        step_label = _truncate(steps[i].get("label", ""), max_label_chars)

        _add_textbox(
            slide, lx, ly, label_w, label_h,
            step_label, font_body, label_font, text_c,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )

    return None
