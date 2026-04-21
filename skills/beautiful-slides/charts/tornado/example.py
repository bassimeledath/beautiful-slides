"""Render the tornado chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_SENSITIVITY = {
    "title": "Sensitivity analysis: NPV drivers",
    "categories": ["Discount rate", "Revenue growth", "COGS %", "CapEx", "Terminal value"],
    "left":  {"name": "Downside", "values": [38, 25, 22, 18, 15]},
    "right": {"name": "Upside",   "values": [42, 30, 19, 14, 12]},
    "value_suffix": "M",
    "show_values": True,
}

DATA_POPULATION = {
    "title": "Population by age group, 2024",
    "categories": ["0-14", "15-24", "25-34", "35-44", "45-54", "55-64", "65+"],
    "left":  {"name": "Male",   "values": [9.2, 8.4, 10.1, 9.8, 8.6, 7.2, 5.8]},
    "right": {"name": "Female", "values": [8.8, 8.1, 10.4, 10.0, 8.9, 7.6, 6.9]},
    "value_suffix": "%",
    "show_values": True,
}

DATA_COMPARISON = {
    "title": "Feature satisfaction: Product A vs. B",
    "categories": ["Performance", "Reliability", "UX", "Support", "Price"],
    "left":  {"name": "Product A", "values": [82, 75, 90, 68, 72]},
    "right": {"name": "Product B", "values": [78, 88, 65, 81, 70]},
    "value_suffix": "",
    "show_values": True,
}


MODE_DATA = {
    "sv-keynote":           DATA_SENSITIVITY,
    "editorial-magazine":   DATA_POPULATION,
    "playful-marketing":    DATA_COMPARISON,
    "consulting-boardroom": DATA_SENSITIVITY,
    "craft-minimal":        DATA_POPULATION,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
