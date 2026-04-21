"""Tornado (butterfly) chart renderer — native python-pptx shapes only.

Two horizontal bar charts back-to-back sharing a common vertical axis,
extending left and right.  Useful for sensitivity analysis, population
pyramids, and bilateral comparisons.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt
import math


_EMU_PER_PX = 9525


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _set_fill(shape, hex_):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_)


def _no_line(shape):
    shape.line.fill.background()


def _add_text(slide, x, y, w, h, text, font_name, size_pt, hex_color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold = bold
    return tb


def _add_bar(slide, x, y, w, h, fill_hex, radius_px):
    w_i = max(1, int(w))
    h_i = max(1, int(h))
    if radius_px and radius_px > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i),
        )
        short_emu = min(w_i, h_i)
        radius_emu = radius_px * _EMU_PER_PX
        ratio = max(0.0, min(0.5, radius_emu / short_emu / 2.0))
        try:
            shape.adjustments[0] = ratio
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i),
        )
    _set_fill(shape, fill_hex)
    _no_line(shape)
    return shape


def _add_line(slide, x1, y1, x2, y2, hex_color, width_emu):
    ln = slide.shapes.add_connector(
        1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)),
    )
    ln.line.color.rgb = _rgb(hex_color)
    ln.line.width = Emu(int(width_emu))
    return ln


def _fmt(v, suffix=""):
    if abs(v - round(v)) < 1e-9:
        s = f"{int(round(v))}"
    else:
        s = f"{v:.1f}"
    return f"{s}{suffix}" if suffix else s


def render(slide, data, tokens, bounds):
    """Render a tornado chart onto *slide* inside *bounds*.

    ``data`` keys:
      categories : list[str]
          Row labels along the shared center axis.
      left : dict
          ``{"name": str, "values": list[number]}`` -- series extending left.
      right : dict
          ``{"name": str, "values": list[number]}`` -- series extending right.
      title : str or None
      value_suffix : str
      show_values : bool (default True)
    """
    x, y, w, h = bounds

    categories = list(data.get("categories") or [])
    left = data.get("left") or {}
    right = data.get("right") or {}
    left_vals = list(left.get("values") or [])
    right_vals = list(right.get("values") or [])
    title = data.get("title")
    value_suffix = data.get("value_suffix") or ""
    show_values = bool(data.get("show_values", True))

    if not categories or not left_vals or not right_vals:
        return

    n = len(categories)

    # Theme tokens
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    _set_fill(bg_shape, bg)
    _no_line(bg_shape)

    # Outer padding
    pad = int(min(w, h) * 0.04)
    ix = x + pad
    iy = y + pad
    iw = w - 2 * pad
    ih = h - 2 * pad

    # Title
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.6))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_text(
            slide, ix, iy, iw, title_h, title,
            font_name=font_display, size_pt=title_pt, hex_color=text_c,
            bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        iy += title_h + int(Pt(base_pt).emu * 0.6)
        ih = (y + h - pad) - iy

    # Legend at top-right
    legend_pt = max(int(base_pt * 0.85), 8)
    legend_h = int(Pt(legend_pt).emu * 2.0)
    _draw_legend(
        slide, ix, iy, iw, legend_h,
        left.get("name", "Left"), right.get("name", "Right"),
        primary, accent,
        font_body, legend_pt, text_c, muted, radius_px,
    )
    iy += legend_h
    ih -= legend_h

    # Layout
    cat_pt = max(int(base_pt * 0.92), 9)
    val_pt = max(int(base_pt * 0.78), 8)

    # Center column for category labels
    max_cat_chars = max(len(str(c)) for c in categories)
    center_w = int(Pt(cat_pt).emu * 0.55 * min(max_cat_chars + 2, 16))
    center_w = max(center_w, int(iw * 0.12))
    center_w = min(center_w, int(iw * 0.25))

    # Value label space
    val_margin = int(Pt(val_pt).emu * 3.5) if show_values else int(Pt(val_pt).emu * 0.5)

    # Bar regions
    side_w = max(1, (iw - center_w) // 2 - val_margin)
    left_region_x = ix + val_margin
    right_region_x = ix + val_margin + side_w + center_w
    center_x = ix + val_margin + side_w

    # Determine scale (shared for both sides)
    all_vals = left_vals + right_vals
    vmax = max(abs(v) for v in all_vals) if all_vals else 1
    if vmax <= 0:
        vmax = 1

    # Vertical layout
    top_pad = int(Pt(base_pt).emu * 0.3)
    row_gap = int(ih * 0.04)
    total_gaps = row_gap * (n - 1)
    row_h = max(int((ih - total_gaps - top_pad) / n), int(Pt(base_pt * 1.5)))

    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    # Center axis line
    center_line_x = center_x + center_w // 2
    axis_top = iy + top_pad
    axis_bottom = iy + top_pad + n * row_h + (n - 1) * row_gap
    _add_line(slide, center_line_x, axis_top, center_line_x, axis_bottom,
              muted, hairline * 2)

    bar_fraction = 0.65  # fraction of row_h used for the bar
    for ci in range(n):
        row_y = iy + top_pad + ci * (row_h + row_gap)

        # Opaque background behind category label so the center axis
        # line doesn't bleed through the text.
        cat_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(center_x)), Emu(int(row_y)),
            Emu(int(center_w)), Emu(int(row_h)),
        )
        _set_fill(cat_bg, bg)
        _no_line(cat_bg)

        # Category label (centered)
        _add_text(
            slide, center_x, row_y, center_w, row_h,
            str(categories[ci]),
            font_name=font_body, size_pt=cat_pt, hex_color=text_c,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        bar_h = int(row_h * bar_fraction)
        bar_y = row_y + (row_h - bar_h) // 2

        # Left bar (extends from center to left)
        lv = abs(left_vals[ci]) if ci < len(left_vals) else 0
        lw = int(side_w * (lv / vmax)) if vmax > 0 else 0
        if lw > 0:
            lx = center_x - lw
            _add_bar(slide, lx, bar_y, lw, bar_h, primary, radius_px)

        # Right bar (extends from center to right)
        rv = abs(right_vals[ci]) if ci < len(right_vals) else 0
        rw = int(side_w * (rv / vmax)) if vmax > 0 else 0
        if rw > 0:
            rx = center_x + center_w
            _add_bar(slide, rx, bar_y, rw, bar_h, accent, radius_px)

        # Value labels
        if show_values:
            vh = int(Pt(val_pt).emu * 1.3)
            # Left value: to the left of the left bar
            if lw > 0:
                lv_x = center_x - lw - int(Pt(val_pt).emu * 3.5)
                lv_w = int(Pt(val_pt).emu * 3.2)
                _add_text(
                    slide, lv_x, row_y, lv_w, row_h,
                    _fmt(lv, value_suffix),
                    font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                )
            # Right value: to the right of the right bar
            if rw > 0:
                rv_x = center_x + center_w + rw + int(Pt(val_pt).emu * 0.3)
                rv_w = int(Pt(val_pt).emu * 3.2)
                _add_text(
                    slide, rv_x, row_y, rv_w, row_h,
                    _fmt(rv, value_suffix),
                    font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                )


def _draw_legend(slide, x, y, w, h, left_name, right_name,
                 left_color, right_color, font_body, pt, text_c, muted, radius_px):
    """Right-aligned legend with two swatches."""
    swatch = int(Pt(pt).emu * 0.9)
    gap = int(Pt(pt).emu * 0.4)
    item_gap = int(Pt(pt).emu * 1.2)

    items = [(left_name, left_color), (right_name, right_color)]
    widths = []
    for name, _ in items:
        est_text_w = int(Pt(pt).emu * 0.55 * max(len(name), 1))
        widths.append(swatch + gap + est_text_w)
    total = sum(widths) + item_gap * (len(widths) - 1)
    cursor = x + max(0, w - total)
    cy = y + (h - swatch) // 2

    for (name, color), wd in zip(items, widths):
        sw = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(int(cursor)), Emu(int(cy)),
            Emu(swatch), Emu(swatch),
        )
        _set_fill(sw, color)
        _no_line(sw)
        tx = cursor + swatch + gap
        tw = wd - swatch - gap
        _add_text(
            slide, tx, y, tw, h, name,
            font_name=font_body, size_pt=pt, hex_color=text_c,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )
        cursor += wd + item_gap
