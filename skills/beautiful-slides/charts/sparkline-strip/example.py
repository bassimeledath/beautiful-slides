"""Run `render()` 5 times -- once per mode -- writing example-<mode>.pptx per run."""

import os
import sys

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module
from tokens import MODES  # noqa: E402
from render import render  # noqa: E402


DATA = {
    "title": "Key metrics dashboard",
    "metrics": [
        {
            "name": "Revenue",
            "value": 2.4,
            "prefix": "$",
            "suffix": "M",
            "delta": 0.3,
            "delta_label": "vs prev quarter",
            "sparkline": [1.8, 1.9, 2.0, 2.1, 2.0, 2.2, 2.3, 2.4],
        },
        {
            "name": "Active Users",
            "value": 14200,
            "delta": 1200,
            "delta_label": "MoM",
            "sparkline": [11000, 11500, 12000, 12800, 13200, 13500, 14000, 14200],
        },
        {
            "name": "Churn Rate",
            "value": 3.2,
            "suffix": "%",
            "delta": -0.5,
            "delta_label": "vs last month",
            "sparkline": [4.1, 3.9, 3.8, 3.7, 3.5, 3.4, 3.3, 3.2],
        },
        {
            "name": "NPS Score",
            "value": 72,
            "delta": 4,
            "sparkline": [60, 62, 65, 67, 68, 70, 71, 72],
        },
        {
            "name": "Avg Deal Size",
            "value": 48.5,
            "prefix": "$",
            "suffix": "K",
            "delta": 3.2,
            "delta_label": "QoQ",
            "sparkline": [38, 40, 42, 44, 45, 46, 47, 48.5],
        },
        {
            "name": "Support CSAT",
            "value": 94,
            "suffix": "%",
            "delta": 2,
            "delta_label": "vs target",
            "sparkline": [88, 89, 90, 91, 92, 93, 93, 94],
        },
    ],
}


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    for mode_name, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        margin = Inches(0.4)
        bounds = (
            int(margin),
            int(margin),
            int(prs.slide_width - 2 * margin),
            int(prs.slide_height - 2 * margin),
        )
        render(slide, DATA, tokens, bounds)

        out_path = os.path.join(out_dir, f"example-{mode_name}.pptx")
        prs.save(out_path)
        print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
