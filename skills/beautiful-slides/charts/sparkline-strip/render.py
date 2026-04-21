"""Sparkline strip — native python-pptx shapes only.

Row of tiny inline line charts (sparklines) paired with KPI labels and current
values. 4-8 metrics, each sparkline ~60-80px tall, no axes.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Emu(int(x1)), Emu(int(y1)),
                                      Emu(int(x2)), Emu(int(y2)))
    line = conn.line
    line.color.rgb = _rgb(color_hex)
    line.width = Pt(weight_pt)
    return conn


def _add_filled_circle(slide, cx, cy, r_emu, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
                                 Emu(int(r_emu * 2)), Emu(int(r_emu * 2)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y)),
                                 Emu(int(w)), Emu(int(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _fmt_value(v, prefix="", suffix=""):
    if isinstance(v, float):
        if abs(v - round(v)) < 1e-6:
            s = f"{int(round(v))}"
        elif abs(v) >= 100:
            s = f"{v:.0f}"
        elif abs(v) >= 10:
            s = f"{v:.1f}"
        else:
            s = f"{v:.2f}"
    else:
        s = str(v)
    return f"{prefix}{s}{suffix}"


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 1].rstrip() + "\u2026"


def render(slide, data, tokens, bounds):
    x0, y0, w0, h0 = bounds
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # Background
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    title = data.get("title")
    metrics = list(data.get("metrics", []))

    if not metrics:
        return

    # Clamp to 4-8
    metrics = metrics[:8]

    pad = int(min(w0, h0) * 0.035)

    # Title
    cursor_y = y0 + pad
    if title:
        title_pt = int(base_pt * 1.55)
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, x0 + pad, cursor_y, w0 - 2 * pad, title_h,
                     title, font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += title_h + int(pad * 0.5)

    # Layout: grid of metric cards
    n = len(metrics)
    # Determine grid: prefer 4 columns for 4-8 items
    if n <= 4:
        cols = n
        rows = 1
    elif n <= 6:
        cols = 3
        rows = 2
    else:
        cols = 4
        rows = 2

    avail_w = w0 - 2 * pad
    avail_h = y0 + h0 - pad - cursor_y
    gap = int(min(avail_w, avail_h) * 0.025)

    cell_w = (avail_w - gap * (cols - 1)) // cols
    cell_h = (avail_h - gap * (rows - 1)) // rows

    if cell_w <= 0 or cell_h <= 0:
        return

    # Sparkline target height: ~60-80px equivalent in EMU (9525 EMU/px)
    sparkline_h = min(int(cell_h * 0.35), 80 * 9525)
    sparkline_h = max(sparkline_h, int(cell_h * 0.2))

    # Adaptive font sizes — reduce more aggressively for dense grids
    if n >= 7:
        label_pt = max(7, int(base_pt * 0.55))
        value_pt = max(10, int(base_pt * 1.0))
        delta_pt = max(6, int(base_pt * 0.5))
    elif n >= 5:
        label_pt = max(7, int(base_pt * 0.7))
        value_pt = int(base_pt * 1.2)
        delta_pt = max(7, int(base_pt * 0.65))
    else:
        label_pt = max(8, int(base_pt * 0.8))
        value_pt = int(base_pt * 1.4)
        delta_pt = max(8, int(base_pt * 0.75))

    # Estimate max chars for metric name truncation
    label_char_w = Pt(label_pt).emu * 0.55
    max_label_chars = max(10, int(cell_w / label_char_w))

    for idx, metric in enumerate(metrics):
        row = idx // cols
        col = idx % cols

        cx = x0 + pad + col * (cell_w + gap)
        cy = cursor_y + row * (cell_h + gap)

        name = metric.get("name", "")
        value = metric.get("value", 0)
        delta = metric.get("delta")
        delta_label = metric.get("delta_label", "")
        sparkline_data = metric.get("sparkline", [])
        prefix = metric.get("prefix", "")
        suffix = metric.get("suffix", "")

        # Card internal layout: label, value + delta, sparkline
        inner_pad = int(cell_w * 0.06)
        content_x = cx + inner_pad
        content_w = cell_w - 2 * inner_pad
        content_y = cy + inner_pad

        # Truncate long metric names to prevent overflow
        display_name = _truncate(name, max_label_chars)

        # 1. Metric label — single line, no wrap
        label_h = int(Pt(label_pt).emu * 1.3)
        _add_textbox(slide, content_x, content_y, content_w, label_h,
                     display_name, font_body, label_pt, text_c,
                     bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        content_y += label_h + int(inner_pad * 0.15)

        # 2. Value (hero number)
        value_h = int(Pt(value_pt).emu * 1.3)
        value_text = _fmt_value(value, prefix, suffix)
        _add_textbox(slide, content_x, content_y, content_w, value_h,
                     value_text, font_display, value_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        content_y += value_h

        # 3. Delta (optional)
        if delta is not None:
            delta_h = int(Pt(delta_pt).emu * 1.3)
            if isinstance(delta, (int, float)):
                delta_positive = delta >= 0
                delta_str = f"+{_fmt_value(delta)}" if delta >= 0 else _fmt_value(delta)
            else:
                delta_str = str(delta)
                delta_positive = not delta_str.startswith("-")
            if delta_label:
                delta_str = f"{delta_str} {delta_label}"
            delta_color = accent if delta_positive else primary
            _add_textbox(slide, content_x, content_y, content_w, delta_h,
                         delta_str, font_mono, delta_pt, delta_color,
                         bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
            content_y += delta_h + int(inner_pad * 0.1)
        else:
            content_y += int(inner_pad * 0.2)

        # 4. Sparkline
        if sparkline_data and len(sparkline_data) >= 2:
            spark_y = content_y
            spark_h = min(sparkline_h, cy + cell_h - inner_pad - spark_y)
            if spark_h < int(Pt(6).emu):
                continue

            spark_x = content_x
            spark_w = content_w

            # Filter valid values
            valid = [(i, float(v)) for i, v in enumerate(sparkline_data) if v is not None]
            if len(valid) < 2:
                continue

            vmin = min(v for _, v in valid)
            vmax = max(v for _, v in valid)
            if vmax == vmin:
                vmax = vmin + 1

            n_pts = len(sparkline_data)
            step_x = spark_w / max(n_pts - 1, 1)

            def sx(i):
                return spark_x + i * step_x

            def sy(v):
                frac = (v - vmin) / (vmax - vmin)
                # Small inner padding for sparkline
                margin_v = spark_h * 0.08
                return spark_y + spark_h - margin_v - frac * (spark_h - 2 * margin_v)

            # Draw sparkline segments
            spark_color = primary
            for a, b in zip(valid, valid[1:]):
                _add_line(slide, sx(a[0]), sy(a[1]), sx(b[0]), sy(b[1]),
                          spark_color, 1.5)

            # End dot
            last = valid[-1]
            dot_r = int(Pt(2.5).emu)
            _add_filled_circle(slide, sx(last[0]), sy(last[1]), dot_r, accent)
