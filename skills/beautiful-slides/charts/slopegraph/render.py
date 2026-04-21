"""Slopegraph — native python-pptx shapes only.

Two vertical axes (before/after) with lines connecting each item's position,
showing rank or value changes. Direct-labeled, no legend.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Emu(int(x1)), Emu(int(y1)),
                                      Emu(int(x2)), Emu(int(y2)))
    line = conn.line
    line.color.rgb = _rgb(color_hex)
    line.width = Pt(weight_pt)
    return conn


def _add_filled_circle(slide, cx, cy, r_emu, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
                                 Emu(int(r_emu * 2)), Emu(int(r_emu * 2)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y)),
                                 Emu(int(w)), Emu(int(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _fmt_num(v):
    if abs(v - round(v)) < 1e-6:
        return f"{int(round(v))}"
    return f"{v:.1f}"


def _boost_muted(muted_hex, bg_hex, min_contrast=80):
    """Ensure muted color has enough contrast against the background.

    If the muted color is too close to the background, lighten or darken
    it to achieve at least *min_contrast* distance (sum of abs channel diffs).
    """
    m = muted_hex.lstrip("#")
    b = bg_hex.lstrip("#")
    mr, mg, mb = int(m[0:2], 16), int(m[2:4], 16), int(m[4:6], 16)
    br, bg_, bb = int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16)
    dist = abs(mr - br) + abs(mg - bg_) + abs(mb - bb)
    if dist >= min_contrast:
        return muted_hex
    # Determine direction: if bg is dark, lighten muted; if light, darken
    bg_lum = br * 0.299 + bg_ * 0.587 + bb * 0.114
    if bg_lum < 128:
        # Dark background -- push muted toward white
        factor = 0.45
        nr = min(255, mr + int((255 - mr) * factor))
        ng = min(255, mg + int((255 - mg) * factor))
        nb = min(255, mb + int((255 - mb) * factor))
    else:
        # Light background -- push muted toward black
        factor = 0.45
        nr = max(0, mr - int(mr * factor))
        ng = max(0, mg - int(mg * factor))
        nb = max(0, mb - int(mb * factor))
    return f"#{nr:02X}{ng:02X}{nb:02X}"


def _truncate(s, max_chars=30):
    """Truncate string with ellipsis if longer than max_chars."""
    if len(s) > max_chars:
        return s[:max_chars - 1] + "\u2026"
    return s


def render(slide, data, tokens, bounds):
    x0, y0, w0, h0 = bounds
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = int(tokens["font_size_base_pt"])

    # Improve muted contrast on dark backgrounds
    muted = _boost_muted(muted, bg)

    # Background
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    title = data.get("title")
    left_label = data.get("left_label", "Before")
    right_label = data.get("right_label", "After")
    items = list(data.get("items", []))
    highlight = data.get("highlight", [])
    value_suffix = data.get("value_suffix", "")

    if not items:
        return

    # Clamp items to 5-12
    items = items[:12]

    pad = int(min(w0, h0) * 0.035)

    # Title
    cursor_y = y0 + pad
    title_h = 0
    if title:
        title_pt = int(base_pt * 1.55)
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, x0 + pad, cursor_y, w0 - 2 * pad, title_h,
                     title, font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += title_h + int(pad * 0.3)

    # Layout: left label column, slope area, right label column
    label_pt = int(base_pt * 0.9)
    value_pt = int(base_pt * 0.85)
    header_pt = int(base_pt * 0.95)

    # Reserve width for labels on each side
    max_name_len = max((len(item.get("name", "")) for item in items), default=5)
    max_val_len = max(
        max((len(_fmt_num(item.get("left", 0)) + value_suffix) for item in items), default=3),
        max((len(_fmt_num(item.get("right", 0)) + value_suffix) for item in items), default=3),
    )
    label_col_w = int(Pt(label_pt).emu * 0.55 * (max_name_len + max_val_len + 3))
    label_col_w = min(label_col_w, int(w0 * 0.32))
    label_col_w = max(label_col_w, int(w0 * 0.18))

    # Header row for left/right labels
    header_h = int(Pt(header_pt).emu * 2.0)
    _add_textbox(slide, x0 + pad, cursor_y, label_col_w, header_h,
                 left_label, font_body, header_pt, text_c,
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    _add_textbox(slide, x0 + w0 - pad - label_col_w, cursor_y, label_col_w, header_h,
                 right_label, font_body, header_pt, text_c,
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    cursor_y += header_h + int(pad * 0.3)

    # Slope area
    slope_x = x0 + pad + label_col_w
    slope_w = w0 - 2 * pad - 2 * label_col_w
    slope_y = cursor_y
    slope_h = y0 + h0 - pad - cursor_y

    if slope_w <= 0 or slope_h <= 0:
        return

    # Compute value range for positioning
    all_vals = []
    for item in items:
        lv = item.get("left")
        rv = item.get("right")
        if lv is not None:
            all_vals.append(float(lv))
        if rv is not None:
            all_vals.append(float(rv))
    if not all_vals:
        return

    vmin = min(all_vals)
    vmax = max(all_vals)
    span = vmax - vmin if vmax > vmin else max(abs(vmax), 1.0)
    # Add padding to range
    vmin_p = vmin - span * 0.08
    vmax_p = vmax + span * 0.08

    # Item vertical spacing: use value-based positioning
    item_h = int(Pt(label_pt).emu * 1.6)
    dot_r = int(Pt(3.5).emu)

    def val_to_y(v):
        """Map value to y position (higher value = higher on slide = lower y)."""
        if vmax_p == vmin_p:
            return slope_y + slope_h / 2
        frac = (v - vmin_p) / (vmax_p - vmin_p)
        return slope_y + slope_h - frac * slope_h

    # Determine highlight set
    highlight_set = set(highlight) if highlight else set()

    # Compute raw y positions for left and right columns
    left_positions = []  # (raw_y, index)
    right_positions = []
    for idx, item in enumerate(items):
        lv = float(item.get("left", 0))
        rv = float(item.get("right", 0))
        left_positions.append((val_to_y(lv), idx))
        right_positions.append((val_to_y(rv), idx))

    # Label de-collision pass: nudge overlapping y-positions apart
    min_spacing = item_h * 1.1

    def _decollide(positions):
        """Sort by y and push apart any that overlap, return index->y map."""
        sorted_pos = sorted(positions, key=lambda p: p[0])
        for j in range(1, len(sorted_pos)):
            if sorted_pos[j][0] - sorted_pos[j - 1][0] < min_spacing:
                sorted_pos[j] = (sorted_pos[j - 1][0] + min_spacing, sorted_pos[j][1])
        # Clamp: shift all back if they overflow the slope area bottom
        if sorted_pos:
            overflow = sorted_pos[-1][0] - (slope_y + slope_h - item_h / 2)
            if overflow > 0:
                sorted_pos = [(sy - overflow, si) for sy, si in sorted_pos]
        return {si: sy for sy, si in sorted_pos}

    left_y_map = _decollide(left_positions)
    right_y_map = _decollide(right_positions)

    # Draw each item
    for idx, item in enumerate(items):
        name = _truncate(item.get("name", ""), 30)
        lv = float(item.get("left", 0))
        rv = float(item.get("right", 0))
        is_highlight = item.get("name", "") in highlight_set or not highlight_set

        # Use raw y for line/dot positions, de-collided y for labels
        ly_raw = val_to_y(lv)
        ry_raw = val_to_y(rv)
        ly_label = left_y_map[idx]
        ry_label = right_y_map[idx]

        # Choose color based on highlight
        if highlight_set:
            line_color = primary if is_highlight else muted
            text_color = text_c
            line_weight = 2.0 if is_highlight else 1.0
        else:
            # Color by direction: increase = accent, decrease = primary, flat = muted
            if rv > lv:
                line_color = accent
            elif rv < lv:
                line_color = primary
            else:
                line_color = muted
            text_color = text_c
            line_weight = 1.75

        # Draw connecting line (at raw data positions)
        _add_line(slide, slope_x, ly_raw, slope_x + slope_w, ry_raw, line_color, line_weight)

        # Dots at endpoints
        _add_filled_circle(slide, slope_x, ly_raw, dot_r, line_color)
        _add_filled_circle(slide, slope_x + slope_w, ry_raw, dot_r, line_color)

        # Left label: "Name  value" (at de-collided y)
        left_text = f"{name}  {_fmt_num(lv)}{value_suffix}"
        _add_textbox(slide,
                     x0 + pad, ly_label - item_h / 2,
                     label_col_w - int(pad * 0.3), item_h,
                     left_text, font_body, label_pt, text_color,
                     bold=(item.get("name", "") in highlight_set and bool(highlight_set)),
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

        # Right label: "value  Name" (at de-collided y)
        right_text = f"{_fmt_num(rv)}{value_suffix}  {name}"
        _add_textbox(slide,
                     x0 + w0 - pad - label_col_w + int(pad * 0.3), ry_label - item_h / 2,
                     label_col_w - int(pad * 0.3), item_h,
                     right_text, font_body, label_pt, text_color,
                     bold=(item.get("name", "") in highlight_set and bool(highlight_set)),
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
