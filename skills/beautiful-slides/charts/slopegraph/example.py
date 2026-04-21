"""Run `render()` 5 times -- once per mode -- writing example-<mode>.pptx per run."""

import os
import sys

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module
from tokens import MODES  # noqa: E402
from render import render  # noqa: E402


DATA = {
    "title": "Market share shift, 2024 vs 2025",
    "left_label": "2024",
    "right_label": "2025",
    "items": [
        {"name": "Alpha Corp",  "left": 32, "right": 28},
        {"name": "Beta Inc",    "left": 25, "right": 30},
        {"name": "Gamma Ltd",   "left": 18, "right": 19},
        {"name": "Delta Co",    "left": 15, "right": 14},
        {"name": "Epsilon AG",  "left": 10, "right": 9},
        {"name": "Zeta Group",  "left": 8,  "right": 12},
        {"name": "Eta Partners", "left": 6, "right": 5},
    ],
    "highlight": ["Beta Inc", "Zeta Group"],
    "value_suffix": "%",
}


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    for mode_name, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        margin = Inches(0.4)
        bounds = (
            int(margin),
            int(margin),
            int(prs.slide_width - 2 * margin),
            int(prs.slide_height - 2 * margin),
        )
        render(slide, DATA, tokens, bounds)

        out_path = os.path.join(out_dir, f"example-{mode_name}.pptx")
        prs.save(out_path)
        print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
