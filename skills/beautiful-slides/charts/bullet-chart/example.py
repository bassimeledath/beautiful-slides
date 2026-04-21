"""Render the bullet chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_SCORECARD = {
    "title": "Q1 performance scorecard",
    "metrics": [
        {"label": "Revenue",       "actual": 82,  "target": 90, "ranges": [100, 75, 50], "suffix": "M"},
        {"label": "Profit margin", "actual": 22,  "target": 25, "ranges": [30, 20, 10],  "suffix": "%"},
        {"label": "NPS",           "actual": 65,  "target": 70, "ranges": [80, 60, 40],  "suffix": ""},
        {"label": "Retention",     "actual": 91,  "target": 95, "ranges": [100, 85, 70], "suffix": "%"},
    ],
    "show_values": True,
}

DATA_SLA = {
    "title": "SLA dashboard, March",
    "metrics": [
        {"label": "Uptime",        "actual": 99.7, "target": 99.9, "ranges": [100, 99.5, 99], "suffix": "%"},
        {"label": "Avg latency",   "actual": 142,  "target": 120,  "ranges": [200, 150, 100], "suffix": "ms"},
        {"label": "Error rate",    "actual": 0.3,  "target": 0.5,  "ranges": [1.0, 0.5, 0.2], "suffix": "%"},
    ],
    "show_values": True,
}

DATA_GOALS = {
    "title": "Quarterly goal progress",
    "metrics": [
        {"label": "New customers",  "actual": 340,  "target": 400,  "ranges": [500, 350, 200], "suffix": ""},
        {"label": "Pipeline",       "actual": 4.2,  "target": 5.0,  "ranges": [6.0, 4.0, 2.0], "suffix": "M"},
        {"label": "Win rate",       "actual": 28,   "target": 35,   "ranges": [40, 30, 20],    "suffix": "%"},
        {"label": "Avg deal size",  "actual": 48,   "target": 55,   "ranges": [70, 50, 30],    "suffix": "K"},
        {"label": "Churn",          "actual": 2.1,  "target": 1.5,  "ranges": [3.0, 2.0, 1.0], "suffix": "%"},
    ],
    "show_values": True,
}


MODE_DATA = {
    "sv-keynote":           DATA_SCORECARD,
    "editorial-magazine":   DATA_SLA,
    "playful-marketing":    DATA_SCORECARD,
    "consulting-boardroom": DATA_GOALS,
    "craft-minimal":        DATA_SLA,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
