"""Bullet chart renderer — native python-pptx shapes only.

Stephen Few's bullet graph: a thin horizontal performance bar against a
target marker and qualitative range bands (poor / satisfactory / good).
Stacks 3-6 metrics vertically.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


_EMU_PER_PX = 9525


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _set_fill(shape, hex_):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_)


def _no_line(shape):
    shape.line.fill.background()


def _lerp_hex(h1, h2, t):
    """Linearly interpolate between two hex colours."""
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_text(slide, x, y, w, h, text, font_name, size_pt, hex_color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold = bold
    return tb


def _add_rect(slide, x, y, w, h, fill_hex, radius_px=0):
    w_i = max(1, int(w))
    h_i = max(1, int(h))
    if radius_px and radius_px > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i),
        )
        short_emu = min(w_i, h_i)
        radius_emu = radius_px * _EMU_PER_PX
        ratio = max(0.0, min(0.5, radius_emu / short_emu / 2.0))
        try:
            shape.adjustments[0] = ratio
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i),
        )
    _set_fill(shape, fill_hex)
    _no_line(shape)
    return shape


def _fmt(v, suffix=""):
    """Format a value with compact number abbreviation (K, M) and optional suffix."""
    abs_v = abs(v)
    if abs_v >= 1_000_000:
        s = f"{v / 1_000_000:.1f}M"
        # Remove trailing .0 for clean display
        s = s.replace(".0M", "M")
    elif abs_v >= 10_000:
        s = f"{v / 1_000:.0f}K"
    elif abs_v >= 1_000:
        s = f"{v / 1_000:.1f}K"
        s = s.replace(".0K", "K")
    elif abs(v - round(v)) < 1e-9:
        s = f"{int(round(v))}"
    else:
        s = f"{v:.1f}"
    return f"{s}{suffix}" if suffix else s


def render(slide, data, tokens, bounds):
    """Render a bullet chart onto *slide* inside *bounds*.

    ``data`` keys:
      metrics : list[dict]
          Each dict has:
            label  : str      -- metric name
            actual : number   -- the performance bar value
            target : number   -- the target marker value
            ranges : list     -- [poor, satisfactory, good] thresholds
            suffix : str      -- optional value suffix (e.g. "%")
      title : str or None
      show_values : bool (default True)
    """
    x, y, w, h = bounds

    metrics = data.get("metrics") or []
    if not metrics:
        return
    title = data.get("title")
    show_values = bool(data.get("show_values", True))

    # Theme tokens
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    _set_fill(bg_shape, bg)
    _no_line(bg_shape)

    # Outer padding
    pad = int(min(w, h) * 0.04)
    ix = x + pad
    iy = y + pad
    iw = w - 2 * pad
    ih = h - 2 * pad

    # Title -- dynamic height for multi-line wrapping
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.6))
        # Estimate how many lines the title will wrap to
        char_width_emu = Pt(title_pt).emu * 0.55
        chars_per_line = max(1, int(iw / char_width_emu))
        title_lines = max(1, -(-len(title) // chars_per_line))  # ceil division
        title_h = int(Pt(title_pt).emu * 1.4 * title_lines + Pt(title_pt).emu * 0.4)
        title_h = min(title_h, int(ih * 0.25))  # cap at 25% of inner height
        _add_text(
            slide, ix, iy, iw, title_h, title,
            font_name=font_display, size_pt=title_pt, hex_color=text_c,
            bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        iy += title_h + int(Pt(base_pt).emu * 0.6)
        ih = (y + h - pad) - iy

    n = len(metrics)
    label_pt = max(int(base_pt * 0.92), 9)
    val_pt = max(int(base_pt * 0.78), 8)

    # Label column on left
    max_label_chars = max(len(m.get("label", "")) for m in metrics) if metrics else 4
    label_col_w = int(Pt(label_pt).emu * 0.55 * min(max_label_chars + 2, 24))
    label_col_w = min(label_col_w, int(iw * 0.35))
    value_col_w = int(Pt(val_pt).emu * 4.0) if show_values else 0

    bar_region_x = ix + label_col_w + int(Pt(label_pt).emu * 0.4)
    bar_region_w = max(1, iw - label_col_w - int(Pt(label_pt).emu * 0.4) - value_col_w)

    # Vertical layout
    row_gap = int(ih * 0.06)
    total_gaps = row_gap * (n - 1)
    row_h = max(int((ih - total_gaps) / n), int(Pt(base_pt * 2)))

    # Qualitative range band colours: 3 bands from lightest to darkest
    # We lerp from bg toward muted at 3 intensities
    band_colors = [
        _lerp_hex(bg, muted, 0.12),   # good (lightest)
        _lerp_hex(bg, muted, 0.25),   # satisfactory
        _lerp_hex(bg, muted, 0.40),   # poor (darkest)
    ]

    for mi, metric in enumerate(metrics):
        row_y = iy + mi * (row_h + row_gap)

        label = metric.get("label", "")
        actual = metric.get("actual", 0)
        target = metric.get("target", 0)
        ranges = metric.get("ranges") or [0, 0, 0]
        suffix = metric.get("suffix", "")

        # Determine the scale max (the largest range value)
        scale_max = max(ranges) if ranges else max(actual, target, 1)
        if scale_max <= 0:
            scale_max = 1

        # Label on left
        _add_text(
            slide, ix, row_y, label_col_w, row_h, label,
            font_name=font_body, size_pt=label_pt, hex_color=text_c,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Qualitative range bands (drawn widest to narrowest)
        # ranges = [poor, satisfactory, good] where good < satisfactory < poor
        # We sort descending to draw widest first
        sorted_ranges = sorted(enumerate(ranges), key=lambda x: x[1], reverse=True)

        band_h = int(row_h * 0.65)
        band_y = row_y + (row_h - band_h) // 2

        for rank, (orig_idx, rval) in enumerate(sorted_ranges):
            rw = int(bar_region_w * (rval / scale_max))
            if rw < 1:
                continue
            _add_rect(
                slide, bar_region_x, band_y, rw, band_h,
                band_colors[rank], radius_px,
            )

        # Performance bar (thinner, centered)
        perf_h = int(band_h * 0.38)
        perf_y = band_y + (band_h - perf_h) // 2
        perf_w = int(bar_region_w * (max(actual, 0) / scale_max))
        if perf_w > 0:
            _add_rect(
                slide, bar_region_x, perf_y, perf_w, perf_h,
                primary, radius_px,
            )

        # Target marker (vertical line)
        target_x = bar_region_x + int(bar_region_w * (max(target, 0) / scale_max))
        marker_h = int(band_h * 0.85)
        marker_y = band_y + (band_h - marker_h) // 2
        marker_w = max(int(_EMU_PER_PX * 2.5), int(Pt(2).emu))
        _add_rect(
            slide, target_x - marker_w // 2, marker_y, marker_w, marker_h,
            accent, 0,
        )

        # Value label on right
        if show_values:
            val_x = bar_region_x + bar_region_w + int(Pt(val_pt).emu * 0.3)
            val_w = value_col_w
            _add_text(
                slide, val_x, row_y, val_w, row_h,
                _fmt(actual, suffix),
                font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                bold=True,
            )
