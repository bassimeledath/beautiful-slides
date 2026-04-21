import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


DATA = {
    "title": "Website traffic flow",
    "stages": [
        # Stage 0: Traffic sources
        [
            {"label": "Organic",  "value": 5000},
            {"label": "Paid",     "value": 3000},
            {"label": "Referral", "value": 2000},
        ],
        # Stage 1: Landing pages
        [
            {"label": "Homepage", "value": 6000},
            {"label": "Pricing",  "value": 2500},
            {"label": "Blog",     "value": 1500},
        ],
        # Stage 2: Outcomes
        [
            {"label": "Signup",   "value": 3200},
            {"label": "Bounce",   "value": 6800},
        ],
    ],
    "flows": [
        # Organic -> pages
        {"from": [0, 0], "to": [1, 0], "value": 3000},
        {"from": [0, 0], "to": [1, 1], "value": 1500},
        {"from": [0, 0], "to": [1, 2], "value":  500},
        # Paid -> pages
        {"from": [0, 1], "to": [1, 0], "value": 2000},
        {"from": [0, 1], "to": [1, 1], "value":  700},
        {"from": [0, 1], "to": [1, 2], "value":  300},
        # Referral -> pages
        {"from": [0, 2], "to": [1, 0], "value": 1000},
        {"from": [0, 2], "to": [1, 1], "value":  300},
        {"from": [0, 2], "to": [1, 2], "value":  700},
        # Pages -> outcomes
        {"from": [1, 0], "to": [2, 0], "value": 2400},
        {"from": [1, 0], "to": [2, 1], "value": 3600},
        {"from": [1, 1], "to": [2, 0], "value":  600},
        {"from": [1, 1], "to": [2, 1], "value": 1900},
        {"from": [1, 2], "to": [2, 0], "value":  200},
        {"from": [1, 2], "to": [2, 1], "value": 1300},
    ],
}


def _set_slide_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))


def build(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA, tokens, (x, y, w, h))

    prs.save(out_path)


def main():
    for name, tokens in MODES.items():
        path = os.path.join(HERE, f"example-{name}.pptx")
        build(name, tokens, path)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
