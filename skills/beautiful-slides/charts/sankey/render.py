"""Sankey diagram — weighted flow bands between stages where band width = volume.

Bezier curves are approximated with freeform line segments (polyline).
Supports 2-4 stages with 2-12 flows.
"""

import math

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(2)))
    tf.margin_right = Emu(int(Pt(2)))
    tf.margin_top = Emu(int(Pt(1)))
    tf.margin_bottom = Emu(int(Pt(1)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _cubic_bezier(p0, p1, p2, p3, n_segments=20):
    """Evaluate a cubic bezier curve at n_segments+1 points.

    Returns list of (x, y) tuples.
    """
    points = []
    for i in range(n_segments + 1):
        t = i / n_segments
        t2 = t * t
        t3 = t2 * t
        mt = 1 - t
        mt2 = mt * mt
        mt3 = mt2 * mt

        x = mt3 * p0[0] + 3 * mt2 * t * p1[0] + 3 * mt * t2 * p2[0] + t3 * p3[0]
        y = mt3 * p0[1] + 3 * mt2 * t * p1[1] + 3 * mt * t2 * p2[1] + t3 * p3[1]
        points.append((x, y))
    return points


def _draw_flow_band(slide, x_left, y_top_left, y_bot_left,
                    x_right, y_top_right, y_bot_right, fill_hex, n_segments=20):
    """Draw a curved flow band (filled shape) between two vertical edges.

    The band connects:
      left edge:  from (x_left, y_top_left) to (x_left, y_bot_left)
      right edge: from (x_right, y_top_right) to (x_right, y_bot_right)

    Top and bottom edges are cubic bezier curves approximated with line segments.
    """
    x_mid = (x_left + x_right) / 2
    ctrl_offset = (x_right - x_left) * 0.4

    # Top curve: left-top to right-top
    top_pts = _cubic_bezier(
        (x_left, y_top_left),
        (x_left + ctrl_offset, y_top_left),
        (x_right - ctrl_offset, y_top_right),
        (x_right, y_top_right),
        n_segments,
    )

    # Bottom curve: right-bottom to left-bottom (reversed so we trace the outline)
    bot_pts = _cubic_bezier(
        (x_right, y_bot_right),
        (x_right - ctrl_offset, y_bot_right),
        (x_left + ctrl_offset, y_bot_left),
        (x_left, y_bot_left),
        n_segments,
    )

    # Build closed polygon: top curve forward, then right edge down,
    # then bottom curve (already reversed), then left edge up (implicit close)
    all_pts = top_pts + bot_pts

    if len(all_pts) < 3:
        return

    start = all_pts[0]
    ff = slide.shapes.build_freeform(int(start[0]), int(start[1]), scale=1.0)
    segments = [(int(p[0]), int(p[1])) for p in all_pts[1:]]
    ff.add_line_segments(segments, close=True)
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    return shape


def render(slide, data, tokens, bounds):
    """Render a Sankey diagram.

    data:
        title    - optional string
        stages   - list of lists of {"label": str, "value": number}
                   e.g. [
                       [{"label": "Source A", "value": 100}, ...],  # stage 0
                       [{"label": "Middle X", "value": 80}, ...],   # stage 1
                       ...
                   ]
        flows    - list of {"from": [stage_idx, node_idx],
                            "to": [stage_idx, node_idx],
                            "value": number}
                   value determines the width of the flow band.
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]
    radius_px = tokens.get("radius_px", 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    stages = data.get("stages", [])
    flows = data.get("flows", [])
    if not stages or not flows:
        return
    title = data.get("title")
    n_stages = len(stages)

    # --- layout ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_h = (y + h) - cur_y
    avail_w = w

    # Node column width and spacing
    node_w = int(avail_w * 0.08)
    # Space between stage columns
    if n_stages <= 1:
        stage_spacing = avail_w
    else:
        stage_spacing = (avail_w - node_w) // (n_stages - 1)

    # Compute the max total value across any single stage for vertical scaling
    stage_totals = []
    for stage_nodes in stages:
        total = sum(n.get("value", 0) for n in stage_nodes)
        stage_totals.append(total)
    max_total = max(stage_totals) if stage_totals else 1

    # Vertical padding between nodes within a stage
    node_v_gap = int(avail_h * 0.03)

    # Label space below/beside nodes
    label_h = int(Pt(base_pt) * 2.0)

    # Calculate positions for all nodes
    # node_positions[stage_idx][node_idx] = {x, y, w, h, cy_top, cy_bot}
    node_positions = []
    for si, stage_nodes in enumerate(stages):
        n_nodes = len(stage_nodes)
        total_val = stage_totals[si]

        # X position for this stage
        if n_stages == 1:
            nx = x + avail_w // 2 - node_w // 2
        else:
            nx = x + int(si * stage_spacing)

        # Available height for nodes (minus gaps and label space)
        usable_h = avail_h - label_h - node_v_gap * max(n_nodes - 1, 0)
        # Scale factor: pixels per unit value
        if max_total > 0:
            px_per_val = usable_h / max_total
        else:
            px_per_val = 1

        positions = []
        current_y = cur_y
        for ni, node in enumerate(stage_nodes):
            val = node.get("value", 0)
            nh = max(int(val * px_per_val), int(Pt(base_pt * 0.8)))

            positions.append({
                "x": nx,
                "y": current_y,
                "w": node_w,
                "h": nh,
                "value": val,
                "label": node.get("label", ""),
                # Track consumed offsets on left and right edges for flow stacking
                "right_offset": 0,
                "left_offset": 0,
            })
            current_y += nh + node_v_gap

        node_positions.append(positions)

    # --- Draw flow bands first (behind nodes) ---
    for fi, flow in enumerate(flows):
        from_stage, from_node = flow["from"]
        to_stage, to_node = flow["to"]
        flow_val = flow.get("value", 0)
        if flow_val <= 0:
            continue

        if from_stage >= n_stages or to_stage >= n_stages:
            continue
        if from_node >= len(node_positions[from_stage]):
            continue
        if to_node >= len(node_positions[to_stage]):
            continue

        src = node_positions[from_stage][from_node]
        dst = node_positions[to_stage][to_node]

        # Scale flow width proportionally to the node heights
        src_total = stage_totals[from_stage]
        if src_total > 0:
            flow_h_src = max(int(src["h"] * flow_val / src["value"]) if src["value"] > 0 else int(Pt(2)), 1)
        else:
            flow_h_src = int(Pt(2))

        dst_total = stage_totals[to_stage]
        if dst_total > 0:
            flow_h_dst = max(int(dst["h"] * flow_val / dst["value"]) if dst["value"] > 0 else int(Pt(2)), 1)
        else:
            flow_h_dst = int(Pt(2))

        # Source edge: right side of source node
        x_left = src["x"] + src["w"]
        y_top_left = src["y"] + src["right_offset"]
        y_bot_left = y_top_left + flow_h_src

        # Target edge: left side of target node
        x_right = dst["x"]
        y_top_right = dst["y"] + dst["left_offset"]
        y_bot_right = y_top_right + flow_h_dst

        # Update consumed offsets
        src["right_offset"] += flow_h_src
        dst["left_offset"] += flow_h_dst

        # Flow color: blend from source node color
        t_src = from_node / max(len(node_positions[from_stage]) - 1, 1)
        flow_color_base = _lerp_hex(primary, accent, t_src)
        # Make flow slightly transparent by blending toward bg
        flow_color = _lerp_hex(flow_color_base, bg, 0.35)

        _draw_flow_band(
            slide, x_left, y_top_left, y_bot_left,
            x_right, y_top_right, y_bot_right,
            flow_color, n_segments=16,
        )

    # --- Draw node rectangles ---
    for si, stage_nodes in enumerate(node_positions):
        for ni, pos in enumerate(stage_nodes):
            t = ni / max(len(stage_nodes) - 1, 1)
            node_color = _lerp_hex(primary, accent, t)

            # Draw node bar
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE
            bar = slide.shapes.add_shape(
                shape_type,
                Emu(pos["x"]), Emu(pos["y"]),
                Emu(pos["w"]), Emu(pos["h"]),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = _rgb(node_color)
            bar.line.fill.background()

            # Node label
            lbl_w = int(stage_spacing * 0.45) if n_stages > 1 else int(avail_w * 0.3)
            lbl_font_size = max(int(base_pt * 0.7), 7)
            lbl_h = int(Pt(lbl_font_size) * 3)

            # Position label beside or below the node
            if si == 0:
                # First stage: label to the left of the node
                lbl_x = pos["x"] - lbl_w - int(Pt(2))
                lbl_y = pos["y"] + (pos["h"] - lbl_h) // 2
                lbl_align = PP_ALIGN.RIGHT
            elif si == n_stages - 1:
                # Last stage: label to the right of the node
                lbl_x = pos["x"] + pos["w"] + int(Pt(2))
                lbl_y = pos["y"] + (pos["h"] - lbl_h) // 2
                lbl_align = PP_ALIGN.LEFT
            else:
                # Middle stages: label above the node
                lbl_x = pos["x"] - (lbl_w - pos["w"]) // 2
                lbl_y = pos["y"] - lbl_h
                lbl_align = PP_ALIGN.CENTER

            # Clamp label to bounds
            lbl_x = max(x, min(lbl_x, x + w - lbl_w))
            lbl_y = max(cur_y, min(lbl_y, y + h - lbl_h))

            # Compose label with value
            label_text = pos["label"]
            if pos["value"] > 0:
                label_text += f" ({pos['value']:,})"

            _add_textbox(
                slide, lbl_x, lbl_y, lbl_w, lbl_h,
                label_text, font_body, lbl_font_size, text_c,
                align=lbl_align, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )

    return None
