"""Data pipeline — directed left-to-right pipeline diagram.

Shows sources -> transforms -> stores -> sinks as labeled boxes grouped
by stage with connecting arrows between stages.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ---- helpers ---------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    """Linearly interpolate between two hex colors."""
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _lighten_hex(hex_, factor=0.85):
    """Blend a hex color toward white by *factor* (0=original, 1=white)."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02X}{g:02X}{b:02X}"


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "\u2026"


def _luminance(hex_):
    """Compute relative luminance for contrast decisions."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 word_wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Emu(int(Pt(3).emu))
    tf.margin_right = Emu(int(Pt(3).emu))
    tf.margin_top = Emu(int(Pt(2).emu))
    tf.margin_bottom = Emu(int(Pt(2).emu))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_rounded_rect(slide, x, y, w, h, fill_hex, line_hex=None,
                      line_width_pt=0.75, radius_px=0):
    shape_type = (MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0
                  else MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(
        shape_type,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Pt(line_width_pt)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _draw_arrow(slide, x1, y1, x2, y2, color_hex, thickness):
    """Draw a directional arrow from (x1,y1) to (x2,y2) as shaft + triangle head."""
    dx = x2 - x1
    dy = y2 - y1
    length = math.sqrt(dx * dx + dy * dy)
    if length < 1:
        return

    ux = dx / length
    uy = dy / length
    px = -uy
    py = ux

    head_len = min(length * 0.3, thickness * 4)
    head_w = head_len * 0.7

    base_x = x2 - ux * head_len
    base_y = y2 - uy * head_len

    # Triangle head
    p1x, p1y = int(x2), int(y2)
    p2x, p2y = int(base_x + px * head_w / 2), int(base_y + py * head_w / 2)
    p3x, p3y = int(base_x - px * head_w / 2), int(base_y - py * head_w / 2)

    ff = slide.shapes.build_freeform(p1x, p1y, scale=1.0)
    ff.add_line_segments([(p2x, p2y), (p3x, p3y)], close=True)
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_hex)
    shape.line.fill.background()

    # Shaft
    shaft_len = length - head_len
    if shaft_len > 0:
        sw = thickness / 2
        s1x = int(x1 + px * sw)
        s1y = int(y1 + py * sw)
        s2x = int(x1 - px * sw)
        s2y = int(y1 - py * sw)
        s3x = int(base_x - px * sw)
        s3y = int(base_y - py * sw)
        s4x = int(base_x + px * sw)
        s4y = int(base_y + py * sw)

        ff2 = slide.shapes.build_freeform(s1x, s1y, scale=1.0)
        ff2.add_line_segments([(s2x, s2y), (s3x, s3y), (s4x, s4y)], close=True)
        shaft = ff2.convert_to_shape()
        shaft.fill.solid()
        shaft.fill.fore_color.rgb = _rgb(color_hex)
        shaft.line.fill.background()


# ---- main render -----------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a data pipeline diagram onto *slide* inside *bounds*.

    data:
        title  - optional string
        stages - list of {"label": str, "nodes": [str, ...]}
                 Stages render left-to-right. Each stage has a header and
                 vertically stacked node boxes. Arrows connect stage groups.
    """
    x0, y0, w0, h0 = bounds

    # --- unpack tokens -------------------------------------------------------
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px", 0))

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x0)), Emu(int(y0)),
        Emu(int(w0)), Emu(int(h0)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    # --- unpack data ---------------------------------------------------------
    stages = data.get("stages", [])
    if not stages:
        return
    title = data.get("title")
    n_stages = len(stages)

    # --- layout constants ----------------------------------------------------
    pad = int(min(w0, h0) * 0.025)

    ix, iy = x0 + pad, y0 + pad
    iw, ih = w0 - 2 * pad, h0 - 2 * pad

    # --- title ---------------------------------------------------------------
    if title:
        title_pt = int(round(base_pt * 1.5))
        # Estimate line count for long titles
        est_chars_per_line = max(1, int(iw / (Pt(title_pt).emu * 0.6)))
        est_lines = max(1, (len(title) + est_chars_per_line - 1) // est_chars_per_line)
        title_h = int(Pt(title_pt).emu * 1.8 * min(est_lines, 3))
        _add_textbox(slide, ix, iy, iw, title_h, title,
                     font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + pad  # increased gap to prevent collision with headers
        ih = (y0 + h0 - pad) - iy

    # --- stage columns layout ------------------------------------------------
    # Space between stage groups is used for arrows -- reduce for many stages
    arrow_gap = int(iw * max(0.02, 0.06 - 0.008 * n_stages))
    total_arrow_gap = arrow_gap * (n_stages - 1) if n_stages > 1 else 0
    stage_w = max(1, int((iw - total_arrow_gap) / n_stages))

    # Header and node sizing -- adaptive to stage count
    header_pt = max(7, int(base_pt * 0.9 * min(1.0, 5 / n_stages)))
    header_h = int(Pt(header_pt).emu * 2.4)
    node_pt = max(6, int(base_pt * 0.8 * min(1.0, 5 / n_stages)))
    node_h = int(Pt(node_pt).emu * (2.0 if n_stages > 5 else 2.6))
    node_gap = int(Pt(4).emu)
    node_pad = int(Pt(max(3, 6 - n_stages * 0.4)).emu)

    # Maximum nodes that can fit vertically
    avail_node_h = ih - header_h - int(pad * 0.5)
    max_nodes = max(1, int(avail_node_h / (node_h + node_gap)))

    # Track stage bounding boxes for arrows
    stage_rects = []  # (x, y, w, h) of each stage group area

    for si, stage in enumerate(stages):
        sx = ix + si * (stage_w + arrow_gap)
        sy = iy

        # Stage color: interpolate from primary to accent
        t = si / (n_stages - 1) if n_stages > 1 else 0.0
        stage_color = _lerp_hex(primary, accent, t)
        stage_bg_color = _lighten_hex(stage_color, 0.88)

        label = stage.get("label", f"Stage {si + 1}")
        nodes = list(stage.get("nodes", []))[:max_nodes]

        # Compute actual stage height needed
        nodes_total_h = len(nodes) * (node_h + node_gap) if nodes else 0
        stage_total_h = header_h + int(pad * 0.3) + nodes_total_h

        # Stage background column
        _add_rounded_rect(slide, sx, sy, stage_w, ih,
                          stage_bg_color, line_hex=stage_color,
                          line_width_pt=0.75, radius_px=radius_px)

        # Stage header bar
        _add_rounded_rect(slide, sx, sy, stage_w, header_h,
                          stage_color, radius_px=radius_px)
        # Contrast-aware header text color
        header_text_color = bg if _luminance(stage_color) < 140 else text_c
        # Truncate long header labels
        max_header_chars = max(10, int(stage_w / (Pt(header_pt).emu * 0.55)))
        _add_textbox(slide, sx, sy, stage_w, header_h,
                     _truncate(label, max_header_chars),
                     font_display, header_pt, header_text_color,
                     bold=True, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

        # Node boxes within the stage
        node_y = sy + header_h + int(pad * 0.3)
        node_inner_w = stage_w - 2 * node_pad

        for ni, node_label in enumerate(nodes):
            nx = sx + node_pad
            ny = node_y

            # Clamp to bounds
            if ny + node_h > y0 + h0 - pad:
                break

            _add_rounded_rect(slide, nx, ny, node_inner_w, node_h,
                              bg, line_hex=muted,
                              line_width_pt=0.5, radius_px=radius_px)
            # Truncate long node labels
            max_node_chars = max(8, int(node_inner_w / (Pt(node_pt).emu * 0.55)))
            _add_textbox(slide, nx, ny, node_inner_w, node_h,
                         _truncate(node_label, max_node_chars),
                         font_body, node_pt, text_c,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            node_y += node_h + node_gap

        stage_rects.append((sx, sy, stage_w, ih))

    # --- draw arrows between stages ------------------------------------------
    arrow_thickness = max(int(Pt(base_pt * 0.15)), int(Pt(2).emu))
    arrow_color = _lerp_hex(muted, primary, 0.4)

    for si in range(n_stages - 1):
        s1_x, s1_y, s1_w, s1_h = stage_rects[si]
        s2_x, s2_y, s2_w, s2_h = stage_rects[si + 1]

        # Arrow from right edge of stage si to left edge of stage si+1
        # Vertically centered in the node area (below header)
        arrow_y = s1_y + header_h + int(ih - header_h) // 2
        ax1 = s1_x + s1_w + int(arrow_gap * 0.1)
        ax2 = s2_x - int(arrow_gap * 0.1)

        _draw_arrow(slide, ax1, arrow_y, ax2, arrow_y,
                    arrow_color, arrow_thickness)

    return None
