import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


DATA = {
    "title": "Real-time analytics pipeline",
    "stages": [
        {
            "label": "Sources",
            "nodes": ["Clickstream", "App events", "Server logs"],
        },
        {
            "label": "Ingestion",
            "nodes": ["Kafka", "Kinesis"],
        },
        {
            "label": "Transform",
            "nodes": ["Flink", "dbt", "Spark"],
        },
        {
            "label": "Storage",
            "nodes": ["Snowflake", "S3", "Redis"],
        },
        {
            "label": "Serve",
            "nodes": ["API", "Dashboard", "Alerts"],
        },
    ],
}


def _set_slide_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))


def build(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA, tokens, (x, y, w, h))

    prs.save(out_path)


def main():
    for name, tokens in MODES.items():
        path = os.path.join(HERE, f"example-{name}.pptx")
        build(name, tokens, path)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
