"""Decision tree — branching path of choices and outcomes.

Each node has a question or label. Branches have option labels. Tree expands
left-to-right or top-to-bottom from a root question.
"""

import math

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "\u2026"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(3)))
    tf.margin_right = Emu(int(Pt(3)))
    tf.margin_top = Emu(int(Pt(2)))
    tf.margin_bottom = Emu(int(Pt(2)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_arrow(slide, x1, y1, x2, y2, color_hex, thickness):
    """Draw a directional arrow from (x1,y1) to (x2,y2)."""
    dx = x2 - x1
    dy = y2 - y1
    length = math.sqrt(dx * dx + dy * dy)
    if length < 1:
        return

    ux = dx / length
    uy = dy / length
    px = -uy
    py = ux

    head_len = min(length * 0.25, thickness * 5)
    head_w = head_len * 0.7

    base_x = x2 - ux * head_len
    base_y = y2 - uy * head_len

    p1x, p1y = int(x2), int(y2)
    p2x, p2y = int(base_x + px * head_w / 2), int(base_y + py * head_w / 2)
    p3x, p3y = int(base_x - px * head_w / 2), int(base_y - py * head_w / 2)

    ff = slide.shapes.build_freeform(p1x, p1y, scale=1.0)
    ff.add_line_segments([(p2x, p2y), (p3x, p3y)], close=True)
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_hex)
    shape.line.fill.background()

    shaft_len = length - head_len
    if shaft_len > 0:
        sw = thickness / 2
        s1x = int(x1 + px * sw)
        s1y = int(y1 + py * sw)
        s2x = int(x1 - px * sw)
        s2y = int(y1 - py * sw)
        s3x = int(base_x - px * sw)
        s3y = int(base_y - py * sw)
        s4x = int(base_x + px * sw)
        s4y = int(base_y + py * sw)

        ff2 = slide.shapes.build_freeform(s1x, s1y, scale=1.0)
        ff2.add_line_segments([(s2x, s2y), (s3x, s3y), (s4x, s4y)], close=True)
        shaft = ff2.convert_to_shape()
        shaft.fill.solid()
        shaft.fill.fore_color.rgb = _rgb(color_hex)
        shaft.line.fill.background()


def _flatten_tree(node, depth=0):
    """Recursively flatten a nested tree into a list of (node, depth, parent_id, branch_label)."""
    result = []
    nid = node.get("id", f"node_{id(node)}")
    result.append((node, depth, None, None))
    for child_info in node.get("children", []):
        child = child_info.get("node", child_info)
        branch_label = child_info.get("label", "")
        child_id = child.get("id", f"node_{id(child)}")
        sub = _flatten_tree(child, depth + 1)
        # Mark parent
        if sub:
            sub[0] = (sub[0][0], sub[0][1], nid, branch_label)
        result.extend(sub)
    return result


def _count_leaves(node):
    """Count the number of leaf nodes (no children) in the subtree."""
    children = node.get("children", [])
    if not children:
        return 1
    total = 0
    for child_info in children:
        child = child_info.get("node", child_info)
        total += _count_leaves(child)
    return total


def _layout_tree(node, x, y, w, h, depth, max_depth, direction):
    """Recursively assign positions to each node in the tree.

    Returns list of (node_id, cx, cy, depth, children_positions, node_ref)
    where children_positions is list of (child_id, branch_label).
    """
    nid = node.get("id", f"node_{id(node)}")
    children = node.get("children", [])

    if direction == "LR":
        # Left-to-right: depth controls x-position, leaves spread vertically
        col_w = w / (max_depth + 1) if max_depth > 0 else w
        cx = x + int((depth + 0.5) * col_w)
        cy = y + h // 2  # will be adjusted for leaves
    else:
        # Top-to-bottom: depth controls y-position, leaves spread horizontally
        row_h = h / (max_depth + 1) if max_depth > 0 else h
        cx = x + w // 2  # will be adjusted for leaves
        cy = y + int((depth + 0.5) * row_h)

    if not children:
        return [(nid, cx, cy, depth, [], node)]

    results = []
    child_layouts = []
    total_leaves = _count_leaves(node)

    # Distribute children across the available cross-axis space
    leaf_idx = 0
    for child_info in children:
        child = child_info.get("node", child_info)
        branch_label = child_info.get("label", "")
        child_leaves = _count_leaves(child)

        if direction == "LR":
            # Children spread vertically
            child_y = y + int(h * leaf_idx / total_leaves)
            child_h = int(h * child_leaves / total_leaves)
            child_x = x
            child_w = w
        else:
            # Children spread horizontally
            child_x = x + int(w * leaf_idx / total_leaves)
            child_w = int(w * child_leaves / total_leaves)
            child_y = y
            child_h = h

        sub = _layout_tree(child, child_x, child_y, child_w, child_h,
                           depth + 1, max_depth, direction)
        child_id = child.get("id", f"node_{id(child)}")
        child_layouts.append((child_id, branch_label, sub))
        leaf_idx += child_leaves

    # Position this node at the average of its children's cross-axis
    if child_layouts:
        if direction == "LR":
            avg_y = sum(s[0][2] for _, _, s in child_layouts) / len(child_layouts)
            cy = int(avg_y)
        else:
            avg_x = sum(s[0][1] for _, _, s in child_layouts) / len(child_layouts)
            cx = int(avg_x)

    children_refs = [(cid, bl) for cid, bl, _ in child_layouts]
    results.append((nid, cx, cy, depth, children_refs, node))

    for _, _, sub in child_layouts:
        results.extend(sub)

    return results


def _max_depth(node, d=0):
    """Find the maximum depth of the tree."""
    children = node.get("children", [])
    if not children:
        return d
    return max(_max_depth(c.get("node", c), d + 1) for c in children)


def render(slide, data, tokens, bounds):
    """Render a decision tree.

    data:
        title     - optional string
        direction - "LR" (left-to-right, default) or "TB" (top-to-bottom)
        root      - nested tree: {"id": str, "label": str, "children": [
                        {"label": str (branch label), "node": {same structure}}
                    ]}
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    root = data.get("root")
    if not root:
        return
    title = data.get("title")
    direction = data.get("direction", "LR")

    # --- layout title ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_w = w
    avail_h = (y + h) - cur_y
    area_x = x
    area_y = cur_y

    # Compute tree layout
    md = _max_depth(root)
    positions = _layout_tree(root, area_x, area_y, avail_w, avail_h,
                             0, md, direction)

    # Build lookup
    pos_map = {}
    for nid, cx, cy, depth, children_refs, node_ref in positions:
        pos_map[nid] = (cx, cy, depth, children_refs, node_ref)

    num_depths = md + 1
    total_leaves = _count_leaves(root)

    # Node sizing -- adaptive to leaf count and depth
    if direction == "LR":
        node_w = int(avail_w / (num_depths + 0.5) * 0.65)
        node_h = int(Pt(base_pt) * 3)
    else:
        node_w = int(Pt(base_pt) * 7)
        node_h = int(avail_h / (num_depths + 0.5) * 0.45)

    node_w = max(node_w, int(Pt(base_pt) * 5))
    node_h = max(node_h, int(Pt(base_pt) * 2.2))
    # Cap to keep within reasonable cell size
    node_w = min(node_w, int(avail_w * 0.20))
    node_h = min(node_h, int(avail_h * 0.12))

    # Adaptive: shrink node_h when leaves would overlap
    if total_leaves > 0:
        max_leaf_h = int(avail_h / total_leaves * 0.85)
        node_h = min(node_h, max(max_leaf_h, int(Pt(base_pt) * 1.4)))

    # Adaptive font scaling based on leaf count
    font_scale = min(1.0, 5 / max(total_leaves, 1))
    adaptive_pt = max(int(base_pt * 0.7 * font_scale + base_pt * 0.3), 7)

    # --- draw edges ---
    arrow_thick = max(int(Pt(base_pt * 0.12)), int(Pt(1.5)))
    arrow_color = _lerp_hex(muted, primary, 0.3)

    for nid, cx, cy, depth, children_refs, node_ref in positions:
        for child_id, branch_label in children_refs:
            if child_id not in pos_map:
                continue
            ccx, ccy, _, _, _ = pos_map[child_id]

            # Compute arrow start/end offset from node edges
            dx = ccx - cx
            dy = ccy - cy
            dist = math.sqrt(dx * dx + dy * dy)
            if dist < 1:
                continue
            ux = dx / dist
            uy = dy / dist

            if abs(ux) > abs(uy):
                src_off = node_w * 0.5
                dst_off = node_w * 0.5
            else:
                src_off = node_h * 0.5
                dst_off = node_h * 0.5

            a_sx = cx + ux * src_off
            a_sy = cy + uy * src_off
            a_ex = ccx - ux * dst_off
            a_ey = ccy - uy * dst_off

            _draw_arrow(slide, a_sx, a_sy, a_ex, a_ey, arrow_color, arrow_thick)

            # Branch label near the midpoint
            if branch_label:
                mid_x = (a_sx + a_ex) / 2
                mid_y = (a_sy + a_ey) / 2
                lbl_w = int(Pt(base_pt) * 5)
                lbl_h = int(Pt(base_pt) * 1.6)
                # Nudge perpendicular to arrow direction
                perp_x = -uy
                perp_y = ux
                nudge = int(Pt(base_pt) * 0.8)
                lbl_x = int(mid_x + perp_x * nudge) - lbl_w // 2
                lbl_y = int(mid_y + perp_y * nudge) - lbl_h // 2
                # Clamp
                lbl_x = max(x, min(lbl_x, x + w - lbl_w))
                lbl_y = max(y, min(lbl_y, y + h - lbl_h))
                _add_textbox(
                    slide, lbl_x, lbl_y, lbl_w, lbl_h,
                    _truncate(branch_label, 20), font_body,
                    max(int(adaptive_pt * 0.85), 7), text_c,
                    align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE,
                )

    # --- draw nodes ---
    # Estimate max characters that fit based on node_w and font size
    max_label_chars = max(12, int(node_w / (Pt(adaptive_pt) * 0.6)))

    for nid, cx, cy, depth, children_refs, node_ref in positions:
        label = _truncate(node_ref.get("label", ""), max_label_chars)
        is_leaf = len(children_refs) == 0
        is_root = depth == 0

        t = depth / max(md, 1)

        # Clamp node position to stay within bounds
        nx = max(area_x, min(cx - node_w // 2, area_x + avail_w - node_w))
        ny = max(area_y, min(cy - node_h // 2, area_y + avail_h - node_h))

        if is_root:
            # Root: prominent filled rounded rect
            fill = primary
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(nx), Emu(ny), Emu(node_w), Emu(node_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.fill.background()
            _add_textbox(
                slide, nx, ny, node_w, node_h,
                label, font_body, adaptive_pt, bg,
                align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )
        elif is_leaf:
            # Leaf: outcome node — accent-tinted
            fill = _lerp_hex(accent, primary, t * 0.3)
            leaf_font = max(int(adaptive_pt * 0.85), 7)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(nx), Emu(ny), Emu(node_w), Emu(node_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.fill.background()
            _add_textbox(
                slide, nx, ny, node_w, node_h,
                label, font_body, leaf_font, bg,
                align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )
        else:
            # Interior decision node — bordered rectangle
            fill = _lerp_hex(bg, primary, 0.06)
            border = _lerp_hex(primary, muted, 0.4)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(nx), Emu(ny), Emu(node_w), Emu(node_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.color.rgb = _rgb(border)
            shape.line.width = Pt(1.5)
            _add_textbox(
                slide, nx, ny, node_w, node_h,
                label, font_body, adaptive_pt, text_c,
                align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE,
            )

    return None
