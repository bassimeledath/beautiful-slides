import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


DATA = {
    "title": "Pricing strategy",
    "direction": "LR",
    "root": {
        "id": "q1",
        "label": "Enterprise customer?",
        "children": [
            {
                "label": "Yes",
                "node": {
                    "id": "q2",
                    "label": "Annual contract?",
                    "children": [
                        {"label": "Yes", "node": {"id": "a1", "label": "Enterprise Annual"}},
                        {"label": "No",  "node": {"id": "a2", "label": "Enterprise Monthly"}},
                    ],
                },
            },
            {
                "label": "No",
                "node": {
                    "id": "q3",
                    "label": "Team size > 10?",
                    "children": [
                        {"label": "Yes", "node": {"id": "a3", "label": "Business plan"}},
                        {"label": "No",  "node": {"id": "a4", "label": "Starter plan"}},
                    ],
                },
            },
        ],
    },
}


def _set_slide_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))


def build(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA, tokens, (x, y, w, h))

    prs.save(out_path)


def main():
    for name, tokens in MODES.items():
        path = os.path.join(HERE, f"example-{name}.pptx")
        build(name, tokens, path)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
