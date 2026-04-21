"""Render the lollipop chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_HORIZONTAL = {
    "orientation": "horizontal",
    "title": "Top 10 accounts by ARR",
    "items": [
        {"label": "Acme Corp",      "value": 2.4},
        {"label": "Globex Inc",     "value": 1.9},
        {"label": "Initech",        "value": 1.7},
        {"label": "Umbrella Corp",  "value": 1.3},
        {"label": "Wonka Industries", "value": 0.9},
        {"label": "Stark Ent",      "value": 0.85},
        {"label": "Wayne Corp",     "value": 0.72},
        {"label": "Hooli",          "value": 0.65},
    ],
    "value_suffix": "M",
    "show_values": True,
}

DATA_VERTICAL = {
    "orientation": "vertical",
    "title": "Monthly active users",
    "items": [
        {"label": "Jan", "value": 120},
        {"label": "Feb", "value": 135},
        {"label": "Mar", "value": 142},
        {"label": "Apr", "value": 128},
        {"label": "May", "value": 155},
        {"label": "Jun", "value": 170},
    ],
    "value_suffix": "K",
    "show_values": True,
}


MODE_DATA = {
    "sv-keynote": DATA_HORIZONTAL,
    "editorial-magazine": DATA_VERTICAL,
    "playful-marketing": DATA_HORIZONTAL,
    "consulting-boardroom": DATA_HORIZONTAL,
    "craft-minimal": DATA_VERTICAL,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
