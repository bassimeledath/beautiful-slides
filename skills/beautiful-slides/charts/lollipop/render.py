"""Lollipop chart renderer -- native python-pptx shapes only.

A dot on a thin stem for each item, horizontal by default.
High data-ink ratio alternative to a bar chart.  Best for ranking 5-15 items.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# ── constants ──────────────────────────────────────────────────────────
_EMU_PER_PX = 9525


# ── helpers ────────────────────────────────────────────────────────────
def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _set_fill(shape, hex_):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_)


def _no_line(shape):
    shape.line.fill.background()


def _no_fill(shape):
    shape.fill.background()


def _style_run(run, font_name, size_pt, hex_color, bold=False):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold = bold


def _add_text(slide, x, y, w, h, text, font_name, size_pt, hex_color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _style_run(run, font_name, size_pt, hex_color, bold=bold)
    return tb


def _add_line(slide, x1, y1, x2, y2, hex_color, width_emu):
    ln = slide.shapes.add_connector(
        1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)),
    )
    ln.line.color.rgb = _rgb(hex_color)
    ln.line.width = Emu(int(width_emu))
    return ln


def _add_dot(slide, cx, cy, diameter, hex_color):
    r = diameter // 2
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(cx - r)), Emu(int(cy - r)),
        Emu(int(diameter)), Emu(int(diameter)),
    )
    _set_fill(dot, hex_color)
    _no_line(dot)
    return dot


def _nice_ticks(vmax, target=5):
    if vmax <= 0:
        return [0, 1], 1
    import math
    raw = vmax / target
    magnitude = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    residual = raw / magnitude
    if residual < 1.5:
        step = 1 * magnitude
    elif residual < 3:
        step = 2 * magnitude
    elif residual < 7:
        step = 5 * magnitude
    else:
        step = 10 * magnitude
    top = step * math.ceil(vmax / step)
    ticks = []
    v = 0.0
    while v <= top + 1e-9:
        ticks.append(v)
        v += step
    return ticks, top


def _fmt(v, suffix=""):
    if abs(v - round(v)) < 1e-9:
        s = f"{int(round(v))}"
    else:
        s = f"{v:.1f}"
    return f"{s}{suffix}" if suffix else s


# ── public API ─────────────────────────────────────────────────────────
def render(slide, data, tokens, bounds):
    x, y, w, h = bounds

    orientation = (data.get("orientation") or "horizontal").lower()
    title = data.get("title")
    items = list(data.get("items") or [])
    value_suffix = data.get("value_suffix") or ""
    show_values = bool(data.get("show_values", True))

    if not items:
        return

    items = items[:15]  # cap

    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    _set_fill(bg_shape, bg)
    _no_line(bg_shape)

    # Padding
    pad = int(min(w, h) * 0.04)
    ix, iy = x + pad, y + pad
    iw, ih = w - 2 * pad, h - 2 * pad

    # Title
    if title:
        title_pt = int(round(base_pt * 1.6))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_text(slide, ix, iy, iw, title_h, title,
                  font_name=font_display, size_pt=title_pt, hex_color=text_c,
                  bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + int(Pt(base_pt).emu * 0.6)
        ih = (y + h - pad) - iy

    labels = [str(item.get("label", "")) for item in items]
    values = [item.get("value", 0) for item in items]
    vmax_raw = max(values) if values else 1

    if orientation == "horizontal":
        _draw_horizontal(
            slide, ix, iy, iw, ih, labels, values, vmax_raw,
            primary, text_c, muted, font_body, font_mono,
            base_pt, value_suffix, show_values,
        )
    else:
        _draw_vertical(
            slide, ix, iy, iw, ih, labels, values, vmax_raw,
            primary, text_c, muted, font_body, font_mono,
            base_pt, value_suffix, show_values,
        )


def _draw_horizontal(slide, x, y, w, h, labels, values, vmax_raw,
                     primary, text_c, muted, font_body, font_mono,
                     base_pt, suffix, show_values):
    """Horizontal lollipop: categories on the left, stems go right."""
    cat_pt = max(int(base_pt * 0.92), 9)
    tick_pt = max(int(base_pt * 0.78), 8)
    val_pt = max(int(base_pt * 0.78), 8)

    # Adaptive font scaling for dense data
    n = len(labels)
    if n > 10:
        cat_pt = max(int(base_pt * 0.78), 8)

    # Ellipsis truncation for long labels
    max_display_chars = 24
    labels = [
        (l[:max_display_chars - 1] + "\u2026" if len(l) > max_display_chars else l)
        for l in labels
    ]

    max_cat_chars = max((len(l) for l in labels), default=4)
    left_margin = int(Pt(cat_pt).emu * 0.55 * min(max_cat_chars + 2, 24))
    bottom_margin = int(Pt(tick_pt).emu * 2.0)
    top_margin = int(Pt(val_pt).emu * 0.6)
    right_margin = int(Pt(val_pt).emu * 4.0)  # room for value labels

    plot_x = x + left_margin
    plot_y = y + top_margin
    plot_w = max(1, w - left_margin - right_margin)
    plot_h = max(1, h - top_margin - bottom_margin)

    ticks, vmax = _nice_ticks(max(vmax_raw, 0.0))
    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    # X-axis tick labels
    for i, t in enumerate(ticks):
        tx = plot_x + int(plot_w * (t / vmax)) if vmax > 0 else plot_x
        tlh = int(Pt(tick_pt).emu * 1.4)
        label = _fmt(t) + (suffix if t != 0 and i == len(ticks) - 1 else "")
        _add_text(
            slide, tx - int(Pt(tick_pt).emu * 2),
            plot_y + plot_h + int(Pt(tick_pt).emu * 0.3),
            int(Pt(tick_pt).emu * 4), tlh, label,
            font_name=font_mono, size_pt=tick_pt, hex_color=text_c,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )
        if i > 0:
            _add_line(slide, tx, plot_y, tx, plot_y + plot_h, muted, hairline)

    # Baseline
    _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, muted, hairline * 2)

    # Minimum row height guard -- trim items if rows would be too tight
    row_h = plot_h / n
    min_row_h = int(Pt(cat_pt).emu * 1.6)
    if row_h < min_row_h:
        n = max(1, int(plot_h // min_row_h))
        labels = labels[:n]
        values = values[:n]
        row_h = plot_h / n
    dot_diameter = int(min(row_h * 0.35, Pt(base_pt).emu * 1.2))
    stem_width = max(int(_EMU_PER_PX * 1.5), int(dot_diameter * 0.18))

    for i, (label, val) in enumerate(zip(labels, values)):
        row_top = plot_y + i * row_h
        cy = int(row_top + row_h / 2)

        # Category label
        cat_h_em = int(Pt(cat_pt).emu * 1.4)
        _add_text(
            slide, x, row_top + (row_h - cat_h_em) / 2,
            left_margin - int(Pt(cat_pt).emu * 0.4), cat_h_em, label,
            font_name=font_body, size_pt=cat_pt, hex_color=text_c,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Stem
        stem_len = int(plot_w * (max(val, 0) / vmax)) if vmax > 0 else 0
        if stem_len > 0:
            _add_line(slide, plot_x, cy, plot_x + stem_len, cy,
                      muted, stem_width)
            # Dot
            _add_dot(slide, plot_x + stem_len, cy, dot_diameter, primary)

        # Value label
        if show_values:
            vt = _fmt(val, suffix)
            vh = int(Pt(val_pt).emu * 1.3)
            vw = int(Pt(val_pt).emu * 4)
            _add_text(
                slide, plot_x + stem_len + dot_diameter // 2 + int(Pt(val_pt).emu * 0.3),
                cy - vh // 2, vw, vh, vt,
                font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
            )


def _draw_vertical(slide, x, y, w, h, labels, values, vmax_raw,
                   primary, text_c, muted, font_body, font_mono,
                   base_pt, suffix, show_values):
    """Vertical lollipop: categories on the bottom, stems go up."""
    cat_pt = max(int(base_pt * 0.92), 9)
    tick_pt = max(int(base_pt * 0.78), 8)
    val_pt = max(int(base_pt * 0.78), 8)

    tick_labels_raw = _nice_ticks(max(vmax_raw, 0.0))
    ticks, vmax = tick_labels_raw
    tick_labels = [_fmt(t) for t in ticks]
    max_tick_chars = max(len(s) for s in tick_labels) if tick_labels else 1

    left_margin = int(Pt(tick_pt).emu * 0.65 * (max_tick_chars + 1))
    bottom_margin = int(Pt(cat_pt).emu * 2.2)
    top_margin = int(Pt(val_pt).emu * 1.6)
    right_margin = int(Pt(val_pt).emu * 1.2)

    plot_x = x + left_margin
    plot_y = y + top_margin
    plot_w = max(1, w - left_margin - right_margin)
    plot_h = max(1, h - top_margin - bottom_margin)

    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    # Y-axis tick labels
    for i, t in enumerate(ticks):
        ty = plot_y + plot_h - int(plot_h * (t / vmax)) if vmax > 0 else plot_y + plot_h
        lbl_w = left_margin - int(Pt(tick_pt).emu * 0.3)
        lbl_h = int(Pt(tick_pt).emu * 1.4)
        label = _fmt(t) + (suffix if t != 0 and i == len(ticks) - 1 else "")
        _add_text(slide, x, ty - lbl_h // 2, lbl_w, lbl_h, label,
                  font_name=font_mono, size_pt=tick_pt, hex_color=text_c,
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        if i > 0:
            _add_line(slide, plot_x, ty, plot_x + plot_w, ty, muted, hairline)

    # Baseline
    base_y = plot_y + plot_h
    _add_line(slide, plot_x, base_y, plot_x + plot_w, base_y, muted, hairline * 2)

    n = len(labels)
    col_w = plot_w / n
    dot_diameter = int(min(col_w * 0.35, Pt(base_pt).emu * 1.2))
    stem_width = max(int(_EMU_PER_PX * 1.5), int(dot_diameter * 0.18))

    for i, (label, val) in enumerate(zip(labels, values)):
        cx = int(plot_x + i * col_w + col_w / 2)

        # Category label
        cat_y = base_y + int(Pt(cat_pt).emu * 0.4)
        cat_h = int(Pt(cat_pt).emu * 1.6)
        _add_text(slide, plot_x + i * col_w, cat_y, col_w, cat_h, label,
                  font_name=font_body, size_pt=cat_pt, hex_color=text_c,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        # Stem
        stem_h = int(plot_h * (max(val, 0) / vmax)) if vmax > 0 else 0
        dot_y = base_y - stem_h
        if stem_h > 0:
            _add_line(slide, cx, base_y, cx, dot_y, muted, stem_width)
            # Dot
            _add_dot(slide, cx, dot_y, dot_diameter, primary)

        # Value label
        if show_values:
            vt = _fmt(val, suffix)
            vh = int(Pt(val_pt).emu * 1.3)
            vw = int(Pt(val_pt).emu * 4)
            _add_text(
                slide, cx - vw // 2,
                dot_y - vh - int(Pt(val_pt).emu * 0.1),
                vw, vh, vt,
                font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
            )
