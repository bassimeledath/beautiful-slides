"""Roadmap Gantt — time-scaled horizontal bars with optional milestones and today marker."""

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.LEFT, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(2)))
    tf.margin_right = Emu(int(Pt(2)))
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


import math


def _estimate_lines(text, font_pt, avail_w):
    """Estimate how many lines text will wrap to given available width."""
    char_w = Pt(font_pt) * 0.55
    chars_per_line = max(1, int(avail_w / char_w))
    return max(1, math.ceil(len(text) / chars_per_line))


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 1].rstrip() + "\u2026"


def render(slide, data, tokens, bounds):
    """Render a roadmap Gantt chart.

    data:
        title      - optional string
        time_units - list of str (e.g. ["Q1", "Q2", "Q3", "Q4"])
        rows       - list of {"label": str, "start": int, "end": int,
                               "milestones": [int, ...] (optional)}
                     start/end are 0-based indices into time_units
        today      - optional int index for a vertical "today" marker
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]
    radius_px = tokens.get("radius_px", 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    rows = data.get("rows", [])
    if not rows:
        return
    time_units = data.get("time_units", [])
    if not time_units:
        return
    title = data.get("title")
    today = data.get("today")

    n_rows = len(rows)
    n_cols = len(time_units)

    # --- layout ---
    cur_y = y
    if title:
        title_pt = int(base_pt * 1.5)
        title_lines = _estimate_lines(title, title_pt, w)
        title_h = int(Pt(title_pt) * 1.3 * title_lines + Pt(title_pt) * 0.5)
        title_h = min(title_h, int(h * 0.20))
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, title_pt, text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.8))

    avail_h = (y + h) - cur_y
    avail_w = w

    # Label column on the left -- dynamically size based on longest label
    max_label_len = max((len(r.get("label", "")) for r in rows), default=10)
    label_frac = min(0.30, 0.20 + 0.01 * max(0, max_label_len - 20))
    label_col_w = int(avail_w * label_frac)
    gap = int(avail_w * 0.01)
    right_pad = int(avail_w * 0.01)  # small right margin to prevent edge clipping
    chart_x = x + label_col_w + gap
    chart_w = avail_w - label_col_w - gap - right_pad

    # Header row for time units
    header_h = int(Pt(base_pt) * 2.2)
    body_y = cur_y + header_h
    body_h = avail_h - header_h

    # Use exact chart_w for positioning to prevent right-edge overflow.
    # Map time indices [0, n_cols] to pixel range [chart_x, chart_x + chart_w].
    # Each column occupies chart_w / n_cols width.
    def time_to_x(t):
        """Convert a time index (0 to n_cols) to an x EMU coordinate."""
        return int(chart_x + (t / n_cols) * chart_w)

    col_w = chart_w / n_cols  # float for precision

    # Draw header time unit labels
    for ci, tu in enumerate(time_units):
        col_x = time_to_x(ci)
        col_end = time_to_x(ci + 1)
        _add_textbox(
            slide, col_x, cur_y, col_end - col_x, header_h,
            tu, font_body, max(int(base_pt * 0.8), 8), text_c,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.BOTTOM,
        )

    # Draw light vertical gridlines
    for ci in range(1, n_cols):
        col_x = time_to_x(ci)
        gl = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(col_x), Emu(body_y),
            Emu(max(int(Pt(0.5)), 1)), Emu(body_h),
        )
        gl.fill.solid()
        gl.fill.fore_color.rgb = _rgb(muted)
        gl.line.fill.background()

    # Row layout -- adaptive font for many rows
    label_pt = max(7, base_pt - max(0, n_rows - 6)) if n_rows > 8 else base_pt
    # Adaptive row gap — shrink for many rows to fit
    if n_rows >= 10:
        row_gap = int(Pt(base_pt * 0.15))
    elif n_rows >= 8:
        row_gap = int(Pt(base_pt * 0.25))
    else:
        row_gap = int(Pt(base_pt * 0.35))
    total_gaps = row_gap * (n_rows - 1) if n_rows > 1 else 0
    row_h = int((body_h - total_gaps) / n_rows)
    # Enforce minimum row height but not at the expense of overflowing
    min_row_h = int(Pt(base_pt * 1.2))
    row_h = max(row_h, min_row_h)
    # Re-check: if total exceeds body_h, shrink row_h to fit exactly
    total_needed = row_h * n_rows + total_gaps
    if total_needed > body_h:
        row_h = int((body_h - total_gaps) / n_rows)
    bar_h = max(int(row_h * 0.55), int(Pt(base_pt * 0.9)))
    bar_h = min(bar_h, row_h - 2)  # ensure bar fits within row

    # Calculate max chars for label truncation
    label_char_w = Pt(label_pt) * 0.55
    max_label_chars = max(10, int(label_col_w / label_char_w))

    for ri, row in enumerate(rows):
        ry = body_y + ri * (row_h + row_gap)
        bar_y = ry + (row_h - bar_h) // 2

        # Row label -- truncate with ellipsis if too long
        row_label = _truncate(row.get("label", ""), max_label_chars)
        _add_textbox(
            slide, x, ry, label_col_w, row_h,
            row_label, font_body, label_pt, text_c,
            align=PP_ALIGN.RIGHT, bold=False, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Alternating row background stripe
        if ri % 2 == 0:
            stripe_color = _lerp_hex(bg, muted, 0.08)
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(chart_x), Emu(ry),
                Emu(chart_w), Emu(row_h),
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = _rgb(stripe_color)
            stripe.line.fill.background()

        # Bar
        start_idx = row.get("start", 0)
        end_idx = row.get("end", n_cols)
        # Clamp to [0, n_cols]
        start_idx = max(0, min(start_idx, n_cols))
        end_idx = max(start_idx, min(end_idx, n_cols))

        bar_x = time_to_x(start_idx)
        bar_end_x = time_to_x(end_idx)
        bar_w = bar_end_x - bar_x
        if bar_w < 1:
            bar_w = int(col_w * 0.1)

        t = ri / (n_rows - 1) if n_rows > 1 else 0.0
        bar_color = _lerp_hex(primary, accent, t)

        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE,
            Emu(bar_x), Emu(bar_y),
            Emu(bar_w), Emu(bar_h),
        )
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = _rgb(bar_color)
        bar_shape.line.fill.background()

        # Milestone diamonds
        milestones = row.get("milestones", [])
        diamond_size = int(bar_h * 0.55)
        for mi in milestones:
            if mi < 0 or mi > n_cols:
                continue
            mx = time_to_x(mi)
            # Clamp diamond to stay within chart bounds
            mx = max(chart_x + diamond_size // 2, min(mx, chart_x + chart_w - diamond_size // 2))
            my = bar_y + (bar_h - diamond_size) // 2
            d = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Emu(mx - diamond_size // 2), Emu(my),
                Emu(diamond_size), Emu(diamond_size),
            )
            d.fill.solid()
            d.fill.fore_color.rgb = _rgb(text_c)
            d.line.fill.background()

    # Today marker
    if today is not None and 0 <= today <= n_cols:
        today_x = time_to_x(today)
        # Clamp today marker within chart area
        marker_w = max(int(Pt(2)), 2)
        today_x = min(today_x, chart_x + chart_w - marker_w)
        marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(today_x - marker_w // 2), Emu(cur_y),
            Emu(marker_w), Emu(avail_h),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = _rgb(accent)
        marker.line.fill.background()

        # "Today" label
        today_label_w = int(col_w * 0.6)
        today_label_h = int(Pt(base_pt * 0.75) * 1.8)
        _add_textbox(
            slide,
            today_x - today_label_w // 2, cur_y - today_label_h,
            today_label_w, today_label_h,
            "Today", font_body, max(int(base_pt * 0.65), 7), accent,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.BOTTOM,
        )

    return None
