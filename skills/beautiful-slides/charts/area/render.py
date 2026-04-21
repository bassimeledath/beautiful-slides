"""Area chart — native python-pptx shapes only.

Line chart with filled region between line and x-axis. Supports single-series
and stacked (max 4 series) with semi-transparent fills.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _hex_to_components(hex_):
    """Return (r, g, b) ints from '#RRGGBB'."""
    h = hex_.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 wrap=False):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Emu(int(x1)), Emu(int(y1)),
                                      Emu(int(x2)), Emu(int(y2)))
    line = conn.line
    line.color.rgb = _rgb(color_hex)
    line.width = Pt(weight_pt)
    if dash is not None:
        try:
            line.dash_style = dash
        except Exception:
            pass
    return conn


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y)),
                                 Emu(int(w)), Emu(int(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_freeform_fill(slide, points, fill_hex, alpha_pct=40):
    """Draw a filled polygon (freeform shape) with semi-transparent fill.

    points: list of (x_emu, y_emu) tuples defining the polygon.
    alpha_pct: fill opacity as a percentage (0=transparent, 100=opaque).
    """
    if len(points) < 3:
        return None

    # Use freeform builder
    builder = slide.shapes.build_freeform(
        start_x=Emu(int(points[0][0])),
        start_y=Emu(int(points[0][1]))
    )
    segments = [(Emu(int(px)), Emu(int(py))) for px, py in points[1:]]
    builder.add_line_segments(segments)

    # Close and add the shape
    shp = builder.convert_to_shape()
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    # Set transparency via alpha
    # python-pptx stores alpha as integer 0-100000 (100000 = fully opaque)
    try:
        fill_elem = shp.fill._fill
        solid = fill_elem.find(qn('a:solidFill'))
        if solid is not None:
            srgb = solid.find(qn('a:srgbClr'))
            if srgb is not None:
                from lxml import etree
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', str(int(alpha_pct * 1000)))
    except Exception:
        pass
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _nice_ticks(vmin, vmax, target=5):
    if vmax <= vmin:
        vmax = vmin + 1
    span = vmax - vmin
    raw = span / max(target, 1)
    mag = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    for mult in (1, 2, 2.5, 5, 10):
        step = mult * mag
        if span / step <= target * 1.5:
            break
    lo = math.floor(vmin / step) * step
    hi = math.ceil(vmax / step) * step
    ticks = []
    v = lo
    while v <= hi + 1e-9:
        ticks.append(round(v, 6))
        v += step
    return ticks, lo, hi


def _fmt_num(v):
    """Format number with compact notation (K/M) for large values."""
    av = abs(v)
    if av >= 1_000_000:
        s = f"{v / 1_000_000:.1f}M"
        return s.replace(".0M", "M")
    if av >= 10_000:
        s = f"{v / 1_000:.0f}K"
        return s
    if av >= 1_000:
        s = f"{v / 1_000:.1f}K"
        return s.replace(".0K", "K")
    if abs(v - round(v)) < 1e-6:
        return f"{int(round(v))}"
    return f"{v:.1f}"


def _estimate_lines(text, font_pt, avail_w):
    """Estimate how many lines text will wrap to given available width."""
    char_w = Pt(font_pt).emu * 0.55
    chars_per_line = max(1, int(avail_w / char_w))
    return max(1, math.ceil(len(text) / chars_per_line))


def render(slide, data, tokens, bounds):
    x0, y0, w0, h0 = bounds
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # Background
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    title = data.get("title")
    x_labels = list(data.get("x_labels", []))
    series = list(data.get("series", []))
    x_label = data.get("x_label")
    y_label = data.get("y_label")
    stacked = bool(data.get("stacked", False))

    # Limit to 4 series
    series = series[:4]

    if not series:
        return

    pad = int(min(w0, h0) * 0.035)

    # Series colors: cycle through primary, accent, then muted-blended variants
    series_colors = []
    base_colors = [primary, accent]
    # Generate extra colors by blending primary/accent with muted
    pr, pg, pb = _hex_to_components(primary)
    ar, ag, ab = _hex_to_components(accent)
    mr, mg, mb = _hex_to_components(muted)
    blend3 = f"#{(pr + mr) // 2:02x}{(pg + mg) // 2:02x}{(pb + mb) // 2:02x}"
    blend4 = f"#{(ar + mr) // 2:02x}{(ag + mg) // 2:02x}{(ab + mb) // 2:02x}"
    base_colors.extend([blend3, blend4])
    for i in range(len(series)):
        series_colors.append(base_colors[i % len(base_colors)])

    # Title
    cursor_y = y0 + pad
    title_h = 0
    if title:
        title_pt = int(base_pt * 1.55)
        title_avail_w = w0 - 2 * pad
        title_lines = _estimate_lines(title, title_pt, title_avail_w)
        title_h = int(Pt(title_pt).emu * 1.3 * title_lines + Pt(title_pt).emu * 0.5)
        title_h = min(title_h, int(h0 * 0.25))  # cap at 25% of height
        _add_textbox(slide, x0 + pad, cursor_y, title_avail_w, title_h,
                     title, font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                     wrap=True)
        cursor_y += title_h + int(pad * 0.3)

    if y_label:
        sub_h = int(Pt(base_pt * 0.85).emu * 1.6)
        _add_textbox(slide, x0 + pad, cursor_y, w0 - 2 * pad, sub_h,
                     y_label, font_body, int(base_pt * 0.85), text_c,
                     bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += sub_h

    # Plot area
    tick_pt = max(8, int(base_pt * 0.8))
    x_tick_h = int(Pt(tick_pt).emu * 1.6)
    x_axis_label_h = int(Pt(tick_pt).emu * 1.8)

    left_pad = int(w0 * 0.10)
    right_pad = int(w0 * 0.03)
    bottom_pad = x_tick_h + (x_axis_label_h if x_label else 0) + int(pad * 0.5)

    # Legend space
    legend_h = 0
    n_series = len(series)
    if n_series > 1:
        legend_h = int(Pt(tick_pt).emu * 2.0)

    plot_x = x0 + pad + left_pad
    plot_y = cursor_y + int(pad * 0.2) + legend_h
    plot_w = (x0 + w0 - pad - right_pad) - plot_x
    plot_h = (y0 + h0 - pad - bottom_pad) - plot_y

    if plot_w <= 0 or plot_h <= 0:
        return

    # Compute stacked values if needed
    n_x = max(len(x_labels), 1)
    if stacked:
        # Build cumulative stacks
        stacked_vals = []
        cumulative = [0.0] * n_x
        for s in series:
            vals = s.get("values", [])
            layer = []
            for i in range(n_x):
                v = float(vals[i]) if i < len(vals) and vals[i] is not None else 0.0
                cumulative[i] += v
                layer.append(cumulative[i])
            stacked_vals.append(layer)
        # Compute range from stacked top
        all_vals = [v for layer in stacked_vals for v in layer]
        all_vals.append(0.0)
    else:
        stacked_vals = None
        all_vals = []
        for s in series:
            for v in s.get("values", []):
                if v is not None:
                    all_vals.append(float(v))
        all_vals.append(0.0)  # area always includes zero baseline

    if not all_vals:
        return

    vmin = min(all_vals)
    vmax = max(all_vals)
    span = vmax - vmin if vmax > vmin else max(abs(vmax), 1.0)
    vmin_p = min(0, vmin - span * 0.05)  # area charts typically start at 0
    vmax_p = vmax + span * 0.08
    ticks, lo, hi = _nice_ticks(vmin_p, vmax_p, target=5)
    if hi == lo:
        hi = lo + 1

    def y_to_emu(v):
        frac = (v - lo) / (hi - lo)
        return plot_y + plot_h - frac * plot_h

    if n_x == 1:
        def x_to_emu(i):
            return plot_x + plot_w / 2
    else:
        step_x = plot_w / (n_x - 1)
        def x_to_emu(i):
            return plot_x + i * step_x

    # Gridlines
    for tv in ticks:
        yy = y_to_emu(tv)
        _add_line(slide, plot_x, yy, plot_x + plot_w, yy, muted, 0.5)

    # Y tick labels
    tick_label_w = int(left_pad * 0.95)
    tick_label_h = int(Pt(tick_pt).emu * 1.4)
    for tv in ticks:
        yy = y_to_emu(tv)
        _add_textbox(slide,
                     plot_x - tick_label_w - int(pad * 0.15),
                     yy - tick_label_h / 2,
                     tick_label_w, tick_label_h,
                     _fmt_num(tv), font_mono, tick_pt, text_c,
                     bold=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Left axis hairline
    _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, muted, 0.75)
    # Baseline
    _add_line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h, muted, 0.75)

    # X tick labels
    max_ticks = 13
    skip = max(1, (n_x + max_ticks - 1) // max_ticks)
    for i, lab in enumerate(x_labels):
        if i % skip != 0 and i != n_x - 1:
            continue
        xx = x_to_emu(i)
        label_w = int(plot_w / max(n_x, 1) * 1.6)
        # Constrain label position to canvas bounds
        label_x = xx - label_w / 2
        label_x = max(label_x, x0)  # don't go left of canvas
        label_x = min(label_x, x0 + w0 - label_w)  # don't go right of canvas
        _add_textbox(slide,
                     label_x,
                     plot_y + plot_h + int(pad * 0.15),
                     label_w, x_tick_h,
                     str(lab), font_body, tick_pt, text_c,
                     bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # X-axis label
    if x_label:
        _add_textbox(slide,
                     plot_x, plot_y + plot_h + x_tick_h + int(pad * 0.2),
                     plot_w, x_axis_label_h,
                     x_label, font_body, int(base_pt * 0.85), text_c,
                     bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Draw filled areas and lines
    baseline_y = y_to_emu(0)

    if stacked:
        # Draw from top series to bottom (so bottom series is on top visually layered correctly)
        # Actually, draw from top of stack down so lower fills don't cover upper
        for idx in range(n_series - 1, -1, -1):
            color = series_colors[idx]
            layer_vals = stacked_vals[idx]
            # Bottom boundary: either previous layer or zero
            if idx == 0:
                bottom_vals = [0.0] * n_x
            else:
                bottom_vals = stacked_vals[idx - 1]

            # Build polygon: top line left-to-right, then bottom line right-to-left
            poly_pts = []
            for i in range(n_x):
                poly_pts.append((x_to_emu(i), y_to_emu(layer_vals[i])))
            for i in range(n_x - 1, -1, -1):
                poly_pts.append((x_to_emu(i), y_to_emu(bottom_vals[i])))

            _add_freeform_fill(slide, poly_pts, color, alpha_pct=50)

            # Top line
            for i in range(n_x - 1):
                _add_line(slide,
                          x_to_emu(i), y_to_emu(layer_vals[i]),
                          x_to_emu(i + 1), y_to_emu(layer_vals[i + 1]),
                          color, 1.75)
    else:
        # Single or overlapping series
        for idx in range(n_series - 1, -1, -1):
            s = series[idx]
            values = s.get("values", [])
            color = series_colors[idx]

            # Build valid points
            pts = []
            for i in range(min(len(values), n_x)):
                v = values[i]
                if v is not None:
                    pts.append((i, float(v)))

            if not pts:
                continue

            # Build polygon for fill
            poly_pts = []
            for i, v in pts:
                poly_pts.append((x_to_emu(i), y_to_emu(v)))
            # Close down to baseline
            poly_pts.append((x_to_emu(pts[-1][0]), baseline_y))
            poly_pts.append((x_to_emu(pts[0][0]), baseline_y))

            alpha = 35 if n_series > 1 else 45
            _add_freeform_fill(slide, poly_pts, color, alpha_pct=alpha)

            # Draw line on top
            for a, b in zip(pts, pts[1:]):
                _add_line(slide,
                          x_to_emu(a[0]), y_to_emu(a[1]),
                          x_to_emu(b[0]), y_to_emu(b[1]),
                          color, 2.0)

    # Legend (compact, above plot area)
    if n_series > 1:
        legend_y = plot_y - legend_h
        cx = plot_x + plot_w
        swatch_w = int(Pt(tick_pt).emu * 0.9)
        gap = int(Pt(tick_pt).emu * 0.4)
        left_limit = x0 + pad
        for idx in range(n_series - 1, -1, -1):
            s = series[idx]
            name = s.get("name", f"Series {idx + 1}")
            color = series_colors[idx]
            est_w = int(Pt(tick_pt).emu * 0.55 * max(len(name), 3))
            item_total_w = est_w + gap + swatch_w + gap * 2
            # Bounds check: stop if this item would go off-slide
            if cx - item_total_w < left_limit:
                break
            cx -= est_w
            _add_textbox(slide,
                         cx, legend_y,
                         est_w, int(Pt(tick_pt).emu * 1.4),
                         name, font_body, tick_pt, text_c,
                         bold=(idx == 0), align=PP_ALIGN.LEFT,
                         anchor=MSO_ANCHOR.MIDDLE)
            cx -= gap
            sy = legend_y + int(Pt(tick_pt).emu * 0.7)
            _add_line(slide, cx - swatch_w, sy, cx, sy, color, 2.0)
            cx -= swatch_w + gap * 2
