"""Line chart — native python-pptx shapes only.

Exposes a single public function: render(slide, data, tokens, bounds).
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Emu(int(x1)), Emu(int(y1)),
                                      Emu(int(x2)), Emu(int(y2)))
    line = conn.line
    line.color.rgb = _rgb(color_hex)
    line.width = Pt(weight_pt)
    if dash is not None:
        try:
            line.dash_style = dash
        except Exception:
            pass
    return conn


def _add_filled_circle(slide, cx, cy, r_emu, color_hex, line_color=None, line_weight_pt=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
                                 Emu(int(r_emu * 2)), Emu(int(r_emu * 2)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_color)
        shp.line.width = Pt(line_weight_pt)
    shp.shadow.inherit = False
    return shp


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y)),
                                 Emu(int(w)), Emu(int(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _nice_ticks(vmin, vmax, target=5):
    if vmax <= vmin:
        vmax = vmin + 1
    span = vmax - vmin
    import math
    raw = span / max(target, 1)
    mag = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    for mult in (1, 2, 2.5, 5, 10):
        step = mult * mag
        if span / step <= target * 1.5:
            break
    lo = math.floor(vmin / step) * step
    hi = math.ceil(vmax / step) * step
    ticks = []
    v = lo
    while v <= hi + 1e-9:
        ticks.append(round(v, 6))
        v += step
    return ticks, lo, hi


def _fmt_num(v):
    if abs(v - round(v)) < 1e-6:
        return f"{int(round(v))}"
    return f"{v:.1f}"


def render(slide, data, tokens, bounds):
    x0, y0, w0, h0 = bounds
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # Background
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    title = data.get("title")
    x_labels = list(data.get("x_labels", []))
    series = list(data.get("series", []))
    x_label = data.get("x_label")
    y_label = data.get("y_label")
    emphasize_last = bool(data.get("emphasize_last_series", False))
    end_labels = bool(data.get("end_labels", False))

    # Layout: top title area, optional y-label column on left, x-label row at bottom
    pad = int(min(w0, h0) * 0.035)
    title_h = Emu(0)
    if title:
        title_h = int(Pt(base_pt * 1.55).emu * 1.8)
    sub_h = 0
    if y_label:
        sub_h = int(Pt(base_pt * 0.85).emu * 1.6)

    # Title
    cursor_y = y0 + pad
    if title:
        _add_textbox(slide, x0 + pad, cursor_y, w0 - 2 * pad, title_h,
                     title, font_display, int(base_pt * 1.55), text,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += title_h + int(pad * 0.3)

    if y_label:
        _add_textbox(slide, x0 + pad, cursor_y, w0 - 2 * pad, sub_h,
                     y_label, font_body, int(base_pt * 0.85), text,
                     bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += sub_h

    # Plot area
    tick_pt = max(8, int(base_pt * 0.8))
    x_axis_label_h = int(Pt(tick_pt).emu * 1.8)
    x_tick_h = int(Pt(tick_pt).emu * 1.6)

    # Reserve space on the right for end-of-line labels
    right_pad = int(w0 * 0.14) if end_labels else int(w0 * 0.03)
    # Reserve left for y tick labels
    left_pad = int(w0 * 0.08)
    bottom_pad = x_tick_h + (x_axis_label_h if x_label else 0) + int(pad * 0.5)

    plot_x = x0 + pad + left_pad
    plot_y = cursor_y + int(pad * 0.2)
    plot_w = (x0 + w0 - pad - right_pad) - plot_x
    plot_h = (y0 + h0 - pad - bottom_pad) - plot_y

    if plot_w <= 0 or plot_h <= 0:
        return

    # Compute y range
    all_vals = []
    for s in series:
        for v in s.get("values", []):
            if v is not None:
                all_vals.append(float(v))
    if not all_vals:
        return
    vmin = min(all_vals)
    vmax = max(all_vals)
    # Pad y range slightly
    span = vmax - vmin if vmax > vmin else max(abs(vmax), 1.0)
    vmin_p = vmin - span * 0.08
    vmax_p = vmax + span * 0.08
    ticks, lo, hi = _nice_ticks(vmin_p, vmax_p, target=5)
    if hi == lo:
        hi = lo + 1

    def y_to_emu(v):
        frac = (v - lo) / (hi - lo)
        return plot_y + plot_h - frac * plot_h

    # X positions: evenly spaced across x_labels
    n_x = max(len(x_labels), 1)
    if n_x == 1:
        def x_to_emu(i):
            return plot_x + plot_w / 2
    else:
        step_x = plot_w / (n_x - 1)
        def x_to_emu(i):
            return plot_x + i * step_x

    # Gridlines (horizontal, subtle)
    for tv in ticks:
        yy = y_to_emu(tv)
        _add_line(slide, plot_x, yy, plot_x + plot_w, yy, muted, 0.5)

    # Y tick labels (on left of plot)
    tick_label_w = int(left_pad * 0.95)
    tick_label_h = int(Pt(tick_pt).emu * 1.4)
    for tv in ticks:
        yy = y_to_emu(tv)
        _add_textbox(slide,
                     plot_x - tick_label_w - int(pad * 0.15),
                     yy - tick_label_h / 2,
                     tick_label_w, tick_label_h,
                     _fmt_num(tv), font_mono, tick_pt, text,
                     bold=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Left axis hairline
    _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, muted, 0.75)
    # Baseline (bottom axis)
    _add_line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h, muted, 0.75)

    # X tick labels
    max_ticks = 13
    skip = max(1, (n_x + max_ticks - 1) // max_ticks)
    for i, lab in enumerate(x_labels):
        if i % skip != 0 and i != n_x - 1:
            continue
        xx = x_to_emu(i)
        label_w = int(plot_w / max(n_x, 1) * 1.6)
        _add_textbox(slide,
                     xx - label_w / 2,
                     plot_y + plot_h + int(pad * 0.15),
                     label_w, x_tick_h,
                     str(lab), font_body, tick_pt, text,
                     bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # X-axis label
    if x_label:
        _add_textbox(slide,
                     plot_x, plot_y + plot_h + x_tick_h + int(pad * 0.2),
                     plot_w, x_axis_label_h,
                     x_label, font_body, int(base_pt * 0.85), text,
                     bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Draw series
    n_series = len(series)
    dash_styles = [None, None]
    try:
        from pptx.enum.dml import MSO_LINE
        dash_styles = [MSO_LINE.DASH, MSO_LINE.ROUND_DOT, MSO_LINE.LONG_DASH, MSO_LINE.DASH_DOT]
    except Exception:
        dash_styles = [None, None, None, None]

    for idx, s in enumerate(series):
        values = s.get("values", [])
        name = s.get("name", f"Series {idx+1}")
        is_primary_idx = (idx == n_series - 1) if emphasize_last else (idx == 0)

        if n_series == 1:
            color = primary
            weight = 2.25
            dash = None
        elif emphasize_last:
            if is_primary_idx:
                color = primary
                weight = 2.5
                dash = None
            else:
                color = muted
                weight = 1.25
                # dash pattern for variety on non-emphasized series
                ds_idx = (n_series - 2 - idx) % max(len(dash_styles) - 1, 1)
                dash = dash_styles[ds_idx + 1] if len(dash_styles) > 1 else None
        else:
            if idx == 0:
                color = primary
                weight = 2.25
                dash = None
            else:
                color = muted
                weight = 1.25
                ds_idx = (idx - 1) % max(len(dash_styles) - 1, 1)
                dash = dash_styles[ds_idx + 1] if len(dash_styles) > 1 else None

        # Build list of valid points
        pts = []
        for i, v in enumerate(values):
            if v is None or i >= n_x:
                continue
            pts.append((x_to_emu(i), y_to_emu(float(v)), i, float(v)))

        # Draw segments
        for a, b in zip(pts, pts[1:]):
            _add_line(slide, a[0], a[1], b[0], b[1], color, weight, dash=dash)

        # Highlight last data point of primary/emphasized series with accent
        if is_primary_idx and pts:
            last = pts[-1]
            r = max(Emu(Pt(3).emu).emu, int(Pt(3.5).emu))
            # Outer dot in accent
            _add_filled_circle(slide, last[0], last[1], int(Pt(4).emu), accent)
            # Inner white-ish in bg for contrast
            _add_filled_circle(slide, last[0], last[1], int(Pt(1.6).emu), bg)

        # End-of-line labels
        if end_labels and pts:
            last = pts[-1]
            lbl_w = int(right_pad * 0.95)
            lbl_h = int(Pt(tick_pt).emu * 1.6)
            _add_textbox(slide,
                         last[0] + int(Pt(4).emu),
                         last[1] - lbl_h / 2,
                         lbl_w, lbl_h,
                         str(name), font_body, tick_pt,
                         text,
                         bold=is_primary_idx, align=PP_ALIGN.LEFT,
                         anchor=MSO_ANCHOR.MIDDLE)

    # If not end_labels, render a compact legend at top-right of plot area
    if not end_labels and n_series > 1:
        legend_y = plot_y - int(Pt(tick_pt).emu * 1.5)
        if legend_y < y0 + pad:
            legend_y = plot_y + int(Pt(tick_pt).emu * 0.2)
        cx = plot_x + plot_w
        swatch_w = int(Pt(tick_pt).emu * 0.9)
        gap = int(Pt(tick_pt).emu * 0.4)
        # Build legend right-to-left
        for idx in range(n_series - 1, -1, -1):
            s = series[idx]
            name = s.get("name", f"Series {idx+1}")
            is_primary_idx = (idx == n_series - 1) if emphasize_last else (idx == 0)
            color = primary if (is_primary_idx or n_series == 1) else muted
            # Estimate label width
            est_w = int(Pt(tick_pt).emu * 0.55 * max(len(name), 3))
            cx -= est_w
            _add_textbox(slide,
                         cx, legend_y,
                         est_w, int(Pt(tick_pt).emu * 1.4),
                         name, font_body, tick_pt, text,
                         bold=is_primary_idx, align=PP_ALIGN.LEFT,
                         anchor=MSO_ANCHOR.MIDDLE)
            cx -= gap
            # Swatch line
            sy = legend_y + int(Pt(tick_pt).emu * 0.7)
            _add_line(slide, cx - swatch_w, sy, cx, sy, color, 2.0)
            cx -= swatch_w + gap * 2
