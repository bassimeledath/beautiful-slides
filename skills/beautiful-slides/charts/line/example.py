"""Run `render()` 5 times — once per mode — writing example-<mode>.pptx per run."""

import os
import sys

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module
from tokens import MODES  # noqa: E402
from render import render  # noqa: E402


DATA = {
    "title": "Revenue retention by cohort, first 12 months",
    "x_labels": ["M0", "M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8", "M9", "M10", "M11", "M12"],
    "series": [
        {"name": "Q1 '25", "values": [100, 98, 97, 98, 100, 101, 102, 103, 104, 104, 105, 105, 106]},
        {"name": "Q2 '25", "values": [100, 99, 100, 102, 104, 105, 106, 106, 107, 108, 109, 110, 110]},
        {"name": "Q3 '25", "values": [100, 102, 106, 109, 110, 112, 112, 113, 113]},
        {"name": "Q4 '25", "values": [100, 104, 111, 118]},
    ],
    "x_label": "Months from cohort start",
    "y_label": "Revenue retention (%)",
    "emphasize_last_series": True,
    "end_labels": True,
}


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    for mode_name, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        margin = Inches(0.4)
        bounds = (
            int(margin),
            int(margin),
            int(prs.slide_width - 2 * margin),
            int(prs.slide_height - 2 * margin),
        )
        render(slide, DATA, tokens, bounds)

        out_path = os.path.join(out_dir, f"example-{mode_name}.pptx")
        prs.save(out_path)
        print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
