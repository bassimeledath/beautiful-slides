"""Donut chart -- ring with proportional arc segments and a hero center label.

Strategy: draw each segment as a filled pie wedge (freeform polygon: center ->
arc points -> center, all in EMU).  Then overlay an oval in bg color to punch
out the donut hole.  Finally add centered hero text and a legend below.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _luminance(hex_):
    h = hex_.lstrip("#")
    r = int(h[0:2], 16) / 255.0
    g = int(h[2:4], 16) / 255.0
    b = int(h[4:6], 16) / 255.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _darken_hex(hex_, factor=0.3):
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r * (1 - factor))
    g = int(g * (1 - factor))
    b = int(b * (1 - factor))
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


# ---------------------------------------------------------------------------
# Pie-wedge drawing (freeform polygon)
# ---------------------------------------------------------------------------

def _draw_pie_wedge(slide, cx, cy, r, start_rad, end_rad, fill_hex,
                    n_points=72):
    """Draw a pie-wedge as a freeform polygon.

    All coordinates are in absolute EMU.  python-pptx's build_freeform
    auto-computes the bounding box and converts to relative path coords.

    CRITICAL: all points must be passed in a **single** add_line_segments
    call.  Multiple calls each inject a spurious <a:close/> element into
    the path XML, which splits the polygon into disconnected sub-paths and
    causes the wedge to render incorrectly (the bug that plagued earlier
    versions of this file).
    """
    sweep = end_rad - start_rad
    if abs(sweep) < 1e-9:
        return None

    steps = max(12, int(abs(sweep) / (2 * math.pi) * n_points))

    # Build arc points along the outer edge
    arc_pts = []
    for i in range(steps + 1):
        a = start_rad + sweep * i / steps
        px = int(round(cx + r * math.cos(a)))
        py = int(round(cy + r * math.sin(a)))
        arc_pts.append((px, py))

    icx, icy = int(round(cx)), int(round(cy))

    # Single list: arc points followed by return to center
    all_pts = arc_pts + [(icx, icy)]

    ff = slide.shapes.build_freeform(icx, icy, scale=1.0)
    ff.add_line_segments(all_pts, close=True)

    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape


# ---------------------------------------------------------------------------
# Main render
# ---------------------------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a donut chart with 1-3 segments and a hero center label.

    Parameters
    ----------
    slide : pptx.slide.Slide
    data  : dict
        segments     : list of {label, value}
        center_value : str -- hero number in center
        center_label : str -- subtitle below hero number
        title        : str -- optional title above chart
    tokens : dict
    bounds : tuple (x, y, w, h) in EMU
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = int(tokens["font_size_base_pt"])

    # --- Background rectangle ---
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()

    # --- Contrast safety ---
    if _luminance(bg) > 0.7 and _luminance(text_c) > 0.7:
        text_c = "#333333"
    if _luminance(bg) > 0.7 and _luminance(muted) > 0.7:
        muted = _darken_hex(muted, 0.45)

    segments = list(data.get("segments", []))[:3]
    center_value = data.get("center_value", "")
    center_label = data.get("center_label", "")
    title = data.get("title")

    if not segments:
        return

    seg_colors = [primary, accent, _darken_hex(muted, 0.3)]

    total = sum(s.get("value", 0) for s in segments)
    if total <= 0:
        return

    # --- Layout ---
    pad = int(min(w, h) * 0.04)
    cursor_y = y + pad

    # Title
    title_h = 0
    if title:
        title_pt = int(base_pt * 1.5)
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, x + pad, cursor_y, w - 2 * pad, title_h,
                     title, font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += title_h + int(Pt(base_pt * 0.4).emu)

    avail_h = (y + h - pad) - cursor_y
    avail_w = w - 2 * pad

    # Legend reservation
    legend_line_h = int(Pt(base_pt).emu * 1.8)
    n_legend = sum(1 for s in segments if s.get("value", 0) > 0)
    legend_total_h = legend_line_h * n_legend + int(Pt(base_pt * 0.3).emu)

    donut_area_h = avail_h - legend_total_h
    donut_area_w = avail_w

    # Ring sizing
    ring_diameter = min(donut_area_w, donut_area_h) * 0.80
    outer_r = ring_diameter / 2
    inner_r = outer_r * 0.58

    ring_cx = x + w / 2
    ring_cy = cursor_y + donut_area_h / 2

    # Gap between segments (small visual separator)
    gap_rad = 0.03 if len(segments) > 1 else 0.0

    # --- Draw pie-wedge segments (start from top, -pi/2) ---
    current_angle = -math.pi / 2
    for i, seg in enumerate(segments):
        val = seg.get("value", 0)
        if val <= 0:
            continue
        sweep = (val / total) * 2 * math.pi
        start = current_angle + gap_rad / 2
        end = current_angle + sweep - gap_rad / 2
        if end <= start:
            end = start + 0.01
        color = seg_colors[i % len(seg_colors)]
        _draw_pie_wedge(slide, ring_cx, ring_cy, outer_r, start, end, color)
        current_angle += sweep

    # --- Donut hole (bg-colored circle) ---
    hole_d = int(inner_r * 2)
    hole_left = int(ring_cx - inner_r)
    hole_top = int(ring_cy - inner_r)
    hole_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(hole_left), Emu(hole_top),
        Emu(hole_d), Emu(hole_d),
    )
    hole_shape.fill.solid()
    hole_shape.fill.fore_color.rgb = _rgb(bg)
    hole_shape.line.fill.background()
    try:
        hole_shape.shadow.inherit = False
    except Exception:
        pass

    # --- Center hero text ---
    center_box_w = inner_r * 2 * 0.85
    center_box_h = inner_r * 2 * 0.75
    center_x = ring_cx - center_box_w / 2
    center_y_pos = ring_cy - center_box_h / 2

    if center_value:
        hero_pt = max(24, min(int((inner_r * 2 / Pt(1).emu) * 0.28), 72))
        est_w = Pt(hero_pt).emu * 0.55 * max(1, len(center_value))
        if est_w > center_box_w:
            hero_pt = max(18, int(hero_pt * center_box_w / est_w))

        hero_h = Pt(hero_pt).emu * 1.3
        if center_label:
            hero_y = center_y_pos + center_box_h * 0.1
            _add_textbox(slide, center_x, hero_y, center_box_w, hero_h,
                         center_value, font_display, hero_pt, text_c,
                         bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.BOTTOM)

            label_pt = max(9, int(base_pt * 0.85))
            label_h = Pt(label_pt).emu * 1.4
            _add_textbox(slide, center_x, hero_y + hero_h,
                         center_box_w, label_h,
                         center_label, font_body, label_pt, text_c,
                         bold=False, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.TOP)
        else:
            _add_textbox(slide, center_x, center_y_pos,
                         center_box_w, center_box_h,
                         center_value, font_display, hero_pt, text_c,
                         bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)

    # --- Legend ---
    legend_y = cursor_y + donut_area_h + int(Pt(base_pt * 0.2).emu)
    swatch_size = int(Pt(base_pt * 0.7).emu)
    swatch_gap = int(Pt(base_pt * 0.4).emu)
    legend_x = x + pad

    legend_idx = 0
    for i, seg in enumerate(segments):
        val = seg.get("value", 0)
        if val <= 0:
            continue
        item_y = legend_y + legend_idx * legend_line_h
        legend_idx += 1
        color = seg_colors[i % len(seg_colors)]

        swatch_y = item_y + (legend_line_h - swatch_size) // 2
        sw = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(int(legend_x)), Emu(int(swatch_y)),
            Emu(swatch_size), Emu(swatch_size),
        )
        sw.fill.solid()
        sw.fill.fore_color.rgb = _rgb(color)
        sw.line.fill.background()
        try:
            sw.shadow.inherit = False
        except Exception:
            pass

        lbl_x = legend_x + swatch_size + swatch_gap
        lbl_w = avail_w - swatch_size - swatch_gap
        label_text = seg.get("label", f"Segment {i + 1}")
        pct = (val / total * 100) if total > 0 else 0
        display = f"{label_text}  {pct:.0f}%"
        _add_textbox(slide, lbl_x, item_y, lbl_w, legend_line_h,
                     display, font_body, base_pt, text_c,
                     bold=False, align=PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)

    return None
