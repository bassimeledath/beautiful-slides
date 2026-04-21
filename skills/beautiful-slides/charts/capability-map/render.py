"""Capability Map — tiled/grouped view of business capabilities by domain."""

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(3)))
    tf.margin_right = Emu(int(Pt(3)))
    tf.margin_top = Emu(int(Pt(2)))
    tf.margin_bottom = Emu(int(Pt(2)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def render(slide, data, tokens, bounds):
    """Render a capability map (grouped tiles by domain).

    data:
        title    - optional string
        domains  - list of {"name": str, "capabilities": [str, ...]}
                   Each domain becomes a column with a colored header and
                   capability tiles stacked below.
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]
    radius_px = tokens.get("radius_px", 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    domains = data.get("domains", [])
    if not domains:
        return
    title = data.get("title")

    n_domains = len(domains)

    # --- layout ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_h = (y + h) - cur_y
    avail_w = w

    # Column layout: equal-width columns with gaps
    col_gap = max(int(avail_w * 0.015), int(Pt(base_pt * 0.4)))
    total_gaps = col_gap * (n_domains - 1)
    col_w = (avail_w - total_gaps) // n_domains

    # Header height
    header_h = int(Pt(base_pt) * 2.5)
    tile_area_y = cur_y + header_h + int(Pt(base_pt * 0.3))
    tile_area_h = avail_h - header_h - int(Pt(base_pt * 0.3))

    # Find max capabilities across domains to compute tile height
    max_caps = max(len(d.get("capabilities", [])) for d in domains)
    tile_gap = max(int(Pt(base_pt * 0.3)), int(avail_w * 0.005))
    tile_h = max(
        int((tile_area_h - tile_gap * (max_caps - 1)) / max_caps) if max_caps > 0 else int(Pt(base_pt * 2)),
        int(Pt(base_pt * 1.6)),
    )

    for di, domain in enumerate(domains):
        col_x = x + di * (col_w + col_gap)
        capabilities = domain.get("capabilities", [])
        domain_name = domain.get("name", "")

        # Color for this domain: interpolate across primary -> accent
        t = di / (n_domains - 1) if n_domains > 1 else 0.0
        domain_color = _lerp_hex(primary, accent, t)
        tile_bg = _lerp_hex(bg, domain_color, 0.12)
        tile_border_color = _lerp_hex(bg, domain_color, 0.35)

        # Domain header
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE
        header = slide.shapes.add_shape(
            shape_type,
            Emu(col_x), Emu(cur_y),
            Emu(col_w), Emu(header_h),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _rgb(domain_color)
        header.line.fill.background()

        # Header text
        header_font_size = max(int(base_pt * 0.85), 8)
        _add_textbox(
            slide, col_x, cur_y, col_w, header_h,
            domain_name, font_display, header_font_size, bg,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Capability tiles
        for ci, cap in enumerate(capabilities):
            tile_y = tile_area_y + ci * (tile_h + tile_gap)

            # Clamp tile to bounds
            if tile_y + tile_h > y + h:
                tile_h_actual = max((y + h) - tile_y, 0)
                if tile_h_actual <= 0:
                    break
            else:
                tile_h_actual = tile_h

            # Tile background
            tile_shape = slide.shapes.add_shape(
                shape_type,
                Emu(col_x), Emu(tile_y),
                Emu(col_w), Emu(tile_h_actual),
            )
            tile_shape.fill.solid()
            tile_shape.fill.fore_color.rgb = _rgb(tile_bg)
            tile_shape.line.color.rgb = _rgb(tile_border_color)
            tile_shape.line.width = Pt(0.75)

            # Tile text
            tile_font_size = max(int(base_pt * 0.7), 7)
            _add_textbox(
                slide, col_x, tile_y, col_w, tile_h_actual,
                cap, font_body, tile_font_size, text_c,
                align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE,
            )

    return None
