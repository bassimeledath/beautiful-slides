from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    t = max(0.0, min(1.0, t))
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _luminance(hex_):
    h = hex_.lstrip("#")
    r = int(h[0:2], 16) / 255.0
    g = int(h[2:4], 16) / 255.0
    b = int(h[4:6], 16) / 255.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _set_text(tf, text, font_name, size_pt, color_hex, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        for r in list(p.runs):
            r.text = ""
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)


def _add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_w_emu=0, radius_emu=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_emu > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    if radius_emu > 0:
        try:
            short = min(w, h)
            adj = max(0.0, min(0.5, (radius_emu / short)))
            shp.adjustments[0] = adj
        except Exception:
            pass
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None or line_w_emu <= 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Emu(int(line_w_emu))
    shp.text_frame.text = ""
    tf = shp.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return shp


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    _set_text(tb.text_frame, text, font_name, size_pt, color_hex, bold=bold, align=align, anchor=anchor)
    return tb


def render(slide, data, tokens, bounds):
    """Render a risk heat matrix (impact vs likelihood).

    A grid where rows = impact levels (high at top), columns = likelihood levels
    (low at left). Cells are colored by severity (green -> yellow -> red via
    token tints). Named risks are placed as labels in their respective cells.
    """
    x, y, w, h = bounds

    # Grid dimensions (default 5x5)
    grid_size = int(data.get("grid_size", 5))
    impact_labels = list(data.get("impact_labels", []))
    likelihood_labels = list(data.get("likelihood_labels", []))
    risks = list(data.get("risks", []))
    title = data.get("title")
    x_axis_label = data.get("x_axis_label", "Likelihood")
    y_axis_label = data.get("y_axis_label", "Impact")

    # Default labels if not provided
    if not impact_labels:
        if grid_size == 3:
            impact_labels = ["High", "Medium", "Low"]
        else:
            impact_labels = ["Critical", "High", "Medium", "Low", "Negligible"][:grid_size]
    if not likelihood_labels:
        if grid_size == 3:
            likelihood_labels = ["Low", "Medium", "High"]
        else:
            likelihood_labels = ["Rare", "Unlikely", "Possible", "Likely", "Almost Certain"][:grid_size]

    n_rows = len(impact_labels)
    n_cols = len(likelihood_labels)

    # Extract tokens
    bg_hex = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(bg_hex.lstrip("#"))
    primary_hex = tokens["primary"]
    accent_hex = tokens["accent"]
    text_hex = tokens["text"]
    muted_hex = tokens["muted"]
    font_body = tokens["font_body"]
    font_display = tokens["font_display"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg_hex)
    bg_shape.line.fill.background()

    radius_emu = radius_px * 9525

    # Font sizes
    label_pt = max(7, int(round(base_pt * 0.72)))
    title_pt = max(base_pt + 2, int(round(base_pt * 1.15)))
    risk_label_pt = max(6, int(round(base_pt * 0.50)))
    axis_label_pt = max(7, int(round(base_pt * 0.65)))

    # Build severity color for each cell.
    # Severity = (row_from_top_inverted + col) normalized.
    # Row 0 = highest impact. Col n-1 = highest likelihood.
    # severity_score = impact_rank + likelihood_rank, both 0-based from low.
    # impact_rank: row 0 = n_rows-1 (highest), row n_rows-1 = 0 (lowest)
    # likelihood_rank: col 0 = 0 (lowest), col n_cols-1 = n_cols-1 (highest)
    max_severity = (n_rows - 1) + (n_cols - 1)

    # Three severity tiers derived from tokens:
    #   low  = accent tinted toward bg (green-ish feel)
    #   med  = muted lightened (amber/yellow feel via lerp primary+accent)
    #   high = primary saturated (red/intense feel)
    # We create a smooth gradient across severity levels.
    low_hex = accent_hex       # green end
    mid_hex = _lerp_hex(primary_hex, accent_hex, 0.5)  # blended mid
    high_hex = primary_hex     # intense end

    def _severity_color(score):
        """Map 0..max_severity to a color from low(green) -> mid -> high(red)."""
        if max_severity == 0:
            return mid_hex
        t = score / max_severity
        if t <= 0.5:
            return _lerp_hex(low_hex, mid_hex, t * 2)
        else:
            return _lerp_hex(mid_hex, high_hex, (t - 0.5) * 2)

    # Layout
    title_h = 0
    if title:
        title_h = int(Pt(title_pt * 1.6))

    # Axis label space — use a horizontal subtitle row instead of a narrow
    # vertical strip (which produces an unreadable character-per-line stack).
    y_axis_w = int(Pt(base_pt * 5.0))  # wide enough for horizontal text
    x_axis_h = int(Pt(base_pt * 1.2))

    # Row/col label space
    row_label_w = int(Pt(base_pt * 3.0))
    for lbl in impact_labels:
        est = int(Pt(base_pt * 0.55 * max(3, len(str(lbl)))))
        if est > row_label_w:
            row_label_w = est
    row_label_w = min(row_label_w, int(w * 0.22))

    col_label_h = int(Pt(base_pt * 1.4))

    inner_x = x + y_axis_w
    inner_y = y + title_h
    inner_w = w - y_axis_w
    inner_h = h - title_h - x_axis_h

    grid_x = inner_x + row_label_w
    grid_y = inner_y + col_label_h
    grid_w = inner_w - row_label_w
    grid_h = inner_h - col_label_h

    gap = max(1, int(Emu(0.5 * 9525)))
    cell_w = (grid_w - gap * (n_cols - 1)) / n_cols
    cell_h = (grid_h - gap * (n_rows - 1)) / n_rows

    # Title
    if title:
        _add_textbox(
            slide,
            x, y, w, title_h,
            title,
            font_display,
            title_pt,
            text_hex,
            bold=True,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )

    # Y-axis label — rendered as a rotated textbox so the text reads
    # bottom-to-top instead of stacking one character per line.
    yaxis_tb = _add_textbox(
        slide,
        int(x), int(inner_y + col_label_h), int(y_axis_w), int(grid_h),
        y_axis_label,
        font_body,
        axis_label_pt,
        text_hex,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    yaxis_tb.rotation = 270.0

    # X-axis label
    _add_textbox(
        slide,
        int(grid_x), int(inner_y + inner_h), int(grid_w), int(x_axis_h),
        x_axis_label,
        font_body,
        axis_label_pt,
        text_hex,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Column headers (likelihood labels)
    for ci, cl in enumerate(likelihood_labels):
        cx = grid_x + ci * (cell_w + gap)
        _add_textbox(
            slide,
            int(cx), int(inner_y), int(cell_w), int(col_label_h),
            cl,
            font_body,
            label_pt,
            text_hex,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Row labels (impact labels — high at top)
    for ri, rl in enumerate(impact_labels):
        ry = grid_y + ri * (cell_h + gap)
        pad_right = int(Pt(base_pt * 0.3))
        _add_textbox(
            slide,
            int(inner_x), int(ry), int(row_label_w - pad_right), int(cell_h),
            rl,
            font_body,
            label_pt,
            text_hex,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Build a lookup: (row, col) -> list of risk names
    cell_risks = {}
    for risk in risks:
        ri = int(risk.get("impact", 0))
        ci = int(risk.get("likelihood", 0))
        ri = max(0, min(n_rows - 1, ri))
        ci = max(0, min(n_cols - 1, ci))
        cell_risks.setdefault((ri, ci), []).append(risk.get("name", ""))

    # Draw cells
    for ri in range(n_rows):
        impact_rank = (n_rows - 1) - ri  # row 0 = highest impact
        for ci in range(n_cols):
            likelihood_rank = ci
            severity = impact_rank + likelihood_rank
            fill_hex = _severity_color(severity)

            cx = grid_x + ci * (cell_w + gap)
            cy = grid_y + ri * (cell_h + gap)

            # Lighten the fill by lerping toward bg for a tinted look
            fill_hex = _lerp_hex(bg_hex, fill_hex, 0.55)

            _add_rect(
                slide,
                cx, cy, cell_w, cell_h,
                fill_hex,
                line_hex=muted_hex,
                line_w_emu=int(Emu(0.25 * 9525)),
                radius_emu=radius_emu,
            )

            # Place risk labels in the cell (truncate to fit)
            names = cell_risks.get((ri, ci), [])
            if names:
                # Estimate max chars that fit per line in the cell
                max_chars_per_line = max(8, int(cell_w / (Pt(risk_label_pt) * 0.6)))
                max_lines = max(1, int(cell_h / (Pt(risk_label_pt) * 1.6)))
                truncated = []
                for nm in names[:max_lines]:
                    if len(nm) > max_chars_per_line:
                        nm = nm[:max_chars_per_line - 1] + "\u2026"
                    truncated.append(nm)
                label_text = "\n".join(truncated)
                cell_lum = _luminance(fill_hex)
                bg_lum = _luminance(bg_hex)
                if bg_lum < 0.5:
                    label_color = text_hex if cell_lum < 0.55 else bg_hex
                else:
                    label_color = text_hex if cell_lum > 0.35 else bg_hex
                _add_textbox(
                    slide,
                    int(cx), int(cy), int(cell_w), int(cell_h),
                    label_text,
                    font_body,
                    risk_label_pt,
                    label_color,
                    align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE,
                )

    return None
