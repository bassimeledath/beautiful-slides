"""Layered architecture — horizontal stack of layers (top-to-bottom).

Each layer is a full-width rounded rectangle with a label.
Optional items/components listed within each layer.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ---- helpers ---------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    """Linearly interpolate between two hex colors."""
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _lighten_hex(hex_, factor=0.85):
    """Blend a hex color toward white by *factor* (0=original, 1=white)."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02X}{g:02X}{b:02X}"


def _luminance(hex_):
    """Return relative luminance (0-1) of a hex color."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255.0, int(h[2:4], 16) / 255.0, int(h[4:6], 16) / 255.0
    r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
    g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
    b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_text(fill_hex, light="#FFFFFF", dark="#1A1A1A"):
    """Return dark text if fill is light, light text if fill is dark."""
    return dark if _luminance(fill_hex) > 0.35 else light


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "\u2026"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                 word_wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Emu(int(Pt(4).emu))
    tf.margin_right = Emu(int(Pt(4).emu))
    tf.margin_top = Emu(int(Pt(2).emu))
    tf.margin_bottom = Emu(int(Pt(2).emu))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_rounded_rect(slide, x, y, w, h, fill_hex, line_hex=None,
                      line_width_pt=0.75, radius_px=0):
    shape_type = (MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0
                  else MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(
        shape_type,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Pt(line_width_pt)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# ---- main render -----------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a layered architecture diagram onto *slide* inside *bounds*.

    data:
        title  - optional string
        layers - list of {"label": str, "items": [str, ...] (optional)}
                 Layers render top-to-bottom (first = topmost layer).
    """
    x0, y0, w0, h0 = bounds

    # --- unpack tokens -------------------------------------------------------
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px", 0))

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x0)), Emu(int(y0)),
        Emu(int(w0)), Emu(int(h0)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    # --- unpack data ---------------------------------------------------------
    layers = data.get("layers", [])
    if not layers:
        return
    title = data.get("title")
    n = len(layers)

    # --- layout constants ----------------------------------------------------
    pad = int(min(w0, h0) * 0.025)

    # Working area
    ix, iy = x0 + pad, y0 + pad
    iw, ih = w0 - 2 * pad, h0 - 2 * pad

    # --- title ---------------------------------------------------------------
    if title:
        title_pt = int(round(base_pt * 1.5))
        # Estimate line count for long titles and scale title_h accordingly
        est_chars_per_line = max(1, int(iw / (Pt(title_pt).emu * 0.6)))
        est_lines = max(1, (len(title) + est_chars_per_line - 1) // est_chars_per_line)
        title_h = int(Pt(title_pt).emu * 1.8 * min(est_lines, 3))
        _add_textbox(slide, ix, iy, iw, title_h, title,
                     font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + pad
        ih = (y0 + h0 - pad) - iy

    # --- compute layer heights -----------------------------------------------
    layer_gap = int(min(ih * 0.025, Pt(8).emu))
    total_gap = layer_gap * (n - 1)
    layer_h = max(int((ih - total_gap) / n), int(Pt(base_pt * 2)))

    # Ensure we don't overflow bounds
    needed = layer_h * n + total_gap
    if needed > ih:
        layer_h = max(int((ih - total_gap) / n), int(Pt(base_pt * 1.5)))

    # Label column on the left; items area on the right
    # Adaptive label column width based on longest layer name
    max_label_len = max((len(layer.get("label", "")) for layer in layers), default=10)
    label_col_pct = 0.25 if max_label_len <= 20 else min(0.32, 0.25 + 0.005 * (max_label_len - 20))
    label_col_w = int(iw * label_col_pct)
    items_gap = int(iw * 0.02)
    items_area_x = ix + label_col_w + items_gap
    items_area_w = iw - label_col_w - items_gap

    # --- draw layers ---------------------------------------------------------
    cur_y = iy

    for i, layer in enumerate(layers):
        # Color: interpolate from primary (top) to accent (bottom)
        t = i / (n - 1) if n > 1 else 0.0
        layer_color = _lerp_hex(primary, accent, t)
        # Blend layer_color toward bg at 25% — keeps the tint visible on both
        # light and dark backgrounds without washing out.
        layer_bg = _lerp_hex(layer_color, bg, 0.25)
        label_text = layer.get("label", f"Layer {i + 1}")
        items = layer.get("items", [])

        # --- layer background rectangle (full width) -------------------------
        _add_rounded_rect(slide, ix, cur_y, iw, layer_h,
                          layer_bg, line_hex=layer_color,
                          line_width_pt=1.0, radius_px=radius_px)

        # --- left label stripe -----------------------------------------------
        stripe_w = int(Pt(4).emu)
        _add_rounded_rect(slide, ix, cur_y, stripe_w, layer_h,
                          layer_color, radius_px=0)

        # --- layer label text ------------------------------------------------
        label_pt = max(int(base_pt * 1.0), 10)
        label_avail_w = label_col_w - stripe_w - int(Pt(8).emu)
        max_label_chars = max(12, int(label_avail_w / (Pt(label_pt).emu * 0.55)))
        # Luminance check: pick text color that contrasts with layer bg
        layer_label_color = _contrast_text(layer_bg)
        _add_textbox(slide, ix + stripe_w + int(Pt(4).emu), cur_y,
                     label_avail_w, layer_h,
                     _truncate(label_text, max_label_chars),
                     font_display, label_pt, layer_label_color,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        # --- items / components within the layer -----------------------------
        if items:
            n_items = len(items)
            item_gap = int(min(items_area_w * 0.015, Pt(6).emu))
            item_w = max(1, int((items_area_w - item_gap * (n_items - 1))
                                / n_items))
            # Reduce font more aggressively for dense layers
            if n_items >= 6:
                item_pt = max(7, int(base_pt * 0.65))
            else:
                item_pt = max(8, int(base_pt * 0.75))
            item_h_inner = max(int(layer_h * 0.6), int(Pt(item_pt * 2)))
            item_y = cur_y + (layer_h - item_h_inner) // 2

            # Estimate max chars for item labels
            max_item_chars = max(10, int(item_w / (Pt(item_pt).emu * 0.55)))

            for j, item_label in enumerate(items):
                item_x = items_area_x + j * (item_w + item_gap)
                # Clamp to bounds
                if item_x + item_w > x0 + w0 - pad:
                    item_w = max(1, (x0 + w0 - pad) - item_x)
                if item_x + item_w > x0 + w0 - pad:
                    break

                _add_rounded_rect(slide, item_x, item_y, item_w, item_h_inner,
                                  bg, line_hex=muted,
                                  line_width_pt=0.5, radius_px=radius_px)
                item_text_color = _contrast_text(bg)
                _add_textbox(slide, item_x, item_y, item_w, item_h_inner,
                             _truncate(item_label, max_item_chars),
                             font_body, item_pt, item_text_color,
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        cur_y += layer_h + layer_gap

    return None
