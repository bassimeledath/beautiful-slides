"""Market map renderer — native python-pptx shapes only.

Positioned clusters or labeled boxes arranged by segment/category on a grid.
For vendor landscapes, partner ecosystems, whitespace analysis.

Supports two layouts:
- "grid": Items grouped into named categories arranged in a column grid.
- "positioned": Items placed at explicit (x, y) positions with optional
  category backgrounds.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn


# ---- helpers ---------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Emu(int(Pt(3).emu))
    tf.margin_right = Emu(int(Pt(3).emu))
    tf.margin_top = Emu(int(Pt(1).emu))
    tf.margin_bottom = Emu(int(Pt(1).emu))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_rounded_rect(slide, x, y, w, h, fill_hex, line_hex=None,
                      line_width_pt=0.75, radius_px=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Pt(line_width_pt)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(int(x1)), Emu(int(y1)),
        Emu(int(x2)), Emu(int(y2)),
    )
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(weight_pt)
    return conn


def _lighten_hex(hex_, factor=0.85):
    """Blend a hex color toward white by `factor` (0=original, 1=white)."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02X}{g:02X}{b:02X}"


# ---- main render -----------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a market map onto *slide* inside *bounds*, styled by *tokens*."""
    x0, y0, w0, h0 = bounds

    # --- unpack tokens -------------------------------------------------------
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px", 0))

    # --- unpack data ---------------------------------------------------------
    title = data.get("title")
    subtitle = data.get("subtitle")
    categories = list(data.get("categories") or [])
    # Each category: {name: str, color: str (opt), items: [str or {label, color}]}

    if not categories:
        return

    # --- background ----------------------------------------------------------
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    # --- outer padding -------------------------------------------------------
    pad = int(min(w0, h0) * 0.035)
    ix, iy = x0 + pad, y0 + pad
    iw, ih = w0 - 2 * pad, h0 - 2 * pad

    # --- title ---------------------------------------------------------------
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.55))
        # Estimate wrapped lines based on title length vs available width
        char_w_emu = Pt(title_pt).emu * 0.55
        chars_per_line = max(1, int(iw / char_w_emu))
        title_lines = max(1, -(-len(title) // chars_per_line))  # ceil division
        title_h = int(Pt(title_pt).emu * 1.8 * title_lines)
        _add_textbox(slide, ix, iy, iw, title_h, title,
                     font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h
        ih = (y0 + h0 - pad) - iy

    if subtitle:
        sub_pt = max(8, int(base_pt * 0.85))
        sub_h = int(Pt(sub_pt).emu * 1.6)
        _add_textbox(slide, ix, iy, iw, sub_h, subtitle,
                     font_body, sub_pt, text_c,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += sub_h + int(pad * 0.3)
        ih = (y0 + h0 - pad) - iy

    # --- grid layout ---------------------------------------------------------
    # Categories are arranged as columns. Each column has a header + stacked items.
    n_cats = len(categories)
    col_gap = int(min(iw * 0.015, Pt(12).emu))

    # Adaptive header font: scale down when many categories
    if n_cats > 5:
        cat_header_pt = max(7, int(base_pt * 0.9 * 5 / n_cats))
    else:
        cat_header_pt = max(9, int(base_pt * 0.9))

    # Adaptive item font: slightly smaller when many categories
    if n_cats > 5:
        item_pt = max(7, int(base_pt * 0.78) - (n_cats - 5))
    else:
        item_pt = max(8, int(base_pt * 0.78))

    cat_header_h = int(Pt(cat_header_pt).emu * 2.2)
    item_h = int(Pt(item_pt).emu * 2.4)
    item_gap = int(Pt(3).emu)
    item_pad = int(Pt(4).emu)

    total_gap_w = col_gap * (n_cats - 1) if n_cats > 1 else 0
    col_w = max(1, (iw - total_gap_w) / n_cats)

    # Default category colors: cycle through primary and accent
    default_cat_colors = [primary, accent]

    for ci, cat in enumerate(categories):
        cx = ix + ci * (col_w + col_gap)
        cy = iy
        cat_color = cat.get("color") or default_cat_colors[ci % len(default_cat_colors)]
        items = list(cat.get("items") or [])

        # Compute how many items we can fit
        available_h = ih - cat_header_h - int(pad * 0.3)
        max_items = max(1, int(available_h / (item_h + item_gap)))
        # Cap more aggressively when many categories (columns are narrow)
        if n_cats > 5:
            max_items = min(max_items, max(3, 8 - (n_cats - 5)))
        items = items[:max_items]

        # Category header bar
        _add_rounded_rect(slide, cx, cy, col_w, cat_header_h,
                          cat_color, radius_px=radius_px)
        _add_textbox(slide, cx, cy, col_w, cat_header_h,
                     cat.get("name", f"Category {ci+1}"),
                     font_display, cat_header_pt, bg,
                     bold=True, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

        # Category column background (subtle tint)
        col_bg_y = cy + cat_header_h
        col_bg_h = ih - cat_header_h
        tint_color = _lighten_hex(cat_color, 0.88)
        _add_rounded_rect(slide, cx, col_bg_y, col_w, col_bg_h,
                          tint_color, line_hex=muted, line_width_pt=0.5,
                          radius_px=0)

        # Items within the column
        item_y = col_bg_y + item_pad
        item_inner_w = col_w - 2 * item_pad

        # Estimate max chars that fit in the column
        char_w_emu = Pt(item_pt).emu * 0.55
        max_label_chars = max(5, int(item_inner_w / char_w_emu))

        for item in items:
            if isinstance(item, str):
                item_label = item
                item_color = None
            else:
                item_label = item.get("label", "")
                item_color = item.get("color")

            # Truncate long labels that would overflow narrow columns
            if len(item_label) > max_label_chars + 3:
                item_label = item_label[:max_label_chars] + "..."

            # Item box
            box_fill = bg
            box_line = muted
            if item_color:
                box_line = item_color

            _add_rounded_rect(slide,
                              cx + item_pad, item_y,
                              item_inner_w, item_h,
                              box_fill, line_hex=box_line,
                              line_width_pt=0.75, radius_px=radius_px)
            _add_textbox(slide,
                         cx + item_pad, item_y,
                         item_inner_w, item_h,
                         item_label, font_body, item_pt, text_c,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                         word_wrap=True)

            item_y += item_h + item_gap

            # Stop if we'd exceed bounds
            if item_y + item_h > y0 + h0 - pad:
                break
