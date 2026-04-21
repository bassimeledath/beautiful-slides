"""Render the market-map chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_COMPETITIVE = {
    "title": "Competitive landscape",
    "subtitle": "Enterprise SaaS market Q4 2024",
    "categories": [
        {
            "name": "CRM",
            "items": ["Salesforce", "HubSpot", "Zoho", "Pipedrive"],
        },
        {
            "name": "Analytics",
            "items": ["Tableau", "Looker", "Power BI", "Metabase"],
        },
        {
            "name": "DevOps",
            "items": ["GitHub", "GitLab", "Jira", "Linear"],
        },
        {
            "name": "Communication",
            "items": ["Slack", "Teams", "Discord", "Zoom"],
        },
    ],
}

DATA_TECH_STACK = {
    "title": "Modern data stack",
    "subtitle": "Recommended tooling by layer",
    "categories": [
        {
            "name": "Ingestion",
            "items": ["Fivetran", "Airbyte", "Stitch", "Meltano"],
        },
        {
            "name": "Storage",
            "items": ["Snowflake", "BigQuery", "Redshift", "Databricks"],
        },
        {
            "name": "Transform",
            "items": ["dbt", "Spark", "Dataform"],
        },
        {
            "name": "Visualization",
            "items": ["Looker", "Tableau", "Metabase", "Superset"],
        },
        {
            "name": "Orchestration",
            "items": ["Airflow", "Dagster", "Prefect"],
        },
    ],
}

DATA_PARTNER_ECO = {
    "title": "Partner ecosystem",
    "categories": [
        {
            "name": "Technology",
            "items": ["AWS", "Azure", "GCP"],
        },
        {
            "name": "System integrators",
            "items": ["Accenture", "Deloitte", "Cognizant", "Wipro"],
        },
        {
            "name": "ISV partners",
            "items": ["Snowflake", "Databricks", "Confluent"],
        },
    ],
}


MODE_DATA = {
    "sv-keynote": DATA_TECH_STACK,
    "editorial-magazine": DATA_COMPETITIVE,
    "playful-marketing": DATA_PARTNER_ECO,
    "consulting-boardroom": DATA_COMPETITIVE,
    "craft-minimal": DATA_TECH_STACK,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
