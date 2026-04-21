"""Forecast-band chart -- native python-pptx shapes only.

Central trend line with upper/lower confidence or scenario bands.
Shaded band region. Solid line for actuals, dashed for forecast portion.

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Emu(int(x1)), Emu(int(y1)),
                                      Emu(int(x2)), Emu(int(y2)))
    line = conn.line
    line.color.rgb = _rgb(color_hex)
    line.width = Pt(weight_pt)
    if dash is not None:
        try:
            line.dash_style = dash
        except Exception:
            pass
    return conn


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y)),
                                 Emu(int(w)), Emu(int(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_band_slice(slide, x, y_top, y_bot, width, color_hex, opacity_pct=20):
    """Draw one vertical slice of the band as a semi-transparent rectangle."""
    h = max(1, int(y_bot - y_top))
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y_top)),
                                 Emu(max(1, int(width))), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    # Set transparency via XML (0-100000 scale, where 100000 = fully transparent)
    try:
        from lxml import etree
        _ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_elem = shp._element
        srgb = sp_elem.find(f".//{{{_ns}}}srgbClr")
        if srgb is not None:
            alpha = etree.SubElement(srgb, f"{{{_ns}}}alpha")
            alpha.set("val", str(int(opacity_pct * 1000)))
    except Exception:
        pass
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_filled_circle(slide, cx, cy, r_emu, color_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
                                 Emu(int(r_emu * 2)), Emu(int(r_emu * 2)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _nice_ticks(vmin, vmax, target=5):
    if vmax <= vmin:
        vmax = vmin + 1
    span = vmax - vmin
    raw = span / max(target, 1)
    mag = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    for mult in (1, 2, 2.5, 5, 10):
        step = mult * mag
        if span / step <= target * 1.5:
            break
    lo = math.floor(vmin / step) * step
    hi = math.ceil(vmax / step) * step
    ticks = []
    v = lo
    while v <= hi + 1e-9:
        ticks.append(round(v, 6))
        v += step
    return ticks, lo, hi


def _fmt_num(v):
    if abs(v - round(v)) < 1e-6:
        return f"{int(round(v))}"
    return f"{v:.1f}"


def render(slide, data, tokens, bounds):
    """Render a forecast-band chart into *slide* within *bounds*.

    data keys
    ---------
    x_labels : list[str]       Labels along the x-axis.
    actuals : list[float]      Observed values (solid line). May be shorter than x_labels.
    forecast : list[float]     Forecast values (dashed line). Same length as x_labels.
    upper : list[float]        Upper bound of band. Same length as x_labels.
    lower : list[float]        Lower bound of band. Same length as x_labels.
    title : str | None         Optional chart title.
    x_label : str | None       Optional x-axis label.
    y_label : str | None       Optional y-axis subtitle.
    forecast_start : int | None  Index where forecast begins (default: len(actuals)).
    band_label : str | None    Label for the band (e.g. "90% CI").
    """
    x0, y0, w0, h0 = bounds

    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    x_labels = list(data.get("x_labels") or [])
    actuals = list(data.get("actuals") or [])
    forecast = list(data.get("forecast") or [])
    upper = list(data.get("upper") or [])
    lower = list(data.get("lower") or [])
    title = data.get("title")
    x_label = data.get("x_label")
    y_label = data.get("y_label")
    band_label = data.get("band_label")
    forecast_start = data.get("forecast_start")

    if not x_labels:
        return
    n_x = len(x_labels)

    # Default forecast_start to length of actuals
    if forecast_start is None:
        forecast_start = len(actuals)

    # Pad arrays to n_x (use None for missing)
    def _pad(arr):
        return arr + [None] * max(0, n_x - len(arr))
    actuals = _pad(actuals)
    forecast = _pad(forecast)
    upper = _pad(upper)
    lower = _pad(lower)

    # Background
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    pad = int(min(w0, h0) * 0.035)
    cursor_y = y0 + pad

    # Title
    if title:
        title_pt = int(base_pt * 1.55)
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, x0 + pad, cursor_y, w0 - 2 * pad, title_h,
                     title, font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += title_h + int(pad * 0.3)

    if y_label:
        sub_pt = max(int(base_pt * 0.85), 8)
        sub_h = int(Pt(sub_pt).emu * 1.6)
        _add_textbox(slide, x0 + pad, cursor_y, w0 - 2 * pad, sub_h,
                     y_label, font_body, sub_pt, text_c,
                     bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += sub_h

    # Layout for plot area
    tick_pt = max(8, int(base_pt * 0.8))
    x_tick_h = int(Pt(tick_pt).emu * 1.6)
    x_axis_label_h = int(Pt(tick_pt).emu * 1.8) if x_label else 0
    left_pad = int(w0 * 0.08)
    right_pad = int(w0 * 0.04)
    bottom_pad = x_tick_h + x_axis_label_h + int(pad * 0.5)

    plot_x = x0 + pad + left_pad
    plot_y = cursor_y + int(pad * 0.2)
    plot_w = (x0 + w0 - pad - right_pad) - plot_x
    plot_h = (y0 + h0 - pad - bottom_pad) - plot_y

    if plot_w <= 0 or plot_h <= 0:
        return

    # Compute y range from all data
    all_vals = []
    for arr in (actuals, forecast, upper, lower):
        for v in arr:
            if v is not None:
                all_vals.append(float(v))
    if not all_vals:
        return
    vmin = min(all_vals)
    vmax = max(all_vals)
    span = vmax - vmin if vmax > vmin else max(abs(vmax), 1.0)
    vmin_p = vmin - span * 0.08
    vmax_p = vmax + span * 0.08
    ticks, lo, hi = _nice_ticks(vmin_p, vmax_p, target=5)
    if hi == lo:
        hi = lo + 1

    def y_to_emu(v):
        frac = (v - lo) / (hi - lo)
        return plot_y + plot_h - frac * plot_h

    # X positions
    if n_x == 1:
        def x_to_emu(i):
            return plot_x + plot_w / 2
    else:
        step_x = plot_w / (n_x - 1)
        def x_to_emu(i):
            return plot_x + i * step_x

    # Gridlines
    for tv in ticks:
        yy = y_to_emu(tv)
        _add_line(slide, plot_x, yy, plot_x + plot_w, yy, muted, 0.5)

    # Y tick labels
    tick_label_w = int(left_pad * 0.95)
    tick_label_h = int(Pt(tick_pt).emu * 1.4)
    for tv in ticks:
        yy = y_to_emu(tv)
        _add_textbox(slide,
                     plot_x - tick_label_w - int(pad * 0.15),
                     yy - tick_label_h / 2,
                     tick_label_w, tick_label_h,
                     _fmt_num(tv), font_mono, tick_pt, text_c,
                     bold=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Left axis
    _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, muted, 0.75)
    # Bottom axis
    _add_line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h, muted, 0.75)

    # X tick labels
    max_ticks = 13
    skip = max(1, (n_x + max_ticks - 1) // max_ticks)
    for i, lab in enumerate(x_labels):
        if i % skip != 0 and i != n_x - 1:
            continue
        xx = x_to_emu(i)
        label_w = int(plot_w / max(n_x, 1) * 1.6)
        _add_textbox(slide,
                     xx - label_w / 2,
                     plot_y + plot_h + int(pad * 0.15),
                     label_w, x_tick_h,
                     str(lab), font_body, tick_pt, text_c,
                     bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # X-axis label
    if x_label:
        _add_textbox(slide,
                     plot_x, plot_y + plot_h + x_tick_h + int(pad * 0.2),
                     plot_w, x_axis_label_h,
                     x_label, font_body, int(base_pt * 0.85), text_c,
                     bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Draw forecast start divider (vertical dashed line)
    try:
        from pptx.enum.dml import MSO_LINE
        dash_style = MSO_LINE.DASH
    except Exception:
        dash_style = None

    if 0 < forecast_start < n_x:
        fx = x_to_emu(forecast_start)
        _add_line(slide, fx, plot_y, fx, plot_y + plot_h, muted, 0.75, dash=dash_style)

    # Draw band (upper-lower shaded region) as vertical slices
    band_opacity = 18
    slice_count = n_x - 1
    if slice_count > 0 and any(u is not None for u in upper) and any(l is not None for l in lower):
        for i in range(slice_count):
            # Interpolate: use endpoints of each slice
            u1, u2 = upper[i], upper[i + 1]
            l1, l2 = lower[i], lower[i + 1]
            if u1 is None or u2 is None or l1 is None or l2 is None:
                continue
            x1 = x_to_emu(i)
            x2 = x_to_emu(i + 1)
            # Top of slice: min y (max value) = upper
            y_top = min(y_to_emu(u1), y_to_emu(u2))
            y_bot = max(y_to_emu(l1), y_to_emu(l2))
            if y_bot > y_top:
                _add_band_slice(slide, x1, y_top, y_bot, x2 - x1, primary, band_opacity)

    # Draw forecast line (dashed) -- full forecast array
    forecast_pts = []
    for i in range(n_x):
        v = forecast[i]
        if v is not None:
            forecast_pts.append((x_to_emu(i), y_to_emu(float(v)), i))

    for a, b in zip(forecast_pts, forecast_pts[1:]):
        # Only draw forecast segments from forecast_start onward
        if a[2] >= forecast_start - 1:
            _add_line(slide, a[0], a[1], b[0], b[1], primary, 2.0, dash=dash_style)

    # Draw actuals line (solid)
    actual_pts = []
    for i in range(n_x):
        v = actuals[i]
        if v is not None:
            actual_pts.append((x_to_emu(i), y_to_emu(float(v)), i))

    for a, b in zip(actual_pts, actual_pts[1:]):
        _add_line(slide, a[0], a[1], b[0], b[1], primary, 2.5)

    # Highlight junction point (where actuals meet forecast)
    if actual_pts:
        last = actual_pts[-1]
        _add_filled_circle(slide, last[0], last[1], int(Pt(4).emu), accent)
        # Inner dot
        _add_filled_circle(slide, last[0], last[1], int(Pt(1.6).emu), bg)

    # Legend: actuals solid, forecast dashed, band shaded
    legend_y = plot_y - int(Pt(tick_pt).emu * 1.5)
    if legend_y < y0 + pad:
        legend_y = plot_y + int(Pt(tick_pt).emu * 0.2)
    cx = plot_x + plot_w
    swatch_w = int(Pt(tick_pt).emu * 1.2)
    gap = int(Pt(tick_pt).emu * 0.4)
    lbl_h = int(Pt(tick_pt).emu * 1.4)
    sy = legend_y + int(Pt(tick_pt).emu * 0.7)

    # Band label
    if band_label:
        est_w = int(Pt(tick_pt).emu * 0.55 * max(len(band_label), 3))
        cx -= est_w
        _add_textbox(slide, cx, legend_y, est_w, lbl_h,
                     band_label, font_body, tick_pt, text_c,
                     bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        cx -= gap
        # Small swatch rectangle
        sw_h = int(Pt(tick_pt).emu * 0.6)
        _add_band_slice(slide, cx - swatch_w, sy - sw_h // 2, sy + sw_h // 2,
                        swatch_w, primary, 30)
        cx -= swatch_w + gap * 2

    # Forecast label
    est_w = int(Pt(tick_pt).emu * 0.55 * 8)
    cx -= est_w
    _add_textbox(slide, cx, legend_y, est_w, lbl_h,
                 "Forecast", font_body, tick_pt, text_c,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    cx -= gap
    _add_line(slide, cx - swatch_w, sy, cx, sy, primary, 2.0, dash=dash_style)
    cx -= swatch_w + gap * 2

    # Actuals label
    est_w = int(Pt(tick_pt).emu * 0.55 * 7)
    cx -= est_w
    _add_textbox(slide, cx, legend_y, est_w, lbl_h,
                 "Actuals", font_body, tick_pt, text_c,
                 bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    cx -= gap
    _add_line(slide, cx - swatch_w, sy, cx, sy, primary, 2.5)
