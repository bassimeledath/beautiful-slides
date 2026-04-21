"""Render the forecast-band chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render  # noqa: E402


DATA = {
    "title": "ARR forecast with confidence interval",
    "x_labels": ["Q1'24", "Q2'24", "Q3'24", "Q4'24", "Q1'25", "Q2'25", "Q3'25", "Q4'25"],
    "actuals": [42, 48, 53, 61, 68],
    "forecast": [42, 48, 53, 61, 68, 76, 85, 95],
    "upper":    [42, 50, 57, 67, 75, 88, 102, 118],
    "lower":    [42, 46, 49, 55, 61, 64, 68, 72],
    "y_label": "ARR ($M)",
    "x_label": "Quarter",
    "band_label": "90% CI",
}

DATA_ALT = {
    "title": "Monthly active users -- base vs. scenarios",
    "x_labels": ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"],
    "actuals": [1200, 1350, 1500, 1480, 1550, 1700, 1750],
    "forecast": [1200, 1350, 1500, 1480, 1550, 1700, 1750, 1820, 1900, 1980, 2050, 2120],
    "upper":    [1200, 1350, 1520, 1530, 1620, 1800, 1900, 2050, 2220, 2400, 2580, 2760],
    "lower":    [1200, 1350, 1480, 1430, 1480, 1600, 1600, 1590, 1580, 1560, 1520, 1480],
    "y_label": "MAU (thousands)",
    "band_label": "Scenario range",
}

MODE_DATA = {
    "sv-keynote": DATA,
    "editorial-magazine": DATA_ALT,
    "playful-marketing": DATA,
    "consulting-boardroom": DATA_ALT,
    "craft-minimal": DATA,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
