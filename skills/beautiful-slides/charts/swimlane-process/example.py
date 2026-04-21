import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


DATA = {
    "title": "Order fulfilment flow",
    "lanes": [
        {
            "label": "Customer",
            "steps": [
                {"label": "Place order"},
                {"label": "Receive confirmation"},
                {"label": "Track shipment"},
            ],
        },
        {
            "label": "Sales",
            "steps": [
                {"label": "Validate order"},
                {"label": "Process payment"},
                {"label": "Send confirmation"},
            ],
        },
        {
            "label": "Warehouse",
            "steps": [
                {"label": "Pick & pack"},
                {"label": "Ship"},
            ],
        },
    ],
    "connections": [
        {"from": [0, 0], "to": [1, 0]},   # Customer places -> Sales validates
        {"from": [1, 0], "to": [1, 1]},   # Sales validates -> processes payment
        {"from": [1, 1], "to": [1, 2]},   # Sales processes -> sends confirmation
        {"from": [1, 2], "to": [0, 1]},   # Sales sends -> Customer receives
        {"from": [1, 1], "to": [2, 0]},   # Sales payment -> Warehouse picks
        {"from": [2, 0], "to": [2, 1]},   # Warehouse picks -> ships
        {"from": [2, 1], "to": [0, 2]},   # Warehouse ships -> Customer tracks
    ],
}


def _set_slide_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))


def build(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA, tokens, (x, y, w, h))

    prs.save(out_path)


def main():
    for name, tokens in MODES.items():
        path = os.path.join(HERE, f"example-{name}.pptx")
        build(name, tokens, path)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
