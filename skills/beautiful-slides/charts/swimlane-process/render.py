"""Swimlane process — flowchart split into horizontal lanes by team/system/owner.

Steps flow left-to-right within lanes, arrows cross lanes for handoffs.
Lane labels on the left.
"""

import math

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(3)))
    tf.margin_right = Emu(int(Pt(3)))
    tf.margin_top = Emu(int(Pt(2)))
    tf.margin_bottom = Emu(int(Pt(2)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_arrow(slide, x1, y1, x2, y2, color_hex, thickness):
    """Draw a directional arrow from (x1,y1) to (x2,y2)."""
    dx = x2 - x1
    dy = y2 - y1
    length = math.sqrt(dx * dx + dy * dy)
    if length < 1:
        return

    ux = dx / length
    uy = dy / length
    px = -uy
    py = ux

    head_len = min(length * 0.25, thickness * 5)
    head_w = head_len * 0.7

    base_x = x2 - ux * head_len
    base_y = y2 - uy * head_len

    # Triangle head
    p1x, p1y = int(x2), int(y2)
    p2x, p2y = int(base_x + px * head_w / 2), int(base_y + py * head_w / 2)
    p3x, p3y = int(base_x - px * head_w / 2), int(base_y - py * head_w / 2)

    ff = slide.shapes.build_freeform(p1x, p1y, scale=1.0)
    ff.add_line_segments([(p2x, p2y), (p3x, p3y)], close=True)
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_hex)
    shape.line.fill.background()

    # Shaft
    shaft_len = length - head_len
    if shaft_len > 0:
        sw = thickness / 2
        s1x = int(x1 + px * sw)
        s1y = int(y1 + py * sw)
        s2x = int(x1 - px * sw)
        s2y = int(y1 - py * sw)
        s3x = int(base_x - px * sw)
        s3y = int(base_y - py * sw)
        s4x = int(base_x + px * sw)
        s4y = int(base_y + py * sw)

        ff2 = slide.shapes.build_freeform(s1x, s1y, scale=1.0)
        ff2.add_line_segments([(s2x, s2y), (s3x, s3y), (s4x, s4y)], close=True)
        shaft = ff2.convert_to_shape()
        shaft.fill.solid()
        shaft.fill.fore_color.rgb = _rgb(color_hex)
        shaft.line.fill.background()


def render(slide, data, tokens, bounds):
    """Render a swimlane process diagram.

    data:
        title    - optional string
        lanes    - list of {"label": str, "steps": [{"label": str}, ...]}
        connections - list of {"from": [lane_idx, step_idx],
                               "to": [lane_idx, step_idx]}
                      Arrows connecting steps, including cross-lane handoffs.
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]
    radius_px = tokens.get("radius_px", 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    lanes = data.get("lanes", [])
    if not lanes:
        return
    title = data.get("title")
    connections = data.get("connections", [])

    n_lanes = len(lanes)
    # Find max steps in any lane for column sizing
    max_steps = max(len(lane.get("steps", [])) for lane in lanes) if lanes else 1
    max_steps = max(max_steps, 1)

    # --- layout ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_h = (y + h) - cur_y
    avail_w = w

    # Label column on the left
    label_col_w = int(avail_w * 0.15)
    gap = int(avail_w * 0.01)
    chart_x = x + label_col_w + gap
    chart_w = avail_w - label_col_w - gap

    # Lane layout
    lane_gap = int(Pt(base_pt * 0.3))
    total_lane_gaps = lane_gap * (n_lanes - 1) if n_lanes > 1 else 0
    lane_h = max(int((avail_h - total_lane_gaps) / n_lanes), int(Pt(base_pt * 3)))

    # Step box sizing — step_w is now computed per-lane (below).
    # Only step_h and padding are global.
    step_pad_x = int(chart_w * 0.02)
    step_pad_y = int(lane_h * 0.15)
    step_h = max(int(lane_h * 0.50), int(Pt(base_pt * 2.2)))

    # Store step center positions for arrows
    step_centers = {}

    for li, lane in enumerate(lanes):
        lane_y = cur_y + li * (lane_h + lane_gap)
        steps = lane.get("steps", [])
        n_steps = len(steps)

        # Compute step width per-lane based on this lane's step count
        lane_n = max(n_steps, 1)
        step_w = max(int((chart_w - step_pad_x * (lane_n + 1)) / lane_n), int(Pt(base_pt * 3)))

        # Lane background stripe (alternating subtle tint)
        if li % 2 == 0:
            stripe_color = _lerp_hex(bg, muted, 0.06)
        else:
            stripe_color = _lerp_hex(bg, muted, 0.12)

        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(chart_x), Emu(lane_y),
            Emu(chart_w), Emu(lane_h),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _rgb(stripe_color)
        stripe.line.fill.background()

        # Lane divider line at bottom (except last lane)
        if li < n_lanes - 1:
            div_h = max(int(Pt(0.75)), 1)
            div_y = lane_y + lane_h + lane_gap // 2
            div = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(x), Emu(div_y),
                Emu(w), Emu(div_h),
            )
            div.fill.solid()
            div.fill.fore_color.rgb = _rgb(_lerp_hex(bg, muted, 0.25))
            div.line.fill.background()

        # Lane label on the left
        _add_textbox(
            slide, x, lane_y, label_col_w, lane_h,
            lane["label"], font_display, max(int(base_pt * 0.9), 8), text_c,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Draw steps as rounded rectangles within the lane
        lane_color_t = li / (n_lanes - 1) if n_lanes > 1 else 0.0
        step_color = _lerp_hex(primary, accent, lane_color_t)

        for si, step in enumerate(steps):
            # Position step evenly across chart_w
            if n_steps == 1:
                sx = chart_x + step_pad_x
            else:
                total_step_space = chart_w - step_pad_x * 2
                sx = chart_x + step_pad_x + int(si * (total_step_space - step_w) / (n_steps - 1))

            sy = lane_y + (lane_h - step_h) // 2

            # Clamp to bounds
            sx = max(chart_x, min(sx, chart_x + chart_w - step_w))
            sy = max(lane_y, min(sy, lane_y + lane_h - step_h))

            # Step box
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE
            box = slide.shapes.add_shape(
                shape_type,
                Emu(sx), Emu(sy),
                Emu(step_w), Emu(step_h),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = _rgb(step_color)
            box.line.fill.background()

            # Step label — adaptive font + ellipsis truncation
            label_font_size = max(int(base_pt * 0.75), 7)
            if n_steps > 3:
                label_font_size = max(label_font_size - 1, 7)
            step_label = step.get("label", "")
            max_step_chars = max(10, int(step_w / (Pt(label_font_size) * 0.5)))
            if len(step_label) > max_step_chars:
                step_label = step_label[:max_step_chars - 1] + "\u2026"
            _add_textbox(
                slide, sx, sy, step_w, step_h,
                step_label, font_body, label_font_size, bg,
                align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )

            # Record center for connections
            step_centers[(li, si)] = (sx + step_w // 2, sy + step_h // 2)

    # Draw connections (arrows between steps)
    arrow_thickness = max(int(Pt(base_pt * 0.12)), int(Pt(1.5)))
    arrow_color = _lerp_hex(muted, text_c, 0.4)

    for conn in connections:
        from_key = tuple(conn["from"])
        to_key = tuple(conn["to"])
        if from_key not in step_centers or to_key not in step_centers:
            continue

        fx, fy = step_centers[from_key]
        tx, ty = step_centers[to_key]

        # Determine arrow start/end at the edge of the step boxes
        # For same-lane (horizontal): arrow from right edge to left edge
        # For cross-lane (vertical): arrow from bottom/top edge
        from_lane, from_step = from_key
        to_lane, to_step = to_key

        if from_lane == to_lane:
            # Horizontal: start from right edge of source, end at left edge of target
            ax1 = fx + step_w // 2
            ay1 = fy
            ax2 = tx - step_w // 2
            ay2 = ty
        else:
            # Cross-lane: start from bottom/top edge, end at top/bottom
            if to_lane > from_lane:
                ay1 = fy + step_h // 2
                ay2 = ty - step_h // 2
            else:
                ay1 = fy - step_h // 2
                ay2 = ty + step_h // 2
            ax1 = fx
            ax2 = tx

        _draw_arrow(slide, ax1, ay1, ax2, ay2, arrow_color, arrow_thickness)

    return None
