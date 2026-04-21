"""Waterfall chart renderer — native python-pptx shapes only.

Shows how an initial value is affected by sequential positive/negative
changes to reach a final value.  Start/end columns are anchored to the
baseline; intermediates float.  Positive deltas use ``tokens["primary"]``,
negative deltas use ``tokens["accent"]``.  Thin dotted connector lines
bridge successive columns.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
import math


# EMU constants
_EMU_PER_PX = 9525   # 96 DPI reference


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _set_fill(shape, hex_):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_)


def _no_line(shape):
    shape.line.fill.background()


def _add_text(slide, x, y, w, h, text, font_name, size_pt, hex_color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold = bold
    return tb


def _add_bar(slide, x, y, w, h, fill_hex, radius_px):
    w_i = max(1, int(w))
    h_i = max(1, int(h))
    if radius_px and radius_px > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i),
        )
        short_emu = min(w_i, h_i)
        radius_emu = radius_px * _EMU_PER_PX
        ratio = max(0.0, min(0.5, radius_emu / short_emu / 2.0))
        try:
            shape.adjustments[0] = ratio
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i),
        )
    _set_fill(shape, fill_hex)
    _no_line(shape)
    return shape


def _add_dashed_line(slide, x1, y1, x2, y2, hex_color, width_emu):
    """Add a connector styled as a dashed line."""
    ln = slide.shapes.add_connector(
        1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)),
    )
    ln.line.color.rgb = _rgb(hex_color)
    ln.line.width = Emu(int(width_emu))
    # Set dash style to dot via XML (pptx doesn't expose dash style directly)
    ln_elem = ln.line._ln
    ln_elem.set(qn("a:prstDash"), "")
    prstDash = ln_elem.makeelement(qn("a:prstDash"), {qn("a:val"): "dot"})
    # Remove existing prstDash if any
    for old in ln_elem.findall(qn("a:prstDash")):
        ln_elem.remove(old)
    ln_elem.append(prstDash)
    return ln


def _add_line(slide, x1, y1, x2, y2, hex_color, width_emu):
    ln = slide.shapes.add_connector(
        1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)),
    )
    ln.line.color.rgb = _rgb(hex_color)
    ln.line.width = Emu(int(width_emu))
    return ln


def _nice_ticks(vmin, vmax, target=5):
    """Return tick values and effective range for a min/max range."""
    if vmin >= vmax:
        return [0, 1], 0, 1
    span = vmax - vmin
    raw = span / target
    magnitude = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    residual = raw / magnitude
    if residual < 1.5:
        step = 1 * magnitude
    elif residual < 3:
        step = 2 * magnitude
    elif residual < 7:
        step = 5 * magnitude
    else:
        step = 10 * magnitude
    tick_min = step * math.floor(vmin / step)
    tick_max = step * math.ceil(vmax / step)
    ticks = []
    v = tick_min
    while v <= tick_max + 1e-9:
        ticks.append(v)
        v += step
    return ticks, tick_min, tick_max


def _fmt(v, suffix=""):
    if abs(v - round(v)) < 1e-9:
        s = f"{int(round(v))}"
    else:
        s = f"{v:.1f}"
    return f"{s}{suffix}" if suffix else s


def render(slide, data, tokens, bounds):
    """Render a waterfall chart onto *slide* inside *bounds*.

    ``data`` keys:
      steps : list[dict]
          Each dict has ``label`` (str) and ``value`` (number).
          First step is the starting total; last step is the ending total.
          Intermediates are deltas (+/-).
      title : str or None
      value_suffix : str
      show_values : bool (default True)
    """
    x, y, w, h = bounds

    steps = data.get("steps") or []
    if len(steps) < 2:
        return
    title = data.get("title")
    value_suffix = data.get("value_suffix") or ""
    show_values = bool(data.get("show_values", True))

    # Theme tokens
    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    _set_fill(bg_shape, bg)
    _no_line(bg_shape)

    # Outer padding
    pad = int(min(w, h) * 0.04)
    ix = x + pad
    iy = y + pad
    iw = w - 2 * pad
    ih = h - 2 * pad

    # Title
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.6))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_text(
            slide, ix, iy, iw, title_h, title,
            font_name=font_display, size_pt=title_pt, hex_color=text_c,
            bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        iy += title_h + int(Pt(base_pt).emu * 0.6)
        ih = (y + h - pad) - iy

    # Compute running totals
    n = len(steps)
    # First and last are totals; intermediates are deltas
    running = []
    cumulative = steps[0]["value"]
    running.append(cumulative)
    for i in range(1, n - 1):
        cumulative += steps[i]["value"]
        running.append(cumulative)
    running.append(steps[-1]["value"])  # final total (should match cumulative)

    # Determine min/max of all bar extents
    bar_tops = []
    bar_bottoms = []
    # Start bar: from 0 to running[0]
    bar_tops.append(max(0, running[0]))
    bar_bottoms.append(min(0, running[0]))
    # Intermediate bars: float between prev running and current running
    prev = running[0]
    for i in range(1, n - 1):
        cur = running[i]
        bar_tops.append(max(prev, cur))
        bar_bottoms.append(min(prev, cur))
        prev = cur
    # End bar: from 0 to running[-1]
    bar_tops.append(max(0, running[-1]))
    bar_bottoms.append(min(0, running[-1]))

    global_max = max(bar_tops)
    global_min = min(bar_bottoms)
    # Add small padding
    val_range = global_max - global_min
    if val_range == 0:
        val_range = 1
    padding_v = val_range * 0.05
    axis_min = global_min - padding_v
    axis_max = global_max + padding_v

    ticks, tick_min, tick_max = _nice_ticks(axis_min, axis_max)

    # Layout regions
    tick_pt = max(int(base_pt * 0.78), 8)
    cat_pt = max(int(base_pt * 0.85), 9)
    val_pt = max(int(base_pt * 0.78), 8)

    # Reduce cat_pt when steps > 8 to fit more labels
    if n > 8:
        cat_pt = max(int(base_pt * 0.72), 7)

    tick_labels = [_fmt(t) for t in ticks]
    max_tick_chars = max(len(s) for s in tick_labels) if tick_labels else 1
    left_margin = int(Pt(tick_pt).emu * 0.65 * (max_tick_chars + 1))

    # Scale bottom_margin based on longest label word count
    max_label_words = max(len(str(s["label"]).split()) for s in steps)
    lines_needed = min(max_label_words, 4)
    bottom_margin = int(Pt(cat_pt).emu * (1.4 * lines_needed + 0.5))
    top_margin = int(Pt(val_pt).emu * 1.6)
    right_margin = int(Pt(val_pt).emu * 1.2)

    plot_x = ix + left_margin
    plot_y = iy + top_margin
    plot_w = max(1, iw - left_margin - right_margin)
    plot_h = max(1, ih - top_margin - bottom_margin)

    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    def val_to_y(v):
        """Convert a data value to a y coordinate in EMU."""
        if tick_max == tick_min:
            return plot_y + plot_h
        frac = (v - tick_min) / (tick_max - tick_min)
        return plot_y + plot_h - int(plot_h * frac)

    # Gridlines + tick labels
    for i, t in enumerate(ticks):
        ty = val_to_y(t)
        lbl_w = left_margin - int(Pt(tick_pt).emu * 0.3)
        lbl_h = int(Pt(tick_pt).emu * 1.4)
        _add_text(
            slide, ix, ty - lbl_h // 2, lbl_w, lbl_h,
            _fmt(t) + (value_suffix if t != 0 and i == len(ticks) - 1 else ""),
            font_name=font_mono, size_pt=tick_pt, hex_color=text_c,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )
        if i > 0:
            _add_line(slide, plot_x, ty, plot_x + plot_w, ty, muted, hairline)

    # Baseline at zero
    zero_y = val_to_y(0)
    _add_line(slide, plot_x, zero_y, plot_x + plot_w, zero_y, muted, hairline * 2)

    # Bars
    col_w = plot_w / n
    inner_pad = col_w * 0.12
    bar_w = col_w - 2 * inner_pad

    prev_end_y = None  # y-coordinate of the "end" of the previous bar for connector

    for i in range(n):
        col_left = plot_x + i * col_w
        bx = col_left + inner_pad

        if i == 0:
            # Start total: anchored to zero
            top_val = max(0, running[0])
            bot_val = min(0, running[0])
            by = val_to_y(top_val)
            bh = val_to_y(bot_val) - by
            color = muted
            connector_end_y = val_to_y(running[0])
        elif i == n - 1:
            # End total: anchored to zero
            top_val = max(0, running[-1])
            bot_val = min(0, running[-1])
            by = val_to_y(top_val)
            bh = val_to_y(bot_val) - by
            color = muted
            connector_end_y = val_to_y(running[-1])
        else:
            # Intermediate: floating between prev running and current running
            delta = steps[i]["value"]
            prev_running = running[i - 1]
            cur_running = running[i]
            top_val = max(prev_running, cur_running)
            bot_val = min(prev_running, cur_running)
            by = val_to_y(top_val)
            bh = val_to_y(bot_val) - by
            color = primary if delta >= 0 else accent
            connector_end_y = val_to_y(cur_running)

        bh = max(bh, 1)
        _add_bar(slide, bx, by, bar_w, bh, color, radius_px)

        # Dotted connector line from previous bar end to this bar start
        if prev_end_y is not None:
            prev_col_right = plot_x + (i - 1) * col_w + inner_pad + bar_w
            _add_dashed_line(
                slide, prev_col_right, prev_end_y,
                bx, prev_end_y,
                muted, hairline,
            )

        prev_end_y = connector_end_y

        # Category label below -- with ellipsis truncation for long labels
        cat_y = plot_y + plot_h + int(Pt(cat_pt).emu * 0.4)
        cat_h = int(Pt(cat_pt).emu * 3.5)  # allow more wrapping room
        step_label = str(steps[i]["label"])
        if len(step_label) > 20 and n > 6:
            step_label = step_label[:19] + "\u2026"
        _add_text(
            slide, col_left, cat_y, col_w, cat_h,
            step_label,
            font_name=font_body, size_pt=cat_pt, hex_color=text_c,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )

        # Value label above/below the bar
        if show_values:
            if i == 0 or i == n - 1:
                vt = _fmt(steps[i]["value"], value_suffix)
            else:
                delta = steps[i]["value"]
                sign = "+" if delta >= 0 else ""
                vt = f"{sign}{_fmt(delta, value_suffix)}"

            vh = int(Pt(val_pt).emu * 1.3)
            # Place above the bar if positive or total, below if negative delta
            if i == 0 or i == n - 1 or steps[i]["value"] >= 0:
                vy = by - vh - int(Pt(val_pt).emu * 0.1)
            else:
                vy = by + bh + int(Pt(val_pt).emu * 0.1)
            _add_text(
                slide, bx, vy, bar_w, vh, vt,
                font_name=font_mono, size_pt=val_pt, hex_color=text_c,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
            )
