"""Render the waterfall chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_BRIDGE = {
    "title": "Revenue bridge, Q4 to Q1",
    "steps": [
        {"label": "Q4 Revenue", "value": 100},
        {"label": "New logos",   "value": 18},
        {"label": "Expansion",   "value": 12},
        {"label": "Churn",       "value": -8},
        {"label": "Contraction", "value": -5},
        {"label": "Q1 Revenue", "value": 117},
    ],
    "value_suffix": "M",
    "show_values": True,
}

DATA_COST = {
    "title": "Cost buildup per unit",
    "steps": [
        {"label": "Materials",   "value": 42},
        {"label": "Labor",       "value": 28},
        {"label": "Overhead",    "value": 15},
        {"label": "Rebates",     "value": -6},
        {"label": "Total cost", "value": 79},
    ],
    "value_suffix": "",
    "show_values": True,
}

DATA_PROFIT = {
    "title": "Net income walk, FY23 to FY24",
    "steps": [
        {"label": "FY23 Net",    "value": 54},
        {"label": "Rev growth",  "value": 22},
        {"label": "COGS",        "value": -9},
        {"label": "OpEx",        "value": -7},
        {"label": "Tax benefit", "value": 3},
        {"label": "FX impact",   "value": -2},
        {"label": "FY24 Net",    "value": 61},
    ],
    "value_suffix": "M",
    "show_values": True,
}


MODE_DATA = {
    "sv-keynote":           DATA_BRIDGE,
    "editorial-magazine":   DATA_COST,
    "playful-marketing":    DATA_BRIDGE,
    "consulting-boardroom": DATA_PROFIT,
    "craft-minimal":        DATA_COST,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
