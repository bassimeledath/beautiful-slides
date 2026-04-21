"""Progress bar — horizontal completion meters stacked vertically.

Each bar shows a label, current value, target value, and a target marker line.
Stack 1-6 bars for a goal dashboard.
Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


_EMU_PER_PX = 9525


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_bar_shape(slide, x, y, w, h, fill_hex, radius_px):
    """Draw a rounded or sharp rect for the bar track/fill."""
    w_i = max(1, int(w))
    h_i = max(1, int(h))
    if radius_px and radius_px > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i))
        short_side = min(w_i, h_i)
        radius_emu = radius_px * _EMU_PER_PX
        ratio = max(0.0, min(0.5, radius_emu / short_side / 2.0))
        try:
            shape.adjustments[0] = ratio
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(x)), Emu(int(y)), Emu(w_i), Emu(h_i))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape


def _lighten_hex(hex_, factor=0.85):
    """Lighten a hex color by blending toward white."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02X}{g:02X}{b:02X}"


def _darken_hex(hex_, factor=0.3):
    """Darken a hex color by blending toward black."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r * (1 - factor))
    g = int(g * (1 - factor))
    b = int(b * (1 - factor))
    return f"#{r:02X}{g:02X}{b:02X}"


_LIGHT_BG_HEXES = {"F6F1E8", "FFF4EB", "FFFFFF", "FCFBF8"}


def _is_dark_mode(bg_hex):
    return bg_hex.lstrip("#").upper() not in _LIGHT_BG_HEXES


def _luminance(hex_):
    """Relative luminance of a hex color (0=black, 1=white)."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255
    def _lin(c):
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)


def _ensure_contrast(fg_hex, bg_hex, min_ratio=3.0):
    """If fg/bg contrast ratio is below min_ratio, return a better fg color."""
    fg_lum = _luminance(fg_hex)
    bg_lum = _luminance(bg_hex)
    lighter = max(fg_lum, bg_lum)
    darker = min(fg_lum, bg_lum)
    ratio = (lighter + 0.05) / (darker + 0.05)
    if ratio >= min_ratio:
        return fg_hex
    # If bg is dark, lighten fg; if bg is light, darken fg
    if bg_lum < 0.5:
        return _lighten_hex(fg_hex, 0.5)
    else:
        return _darken_hex(fg_hex, 0.4)


def _fmt_compact(v):
    """Format number compactly with K/M suffix."""
    av = abs(v)
    if av >= 1_000_000:
        s = f"{v / 1_000_000:.1f}M"
        return s.replace(".0M", "M")
    if av >= 10_000:
        return f"{v / 1_000:.0f}K"
    if av >= 1_000:
        s = f"{v / 1_000:.1f}K"
        return s.replace(".0K", "K")
    if isinstance(v, float) and abs(v - round(v)) > 1e-6:
        return f"{v:.1f}"
    return f"{int(round(v))}"


import math


def _estimate_lines(text, font_pt, avail_w):
    """Estimate how many lines text will wrap to given available width."""
    char_w = Pt(font_pt).emu * 0.55
    chars_per_line = max(1, int(avail_w / char_w))
    return max(1, math.ceil(len(text) / chars_per_line))


def render(slide, data, tokens, bounds):
    """Render a set of progress bars stacked vertically.

    Parameters
    ----------
    slide : pptx.slide.Slide
    data : dict
        bars: list of {label, value, target, [format]}
        title: str — optional title above bars
    tokens : dict
    bounds : tuple (x, y, w, h) in EMU
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()

    dark = _is_dark_mode(bg)

    bars = list(data.get("bars", []))
    title = data.get("title")

    if not bars:
        return

    # Limit to 6 bars
    bars = bars[:6]

    # Layout
    pad = int(min(w, h) * 0.04)
    cursor_y = y + pad

    # Title -- dynamic height for multi-line wrapping
    if title:
        title_pt = int(base_pt * 1.5)
        title_avail_w = w - 2 * pad
        title_lines = _estimate_lines(title, title_pt, title_avail_w)
        title_h = int(Pt(title_pt).emu * 1.3 * title_lines + Pt(title_pt).emu * 0.5)
        title_h = min(title_h, int(h * 0.22))
        _add_textbox(slide, x + pad, cursor_y, title_avail_w, title_h,
                     title, font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        cursor_y += title_h + int(Pt(base_pt * 1.0).emu)

    # Available space for bars
    avail_h = (y + h - pad) - cursor_y
    inner_w = w - 2 * pad

    # Each bar block: label row + bar row + spacing
    n = len(bars)
    label_pt = max(9, int(base_pt * 0.9))
    value_pt = max(8, int(base_pt * 0.85))
    label_row_h = int(Pt(label_pt).emu * 1.6)
    bar_h = max(int(Pt(base_pt * 0.9).emu), int(avail_h * 0.06))
    inter_spacing = int(Pt(base_pt * 0.6).emu)

    block_h = label_row_h + bar_h + inter_spacing
    # Adjust if blocks exceed available space -- more aggressive reduction
    total_needed = block_h * n - inter_spacing
    if total_needed > avail_h:
        block_h = avail_h // n
        label_row_h = int(block_h * 0.30)
        bar_h = int(block_h * 0.30)
        inter_spacing = block_h - label_row_h - bar_h
        # Also reduce font sizes
        label_pt = max(7, label_pt - 2)
        value_pt = max(7, value_pt - 2)

    # Ensure label text has sufficient contrast against bg
    label_color = _ensure_contrast(text_c, bg)
    value_color = _ensure_contrast(text_c, bg)

    # Track color: subtle tint of bg
    if dark:
        track_color = _lighten_hex(bg, 0.15)
    else:
        track_color = _darken_hex(bg, 0.08)

    # Determine max target for scaling (bars can exceed target)
    max_ref = max(b.get("target", b.get("value", 1)) for b in bars)
    # Also consider max value
    max_val = max(b.get("value", 0) for b in bars)
    scale_max = max(max_ref, max_val) * 1.05  # 5% headroom

    bar_colors = [primary, accent]

    for i, bar in enumerate(bars):
        label = bar.get("label", "")
        value = bar.get("value", 0)
        target = bar.get("target")
        fmt = bar.get("format", "")

        color = bar_colors[i % len(bar_colors)]
        block_y = cursor_y + i * block_h

        # Format value string -- use compact notation when no explicit format
        if fmt:
            val_str = fmt.format(value)
            tgt_str = fmt.format(target) if target is not None else ""
        else:
            val_str = _fmt_compact(value) if isinstance(value, (int, float)) else str(value)
            tgt_str = _fmt_compact(target) if target is not None and isinstance(target, (int, float)) else (str(target) if target is not None else "")

        # Label row: "Label" on left, "value / target" on right
        right_text = val_str
        if target is not None:
            right_text = f"{val_str} / {tgt_str}"

        _add_textbox(slide, x + pad, block_y, inner_w * 0.6, label_row_h,
                     label, font_body, label_pt, label_color,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)

        _add_textbox(slide, x + pad + inner_w * 0.6, block_y,
                     inner_w * 0.4, label_row_h,
                     right_text, font_mono, value_pt, value_color,
                     bold=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

        # Bar track (full width, muted bg)
        bar_y = block_y + label_row_h
        _add_bar_shape(slide, x + pad, bar_y, inner_w, bar_h,
                       track_color, radius_px)

        # Filled portion
        if scale_max > 0:
            fill_frac = min(value / scale_max, 1.0)
        else:
            fill_frac = 0
        fill_w = max(1, int(inner_w * fill_frac))
        if fill_w > 0:
            _add_bar_shape(slide, x + pad, bar_y, fill_w, bar_h,
                           color, radius_px)

        # Target marker line (vertical)
        if target is not None and scale_max > 0:
            tgt_frac = min(target / scale_max, 1.0)
            tgt_x = int((x + pad) + inner_w * tgt_frac)
            marker_overshoot = int(bar_h * 0.25)
            marker_top = bar_y - marker_overshoot
            marker_bottom = bar_y + bar_h + marker_overshoot
            marker_w = max(int(Pt(2).emu), int(bar_h * 0.08))

            marker = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(int(tgt_x - marker_w // 2)), Emu(int(marker_top)),
                Emu(marker_w), Emu(int(marker_bottom - marker_top)))
            marker.fill.solid()
            marker.fill.fore_color.rgb = _rgb(text_c)
            marker.line.fill.background()
            try:
                marker.shadow.inherit = False
            except Exception:
                pass

    return None
