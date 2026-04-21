"""Run `render()` 5 times -- once per mode -- writing example-<mode>.pptx per run."""

import os
import sys

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module
from tokens import MODES  # noqa: E402
from render import render  # noqa: E402


DATA = {
    "title": "Q2 Goal Progress",
    "bars": [
        {"label": "Revenue", "value": 8.4, "target": 10.0, "format": "${:,.1f}M"},
        {"label": "New Logos", "value": 34, "target": 40},
        {"label": "NPS", "value": 68, "target": 70},
        {"label": "Retention", "value": 94, "target": 95, "format": "{:.0f}%"},
        {"label": "CSAT", "value": 4.2, "target": 4.5, "format": "{:.1f}"},
    ],
}


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    for mode_name, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        margin = Inches(0.4)
        bounds = (
            int(margin),
            int(margin),
            int(prs.slide_width - 2 * margin),
            int(prs.slide_height - 2 * margin),
        )
        render(slide, DATA, tokens, bounds)

        out_path = os.path.join(out_dir, f"example-{mode_name}.pptx")
        prs.save(out_path)
        print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
