"""Flowchart — boxes and arrows for linear or branching process logic.

Supports process (rectangle), decision (diamond), and start/end (rounded rect)
node types. Auto-layout top-to-bottom or left-to-right.
"""

import math

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _luminance(hex_):
    """Return relative luminance (0-1) of a hex color."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255.0, int(h[2:4], 16) / 255.0, int(h[4:6], 16) / 255.0
    r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
    g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
    b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(3)))
    tf.margin_right = Emu(int(Pt(3)))
    tf.margin_top = Emu(int(Pt(2)))
    tf.margin_bottom = Emu(int(Pt(2)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _draw_arrow(slide, x1, y1, x2, y2, color_hex, thickness):
    """Draw a directional arrow from (x1,y1) to (x2,y2)."""
    dx = x2 - x1
    dy = y2 - y1
    length = math.sqrt(dx * dx + dy * dy)
    if length < 1:
        return

    ux = dx / length
    uy = dy / length
    px = -uy
    py = ux

    head_len = min(length * 0.25, thickness * 5)
    head_w = head_len * 0.7

    base_x = x2 - ux * head_len
    base_y = y2 - uy * head_len

    p1x, p1y = int(x2), int(y2)
    p2x, p2y = int(base_x + px * head_w / 2), int(base_y + py * head_w / 2)
    p3x, p3y = int(base_x - px * head_w / 2), int(base_y - py * head_w / 2)

    ff = slide.shapes.build_freeform(p1x, p1y, scale=1.0)
    ff.add_line_segments([(p2x, p2y), (p3x, p3y)], close=True)
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_hex)
    shape.line.fill.background()

    shaft_len = length - head_len
    if shaft_len > 0:
        sw = thickness / 2
        s1x = int(x1 + px * sw)
        s1y = int(y1 + py * sw)
        s2x = int(x1 - px * sw)
        s2y = int(y1 - py * sw)
        s3x = int(base_x - px * sw)
        s3y = int(base_y - py * sw)
        s4x = int(base_x + px * sw)
        s4y = int(base_y + py * sw)

        ff2 = slide.shapes.build_freeform(s1x, s1y, scale=1.0)
        ff2.add_line_segments([(s2x, s2y), (s3x, s3y), (s4x, s4y)], close=True)
        shaft = ff2.convert_to_shape()
        shaft.fill.solid()
        shaft.fill.fore_color.rgb = _rgb(color_hex)
        shaft.line.fill.background()


def _draw_diamond(slide, cx, cy, half_w, half_h, fill_hex, border_hex):
    """Draw a diamond (rotated square) centred at (cx, cy)."""
    top = (int(cx), int(cy - half_h))
    right = (int(cx + half_w), int(cy))
    bottom = (int(cx), int(cy + half_h))
    left = (int(cx - half_w), int(cy))

    ff = slide.shapes.build_freeform(top[0], top[1], scale=1.0)
    ff.add_line_segments([right, bottom, left], close=True)
    shape = ff.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.color.rgb = _rgb(border_hex)
    shape.line.width = Pt(1.5)
    return shape


def _assign_columns(nodes, edges):
    """Assign each node to a column for layout.

    Returns dict of node_id -> column index.
    For linear flows, everything is column 0.
    For branching (decision nodes with multiple outgoing edges),
    children get spread across columns.
    """
    # Build adjacency
    children = {}
    for e in edges:
        src = e["from"]
        children.setdefault(src, []).append(e["to"])

    # Find root(s) — nodes not targeted by any edge
    targeted = {e["to"] for e in edges}
    roots = [n["id"] for n in nodes if n["id"] not in targeted]
    if not roots:
        roots = [nodes[0]["id"]]

    node_map = {n["id"]: n for n in nodes}
    col = {}
    row = {}

    def assign(nid, c, r):
        if nid in col:
            return
        col[nid] = c
        row[nid] = r
        kids = children.get(nid, [])
        if len(kids) == 0:
            return
        elif len(kids) == 1:
            assign(kids[0], c, r + 1)
        else:
            # Spread children symmetrically around current column
            spread = len(kids)
            start_col = c - (spread - 1) / 2.0
            for i, kid in enumerate(kids):
                assign(kid, start_col + i, r + 1)

    for root in roots:
        assign(root, 0, 0)

    # Assign any orphans
    for n in nodes:
        if n["id"] not in col:
            col[n["id"]] = 0
            row[n["id"]] = len(row)

    return col, row


def _linear_layout(nodes, edges):
    """Two-column linear layout for large graphs (>8 nodes).

    Walk the main chain (following the first outgoing edge of each node).
    Place main-chain nodes in column 0. Branch targets (second edge of
    decision nodes) go in column 1 at the same row as their source.
    Returns (col_map, row_map) with integer column/row indices.
    """
    children = {}
    for e in edges:
        children.setdefault(e["from"], []).append(e["to"])

    targeted = {e["to"] for e in edges}
    roots = [n["id"] for n in nodes if n["id"] not in targeted]
    if not roots:
        roots = [nodes[0]["id"]]

    node_set = {n["id"] for n in nodes}
    col = {}
    row = {}
    current_row = 0

    # Walk main chain from each root (typically one root)
    visited = set()
    branch_queue = []  # (node_id, source_row)

    for root in roots:
        nid = root
        while nid and nid not in visited and nid in node_set:
            visited.add(nid)
            col[nid] = 0
            row[nid] = current_row
            kids = children.get(nid, [])
            if len(kids) >= 2:
                # First child continues the main chain; second is branch
                for bkid in kids[1:]:
                    if bkid not in visited and bkid in node_set:
                        branch_queue.append((bkid, current_row))
                nid = kids[0]
            elif len(kids) == 1:
                nid = kids[0]
            else:
                nid = None
            current_row += 1

    # Place branch nodes in column 1
    for bnode, src_row in branch_queue:
        if bnode not in visited and bnode in node_set:
            visited.add(bnode)
            col[bnode] = 1
            row[bnode] = src_row
            # Follow any single chain from the branch node downward
            nid = bnode
            kids = children.get(nid, [])
            while kids:
                # Follow first unvisited child
                next_nid = None
                for k in kids:
                    if k not in visited and k in node_set:
                        next_nid = k
                        break
                if next_nid is None:
                    break
                visited.add(next_nid)
                src_row += 1
                col[next_nid] = 1
                row[next_nid] = src_row
                nid = next_nid
                kids = children.get(nid, [])

    # Any remaining unvisited nodes
    for n in nodes:
        if n["id"] not in col:
            col[n["id"]] = 0
            row[n["id"]] = current_row
            current_row += 1

    return col, row


def render(slide, data, tokens, bounds):
    """Render a flowchart with boxes, diamonds, and arrows.

    data:
        title     - optional string
        direction - "TB" (top-to-bottom, default) or "LR" (left-to-right)
        nodes     - list of {"id": str, "label": str, "type": "process"|"decision"|"terminal"}
        edges     - list of {"from": str, "to": str, "label": str (optional)}
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    all_nodes = data.get("nodes", [])
    all_edges = data.get("edges", [])
    if not all_nodes:
        return
    title = data.get("title")
    direction = data.get("direction", "TB")

    # --- Cap visible nodes at 10 for auto-layout ---
    MAX_VISIBLE = 10
    overflow_count = max(0, len(all_nodes) - MAX_VISIBLE)
    nodes = all_nodes[:MAX_VISIBLE]
    visible_ids = {n["id"] for n in nodes}
    # Only keep edges where both endpoints are visible
    edges = [e for e in all_edges if e["from"] in visible_ids and e["to"] in visible_ids]

    # --- layout title ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        # Always use the main text token for the title — never muted/light
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_w = w
    avail_h = (y + h) - cur_y

    # Reserve space for "+N more" note at bottom if needed
    overflow_h = 0
    if overflow_count > 0:
        overflow_pt = max(int(base_pt * 0.75), 8)
        overflow_h = int(Pt(overflow_pt) * 2.2)
        avail_h -= overflow_h

    area_x = x
    area_y = cur_y

    # --- compute grid layout ---
    use_linear = len(nodes) > 8
    if use_linear:
        col_map, row_map = _linear_layout(nodes, edges)
    else:
        col_map, row_map = _assign_columns(nodes, edges)

    # Normalise columns to 0-based integers
    all_cols = set(col_map.values())
    all_rows = set(row_map.values())
    min_col = min(all_cols)
    max_col = max(all_cols)
    min_row = min(all_rows)
    max_row = max(all_rows)

    num_cols = max(1, max_col - min_col + 1)
    num_rows = max(1, max_row - min_row + 1)

    if direction == "LR" and not use_linear:
        # Swap: rows become columns and vice versa
        num_cols, num_rows = num_rows, num_cols
        new_col = {}
        new_row = {}
        for nid in col_map:
            new_col[nid] = row_map[nid] - min_row
            new_row[nid] = col_map[nid] - min_col
        col_map = new_col
        row_map = new_row
        min_col = 0
        min_row = 0
    else:
        # Normalise
        for nid in col_map:
            col_map[nid] -= min_col
            row_map[nid] -= min_row
        min_col = 0
        min_row = 0

    # Cell sizing — enforce minimum spacing
    cell_w = int(avail_w / num_cols) if num_cols > 0 else avail_w
    cell_h = int(avail_h / num_rows) if num_rows > 0 else avail_h

    # Adaptive sizing — shrink nodes when there are many
    total_nodes = len(nodes)
    node_scale = min(1.0, 8 / max(total_nodes, 1))
    node_frac_w = 0.55 * (0.6 + 0.4 * node_scale)   # 0.55 -> ~0.36 at many nodes
    node_frac_h = 0.40 * (0.6 + 0.4 * node_scale)   # 0.40 -> ~0.26 at many nodes

    # Node box sizing (fraction of cell)
    node_w = int(cell_w * node_frac_w)
    node_h = int(cell_h * node_frac_h)
    # Ensure minimum readable size (but not too large for small cells)
    node_w = max(node_w, int(Pt(base_pt) * 4))
    node_h = max(node_h, int(Pt(base_pt) * 1.8))
    # Cap nodes within available space so they cannot exceed cell bounds
    node_w = min(node_w, int(avail_w / max(num_cols, 3) * 0.80))
    node_h = min(node_h, int(cell_h * 0.55))

    # For linear layout, further constrain node size
    if use_linear:
        node_w = min(node_w, int(cell_w * 0.70))
        node_h = min(node_h, int(cell_h * 0.50))
        node_w = max(node_w, int(Pt(base_pt) * 5))
        node_h = max(node_h, int(Pt(base_pt) * 1.5))

    # Minimum spacing between node centres (increased for clarity)
    min_centre_spacing_x = node_w + max(int(node_w * 0.30), int(Pt(base_pt) * 1.5))
    min_centre_spacing_y = node_h + max(int(node_h * 0.40), int(Pt(base_pt) * 1.5))

    # Compute node centres — with overlap resolution instead of simple clamping
    centres = {}
    for n in nodes:
        nid = n["id"]
        c = col_map.get(nid, 0)
        r = row_map.get(nid, 0)
        cx = area_x + int((c + 0.5) * cell_w)
        cy = area_y + int((r + 0.5) * cell_h)
        # Soft clamp to bounds (allow slight overflow rather than pile-up)
        cx = max(area_x + node_w // 2, min(cx, area_x + avail_w - node_w // 2))
        cy = max(area_y + node_h // 2, min(cy, area_y + avail_h - node_h // 2))
        centres[nid] = (cx, cy)

    # Overlap resolution: push apart nodes that overlap, grouped by row
    row_groups = {}
    for n in nodes:
        nid = n["id"]
        r = row_map.get(nid, 0)
        row_groups.setdefault(r, []).append(nid)

    for r, nids in row_groups.items():
        if len(nids) <= 1:
            continue
        # Sort by x position
        nids.sort(key=lambda nid: centres[nid][0])
        for i in range(1, len(nids)):
            prev_cx, _ = centres[nids[i - 1]]
            cur_cx, cur_cy = centres[nids[i]]
            if cur_cx - prev_cx < min_centre_spacing_x:
                new_cx = prev_cx + min_centre_spacing_x
                new_cx = min(new_cx, area_x + avail_w - node_w // 2)
                centres[nids[i]] = (new_cx, cur_cy)

    # Also resolve vertical overlap within columns
    col_groups = {}
    for n in nodes:
        nid = n["id"]
        c = col_map.get(nid, 0)
        col_groups.setdefault(c, []).append(nid)

    for c, nids in col_groups.items():
        if len(nids) <= 1:
            continue
        nids.sort(key=lambda nid: centres[nid][1])
        for i in range(1, len(nids)):
            _, prev_cy = centres[nids[i - 1]]
            cur_cx, cur_cy = centres[nids[i]]
            if cur_cy - prev_cy < min_centre_spacing_y:
                new_cy = prev_cy + min_centre_spacing_y
                new_cy = min(new_cy, area_y + avail_h - node_h // 2)
                centres[nids[i]] = (cur_cx, new_cy)

    # --- draw edges first (behind nodes) ---
    arrow_thick = max(int(Pt(base_pt * 0.15)), int(Pt(2)))
    arrow_color = _lerp_hex(muted, primary, 0.35)

    node_map = {n["id"]: n for n in nodes}

    # Bounds for clamping arrow endpoints (with small inset for arrow heads)
    arrow_margin = int(arrow_thick * 6)
    bx0 = x + arrow_margin
    by0 = y + arrow_margin
    bx1 = x + w - arrow_margin
    by1 = y + h - arrow_margin

    def _clamp_pt(px, py):
        return (max(bx0, min(bx1, px)), max(by0, min(by1, py)))

    for edge in edges:
        src_id = edge["from"]
        dst_id = edge["to"]
        if src_id not in centres or dst_id not in centres:
            continue
        sx, sy = centres[src_id]
        ex, ey = centres[dst_id]

        # Offset start/end from node edges
        dx = ex - sx
        dy = ey - sy
        dist = math.sqrt(dx * dx + dy * dy)
        if dist < 1:
            continue

        ux = dx / dist
        uy = dy / dist

        # Start offset: from edge of source node
        src_type = node_map.get(src_id, {}).get("type", "process")
        dst_type = node_map.get(dst_id, {}).get("type", "process")

        if src_type == "decision":
            src_offset = max(node_w, node_h) * 0.5
        else:
            if abs(ux) > abs(uy):
                src_offset = node_w * 0.5
            else:
                src_offset = node_h * 0.5

        if dst_type == "decision":
            dst_offset = max(node_w, node_h) * 0.5
        else:
            if abs(ux) > abs(uy):
                dst_offset = node_w * 0.5
            else:
                dst_offset = node_h * 0.5

        # Detect upward/backward edges (destination is above or behind source)
        goes_up = (ey < sy - node_h * 0.3)

        # For linear layout, use orthogonal routing (horizontal + vertical)
        if use_linear and abs(dx) > node_w * 0.3 and abs(dy) > node_h * 0.3:
            # Orthogonal L-shaped route: go down from source, then across to dest
            a_sy = sy + node_h * 0.5  # bottom of source
            a_ey = ey - node_h * 0.5  # top of dest (if going down)
            if goes_up:
                a_sy = sy - node_h * 0.5
                a_ey = ey + node_h * 0.5

            mid_y = (a_sy + a_ey) / 2

            p1 = _clamp_pt(sx, a_sy)
            p2 = _clamp_pt(sx, mid_y)
            p3 = _clamp_pt(ex, mid_y)
            p4 = _clamp_pt(ex, a_ey)

            _draw_arrow(slide, p1[0], p1[1], p2[0], p2[1], arrow_color, arrow_thick)
            _draw_arrow(slide, p2[0], p2[1], p3[0], p3[1], arrow_color, arrow_thick)
            _draw_arrow(slide, p3[0], p3[1], p4[0], p4[1], arrow_color, arrow_thick)

            # Edge label near the horizontal segment
            edge_label = edge.get("label")
            if edge_label:
                mid_x = (sx + ex) / 2
                lbl_w = int(Pt(base_pt) * 4)
                lbl_h = int(Pt(base_pt) * 1.8)
                lbl_x = int(mid_x) - lbl_w // 2
                lbl_y = int(mid_y) - lbl_h
                lbl_x = max(x, min(lbl_x, x + w - lbl_w))
                lbl_y = max(y, min(lbl_y, y + h - lbl_h))
                _add_textbox(
                    slide, lbl_x, lbl_y, lbl_w, lbl_h,
                    edge_label, font_body, max(int(base_pt * 0.8), 9), text_c,
                    align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                )

        elif goes_up:
            a_sx = sx + node_w * 0.5
            a_sy = sy
            route_x = max(sx, ex) + node_w * 0.8
            route_x = min(route_x, x + w - int(Pt(base_pt)))
            mid_y_up = ey

            p1 = _clamp_pt(a_sx, a_sy)
            p2 = _clamp_pt(route_x, a_sy)
            p3 = _clamp_pt(route_x, mid_y_up)
            a_ex = ex + node_w * 0.5
            a_ey = ey
            p4 = _clamp_pt(a_ex, a_ey)

            _draw_arrow(slide, p1[0], p1[1], p2[0], p2[1], arrow_color, arrow_thick)
            _draw_arrow(slide, p2[0], p2[1], p3[0], p3[1], arrow_color, arrow_thick)
            _draw_arrow(slide, p3[0], p3[1], p4[0], p4[1], arrow_color, arrow_thick)

            # Edge label
            edge_label = edge.get("label")
            if edge_label:
                mid_x = route_x + int(Pt(base_pt) * 0.5)
                mid_y = (sy + ey) / 2
                lbl_w = int(Pt(base_pt) * 5)
                lbl_h = int(Pt(base_pt) * 2.0)
                lbl_x = int(mid_x) - lbl_w // 2
                lbl_y = int(mid_y) - lbl_h // 2
                lbl_x += int(Pt(base_pt) * 0.5)
                lbl_x = max(x, min(lbl_x, x + w - lbl_w))
                lbl_y = max(y, min(lbl_y, y + h - lbl_h))
                _add_textbox(
                    slide, lbl_x, lbl_y, lbl_w, lbl_h,
                    edge_label, font_body, max(int(base_pt * 0.9), 10), text_c,
                    align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                )
        else:
            a_sx = sx + ux * src_offset
            a_sy = sy + uy * src_offset
            a_ex = ex - ux * dst_offset
            a_ey = ey - uy * dst_offset

            p1 = _clamp_pt(a_sx, a_sy)
            p2 = _clamp_pt(a_ex, a_ey)
            _draw_arrow(slide, p1[0], p1[1], p2[0], p2[1], arrow_color, arrow_thick)

            # Edge label
            edge_label = edge.get("label")
            if edge_label:
                mid_x = (a_sx + a_ex) / 2
                mid_y = (a_sy + a_ey) / 2
                lbl_w = int(Pt(base_pt) * 5)
                lbl_h = int(Pt(base_pt) * 2.0)
                lbl_x = int(mid_x) - lbl_w // 2
                lbl_y = int(mid_y) - lbl_h // 2
                lbl_x += int(Pt(base_pt) * 0.5)
                lbl_x = max(x, min(lbl_x, x + w - lbl_w))
                lbl_y = max(y, min(lbl_y, y + h - lbl_h))
                _add_textbox(
                    slide, lbl_x, lbl_y, lbl_w, lbl_h,
                    edge_label, font_body, max(int(base_pt * 0.9), 10), text_c,
                    align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                )

    # Adaptive font size when many nodes
    node_font_pt = base_pt if total_nodes <= 8 else max(int(base_pt * min(1.0, 5 / total_nodes)), 7)

    # --- draw nodes ---
    for i, n in enumerate(nodes):
        nid = n["id"]
        cx, cy = centres[nid]
        ntype = n.get("type", "process")
        label = n.get("label", "")

        # Truncate long labels with ellipsis
        max_label_chars = max(12, int(node_w / (Pt(node_font_pt) * 0.5)))
        if len(label) > max_label_chars:
            label = label[:max_label_chars - 1] + "\u2026"

        t = i / max(len(nodes) - 1, 1)

        if ntype == "terminal":
            # Rounded rectangle (oval-ish ends)
            fill = _lerp_hex(primary, accent, t)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(cx - node_w // 2), Emu(cy - node_h // 2),
                Emu(node_w), Emu(node_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.fill.background()
            # Text on terminal nodes uses bg color for contrast
            _add_textbox(
                slide,
                cx - node_w // 2, cy - node_h // 2, node_w, node_h,
                label, font_body, node_font_pt, bg,
                align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )

        elif ntype == "decision":
            # Diamond
            diamond_half_w = int(node_w * 0.55)
            diamond_half_h = int(node_h * 0.65)
            fill = _lerp_hex(accent, primary, t)
            border = _lerp_hex(primary, accent, 0.5)
            _draw_diamond(slide, cx, cy, diamond_half_w, diamond_half_h, fill, border)
            # Text inside diamond (smaller area)
            text_w = int(diamond_half_w * 1.2)
            text_h = int(diamond_half_h * 1.0)
            _add_textbox(
                slide,
                cx - text_w // 2, cy - text_h // 2, text_w, text_h,
                label, font_body, max(int(node_font_pt * 0.8), 7), bg,
                align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )

        else:
            # Process — standard rectangle
            fill = _lerp_hex(bg, primary, 0.08)
            border = _lerp_hex(primary, muted, 0.4)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(cx - node_w // 2), Emu(cy - node_h // 2),
                Emu(node_w), Emu(node_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.color.rgb = _rgb(border)
            shape.line.width = Pt(1.5)
            _add_textbox(
                slide,
                cx - node_w // 2, cy - node_h // 2, node_w, node_h,
                label, font_body, node_font_pt, text_c,
                align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE,
            )

    # --- "+N more" overflow note ---
    if overflow_count > 0:
        note_text = f"+{overflow_count} more node{'s' if overflow_count > 1 else ''}"
        note_pt = max(int(base_pt * 0.75), 8)
        note_h = int(Pt(note_pt) * 2.2)
        # Place within the reserved overflow_h zone (already subtracted from avail_h)
        note_y = area_y + avail_h
        # Clamp so it never exceeds bounds
        max_bottom = y + h
        if note_y + note_h > max_bottom:
            note_y = max_bottom - note_h
        _add_textbox(
            slide, x, note_y, w, note_h,
            note_text, font_body, note_pt, muted,
            align=PP_ALIGN.RIGHT, bold=False, anchor=MSO_ANCHOR.TOP,
        )

    return None
