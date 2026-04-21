import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


DATA = {
    "title": "User signup flow",
    "direction": "TB",
    "nodes": [
        {"id": "start", "label": "Start",         "type": "terminal"},
        {"id": "input", "label": "Enter details",  "type": "process"},
        {"id": "valid", "label": "Valid?",          "type": "decision"},
        {"id": "save",  "label": "Save to DB",     "type": "process"},
        {"id": "error", "label": "Show error",      "type": "process"},
        {"id": "end",   "label": "Done",            "type": "terminal"},
    ],
    "edges": [
        {"from": "start", "to": "input"},
        {"from": "input", "to": "valid"},
        {"from": "valid", "to": "save",  "label": "Yes"},
        {"from": "valid", "to": "error", "label": "No"},
        {"from": "error", "to": "input", "label": "Retry"},
        {"from": "save",  "to": "end"},
    ],
}

DATA_MAX = {
    "title": "Complex CI/CD pipeline",
    "direction": "TB",
    "nodes": [
        {"id": "n1",  "label": "Push to repo",       "type": "terminal"},
        {"id": "n2",  "label": "Lint code",           "type": "process"},
        {"id": "n3",  "label": "Unit tests",          "type": "process"},
        {"id": "n4",  "label": "Build container",     "type": "process"},
        {"id": "n5",  "label": "Passes?",             "type": "decision"},
        {"id": "n6",  "label": "Deploy staging",      "type": "process"},
        {"id": "n7",  "label": "Integration tests",   "type": "process"},
        {"id": "n8",  "label": "Manual review",       "type": "process"},
        {"id": "n9",  "label": "Approved?",           "type": "decision"},
        {"id": "n10", "label": "Deploy prod",         "type": "process"},
        {"id": "n11", "label": "Rollback",            "type": "process"},
        {"id": "n12", "label": "Notify team",         "type": "process"},
        {"id": "n13", "label": "Done",                "type": "terminal"},
    ],
    "edges": [
        {"from": "n1",  "to": "n2"},
        {"from": "n2",  "to": "n3"},
        {"from": "n3",  "to": "n4"},
        {"from": "n4",  "to": "n5"},
        {"from": "n5",  "to": "n6",  "label": "Yes"},
        {"from": "n5",  "to": "n11", "label": "No"},
        {"from": "n6",  "to": "n7"},
        {"from": "n7",  "to": "n8"},
        {"from": "n8",  "to": "n9"},
        {"from": "n9",  "to": "n10", "label": "Yes"},
        {"from": "n9",  "to": "n11", "label": "No"},
        {"from": "n10", "to": "n12"},
        {"from": "n11", "to": "n12"},
        {"from": "n12", "to": "n13"},
    ],
}


def _set_slide_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))


def build(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA, tokens, (x, y, w, h))

    prs.save(out_path)


def build_max(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA_MAX, tokens, (x, y, w, h))

    prs.save(out_path)


def main():
    for name, tokens in MODES.items():
        path = os.path.join(HERE, f"example-{name}.pptx")
        build(name, tokens, path)
        print(f"wrote {path}")
        path_max = os.path.join(HERE, f"example-max-{name}.pptx")
        build_max(name, tokens, path_max)
        print(f"wrote {path_max}")


if __name__ == "__main__":
    main()
