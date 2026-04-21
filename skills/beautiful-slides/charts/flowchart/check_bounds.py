"""Verify that every shape rendered by flowchart stays within bounds."""

import os
import sys

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))

from tokens import MODES  # noqa: E402
from render import render  # noqa: E402
from example import DATA, DATA_MAX  # noqa: E402


def check(mode_name, tokens, data, label=""):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, data, tokens, (x, y, w, h))

    violations = []
    for shape in slide.shapes:
        sx = shape.left
        sy = shape.top
        sr = sx + shape.width
        sb = sy + shape.height
        tol = 1
        if sx < x - tol or sy < y - tol or sr > x + w + tol or sb > y + h + tol:
            violations.append(
                f"  shape at ({sx},{sy}) size ({shape.width},{shape.height}) "
                f"exceeds bounds ({x},{y},{w},{h}) — right={sr} bottom={sb}"
            )

    tag = f"{label} [{mode_name}]" if label else f"[{mode_name}]"
    if violations:
        print(f"FAIL {tag} {len(violations)} shape(s) out of bounds:")
        for v in violations:
            print(v)
        return False
    else:
        print(f"OK   {tag} all shapes within bounds")
        return True


def main():
    ok = True
    for name, tokens in MODES.items():
        if not check(name, tokens, DATA, "normal"):
            ok = False
        if not check(name, tokens, DATA_MAX, "max"):
            ok = False
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
