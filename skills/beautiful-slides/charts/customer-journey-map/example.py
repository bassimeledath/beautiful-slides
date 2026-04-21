import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


DATA = {
    "title": "SaaS onboarding journey",
    "stages": [
        {
            "label": "Awareness",
            "actions": "Sees ad, reads blog post",
            "touchpoints": "Social media, blog",
            "feelings": "Curious but skeptical",
            "sentiment": 3,
            "pain_points": "Too many competing options",
            "opportunities": "Targeted content marketing",
        },
        {
            "label": "Consideration",
            "actions": "Compares features, reads reviews",
            "touchpoints": "Website, G2, demo",
            "feelings": "Interested, cautiously optimistic",
            "sentiment": 4,
            "pain_points": "Unclear pricing tiers",
            "opportunities": "Interactive ROI calculator",
        },
        {
            "label": "Sign-up",
            "actions": "Creates account, starts trial",
            "touchpoints": "Signup page, email",
            "feelings": "Excited to try",
            "sentiment": 4,
            "pain_points": "Long form, credit card required",
            "opportunities": "Simplified one-click signup",
        },
        {
            "label": "Onboarding",
            "actions": "Follows setup wizard, imports data",
            "touchpoints": "In-app, help docs",
            "feelings": "Overwhelmed by features",
            "sentiment": 2,
            "pain_points": "Complex configuration",
            "opportunities": "Guided setup with templates",
        },
        {
            "label": "Adoption",
            "actions": "Uses daily, invites team",
            "touchpoints": "In-app, Slack integration",
            "feelings": "Productive, sees value",
            "sentiment": 5,
            "pain_points": "Missing integrations",
            "opportunities": "Team collaboration features",
        },
    ],
}


def _set_slide_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))


def build(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA, tokens, (x, y, w, h))

    prs.save(out_path)


def main():
    for name, tokens in MODES.items():
        path = os.path.join(HERE, f"example-{name}.pptx")
        build(name, tokens, path)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
