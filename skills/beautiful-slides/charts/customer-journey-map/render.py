"""Customer journey map — stage-by-stage journey with touchpoints.

Horizontal stages across the top, rows for different aspects
(actions, feelings, pain points, opportunities).
"""

import math

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(int(Pt(3)))
    tf.margin_right = Emu(int(Pt(3)))
    tf.margin_top = Emu(int(Pt(2)))
    tf.margin_bottom = Emu(int(Pt(2)))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


# Sentiment icons using simple unicode
_SENTIMENT_ICONS = {
    5: "+",   # very positive
    4: "+",   # positive
    3: "~",   # neutral
    2: "-",   # negative
    1: "-",   # very negative
}


def render(slide, data, tokens, bounds):
    """Render a customer journey map.

    data:
        title    - optional string
        stages   - list of {
            "label": str,              # stage name (e.g. "Awareness")
            "actions": str,            # what the customer does
            "touchpoints": str,        # channels / touchpoints
            "feelings": str,           # emotional state description
            "sentiment": int,          # 1-5 (1=very negative, 5=very positive)
            "pain_points": str,        # pain points / frustrations
            "opportunities": str,      # improvement opportunities
        }
        rows     - optional list of str, overrides default row labels
                   default: ["Actions", "Touchpoints", "Feelings", "Pain points", "Opportunities"]
    """
    x, y, w, h = bounds

    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]
    radius_px = tokens.get("radius_px", 0)

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()


    stages = data.get("stages", [])
    if not stages:
        return
    title = data.get("title")
    n_stages = len(stages)

    # Row definitions: key in data, default label
    default_row_defs = [
        ("actions", "Actions"),
        ("touchpoints", "Touchpoints"),
        ("feelings", "Feelings"),
        ("pain_points", "Pain points"),
        ("opportunities", "Opportunities"),
    ]
    custom_rows = data.get("rows")
    if custom_rows and len(custom_rows) == len(default_row_defs):
        row_defs = [(default_row_defs[i][0], custom_rows[i]) for i in range(len(default_row_defs))]
    else:
        row_defs = default_row_defs

    n_rows = len(row_defs)

    # --- layout ---
    cur_y = y
    if title:
        title_h = int(Pt(base_pt * 1.6) * 1.8)
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, int(base_pt * 1.5), text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.4))

    avail_h = (y + h) - cur_y
    avail_w = w

    # Row label column on the left
    label_col_w = int(avail_w * 0.12)
    gap = int(avail_w * 0.008)
    grid_x = x + label_col_w + gap
    grid_w = avail_w - label_col_w - gap

    # Stage header row at top
    header_h = int(Pt(base_pt) * 2.5)
    # Sentiment row (dots/bar between header and body)
    sentiment_h = int(Pt(base_pt) * 2.0)

    body_y = cur_y + header_h + sentiment_h
    body_h = avail_h - header_h - sentiment_h

    col_w = grid_w // n_stages
    col_gap = int(col_w * 0.03)

    # Row heights within the body
    row_gap = int(Pt(base_pt * 0.15))
    total_row_gaps = row_gap * (n_rows - 1)
    row_h = max(int((body_h - total_row_gaps) / n_rows), int(Pt(base_pt * 1.8)))

    # --- Draw stage headers ---
    for si, stage in enumerate(stages):
        sx = grid_x + col_w * si
        t = si / (n_stages - 1) if n_stages > 1 else 0.0
        header_color = _lerp_hex(primary, accent, t)

        # Header background pill/rect
        hdr_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE,
            Emu(sx + col_gap), Emu(cur_y),
            Emu(col_w - 2 * col_gap), Emu(header_h),
        )
        hdr_shape.fill.solid()
        hdr_shape.fill.fore_color.rgb = _rgb(header_color)
        hdr_shape.line.fill.background()

        # Header label
        _add_textbox(
            slide, sx + col_gap, cur_y, col_w - 2 * col_gap, header_h,
            stage.get("label", ""), font_display, max(int(base_pt * 0.85), 8), bg,
            align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )

    # --- Draw sentiment row ---
    # Sentiment dots / indicators
    sentiment_y = cur_y + header_h
    # Track line for sentiment
    track_h = max(int(Pt(1.5)), 1)
    track_cy = sentiment_y + sentiment_h // 2
    track_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(grid_x + col_gap), Emu(track_cy - track_h // 2),
        Emu(grid_w - 2 * col_gap), Emu(track_h),
    )
    track_line.fill.solid()
    track_line.fill.fore_color.rgb = _rgb(_lerp_hex(bg, muted, 0.25))
    track_line.line.fill.background()

    # Sentiment label on the left
    _add_textbox(
        slide, x, sentiment_y, label_col_w, sentiment_h,
        "Sentiment", font_body, max(int(base_pt * 0.7), 7), text_c,
        align=PP_ALIGN.RIGHT, bold=True, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Plot sentiment dots
    dot_r = int(Pt(base_pt * 0.4))
    sentiment_range_h = int(sentiment_h * 0.7)
    sentiment_mid_y = track_cy

    for si, stage in enumerate(stages):
        sentiment = stage.get("sentiment", 3)
        sentiment = max(1, min(5, sentiment))

        # Map sentiment 1-5 to vertical position (5=top, 1=bottom)
        norm = (sentiment - 3) / 2.0  # -1 to 1
        dot_cy = sentiment_mid_y - int(norm * sentiment_range_h / 2)

        dot_cx = grid_x + col_w * si + col_w // 2

        # Color: positive=accent, negative=primary, neutral=muted
        if sentiment >= 4:
            dot_color = accent
        elif sentiment <= 2:
            dot_color = primary
        else:
            dot_color = muted

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(dot_cx - dot_r), Emu(dot_cy - dot_r),
            Emu(dot_r * 2), Emu(dot_r * 2),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(dot_color)
        dot.line.fill.background()

        # Connect dots with line segments
        if si < n_stages - 1:
            next_sentiment = stages[si + 1].get("sentiment", 3)
            next_sentiment = max(1, min(5, next_sentiment))
            next_norm = (next_sentiment - 3) / 2.0
            next_dot_cy = sentiment_mid_y - int(next_norm * sentiment_range_h / 2)
            next_dot_cx = grid_x + col_w * (si + 1) + col_w // 2

            line_w_px = max(int(Pt(1.5)), 1)
            # Draw as freeform thin quad
            dx = next_dot_cx - dot_cx
            dy = next_dot_cy - dot_cy
            seg_len = math.sqrt(dx * dx + dy * dy)
            if seg_len > 0:
                px = -dy / seg_len
                py = dx / seg_len
                hw = line_w_px / 2
                p1x = int(dot_cx + px * hw)
                p1y = int(dot_cy + py * hw)
                p2x = int(dot_cx - px * hw)
                p2y = int(dot_cy - py * hw)
                p3x = int(next_dot_cx - px * hw)
                p3y = int(next_dot_cy - py * hw)
                p4x = int(next_dot_cx + px * hw)
                p4y = int(next_dot_cy + py * hw)

                ff = slide.shapes.build_freeform(p1x, p1y, scale=1.0)
                ff.add_line_segments([(p2x, p2y), (p3x, p3y), (p4x, p4y)], close=True)
                line_shape = ff.convert_to_shape()
                line_shape.fill.solid()
                line_shape.fill.fore_color.rgb = _rgb(_lerp_hex(muted, text_c, 0.3))
                line_shape.line.fill.background()

    # --- Draw body grid ---
    for ri, (row_key, row_label) in enumerate(row_defs):
        ry = body_y + ri * (row_h + row_gap)

        # Row label on the left
        _add_textbox(
            slide, x, ry, label_col_w, row_h,
            row_label, font_body, max(int(base_pt * 0.7), 7), text_c,
            align=PP_ALIGN.RIGHT, bold=True, anchor=MSO_ANCHOR.TOP,
        )

        # Horizontal divider above each row (except first)
        if ri == 0:
            div_h = max(int(Pt(0.75)), 1)
            div = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(grid_x), Emu(ry),
                Emu(grid_w), Emu(div_h),
            )
            div.fill.solid()
            div.fill.fore_color.rgb = _rgb(_lerp_hex(bg, muted, 0.2))
            div.line.fill.background()

        # Cell content for each stage
        cell_font_size = max(int(base_pt * 0.65), 7)
        for si, stage in enumerate(stages):
            cx = grid_x + col_w * si + col_gap
            cw = col_w - 2 * col_gap
            content = stage.get(row_key, "")

            _add_textbox(
                slide, cx, ry, cw, row_h,
                content, font_body, cell_font_size, text_c,
                align=PP_ALIGN.LEFT, bold=False, anchor=MSO_ANCHOR.TOP,
            )

    # --- Vertical gridlines between stage columns ---
    gridline_w = max(int(Pt(0.5)), 1)
    for si in range(1, n_stages):
        gx = grid_x + col_w * si
        gl = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(gx), Emu(cur_y + header_h),
            Emu(gridline_w), Emu(avail_h - header_h),
        )
        gl.fill.solid()
        gl.fill.fore_color.rgb = _rgb(_lerp_hex(bg, muted, 0.15))
        gl.line.fill.background()

    return None
