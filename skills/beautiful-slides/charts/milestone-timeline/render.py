"""Milestone timeline — horizontal dated milestones connected by a track line."""

import math

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 align=PP_ALIGN.CENTER, bold=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _lerp_hex(h1, h2, t):
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _estimate_lines(text, font_pt, avail_w):
    """Estimate how many lines text will wrap to given available width."""
    char_w = Pt(font_pt) * 0.55
    chars_per_line = max(1, int(avail_w / char_w))
    return max(1, math.ceil(len(text) / chars_per_line))


def _truncate(text, max_chars):
    """Truncate text with ellipsis if it exceeds max_chars."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 1].rstrip() + "\u2026"


def render(slide, data, tokens, bounds):
    """Render a horizontal milestone timeline.

    data:
        title      - optional string
        milestones - list of {"date": str, "label": str}  (3-7 items)
    """
    x, y, w, h = bounds

    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    base_pt = tokens["font_size_base_pt"]

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(bg)
    bg_shape.line.fill.background()

    milestones = data.get("milestones", [])
    if not milestones:
        return
    title = data.get("title")
    n = len(milestones)

    # --- layout ---
    cur_y = y
    title_h = 0
    if title:
        title_pt = int(base_pt * 1.5)
        title_avail_w = w
        title_lines = _estimate_lines(title, title_pt, title_avail_w)
        title_h = int(Pt(title_pt) * 1.3 * title_lines + Pt(title_pt) * 0.5)
        title_h = min(title_h, int(h * 0.20))  # cap at 20% of height
        _add_textbox(
            slide, x, cur_y, w, title_h,
            title, font_display, title_pt, text_c,
            align=PP_ALIGN.LEFT, bold=True, anchor=MSO_ANCHOR.TOP,
        )
        cur_y += title_h + int(Pt(base_pt * 0.6))

    avail_h = (y + h) - cur_y
    avail_w = w

    # Horizontal margins so first/last nodes aren't at the edge
    h_margin = int(avail_w * 0.06)
    track_x_start = x + h_margin
    track_x_end = x + w - h_margin
    track_len = track_x_end - track_x_start

    # Vertical centre for the track line
    track_cy = cur_y + avail_h // 2

    # Node (circle) radius
    node_r = int(Pt(base_pt * 0.55))

    # Track line (thin rectangle)
    line_h = max(int(Pt(2)), int(node_r * 0.18))
    track_line_y = track_cy - line_h // 2

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(track_x_start), Emu(track_line_y),
        Emu(track_len), Emu(line_h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(muted)
    shape.line.fill.background()

    # Label sizing -- reduce font when many milestones
    if n > 5:
        label_size = max(8, int(base_pt - (n - 5) * 2))
    else:
        label_size = base_pt
    date_size = max(int(base_pt * 0.78), 8)
    label_col_w = max(int(track_len / n), int(Pt(base_pt) * 8))

    # Space above/below the track for text
    text_above_h = int(avail_h * 0.38)
    text_below_h = int(avail_h * 0.38)

    for i, ms in enumerate(milestones):
        # Horizontal position: evenly spaced
        if n == 1:
            cx = track_x_start + track_len // 2
        else:
            cx = track_x_start + int(track_len * i / (n - 1))

        # Alternating above/below placement
        above = (i % 2 == 0)

        # Color interpolation for variety
        t = i / (n - 1) if n > 1 else 0.0
        node_color = _lerp_hex(primary, accent, t)

        # Draw node circle
        node_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(cx - node_r), Emu(track_cy - node_r),
            Emu(node_r * 2), Emu(node_r * 2),
        )
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = _rgb(node_color)
        node_shape.line.fill.background()

        # Draw connector line from node to label area
        connector_h = int(avail_h * 0.12)
        connector_w = max(int(Pt(1.5)), int(node_r * 0.12))

        if above:
            conn_y = track_cy - node_r - connector_h
        else:
            conn_y = track_cy + node_r

        conn_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(cx - connector_w // 2), Emu(conn_y),
            Emu(connector_w), Emu(connector_h),
        )
        conn_shape.fill.solid()
        conn_shape.fill.fore_color.rgb = _rgb(node_color)
        conn_shape.line.fill.background()

        # Text boxes for date and label
        tb_w = min(label_col_w, int(track_len * 0.35))
        tb_x = cx - tb_w // 2
        # Clamp to bounds
        tb_x = max(x, min(tb_x, x + w - tb_w))

        date_h = int(Pt(date_size) * 2)
        label_h = int(Pt(label_size) * 3.5)

        if above:
            # Date right above connector, label above date
            date_y = conn_y - date_h
            label_y = date_y - label_h
            # Clamp
            if label_y < cur_y:
                label_y = cur_y
                date_y = label_y + label_h
            date_anchor = MSO_ANCHOR.BOTTOM
            label_anchor = MSO_ANCHOR.BOTTOM
        else:
            # Date right below connector, label below date
            date_y = conn_y + connector_h
            label_y = date_y + date_h
            # Clamp
            if label_y + label_h > y + h:
                label_y = y + h - label_h
                date_y = label_y - date_h
            date_anchor = MSO_ANCHOR.TOP
            label_anchor = MSO_ANCHOR.TOP

        # Date label
        _add_textbox(
            slide, tb_x, date_y, tb_w, date_h,
            ms.get("date", ""), font_body, date_size, text_c,
            align=PP_ALIGN.CENTER, bold=False, anchor=date_anchor,
        )

        # Event label -- truncate to ~3 lines worth of characters
        label_char_w = Pt(label_size) * 0.55
        chars_per_line = max(5, int(tb_w / label_char_w))
        max_label_chars = chars_per_line * 3
        event_label = _truncate(ms.get("label", ""), max_label_chars)
        _add_textbox(
            slide, tb_x, label_y, tb_w, label_h,
            event_label, font_body, label_size, text_c,
            align=PP_ALIGN.CENTER, bold=True, anchor=label_anchor,
        )

    return None
