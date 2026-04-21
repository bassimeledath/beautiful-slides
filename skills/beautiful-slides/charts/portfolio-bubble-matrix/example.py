"""Render the portfolio-bubble-matrix chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_GE_MCKINSEY = {
    "title": "GE/McKinsey portfolio matrix",
    "x_label": "Competitive strength",
    "y_label": "Market attractiveness",
    "x_segments": ["Low", "Medium", "High"],
    "y_segments": ["High", "Medium", "Low"],
    "quadrant_labels": [
        "Selective growth", "Invest/grow", "Invest/grow",
        "Harvest", "Selective", "Selective growth",
        "Divest", "Harvest", "Selective",
    ],
    "size_label": "Revenue ($M)",
    "show_labels": True,
    "show_size_legend": True,
    "bubbles": [
        {"x": 0.85, "y": 0.90, "size": 200, "label": "Cloud Platform"},
        {"x": 0.70, "y": 0.75, "size": 150, "label": "Analytics"},
        {"x": 0.50, "y": 0.55, "size": 100, "label": "Mobile App"},
        {"x": 0.30, "y": 0.40, "size": 60,  "label": "Desktop Suite"},
        {"x": 0.15, "y": 0.20, "size": 80,  "label": "Legacy ERP"},
        {"x": 0.60, "y": 0.30, "size": 40,  "label": "IoT Sensors"},
    ],
}

DATA_INVESTMENT = {
    "title": "Investment portfolio positioning",
    "x_label": "Risk profile",
    "y_label": "Expected return",
    "x_segments": ["Conservative", "Moderate", "Aggressive"],
    "y_segments": ["High", "Medium", "Low"],
    "size_label": "AUM ($B)",
    "show_labels": True,
    "show_size_legend": True,
    "bubbles": [
        {"x": 0.10, "y": 0.35, "size": 300, "label": "Bonds"},
        {"x": 0.35, "y": 0.55, "size": 180, "label": "Blue chip"},
        {"x": 0.55, "y": 0.70, "size": 120, "label": "Growth equity"},
        {"x": 0.80, "y": 0.85, "size": 60,  "label": "Venture"},
        {"x": 0.90, "y": 0.50, "size": 40,  "label": "Crypto"},
        {"x": 0.45, "y": 0.25, "size": 90,  "label": "Real estate"},
        {"x": 0.70, "y": 0.40, "size": 70,  "label": "Private equity"},
    ],
}

DATA_SIMPLE_2X2 = {
    "title": "Product strategy matrix",
    "x_label": "Technical feasibility",
    "y_label": "Business value",
    "x_segments": ["Low", "High"],
    "y_segments": ["High", "Low"],
    "quadrant_labels": [
        "Invest & research", "Priority builds",
        "Deprioritize", "Quick wins",
    ],
    "size_label": "Team size",
    "show_labels": True,
    "show_size_legend": True,
    "bubbles": [
        {"x": 0.75, "y": 0.80, "size": 50, "label": "Feature A"},
        {"x": 0.25, "y": 0.85, "size": 30, "label": "Feature B"},
        {"x": 0.80, "y": 0.25, "size": 20, "label": "Feature C"},
        {"x": 0.20, "y": 0.20, "size": 15, "label": "Feature D"},
        {"x": 0.55, "y": 0.60, "size": 40, "label": "Feature E"},
    ],
}


MODE_DATA = {
    "sv-keynote": DATA_GE_MCKINSEY,
    "editorial-magazine": DATA_INVESTMENT,
    "playful-marketing": DATA_SIMPLE_2X2,
    "consulting-boardroom": DATA_GE_MCKINSEY,
    "craft-minimal": DATA_INVESTMENT,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
