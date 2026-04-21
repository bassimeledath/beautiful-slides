"""Portfolio bubble matrix renderer — native python-pptx shapes only.

GE/McKinsey-style portfolio view: 2x2 (or 3x3) quadrant grid combined with
area-scaled bubbles. Each bubble sits at an (x, y) position within the grid
and its radius encodes a third variable (e.g., revenue, headcount).

Public API: render(slide, data, tokens, bounds)
"""

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ---- helpers ---------------------------------------------------------------

def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _luminance(hex_):
    """Return relative luminance (0-1) of a hex color."""
    h = hex_.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255.0, int(h[2:4], 16) / 255.0, int(h[4:6], 16) / 255.0
    r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
    g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
    b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_text(fill_hex, light="#FFFFFF", dark="#1A1A1A"):
    """Return dark text if fill is light, light text if fill is dark."""
    return dark if _luminance(fill_hex) > 0.35 else light


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


def _add_circle(slide, cx, cy, r_emu, fill_hex, line_hex=None, line_width_pt=0):
    d = int(r_emu * 2)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(cx - r_emu)), Emu(int(cy - r_emu)),
        Emu(d), Emu(d),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Pt(line_width_pt)
    shp.shadow.inherit = False
    return shp


def _add_line(slide, x1, y1, x2, y2, color_hex, weight_pt):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(int(x1)), Emu(int(y1)),
        Emu(int(x2)), Emu(int(y2)),
    )
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(weight_pt)
    return conn


def _add_rect_bg(slide, x, y, w, h, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_rect_fill(slide, x, y, w, h, color_hex, alpha=None):
    """Rectangle with optional fill transparency (0-100000 scale)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _fmt(v):
    if abs(v - round(v)) < 1e-6:
        return f"{int(round(v))}"
    return f"{v:.1f}"


# ---- main render -----------------------------------------------------------

def render(slide, data, tokens, bounds):
    """Render a GE/McKinsey portfolio bubble matrix onto *slide*."""
    x0, y0, w0, h0 = bounds

    # --- unpack tokens -------------------------------------------------------
    bg = tokens["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # --- unpack data ---------------------------------------------------------
    title = data.get("title")
    x_label = data.get("x_label")           # e.g., "Competitive strength"
    y_label = data.get("y_label")           # e.g., "Market attractiveness"
    x_segments = data.get("x_segments") or ["Low", "Medium", "High"]
    y_segments = data.get("y_segments") or ["Low", "Medium", "High"]
    quadrant_labels = data.get("quadrant_labels")  # optional NxM grid labels
    bubbles = list(data.get("bubbles") or [])[:20]
    size_label = data.get("size_label")
    show_labels = bool(data.get("show_labels", True))
    show_size_legend = bool(data.get("show_size_legend", True))
    min_radius_pt = data.get("min_radius_pt", 6)
    max_radius_pt = data.get("max_radius_pt", 28)

    if not bubbles:
        return

    n_cols = len(x_segments)
    n_rows = len(y_segments)

    # --- background ----------------------------------------------------------
    _add_rect_bg(slide, x0, y0, w0, h0, bg)

    # --- outer padding -------------------------------------------------------
    pad = int(min(w0, h0) * 0.035)
    ix, iy = x0 + pad, y0 + pad
    iw, ih = w0 - 2 * pad, h0 - 2 * pad

    # --- title ---------------------------------------------------------------
    title_h = 0
    if title:
        title_pt = int(round(base_pt * 1.55))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_textbox(slide, ix, iy, iw, title_h, title,
                     font_display, title_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + int(pad * 0.3)
        ih = (y0 + h0 - pad) - iy

    # --- compute grid area ---------------------------------------------------
    label_pt = max(8, int(base_pt * 0.85))
    axis_label_pt = max(9, int(base_pt * 0.9))
    seg_label_pt = max(8, int(base_pt * 0.78))

    # Reserve space for axis labels / segment headers
    # y-axis label now rendered as a subtitle row above the grid, so left
    # margin only needs to accommodate segment labels, not the full axis title.
    left_margin = int(Pt(seg_label_pt).emu * 3.5) if y_label else int(Pt(seg_label_pt).emu * 1.5)
    bottom_margin = int(Pt(seg_label_pt).emu * 2.0) + (
        int(Pt(axis_label_pt).emu * 1.6) if x_label else 0
    )
    top_margin = int(Pt(axis_label_pt).emu * 2.0) if y_label else int(Pt(seg_label_pt).emu * 1.5)
    legend_w = int(iw * 0.18) if (show_size_legend and size_label) else 0
    right_margin = int(Pt(seg_label_pt).emu * 1.0) + legend_w

    grid_x = ix + left_margin
    grid_y = iy + top_margin
    grid_w = max(1, iw - left_margin - right_margin)
    grid_h = max(1, ih - top_margin - bottom_margin)

    if grid_w <= 0 or grid_h <= 0:
        return

    cell_w = grid_w / n_cols
    cell_h = grid_h / n_rows

    # --- draw grid cells with subtle alternating fill -----------------------
    for r in range(n_rows):
        for c in range(n_cols):
            cx = grid_x + c * cell_w
            cy = grid_y + r * cell_h
            # Subtle checker pattern using muted at very low visual weight
            if (r + c) % 2 == 0:
                _add_rect_fill(slide, cx, cy, cell_w, cell_h, bg)
            else:
                _add_rect_fill(slide, cx, cy, cell_w, cell_h, muted)
                # Override with near-bg tint by overlaying bg at high alpha
                _add_rect_fill(slide, cx, cy, cell_w, cell_h, bg)

    # --- grid lines ----------------------------------------------------------
    line_weight = 1.0
    for c in range(n_cols + 1):
        lx = grid_x + c * cell_w
        _add_line(slide, lx, grid_y, lx, grid_y + grid_h, muted, line_weight)
    for r in range(n_rows + 1):
        ly = grid_y + r * cell_h
        _add_line(slide, grid_x, ly, grid_x + grid_w, ly, muted, line_weight)

    # --- quadrant labels (centered in each cell) ----------------------------
    # Collision detection helpers
    occupied_rects = []  # list of (x1, y1, x2, y2)
    lbl_gap = int(Pt(1).emu)

    def _rects_overlap(r1, r2):
        return r1[0] < r2[2] and r1[2] > r2[0] and r1[1] < r2[3] and r1[3] > r2[1]

    def _any_collision(rect):
        for occ in occupied_rects:
            if _rects_overlap(rect, occ):
                return True
        return False

    if quadrant_labels:
        q_label_pt = max(8, int(base_pt * 0.72))
        q_label_h = int(Pt(q_label_pt).emu * 1.8)
        for idx, ql in enumerate(quadrant_labels):
            if not ql:
                continue
            # Row-major: idx = row * n_cols + col
            r = idx // n_cols
            c = idx % n_cols
            if r >= n_rows or c >= n_cols:
                continue
            cx = grid_x + c * cell_w
            cy = grid_y + r * cell_h
            # Use luminance-aware text color against the cell background
            cell_bg_color = bg  # cells are bg-colored
            q_text_color = _contrast_text(cell_bg_color)
            _add_textbox(slide,
                         cx + int(cell_w * 0.05),
                         cy + int(cell_h * 0.05),
                         int(cell_w * 0.9), q_label_h,
                         ql, font_body, q_label_pt, q_text_color,
                         bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
            # Record as exclusion zone for bubble labels
            occupied_rects.append((
                cx + int(cell_w * 0.05), cy + int(cell_h * 0.05),
                cx + int(cell_w * 0.95), cy + int(cell_h * 0.05) + q_label_h,
            ))

    # --- segment labels along axes ------------------------------------------
    seg_h = int(Pt(seg_label_pt).emu * 1.6)

    # X-axis segment labels (below grid)
    for c, seg in enumerate(x_segments):
        cx = grid_x + c * cell_w
        _add_textbox(slide,
                     cx, grid_y + grid_h + int(pad * 0.15),
                     cell_w, seg_h,
                     seg, font_body, seg_label_pt, text_c,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Y-axis segment labels (left of grid, top-to-bottom = high-to-low visually)
    y_seg_w = left_margin - int(pad * 0.3)
    for r, seg in enumerate(y_segments):
        cy = grid_y + r * cell_h
        _add_textbox(slide,
                     ix, cy,
                     y_seg_w, cell_h,
                     seg, font_body, seg_label_pt, text_c,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # --- axis labels ---------------------------------------------------------
    ax_label_h = int(Pt(axis_label_pt).emu * 1.6)

    if x_label:
        _add_textbox(slide,
                     grid_x, grid_y + grid_h + seg_h + int(pad * 0.2),
                     grid_w, ax_label_h,
                     x_label, font_body, axis_label_pt, text_c,
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    if y_label:
        # Place y-axis label as a horizontal subtitle row above the grid
        # so it doesn't wrap illegibly in a narrow left column.
        y_lbl_display = y_label if len(y_label) <= 50 else y_label[:47] + "..."
        _add_textbox(slide,
                     grid_x, grid_y - ax_label_h - int(pad * 0.1),
                     grid_w, ax_label_h,
                     y_lbl_display, font_body, axis_label_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)

    # --- size scaling (area-proportional) ------------------------------------
    sizes = [b.get("size", 1) for b in bubbles]
    s_min = min(sizes)
    s_max = max(sizes)
    if s_max == s_min:
        s_max = s_min + 1

    def size_to_r_emu(s):
        frac = (s - s_min) / (s_max - s_min)
        min_area = math.pi * min_radius_pt ** 2
        max_area = math.pi * max_radius_pt ** 2
        area = min_area + frac * (max_area - min_area)
        r_pt = math.sqrt(area / math.pi)
        return int(Pt(r_pt).emu)

    # --- draw bubbles (largest first) ----------------------------------------
    sorted_bubbles = sorted(bubbles, key=lambda b: b.get("size", 1), reverse=True)
    point_label_pt = max(7, int(base_pt * 0.7))
    point_label_h = int(Pt(point_label_pt).emu * 1.4)

    for b in sorted_bubbles:
        # x, y are 0..1 fractions within the grid (0=left/bottom, 1=right/top)
        fx = float(b.get("x", 0.5))
        fy = float(b.get("y", 0.5))
        fx = max(0.0, min(1.0, fx))
        fy = max(0.0, min(1.0, fy))

        cx = grid_x + fx * grid_w
        cy = grid_y + (1.0 - fy) * grid_h  # invert y: 1 = top

        r = size_to_r_emu(b.get("size", 1))
        color = b.get("color") or primary

        # Clamp so bubble stays inside grid
        cx = max(grid_x + r, min(grid_x + grid_w - r, cx))
        cy = max(grid_y + r, min(grid_y + grid_h - r, cy))

        _add_circle(slide, cx, cy, r, color, line_hex=bg, line_width_pt=1)

        if show_labels and b.get("label"):
            lbl = str(b["label"])
            # Cap label length at 20 chars
            if len(lbl) > 20:
                lbl = lbl[:17] + "..."
            est_w = int(Pt(point_label_pt).emu * 0.55 * len(lbl))
            est_w = max(est_w, int(Pt(point_label_pt).emu * 2))

            # Try to center label inside bubble if it fits
            if est_w < r * 1.6 and point_label_h < r * 1.6:
                lx = cx - est_w / 2
                ly = cy - point_label_h / 2
                rect = (lx - lbl_gap, ly - lbl_gap,
                        lx + est_w + lbl_gap, ly + point_label_h + lbl_gap)
                if not _any_collision(rect):
                    _add_textbox(slide, lx, ly, est_w, point_label_h,
                                 lbl, font_body, point_label_pt, bg,
                                 bold=True,
                                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                                 word_wrap=False)
                    occupied_rects.append(rect)
                    continue

            # Try multiple external positions: right, left, above, below
            spacing = r + int(Pt(2).emu)
            candidates = [
                (cx + spacing, cy - point_label_h / 2),         # right
                (cx - spacing - est_w, cy - point_label_h / 2), # left
                (cx - est_w / 2, cy - spacing - point_label_h), # above
                (cx - est_w / 2, cy + spacing),                 # below
            ]
            for lx, ly in candidates:
                lx = max(grid_x, min(grid_x + grid_w - est_w, lx))
                ly = max(grid_y, min(grid_y + grid_h - point_label_h, ly))
                rect = (lx - lbl_gap, ly - lbl_gap,
                        lx + est_w + lbl_gap, ly + point_label_h + lbl_gap)
                if not _any_collision(rect):
                    _add_textbox(slide, lx, ly, est_w, point_label_h,
                                 lbl, font_body, point_label_pt, text_c,
                                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                                 word_wrap=False)
                    occupied_rects.append(rect)
                    break
            # If no position works, skip this label to avoid overlap

    # --- size legend ---------------------------------------------------------
    if show_size_legend and size_label and legend_w > 0:
        tick_pt = max(8, int(base_pt * 0.8))
        tick_label_h = int(Pt(tick_pt).emu * 1.4)

        legend_x = grid_x + grid_w + int(pad * 0.5)
        legend_y = grid_y + int(grid_h * 0.1)
        legend_inner_w = legend_w - int(pad * 0.5)
        # Hard limit: legend must not extend beyond bounds
        legend_bottom_limit = y0 + h0 - pad

        leg_title_h = int(Pt(tick_pt).emu * 1.6)
        _add_textbox(slide, legend_x, legend_y, legend_inner_w, leg_title_h,
                     size_label, font_body, tick_pt, text_c,
                     bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        legend_y += leg_title_h + int(pad * 0.5)

        ref_vals = [s_min, (s_min + s_max) / 2, s_max]
        max_r_legend = int(legend_inner_w * 0.16)

        for rv in ref_vals:
            r = size_to_r_emu(rv)
            r_draw = min(r, max_r_legend)

            if legend_y + r_draw * 2 > legend_bottom_limit:
                break

            circle_cx = legend_x + r_draw + int(pad * 0.3)
            circle_cy = legend_y + r_draw
            _add_circle(slide, circle_cx, circle_cy, r_draw, muted)

            lbl_x = circle_cx + r_draw + int(pad * 0.4)
            lbl_w = max(1, (legend_x + legend_inner_w) - lbl_x)
            if lbl_w > 0:
                _add_textbox(slide, lbl_x, circle_cy - tick_label_h / 2,
                             lbl_w, tick_label_h,
                             _fmt(rv), font_mono, tick_pt, text_c,
                             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

            # Increase vertical spacing between reference circles
            legend_y += r_draw * 2 + int(pad * 1.2)
