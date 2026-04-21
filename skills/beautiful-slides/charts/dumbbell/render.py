"""Dumbbell chart renderer -- native python-pptx shapes only.

Two dots connected by a line per category, showing the gap between two
values (before/after, plan/actual, male/female, etc.).  Horizontal layout.

Public API: render(slide, data, tokens, bounds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# ── constants ──────────────────────────────────────────────────────────
_EMU_PER_PX = 9525


# ── helpers ────────────────────────────────────────────────────────────
def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _set_fill(shape, hex_):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_)


def _no_line(shape):
    shape.line.fill.background()


def _style_run(run, font_name, size_pt, hex_color, bold=False):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold = bold


def _add_text(slide, x, y, w, h, text, font_name, size_pt, hex_color,
              bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _style_run(run, font_name, size_pt, hex_color, bold=bold)
    return tb


def _add_line(slide, x1, y1, x2, y2, hex_color, width_emu):
    ln = slide.shapes.add_connector(
        1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)),
    )
    ln.line.color.rgb = _rgb(hex_color)
    ln.line.width = Emu(int(width_emu))
    return ln


def _add_dot(slide, cx, cy, diameter, hex_color):
    r = diameter // 2
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(cx - r)), Emu(int(cy - r)),
        Emu(int(diameter)), Emu(int(diameter)),
    )
    _set_fill(dot, hex_color)
    _no_line(dot)
    return dot


def _nice_ticks(vmin, vmax, target=5):
    if vmax <= vmin:
        return [vmin, vmin + 1], vmin, vmin + 1
    import math
    span = vmax - vmin
    raw = span / target
    magnitude = 10 ** math.floor(math.log10(raw)) if raw > 0 else 1
    residual = raw / magnitude
    if residual < 1.5:
        step = 1 * magnitude
    elif residual < 3:
        step = 2 * magnitude
    elif residual < 7:
        step = 5 * magnitude
    else:
        step = 10 * magnitude
    lo = step * math.floor(vmin / step)
    hi = step * math.ceil(vmax / step)
    ticks = []
    v = lo
    while v <= hi + 1e-9:
        ticks.append(v)
        v += step
    return ticks, lo, hi


def _fmt(v, suffix=""):
    if abs(v - round(v)) < 1e-9:
        s = f"{int(round(v))}"
    else:
        s = f"{v:.1f}"
    return f"{s}{suffix}" if suffix else s


# ── public API ─────────────────────────────────────────────────────────
def render(slide, data, tokens, bounds):
    x, y, w, h = bounds

    title = data.get("title")
    items = list(data.get("items") or [])
    series_names = data.get("series_names") or ["A", "B"]
    value_suffix = data.get("value_suffix") or ""
    show_values = bool(data.get("show_values", True))

    if not items:
        return

    items = items[:15]

    bg = tokens["bg"]
    primary = tokens["primary"]
    accent = tokens["accent"]
    text_c = tokens["text"]
    muted = tokens["muted"]
    font_display = tokens["font_display"]
    font_body = tokens["font_body"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])

    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
        Emu(int(w)), Emu(int(h)),
    )
    _set_fill(bg_shape, bg)
    _no_line(bg_shape)

    # Padding
    pad = int(min(w, h) * 0.04)
    ix, iy = x + pad, y + pad
    iw, ih = w - 2 * pad, h - 2 * pad

    # Title
    if title:
        title_pt = int(round(base_pt * 1.6))
        title_h = int(Pt(title_pt).emu * 1.8)
        _add_text(slide, ix, iy, iw, title_h, title,
                  font_name=font_display, size_pt=title_pt, hex_color=text_c,
                  bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        iy += title_h + int(Pt(base_pt).emu * 0.6)
        ih = (y + h - pad) - iy

    # Legend (two dots with labels)
    legend_pt = max(int(base_pt * 0.85), 8)
    legend_h = int(Pt(legend_pt).emu * 2.0)
    _draw_legend(slide, ix, iy, iw, legend_h, series_names,
                 [primary, accent], font_body, legend_pt, text_c)
    iy += legend_h
    ih -= legend_h

    # Gather all values for axis range
    all_vals = []
    for item in items:
        all_vals.append(item.get("value_a", 0))
        all_vals.append(item.get("value_b", 0))

    vmin_raw = min(all_vals) if all_vals else 0
    vmax_raw = max(all_vals) if all_vals else 1

    # If all values are positive and close to zero, start from zero
    if vmin_raw >= 0:
        vmin_raw = 0

    ticks, axis_lo, axis_hi = _nice_ticks(vmin_raw, vmax_raw)
    axis_span = axis_hi - axis_lo if axis_hi > axis_lo else 1

    cat_pt = max(int(base_pt * 0.92), 9)
    tick_pt = max(int(base_pt * 0.78), 8)
    val_pt = max(int(base_pt * 0.78), 8)

    # Adaptive font scaling for dense data
    n = len(items)
    if n > 10:
        cat_pt = max(int(base_pt * 0.78), 8)

    # Layout: horizontal, labels on left
    labels = [str(item.get("label", "")) for item in items]
    # Ellipsis truncation for long labels
    max_display_chars = 24
    labels = [
        (l[:max_display_chars - 1] + "\u2026" if len(l) > max_display_chars else l)
        for l in labels
    ]
    max_cat_chars = max((len(l) for l in labels), default=4)
    left_margin = int(Pt(cat_pt).emu * 0.55 * min(max_cat_chars + 2, 24))
    bottom_margin = int(Pt(tick_pt).emu * 2.0)
    top_margin = int(Pt(val_pt).emu * 0.6)
    right_margin = int(Pt(val_pt).emu * 3.5)

    plot_x = ix + left_margin
    plot_y = iy + top_margin
    plot_w = max(1, iw - left_margin - right_margin)
    plot_h = max(1, ih - top_margin - bottom_margin)

    hairline = max(int(_EMU_PER_PX * 0.5), 3175)

    # X-axis tick labels and gridlines
    for i, t in enumerate(ticks):
        tx = plot_x + int(plot_w * ((t - axis_lo) / axis_span))
        tlh = int(Pt(tick_pt).emu * 1.4)
        label = _fmt(t) + (value_suffix if t != 0 and i == len(ticks) - 1 else "")
        _add_text(
            slide, tx - int(Pt(tick_pt).emu * 2),
            plot_y + plot_h + int(Pt(tick_pt).emu * 0.3),
            int(Pt(tick_pt).emu * 4), tlh, label,
            font_name=font_mono, size_pt=tick_pt, hex_color=text_c,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )
        if i > 0:
            _add_line(slide, tx, plot_y, tx, plot_y + plot_h, muted, hairline)

    # Baseline (left axis)
    base_x = plot_x + int(plot_w * ((0 - axis_lo) / axis_span)) if axis_lo <= 0 <= axis_hi else plot_x
    _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, muted, hairline * 2)

    # Minimum row height guard -- trim items if rows would be too tight
    row_h = plot_h / n
    min_row_h = int(Pt(cat_pt).emu * 1.6)
    if row_h < min_row_h:
        n = max(1, int(plot_h // min_row_h))
        items = items[:n]
        labels = labels[:n]
        row_h = plot_h / n
    dot_diameter = int(min(row_h * 0.32, Pt(base_pt).emu * 1.1))
    connector_width = max(int(_EMU_PER_PX * 2.0), int(dot_diameter * 0.22))

    for i, item in enumerate(items):
        label = labels[i]
        va = item.get("value_a", 0)
        vb = item.get("value_b", 0)

        row_top = plot_y + i * row_h
        cy = int(row_top + row_h / 2)

        # Category label
        cat_h_em = int(Pt(cat_pt).emu * 1.4)
        _add_text(
            slide, ix, row_top + (row_h - cat_h_em) / 2,
            left_margin - int(Pt(cat_pt).emu * 0.4), cat_h_em, label,
            font_name=font_body, size_pt=cat_pt, hex_color=text_c,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Positions
        xa = plot_x + int(plot_w * ((va - axis_lo) / axis_span))
        xb = plot_x + int(plot_w * ((vb - axis_lo) / axis_span))

        # Connecting line (between the two dots)
        x_left = min(xa, xb)
        x_right = max(xa, xb)
        if x_right > x_left:
            _add_line(slide, x_left, cy, x_right, cy, muted, connector_width)

        # Dots (draw after line so they sit on top)
        _add_dot(slide, xa, cy, dot_diameter, primary)
        _add_dot(slide, xb, cy, dot_diameter, accent)

        # Value labels
        if show_values:
            vh = int(Pt(val_pt).emu * 1.3)
            vw = int(Pt(val_pt).emu * 3.5)

            gap_emu = abs(xa - xb)
            close_threshold = dot_diameter * 3

            if gap_emu < close_threshold:
                # Values are close -- place labels above/below to avoid collision
                _add_text(
                    slide, xa - vw // 2, cy - vh - dot_diameter // 2, vw, vh,
                    _fmt(va, value_suffix),
                    font_name=font_mono, size_pt=val_pt, hex_color=primary,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
                )
                _add_text(
                    slide, xb - vw // 2, cy + dot_diameter // 2, vw, vh,
                    _fmt(vb, value_suffix),
                    font_name=font_mono, size_pt=val_pt, hex_color=accent,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                )
            elif xa <= xb:
                # A is left, B is right
                _add_text(
                    slide, xa - vw - dot_diameter // 2, cy - vh // 2, vw, vh,
                    _fmt(va, value_suffix),
                    font_name=font_mono, size_pt=val_pt, hex_color=primary,
                    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                )
                _add_text(
                    slide, xb + dot_diameter // 2 + int(Pt(val_pt).emu * 0.2),
                    cy - vh // 2, vw, vh,
                    _fmt(vb, value_suffix),
                    font_name=font_mono, size_pt=val_pt, hex_color=accent,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                )
            else:
                # B is left, A is right
                _add_text(
                    slide, xb - vw - dot_diameter // 2, cy - vh // 2, vw, vh,
                    _fmt(vb, value_suffix),
                    font_name=font_mono, size_pt=val_pt, hex_color=accent,
                    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                )
                _add_text(
                    slide, xa + dot_diameter // 2 + int(Pt(val_pt).emu * 0.2),
                    cy - vh // 2, vw, vh,
                    _fmt(va, value_suffix),
                    font_name=font_mono, size_pt=val_pt, hex_color=primary,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                )


def _draw_legend(slide, x, y, w, h, names, colors, font_body, pt, text_c):
    swatch = int(Pt(pt).emu * 0.9)
    gap = int(Pt(pt).emu * 0.4)
    item_gap = int(Pt(pt).emu * 1.2)

    widths = []
    for name in names:
        est_text_w = int(Pt(pt).emu * 0.55 * max(len(name), 1))
        widths.append(swatch + gap + est_text_w)
    total = sum(widths) + item_gap * (len(widths) - 1)
    cursor = x + max(0, w - total)
    cy = y + (h - swatch) // 2

    for name, color, wd in zip(names, colors, widths):
        sw = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(int(cursor)), Emu(int(cy)),
            Emu(swatch), Emu(swatch),
        )
        _set_fill(sw, color)
        _no_line(sw)
        tx = cursor + swatch + gap
        tw = wd - swatch - gap
        _add_text(slide, tx, y, tw, h, name,
                  font_name=font_body, size_pt=pt, hex_color=text_c,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        cursor += wd + item_gap
