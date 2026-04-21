"""Render the dumbbell chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_SATISFACTION = {
    "title": "Satisfaction score, 2023 vs 2024",
    "series_names": ["2023", "2024"],
    "items": [
        {"label": "Onboarding", "value_a": 72, "value_b": 88},
        {"label": "Support",    "value_a": 65, "value_b": 79},
        {"label": "Billing",    "value_a": 58, "value_b": 61},
        {"label": "Product",    "value_a": 80, "value_b": 85},
        {"label": "Docs",       "value_a": 55, "value_b": 74},
        {"label": "Pricing",    "value_a": 62, "value_b": 68},
    ],
    "value_suffix": "%",
    "show_values": True,
}

DATA_PLAN_ACTUAL = {
    "title": "Revenue plan vs actual by region",
    "series_names": ["Plan", "Actual"],
    "items": [
        {"label": "North America", "value_a": 12.0, "value_b": 13.4},
        {"label": "EMEA",          "value_a": 8.5,  "value_b": 7.9},
        {"label": "APAC",          "value_a": 5.0,  "value_b": 5.8},
        {"label": "LATAM",         "value_a": 3.0,  "value_b": 2.4},
    ],
    "value_suffix": "M",
    "show_values": True,
}


MODE_DATA = {
    "sv-keynote": DATA_SATISFACTION,
    "editorial-magazine": DATA_PLAN_ACTUAL,
    "playful-marketing": DATA_SATISFACTION,
    "consulting-boardroom": DATA_PLAN_ACTUAL,
    "craft-minimal": DATA_SATISFACTION,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
