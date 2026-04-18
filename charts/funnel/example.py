import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

from render import render, _rgb  # noqa: E402


MODES = {
    "sv-keynote": {
        "primary": "#21D4FD",
        "accent":  "#17B26A",
        "text":    "#F5F7FA",
        "muted":   "#9AA4B2",
        "bg":      "#05070A",
        "font_display": "Manrope",
        "font_body":    "Manrope",
        "font_mono":    "JetBrains Mono",
        "font_size_base_pt": 18,
        "radius_px": 6,
    },
    "editorial-magazine": {
        "primary": "#8C2F39",
        "accent":  "#9C5B00",
        "text":    "#181514",
        "muted":   "#6F675F",
        "bg":      "#F6F1E8",
        "font_display": "Fraunces",
        "font_body":    "Newsreader",
        "font_mono":    "IBM Plex Mono",
        "font_size_base_pt": 16,
        "radius_px": 0,
    },
    "playful-marketing": {
        "primary": "#FF7A00",
        "accent":  "#0AB39C",
        "text":    "#1B1B1F",
        "muted":   "#6E6A73",
        "bg":      "#FFF4EB",
        "font_display": "Bricolage Grotesque",
        "font_body":    "Plus Jakarta Sans",
        "font_mono":    "Recursive Mono",
        "font_size_base_pt": 18,
        "radius_px": 12,
    },
    "consulting-boardroom": {
        "primary": "#0F4C81",
        "accent":  "#05603A",
        "text":    "#101828",
        "muted":   "#475467",
        "bg":      "#FFFFFF",
        "font_display": "Public Sans",
        "font_body":    "Public Sans",
        "font_mono":    "Public Sans",
        "font_size_base_pt": 14,
        "radius_px": 0,
    },
    "craft-minimal": {
        "primary": "#7C8571",
        "accent":  "#9A6B39",
        "text":    "#22201C",
        "muted":   "#7B776F",
        "bg":      "#FCFBF8",
        "font_display": "Instrument Serif",
        "font_body":    "Instrument Sans",
        "font_mono":    "Instrument Sans",
        "font_size_base_pt": 16,
        "radius_px": 2,
    },
}


DATA = {
    "title": "Enterprise pipeline, Q1",
    "stages": [
        {"label": "Leads",       "value": 10000},
        {"label": "Qualified",   "value":  4200},
        {"label": "Demo'd",      "value":  1800},
        {"label": "Proposal",    "value":   620},
        {"label": "Negotiation", "value":   310},
        {"label": "Closed-won",  "value":   180},
    ],
    "show_conversion": True,
    "value_format": "{:,}",
}


def _set_slide_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))


def build(mode_name, tokens, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _set_slide_bg(slide, tokens["bg"])

    margin = Inches(0.6)
    x = int(margin)
    y = int(margin)
    w = int(prs.slide_width) - 2 * int(margin)
    h = int(prs.slide_height) - 2 * int(margin)

    render(slide, DATA, tokens, (x, y, w, h))

    prs.save(out_path)


def main():
    out_dir = os.path.join(HERE, "renders")
    os.makedirs(out_dir, exist_ok=True)
    for name, tokens in MODES.items():
        path = os.path.join(out_dir, f"funnel_{name}.pptx")
        build(name, tokens, path)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
