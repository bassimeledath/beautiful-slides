"""Run `render()` 5 times — once per mode — writing example-<mode>.pptx per run."""

import os
import sys

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from render import render  # noqa: E402


MODES = {
    "sv-keynote": {
        "primary": "#21D4FD",
        "accent":  "#17B26A",
        "text":    "#F5F7FA",
        "muted":   "#9AA4B2",
        "bg":      "#05070A",
        "font_display": "Manrope",
        "font_body":    "Manrope",
        "font_mono":    "JetBrains Mono",
        "font_size_base_pt": 18,
        "radius_px": 6,
    },
    "editorial-magazine": {
        "primary": "#8C2F39",
        "accent":  "#9C5B00",
        "text":    "#181514",
        "muted":   "#6F675F",
        "bg":      "#F6F1E8",
        "font_display": "Fraunces",
        "font_body":    "Newsreader",
        "font_mono":    "IBM Plex Mono",
        "font_size_base_pt": 16,
        "radius_px": 0,
    },
    "playful-marketing": {
        "primary": "#FF7A00",
        "accent":  "#0AB39C",
        "text":    "#1B1B1F",
        "muted":   "#6E6A73",
        "bg":      "#FFF4EB",
        "font_display": "Bricolage Grotesque",
        "font_body":    "Plus Jakarta Sans",
        "font_mono":    "Recursive Mono",
        "font_size_base_pt": 18,
        "radius_px": 12,
    },
    "consulting-boardroom": {
        "primary": "#0F4C81",
        "accent":  "#05603A",
        "text":    "#101828",
        "muted":   "#475467",
        "bg":      "#FFFFFF",
        "font_display": "Public Sans",
        "font_body":    "Public Sans",
        "font_mono":    "Public Sans",
        "font_size_base_pt": 14,
        "radius_px": 0,
    },
    "craft-minimal": {
        "primary": "#7C8571",
        "accent":  "#9A6B39",
        "text":    "#22201C",
        "muted":   "#7B776F",
        "bg":      "#FCFBF8",
        "font_display": "Instrument Serif",
        "font_body":    "Instrument Sans",
        "font_mono":    "Instrument Sans",
        "font_size_base_pt": 16,
        "radius_px": 2,
    },
}


DATA = {
    "title": "Revenue retention by cohort, first 12 months",
    "x_labels": ["M0", "M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8", "M9", "M10", "M11", "M12"],
    "series": [
        {"name": "Q1 '25", "values": [100, 98, 97, 98, 100, 101, 102, 103, 104, 104, 105, 105, 106]},
        {"name": "Q2 '25", "values": [100, 99, 100, 102, 104, 105, 106, 106, 107, 108, 109, 110, 110]},
        {"name": "Q3 '25", "values": [100, 102, 106, 109, 110, 112, 112, 113, 113]},
        {"name": "Q4 '25", "values": [100, 104, 111, 118]},
    ],
    "x_label": "Months from cohort start",
    "y_label": "Revenue retention (%)",
    "emphasize_last_series": True,
    "end_labels": True,
}


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    for mode_name, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        margin = Inches(0.4)
        bounds = (
            int(margin),
            int(margin),
            int(prs.slide_width - 2 * margin),
            int(prs.slide_height - 2 * margin),
        )
        render(slide, DATA, tokens, bounds)

        out_path = os.path.join(out_dir, f"example-{mode_name}.pptx")
        prs.save(out_path)
        print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
