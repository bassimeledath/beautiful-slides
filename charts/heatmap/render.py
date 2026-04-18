from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _rgb(hex_):
    return RGBColor.from_string(hex_.lstrip("#"))


def _lerp_hex(h1, h2, t):
    t = max(0.0, min(1.0, t))
    h1 = h1.lstrip("#")
    h2 = h2.lstrip("#")
    r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
    r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
    r = round(r1 + (r2 - r1) * t)
    g = round(g1 + (g2 - g1) * t)
    b = round(b1 + (b2 - b1) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _luminance(hex_):
    h = hex_.lstrip("#")
    r = int(h[0:2], 16) / 255.0
    g = int(h[2:4], 16) / 255.0
    b = int(h[4:6], 16) / 255.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _set_text(tf, text, font_name, size_pt, color_hex, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        for r in list(p.runs):
            r.text = ""
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)


def _add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_w_emu=0, radius_emu=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_emu > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    if radius_emu > 0:
        try:
            short = min(w, h)
            adj = max(0.0, min(0.5, (radius_emu / short)))
            shp.adjustments[0] = adj
        except Exception:
            pass
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None or line_w_emu <= 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Emu(int(line_w_emu))
    shp.text_frame.text = ""
    tf = shp.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return shp


def _add_textbox(slide, x, y, w, h, text, font_name, size_pt, color_hex, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    _set_text(tb.text_frame, text, font_name, size_pt, color_hex, bold=bold, align=align, anchor=anchor)
    return tb


def render(slide, data, tokens, bounds):
    x, y, w, h = bounds

    row_labels = list(data.get("row_labels", []))
    col_labels = list(data.get("col_labels", []))
    values = data.get("values", [])
    title = data.get("title")
    show_values = bool(data.get("show_values", False))
    value_format = data.get("value_format", "{:.2f}")

    n_rows = len(row_labels)
    n_cols = len(col_labels)
    if n_rows == 0 or n_cols == 0 or not values:
        return None

    flat = [v for row in values for v in row]
    vmin = data.get("value_min")
    vmax = data.get("value_max")
    if vmin is None:
        vmin = min(flat) if flat else 0.0
    if vmax is None:
        vmax = max(flat) if flat else 1.0
    if vmax == vmin:
        vmax = vmin + 1.0

    bg_hex = tokens["bg"]
    primary_hex = tokens["primary"]
    text_hex = tokens["text"]
    muted_hex = tokens["muted"]
    font_body = tokens["font_body"]
    font_display = tokens["font_display"]
    font_mono = tokens["font_mono"]
    base_pt = int(tokens["font_size_base_pt"])
    radius_px = int(tokens.get("radius_px") or 0)
    radius_emu = radius_px * 9525

    label_pt = max(7, int(round(base_pt * 0.72)))
    title_pt = max(base_pt + 2, int(round(base_pt * 1.15)))
    cell_value_pt = max(6, int(round(base_pt * 0.55)))
    legend_pt = max(6, int(round(base_pt * 0.6)))

    title_h = 0
    if title:
        title_h = int(Pt(title_pt * 1.6))

    legend_h = int(Pt(base_pt * 2.2))

    row_label_w = int(Pt(base_pt * 2.6))
    col_label_h = int(Pt(base_pt * 1.4))

    for lbl in row_labels:
        est = int(Pt(base_pt * 0.55 * max(3, len(str(lbl)))))
        if est > row_label_w:
            row_label_w = est
    row_label_w = min(row_label_w, int(w * 0.22))

    inner_x = x
    inner_y = y + title_h
    inner_w = w
    inner_h = h - title_h - legend_h

    grid_x = inner_x + row_label_w
    grid_y = inner_y + col_label_h
    grid_w = inner_w - row_label_w
    grid_h = inner_h - col_label_h

    gap = max(1, int(Emu(0.5 * 9525)))
    cell_w = (grid_w - gap * (n_cols - 1)) / n_cols
    cell_h = (grid_h - gap * (n_rows - 1)) / n_rows

    if title:
        _add_textbox(
            slide,
            x, y, w, title_h,
            title,
            font_display,
            title_pt,
            text_hex,
            bold=True,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )

    for ci, cl in enumerate(col_labels):
        cx = grid_x + ci * (cell_w + gap)
        _add_textbox(
            slide,
            int(cx), int(inner_y), int(cell_w), int(col_label_h),
            cl,
            font_body,
            label_pt,
            muted_hex,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    for ri, rl in enumerate(row_labels):
        ry = grid_y + ri * (cell_h + gap)
        pad_right = int(Pt(base_pt * 0.3))
        _add_textbox(
            slide,
            int(inner_x), int(ry), int(row_label_w - pad_right), int(cell_h),
            rl,
            font_body,
            label_pt,
            muted_hex,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    for ri in range(n_rows):
        row = values[ri] if ri < len(values) else []
        for ci in range(n_cols):
            v = row[ci] if ci < len(row) else vmin
            try:
                t = (float(v) - float(vmin)) / (float(vmax) - float(vmin))
            except Exception:
                t = 0.0
            t = max(0.0, min(1.0, t))
            fill_hex = _lerp_hex(bg_hex, primary_hex, t)

            cx = grid_x + ci * (cell_w + gap)
            cy = grid_y + ri * (cell_h + gap)

            _add_rect(
                slide,
                cx, cy, cell_w, cell_h,
                fill_hex,
                line_hex=muted_hex,
                line_w_emu=int(Emu(0.25 * 9525)),
                radius_emu=radius_emu,
            )

            if show_values:
                try:
                    vtxt = value_format.format(v)
                except Exception:
                    vtxt = str(v)
                cell_lum = _luminance(fill_hex)
                bg_lum = _luminance(bg_hex)
                if bg_lum < 0.5:
                    label_color = text_hex if cell_lum < 0.55 else bg_hex
                else:
                    label_color = bg_hex if cell_lum < 0.45 else text_hex
                _add_textbox(
                    slide,
                    int(cx), int(cy), int(cell_w), int(cell_h),
                    vtxt,
                    font_mono,
                    cell_value_pt,
                    label_color,
                    align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE,
                )

    legend_y = inner_y + inner_h
    legend_w = int(min(inner_w * 0.35, Pt(base_pt * 14)))
    legend_x = x + w - legend_w
    legend_bar_h = int(Pt(base_pt * 0.55))
    legend_bar_y = legend_y + int(Pt(base_pt * 0.3))

    steps = 24
    step_w = legend_w / steps
    for i in range(steps):
        t = i / (steps - 1) if steps > 1 else 0.0
        fill_hex = _lerp_hex(bg_hex, primary_hex, t)
        sx = legend_x + i * step_w
        _add_rect(
            slide,
            sx, legend_bar_y, step_w + Emu(1), legend_bar_h,
            fill_hex,
            line_hex=None,
            line_w_emu=0,
            radius_emu=0,
        )

    try:
        min_lbl = value_format.format(vmin)
        max_lbl = value_format.format(vmax)
    except Exception:
        min_lbl = f"{vmin}"
        max_lbl = f"{vmax}"

    lbl_y = legend_bar_y + legend_bar_h + int(Pt(base_pt * 0.15))
    lbl_h = int(Pt(legend_pt * 1.4))
    _add_textbox(
        slide,
        int(legend_x), int(lbl_y), int(legend_w / 2), lbl_h,
        min_lbl,
        font_mono,
        legend_pt,
        muted_hex,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
    )
    _add_textbox(
        slide,
        int(legend_x + legend_w / 2), int(lbl_y), int(legend_w / 2), lbl_h,
        max_lbl,
        font_mono,
        legend_pt,
        muted_hex,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.TOP,
    )

    return None
