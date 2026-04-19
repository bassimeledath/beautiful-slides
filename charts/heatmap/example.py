import os
import sys
from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.dirname(HERE))  # charts/ for shared tokens module

from tokens import MODES  # noqa: E402
from render import render, _rgb  # noqa: E402


def _data():
    rows = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    cols = ["00", "02", "04", "06", "08", "10", "12", "14", "16", "18", "20", "22"]
    base = [
        [0.08, 0.05, 0.03, 0.03, 0.18, 0.55, 0.72, 0.85, 0.88, 0.78, 0.60, 0.35],
        [0.07, 0.04, 0.03, 0.04, 0.22, 0.58, 0.76, 0.90, 0.92, 0.82, 0.62, 0.40],
        [0.09, 0.06, 0.04, 0.05, 0.24, 0.60, 0.80, 0.94, 0.96, 0.86, 0.66, 0.45],
        [0.10, 0.06, 0.04, 0.04, 0.26, 0.62, 0.82, 0.95, 0.97, 0.88, 0.70, 0.48],
        [0.14, 0.10, 0.06, 0.06, 0.28, 0.58, 0.78, 0.88, 0.84, 0.80, 0.72, 0.62],
        [0.22, 0.18, 0.12, 0.08, 0.16, 0.32, 0.50, 0.58, 0.64, 0.72, 0.80, 0.78],
        [0.25, 0.20, 0.14, 0.09, 0.14, 0.28, 0.42, 0.50, 0.54, 0.60, 0.66, 0.58],
    ]
    return {
        "title": "Usage intensity — weekday × hour",
        "row_labels": rows,
        "col_labels": cols,
        "values": base,
        "value_min": 0.0,
        "value_max": 1.0,
        "show_values": True,
        "value_format": "{:.0%}",
    }


def _make_pptx(mode_name, tokens, data, out_dir):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg_shape = slide.shapes.add_shape(
        1,
        0, 0,
        prs.slide_width, prs.slide_height,
    )
    bg_shape.shadow.inherit = False
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(tokens["bg"])
    bg_shape.line.fill.background()
    bg_shape.text_frame.text = ""

    margin = Inches(0.6)
    bounds = (
        margin,
        margin,
        prs.slide_width - 2 * margin,
        prs.slide_height - 2 * margin,
    )

    render(slide, data, tokens, bounds)

    out_path = os.path.join(out_dir, f"example-{mode_name}.pptx")
    prs.save(out_path)
    return out_path


def main():
    data = _data()
    written = []
    for name, tokens in MODES.items():
        path = _make_pptx(name, tokens, data, HERE)
        written.append(path)
        print(f"wrote {path}")
    print(f"OK: {len(written)} files")


if __name__ == "__main__":
    main()
