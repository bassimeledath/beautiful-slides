"""Render the bar chart in all 5 modes as separate 16:9 slides."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tokens import MODES  # noqa: E402

from render import render


DATA_SINGLE = {
    "orientation": "vertical",
    "title": "Revenue by segment, Q1",
    "categories": ["Enterprise", "Mid-market", "SMB", "Startup"],
    "series": [
        {"name": "Q1 Actual", "values": [12.4, 8.1, 4.3, 1.8]},
    ],
    "value_suffix": "M",
    "show_values": True,
}

DATA_GROUPED = {
    "orientation": "vertical",
    "title": "Q1 actual vs. plan, by segment",
    "categories": ["Enterprise", "Mid-market", "SMB", "Startup"],
    "series": [
        {"name": "Actual", "values": [12.4, 8.1, 4.3, 1.8]},
        {"name": "Plan",   "values": [11.0, 8.5, 5.0, 2.0]},
    ],
    "value_suffix": "M",
    "show_values": True,
}

DATA_HORIZONTAL = {
    "orientation": "horizontal",
    "title": "Revenue by segment, Q1",
    "categories": ["Enterprise", "Mid-market", "SMB", "Startup"],
    "series": [
        {"name": "Q1 Actual", "values": [12.4, 8.1, 4.3, 1.8]},
    ],
    "value_suffix": "M",
    "show_values": True,
}


# Pair each mode with a data variant so we exercise grouped + horizontal forms.
MODE_DATA = {
    "sv-keynote": DATA_SINGLE,
    "editorial-magazine": DATA_HORIZONTAL,
    "playful-marketing": DATA_GROUPED,
    "consulting-boardroom": DATA_GROUPED,
    "craft-minimal": DATA_SINGLE,
}


def main():
    out_dir = Path(__file__).parent
    for mode, tokens in MODES.items():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.5)
        bounds = (
            margin,
            margin,
            prs.slide_width - 2 * margin,
            prs.slide_height - 2 * margin,
        )

        data = MODE_DATA[mode]
        render(slide, data, tokens, bounds)

        out = out_dir / f"example-{mode}.pptx"
        prs.save(str(out))
        print(f"wrote {out.name}")


if __name__ == "__main__":
    main()
