"""Render KPI tiles in all 5 modes. Generates one .pptx per mode."""

import os
import sys

from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

from render import render, _rgb  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402


MODES = {
    "sv-keynote": {
        "primary": "#21D4FD",
        "accent":  "#17B26A",
        "text":    "#F5F7FA",
        "muted":   "#9AA4B2",
        "bg":      "#05070A",
        "font_display": "Manrope",
        "font_body":    "Manrope",
        "font_mono":    "JetBrains Mono",
        "font_size_base_pt": 18,
        "radius_px": 6,
    },
    "editorial-magazine": {
        "primary": "#8C2F39",
        "accent":  "#9C5B00",
        "text":    "#181514",
        "muted":   "#6F675F",
        "bg":      "#F6F1E8",
        "font_display": "Fraunces",
        "font_body":    "Newsreader",
        "font_mono":    "IBM Plex Mono",
        "font_size_base_pt": 16,
        "radius_px": 0,
    },
    "playful-marketing": {
        "primary": "#FF7A00",
        "accent":  "#0AB39C",
        "text":    "#1B1B1F",
        "muted":   "#6E6A73",
        "bg":      "#FFF4EB",
        "font_display": "Bricolage Grotesque",
        "font_body":    "Plus Jakarta Sans",
        "font_mono":    "Recursive Mono",
        "font_size_base_pt": 18,
        "radius_px": 12,
    },
    "consulting-boardroom": {
        "primary": "#0F4C81",
        "accent":  "#05603A",
        "text":    "#101828",
        "muted":   "#475467",
        "bg":      "#FFFFFF",
        "font_display": "Public Sans",
        "font_body":    "Public Sans",
        "font_mono":    "Public Sans",
        "font_size_base_pt": 14,
        "radius_px": 0,
    },
    "craft-minimal": {
        "primary": "#7C8571",
        "accent":  "#9A6B39",
        "text":    "#22201C",
        "muted":   "#7B776F",
        "bg":      "#FCFBF8",
        "font_display": "Instrument Serif",
        "font_body":    "Instrument Sans",
        "font_mono":    "Instrument Sans",
        "font_size_base_pt": 16,
        "radius_px": 2,
    },
}


TILES = [
    {
        "label": "ARR",
        "value": "$47.2M",
        "delta": "+12.4% vs plan",
        "delta_direction": "up",
        "footnote": "Source: NetSuite, Apr 10",
    },
    {
        "label": "Gross Margin",
        "value": "72.8%",
        "delta": "-1.1 pts QoQ",
        "delta_direction": "down",
        "footnote": "Non-GAAP, ex. one-time",
    },
    {
        "label": "NPS",
        "value": "64",
        "delta": "+6 vs Q4",
        "delta_direction": "up",
        "footnote": "n=1,842, mixed segment",
    },
]


def _paint_slide_bg(slide, hex_):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        slide.part.package.presentation_part.presentation.slide_width,
        slide.part.package.presentation_part.presentation.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(hex_)
    bg.line.fill.background()
    # Send to back.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def build(mode, tokens):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _paint_slide_bg(slide, tokens["bg"])

    sw = prs.slide_width
    sh = prs.slide_height

    # Layout: 3 tiles side-by-side.
    margin_x = int(sw * 0.06)
    margin_y_top = int(sh * 0.22)
    gutter = int(sw * 0.03)
    tile_area_w = sw - 2 * margin_x
    tile_w = (tile_area_w - 2 * gutter) // 3
    tile_h = int(sh * 0.48)

    for i, tile in enumerate(TILES):
        bx = margin_x + i * (tile_w + gutter)
        by = margin_y_top
        render(slide, tile, tokens, (bx, by, tile_w, tile_h))

    out = os.path.join(HERE, f"example-{mode}.pptx")
    prs.save(out)
    return out


def main():
    for mode, tokens in MODES.items():
        path = build(mode, tokens)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
